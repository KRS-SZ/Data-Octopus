from pptx import Presentation
from pptx.util import Inches, Pt, Emu
import sys
sys.stdout.reconfigure(encoding='utf-8')

pptx_path = r'c:\Users\szenklarz\Desktop\VS_Folder\PPT\Wafermap_Analysis_Report101.pptx'
prs = Presentation(pptx_path)

slide_w = prs.slide_width
slide_h = prs.slide_height
print(f'Slide size: {slide_w/914400:.2f}" x {slide_h/914400:.2f}"')
print(f'Total slides: {len(prs.slides)}')
print()

for i, slide in enumerate(prs.slides):
    title_text = ''
    for shape in slide.shapes:
        if shape.has_text_frame:
            txt = shape.text_frame.text[:80]
            if txt:
                title_text = txt
                break
    print(f'=== Slide {i+1}: {title_text} ===')
    for shape in slide.shapes:
        x = shape.left / 914400 if shape.left else 0
        y = shape.top / 914400 if shape.top else 0
        w = shape.width / 914400 if shape.width else 0
        h = shape.height / 914400 if shape.height else 0
        if shape.has_text_frame:
            txt = shape.text_frame.text[:60].replace('\n', ' | ')
            print(f'  TextBox: x={x:.2f}" y={y:.2f}" w={w:.2f}" h={h:.2f}"  "{txt}"')
        elif hasattr(shape, 'image'):
            try:
                _ = shape.image
                print(f'  Image:   x={x:.2f}" y={y:.2f}" w={w:.2f}" h={h:.2f}"')
            except:
                print(f'  Shape:   x={x:.2f}" y={y:.2f}" w={w:.2f}" h={h:.2f}"')
        else:
            print(f'  Shape:   x={x:.2f}" y={y:.2f}" w={w:.2f}" h={h:.2f}"  [{shape.name}]')
    print()
