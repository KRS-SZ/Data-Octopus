"""Extract images from Green Gage R&R PPT slides 71-95"""
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os

pptx_path = r'c:\Users\szenklarz\Desktop\VS_Folder\Data Octopus\PPT\Gage_R&R Green V1.0.pptx'
out_dir = r'c:\Users\szenklarz\Desktop\VS_Folder\Data Octopus\PPT\extracted_slides'
os.makedirs(out_dir, exist_ok=True)

prs = Presentation(pptx_path)

# Extract images from slides 71, 72, 73, 95 (most informative)
target_slides = [71, 72, 73, 74, 95]

for i, slide in enumerate(prs.slides):
    slide_num = i + 1
    if slide_num not in target_slides:
        continue
    img_idx = 0
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            img_idx += 1
            image = shape.image
            ext = image.content_type.split('/')[-1]
            if ext == 'jpeg':
                ext = 'jpg'
            fname = f'slide_{slide_num}_img{img_idx}.{ext}'
            fpath = os.path.join(out_dir, fname)
            with open(fpath, 'wb') as f:
                f.write(image.blob)
            print(f'Saved: {fname} ({len(image.blob)} bytes, {image.content_type})')

print(f'\nDone! Images in: {out_dir}')
