"""
PowerPoint Export Module for Data Octopus

Contains utility functions and classes for generating PowerPoint presentations
with wafermap, statistics, GRR analysis, and comparison slides.

Usage:
    from src.stdf_analyzer.services.ppt_export import PPTExporter
    exporter = PPTExporter()
    exporter.create_presentation(output_path)
"""

import os
import io
import tempfile
from typing import Optional, Dict, List, Any, Tuple, Union
from dataclasses import dataclass, field
from datetime import datetime

import numpy as np
import pandas as pd
import matplotlib.pyplot as plt
from matplotlib.figure import Figure

# python-pptx imports
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RgbColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.dml import MSO_THEME_COLOR
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    print("Warning: python-pptx not installed. PPT export disabled.")

from src.stdf_analyzer.core.parameter_utils import simplify_param_name


# ============================================================
# CONSTANTS
# ============================================================

# Slide dimensions (16:9 widescreen)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Default colors
COLORS = {
    'title_bg': RgbColor(0x00, 0x70, 0xC0) if PPTX_AVAILABLE else None,  # Blue
    'header_bg': RgbColor(0x40, 0x40, 0x40) if PPTX_AVAILABLE else None,  # Dark gray
    'pass': RgbColor(0x00, 0xB0, 0x50) if PPTX_AVAILABLE else None,  # Green
    'fail': RgbColor(0xFF, 0x00, 0x00) if PPTX_AVAILABLE else None,  # Red
    'warning': RgbColor(0xFF, 0xA5, 0x00) if PPTX_AVAILABLE else None,  # Orange
}

# GRR thresholds
GRR_EXCELLENT = 10  # %GRR < 10% = Excellent
GRR_ACCEPTABLE = 30  # %GRR < 30% = Acceptable
NDC_MIN = 5  # ndc >= 5 = Good


@dataclass
class SlideConfig:
    """Configuration for a slide."""
    title: str = ""
    subtitle: str = ""
    layout_index: int = 5  # Default: blank layout
    include_date: bool = True
    include_page_number: bool = True


@dataclass
class WafermapSlideConfig(SlideConfig):
    """Configuration for a wafermap slide."""
    wafer_id: str = ""
    parameter: str = ""
    display_name: str = ""
    image_path: Optional[str] = None
    statistics: Dict[str, float] = field(default_factory=dict)
    limits: Optional[Dict[str, float]] = None


@dataclass
class GRRSlideConfig(SlideConfig):
    """Configuration for a GRR slide."""
    parameter: str = ""
    grr_pct: float = 0.0
    ndc: float = 0.0
    repeatability: float = 0.0
    reproducibility: float = 0.0
    wafer_count: int = 0
    chart_images: List[str] = field(default_factory=list)


class PPTExporter:
    """
    PowerPoint Exporter for Data Octopus.

    Creates professional PowerPoint presentations with:
    - Title slide
    - Agenda slide
    - Wafermap slides (heatmap, binmap)
    - Statistics slides
    - GRR analysis slides
    - Multi-wafer comparison slides
    - Diffmap slides
    """

    def __init__(self, template_path: Optional[str] = None):
        """
        Initialize the PPT Exporter.

        Args:
            template_path: Optional path to a PPTX template file
        """
        if not PPTX_AVAILABLE:
            raise ImportError("python-pptx is required for PPT export")

        self.template_path = template_path
        self.presentation: Optional[Presentation] = None
        self.slide_count = 0
        self.agenda_items: List[Tuple[str, int]] = []

    def create_presentation(self, output_path: Optional[str] = None) -> Presentation:
        """
        Create a new presentation.

        Args:
            output_path: Optional path to save the presentation

        Returns:
            The created Presentation object
        """
        if self.template_path and os.path.exists(self.template_path):
            self.presentation = Presentation(self.template_path)
        else:
            self.presentation = Presentation()
            # Set slide dimensions to 16:9
            self.presentation.slide_width = SLIDE_WIDTH
            self.presentation.slide_height = SLIDE_HEIGHT

        self.slide_count = 0
        self.agenda_items = []

        return self.presentation

    def add_title_slide(self, title: str, subtitle: str = "",
                        author: str = "", date: Optional[str] = None) -> int:
        """
        Add a title slide.

        Args:
            title: Main title
            subtitle: Subtitle or description
            author: Author name
            date: Date string (defaults to today)

        Returns:
            Slide index
        """
        if self.presentation is None:
            self.create_presentation()

        # Get title layout (index 0)
        layout = self.presentation.slide_layouts[0]
        slide = self.presentation.slides.add_slide(layout)
        self.slide_count += 1

        # Set title
        if slide.shapes.title:
            slide.shapes.title.text = title

        # Set subtitle
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 1:  # Subtitle placeholder
                shape.text = subtitle
                if author:
                    shape.text += f"\n{author}"
                if date is None:
                    date = datetime.now().strftime("%Y-%m-%d")
                shape.text += f"\n{date}"
                break

        return self.slide_count - 1

    def add_agenda_slide(self, items: List[Tuple[str, int]]) -> int:
        """
        Add an agenda slide with section names and page numbers.

        Args:
            items: List of (section_name, page_number) tuples

        Returns:
            Slide index
        """
        if self.presentation is None:
            self.create_presentation()

        # Get blank layout
        layout = self.presentation.slide_layouts[5]
        slide = self.presentation.slides.add_slide(layout)
        self.slide_count += 1

        # Add title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(12), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = "Agenda"
        title_para.font.size = Pt(32)
        title_para.font.bold = True

        # Add agenda items
        y_pos = Inches(1.2)
        for section_name, page_num in items:
            item_box = slide.shapes.add_textbox(
                Inches(0.7), y_pos, Inches(10), Inches(0.4)
            )
            item_frame = item_box.text_frame
            item_para = item_frame.paragraphs[0]
            item_para.text = f"• {section_name}"
            item_para.font.size = Pt(18)

            # Page number on the right
            page_box = slide.shapes.add_textbox(
                Inches(11), y_pos, Inches(1.5), Inches(0.4)
            )
            page_frame = page_box.text_frame
            page_para = page_frame.paragraphs[0]
            page_para.text = f"Page {page_num}"
            page_para.font.size = Pt(18)
            page_para.alignment = PP_ALIGN.RIGHT

            y_pos += Inches(0.5)

        return self.slide_count - 1

    def add_wafermap_slide(self, config: WafermapSlideConfig,
                           fig: Optional[Figure] = None) -> int:
        """
        Add a wafermap slide with heatmap image and statistics.

        Args:
            config: WafermapSlideConfig with slide settings
            fig: Optional matplotlib Figure to embed

        Returns:
            Slide index
        """
        if self.presentation is None:
            self.create_presentation()

        layout = self.presentation.slide_layouts[5]  # Blank
        slide = self.presentation.slides.add_slide(layout)
        self.slide_count += 1

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.3), Inches(0.2), Inches(12), Inches(0.6)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = f"{config.display_name or config.parameter} - {config.wafer_id}"
        title_para.font.size = Pt(24)
        title_para.font.bold = True

        # Add image
        if fig is not None:
            # Save figure to bytes
            img_bytes = io.BytesIO()
            fig.savefig(img_bytes, format='png', dpi=150, bbox_inches='tight')
            img_bytes.seek(0)

            # Add to slide
            slide.shapes.add_picture(
                img_bytes, Inches(0.5), Inches(1), Inches(8), Inches(6)
            )
        elif config.image_path and os.path.exists(config.image_path):
            slide.shapes.add_picture(
                config.image_path, Inches(0.5), Inches(1), Inches(8), Inches(6)
            )

        # Statistics box
        if config.statistics:
            stats_box = slide.shapes.add_textbox(
                Inches(9), Inches(1), Inches(4), Inches(5)
            )
            stats_frame = stats_box.text_frame
            stats_frame.word_wrap = True

            # Header
            p = stats_frame.paragraphs[0]
            p.text = "Statistics"
            p.font.size = Pt(16)
            p.font.bold = True

            # Stats
            for key, value in config.statistics.items():
                p = stats_frame.add_paragraph()
                if isinstance(value, float):
                    p.text = f"{key}: {value:.4g}"
                else:
                    p.text = f"{key}: {value}"
                p.font.size = Pt(12)

        # Add to agenda
        self.agenda_items.append((f"Wafermap: {config.display_name}", self.slide_count))

        return self.slide_count - 1

    def add_grr_slide(self, config: GRRSlideConfig,
                      charts: Optional[List[Figure]] = None) -> int:
        """
        Add a GRR analysis slide.

        Args:
            config: GRRSlideConfig with GRR results
            charts: Optional list of matplotlib Figures (%GRR, ndc, Repeat/Reprod)

        Returns:
            Slide index
        """
        if self.presentation is None:
            self.create_presentation()

        layout = self.presentation.slide_layouts[5]  # Blank
        slide = self.presentation.slides.add_slide(layout)
        self.slide_count += 1

        param_label = simplify_param_name(config.parameter)

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.3), Inches(0.2), Inches(12), Inches(0.6)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = f"GRR Analysis: {param_label}"
        title_para.font.size = Pt(24)
        title_para.font.bold = True

        # Status indicator
        if config.grr_pct < GRR_EXCELLENT:
            status = "✅ EXCELLENT"
            status_color = COLORS['pass']
        elif config.grr_pct < GRR_ACCEPTABLE:
            status = "⚠️ ACCEPTABLE"
            status_color = COLORS['warning']
        else:
            status = "❌ UNACCEPTABLE"
            status_color = COLORS['fail']

        status_box = slide.shapes.add_textbox(
            Inches(10), Inches(0.2), Inches(3), Inches(0.5)
        )
        status_frame = status_box.text_frame
        status_para = status_frame.paragraphs[0]
        status_para.text = status
        status_para.font.size = Pt(18)
        status_para.font.bold = True
        status_para.alignment = PP_ALIGN.RIGHT

        # Results table
        results_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1), Inches(4), Inches(2.5)
        )
        results_frame = results_box.text_frame
        results_frame.word_wrap = True

        p = results_frame.paragraphs[0]
        p.text = "GRR Results"
        p.font.size = Pt(16)
        p.font.bold = True

        results_data = [
            (f"%GRR: {config.grr_pct:.2f}%", config.grr_pct < GRR_ACCEPTABLE),
            (f"ndc: {config.ndc:.2f}", config.ndc >= NDC_MIN),
            (f"Repeatability: {config.repeatability:.4g}", True),
            (f"Reproducibility: {config.reproducibility:.4g}", True),
            (f"Wafers: {config.wafer_count}", True),
        ]

        for text, is_good in results_data:
            p = results_frame.add_paragraph()
            p.text = text
            p.font.size = Pt(14)
            if not is_good:
                p.font.color.rgb = COLORS['fail']

        # Add charts
        if charts:
            chart_width = Inches(4)
            chart_height = Inches(3)
            x_positions = [Inches(0.5), Inches(4.5), Inches(8.5)]

            for i, (fig, x_pos) in enumerate(zip(charts, x_positions)):
                img_bytes = io.BytesIO()
                fig.savefig(img_bytes, format='png', dpi=100, bbox_inches='tight')
                img_bytes.seek(0)

                slide.shapes.add_picture(
                    img_bytes, x_pos, Inches(4), chart_width, chart_height
                )

        # Add to agenda
        self.agenda_items.append((f"GRR: {param_label}", self.slide_count))

        return self.slide_count - 1

    def add_comparison_slide(self, title: str,
                             wafer_ids: List[str],
                             fig: Optional[Figure] = None) -> int:
        """
        Add a multi-wafer comparison slide.

        Args:
            title: Slide title
            wafer_ids: List of wafer IDs being compared
            fig: Optional matplotlib Figure with comparison chart

        Returns:
            Slide index
        """
        if self.presentation is None:
            self.create_presentation()

        layout = self.presentation.slide_layouts[5]
        slide = self.presentation.slides.add_slide(layout)
        self.slide_count += 1

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.3), Inches(0.2), Inches(12), Inches(0.6)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = title
        title_para.font.size = Pt(24)
        title_para.font.bold = True

        # Subtitle with wafer IDs
        subtitle_box = slide.shapes.add_textbox(
            Inches(0.3), Inches(0.7), Inches(12), Inches(0.4)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.text = f"Wafers: {', '.join(wafer_ids[:5])}"
        if len(wafer_ids) > 5:
            subtitle_para.text += f" (+{len(wafer_ids) - 5} more)"
        subtitle_para.font.size = Pt(14)
        subtitle_para.font.color.rgb = RgbColor(0x80, 0x80, 0x80)

        # Add chart
        if fig is not None:
            img_bytes = io.BytesIO()
            fig.savefig(img_bytes, format='png', dpi=150, bbox_inches='tight')
            img_bytes.seek(0)

            slide.shapes.add_picture(
                img_bytes, Inches(0.5), Inches(1.2), Inches(12), Inches(6)
            )

        self.agenda_items.append((title, self.slide_count))

        return self.slide_count - 1

    def add_table_slide(self, title: str,
                        headers: List[str],
                        rows: List[List[Any]],
                        col_widths: Optional[List[float]] = None) -> int:
        """
        Add a slide with a data table.

        Args:
            title: Slide title
            headers: Column headers
            rows: Data rows
            col_widths: Optional column widths in inches

        Returns:
            Slide index
        """
        if self.presentation is None:
            self.create_presentation()

        layout = self.presentation.slide_layouts[5]
        slide = self.presentation.slides.add_slide(layout)
        self.slide_count += 1

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.3), Inches(0.2), Inches(12), Inches(0.6)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = title
        title_para.font.size = Pt(24)
        title_para.font.bold = True

        # Calculate table dimensions
        n_cols = len(headers)
        n_rows = len(rows) + 1  # +1 for header

        if col_widths is None:
            total_width = 12.0
            col_widths = [total_width / n_cols] * n_cols

        # Add table
        table = slide.shapes.add_table(
            n_rows, n_cols,
            Inches(0.5), Inches(1),
            Inches(sum(col_widths)), Inches(min(n_rows * 0.4, 6))
        ).table

        # Set column widths
        for i, width in enumerate(col_widths):
            table.columns[i].width = Inches(width)

        # Header row
        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = str(header)
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLORS['header_bg']
            para = cell.text_frame.paragraphs[0]
            para.font.bold = True
            para.font.size = Pt(12)
            para.font.color.rgb = RgbColor(0xFF, 0xFF, 0xFF)

        # Data rows
        for row_idx, row_data in enumerate(rows):
            for col_idx, value in enumerate(row_data):
                cell = table.cell(row_idx + 1, col_idx)
                if isinstance(value, float):
                    cell.text = f"{value:.4g}"
                else:
                    cell.text = str(value)
                para = cell.text_frame.paragraphs[0]
                para.font.size = Pt(11)

        return self.slide_count - 1

    def save(self, output_path: str) -> str:
        """
        Save the presentation to a file.

        Args:
            output_path: Path to save the PPTX file

        Returns:
            The saved file path
        """
        if self.presentation is None:
            raise ValueError("No presentation created")

        # Ensure directory exists
        os.makedirs(os.path.dirname(output_path) or '.', exist_ok=True)

        self.presentation.save(output_path)
        return output_path

    def get_bytes(self) -> bytes:
        """
        Get the presentation as bytes (for upload).

        Returns:
            PPTX file as bytes
        """
        if self.presentation is None:
            raise ValueError("No presentation created")

        pptx_bytes = io.BytesIO()
        self.presentation.save(pptx_bytes)
        pptx_bytes.seek(0)
        return pptx_bytes.read()


# ============================================================
# HELPER FUNCTIONS
# ============================================================

def create_wafermap_figure(df: pd.DataFrame,
                          parameter: str,
                          wafer_id: str = "",
                          figsize: Tuple[float, float] = (8, 8),
                          dpi: int = 100) -> Figure:
    """
    Create a wafermap figure for embedding in PPT.

    Args:
        df: DataFrame with x, y, and parameter columns
        parameter: Column name to visualize
        wafer_id: Wafer identifier for title
        figsize: Figure size in inches
        dpi: Resolution

    Returns:
        Matplotlib Figure
    """
    fig, ax = plt.subplots(figsize=figsize, dpi=dpi)

    x = df['x'].values
    y = df['y'].values
    values = pd.to_numeric(df[parameter], errors='coerce').values

    scatter = ax.scatter(x, y, c=values, cmap='viridis', s=50, marker='s')
    fig.colorbar(scatter, ax=ax, label=simplify_param_name(parameter))

    ax.set_title(f"{simplify_param_name(parameter)} - {wafer_id}")
    ax.set_xlabel("X")
    ax.set_ylabel("Y")
    ax.set_aspect('equal')

    fig.tight_layout()
    return fig


def create_grr_charts(grr_result: Dict,
                     param_label: str) -> List[Figure]:
    """
    Create GRR analysis charts.

    Args:
        grr_result: Dict with GRR analysis results
        param_label: Parameter display name

    Returns:
        List of 3 Figures: [%GRR, ndc, Repeat/Reprod]
    """
    figures = []

    # %GRR chart
    fig1, ax1 = plt.subplots(figsize=(4, 3), dpi=100)
    grr_pct = grr_result.get('grr_pct', 0)
    color = 'green' if grr_pct < 10 else ('orange' if grr_pct < 30 else 'red')
    ax1.bar([param_label], [grr_pct], color=color)
    ax1.axhline(y=10, color='g', linestyle='--', alpha=0.7)
    ax1.axhline(y=30, color='orange', linestyle='--', alpha=0.7)
    ax1.set_title("%GRR")
    ax1.set_ylabel("%GRR")
    ax1.set_ylim(0, max(grr_pct * 1.2, 35))
    fig1.tight_layout()
    figures.append(fig1)

    # ndc chart
    fig2, ax2 = plt.subplots(figsize=(4, 3), dpi=100)
    ndc = grr_result.get('ndc', 0)
    ndc_color = 'green' if ndc >= 5 else 'red'
    ax2.bar([param_label], [ndc], color=ndc_color)
    ax2.axhline(y=5, color='g', linestyle='--', alpha=0.7)
    ax2.set_title("ndc")
    ax2.set_ylabel("ndc")
    ax2.set_ylim(0, max(ndc * 1.2, 6))
    fig2.tight_layout()
    figures.append(fig2)

    # Repeatability/Reproducibility chart
    fig3, ax3 = plt.subplots(figsize=(4, 3), dpi=100)
    total_var = grr_result.get('total_variation', 1)
    repeat = grr_result.get('repeatability', 0)
    reprod = grr_result.get('reproducibility', 0)

    if total_var > 0:
        repeat_pct = 100 * (repeat**2) / (total_var**2)
        reprod_pct = 100 * (reprod**2) / (total_var**2)
    else:
        repeat_pct = 0
        reprod_pct = 0

    ax3.bar([param_label], [repeat_pct], label='Repeatability', color='#2196F3')
    ax3.bar([param_label], [reprod_pct], bottom=[repeat_pct],
           label='Reproducibility', color='#FF9800')
    ax3.set_title("Repeat. & Reprod.")
    ax3.set_ylabel("%")
    ax3.legend(fontsize=8)
    ax3.set_ylim(0, 100)
    fig3.tight_layout()
    figures.append(fig3)

    return figures
