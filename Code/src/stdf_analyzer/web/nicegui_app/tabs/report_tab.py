"""
Report Tab - Web Version.

PowerPoint Report Generation.
"""

from nicegui import ui
from typing import Any, Dict, List, Optional
import io
from datetime import datetime


class ReportTab:
    """
    Report Tab für PowerPoint-Export.
    """

    def __init__(self, app_state: Any):
        self.app_state = app_state
        self.container = None

    def build(self, parent: Any = None) -> None:
        """Tab UI aufbauen."""
        with ui.column().classes('w-full h-full p-2') as self.container:
            ui.label('PowerPoint Report Generator').classes('text-xl font-bold')

            with ui.row().classes('w-full gap-4'):
                # Left: Options
                with ui.card().classes('w-80'):
                    ui.label('Report Options').classes('font-bold')

                    # Wafer Selection
                    ui.label('Wafers:').classes('mt-2')
                    self.wafer_select = ui.select(
                        options=[],
                        multiple=True,
                        label='Select Wafers'
                    ).classes('w-full')

                    # Group Selection
                    ui.label('Groups:').classes('mt-2')
                    self.group_select = ui.select(
                        options=['All Groups', 'DC', 'ANLG', 'FUNC', 'OPTIC'],
                        value='All Groups',
                        multiple=True
                    ).classes('w-full')

                    # Slide Types
                    ui.label('Include Slides:').classes('mt-2 font-bold')
                    self.include_title = ui.checkbox('Title Slide', value=True)
                    self.include_agenda = ui.checkbox('Agenda', value=True)
                    self.include_wafermap = ui.checkbox('Wafermaps', value=True)
                    self.include_stats = ui.checkbox('Statistics', value=True)
                    self.include_boxplot = ui.checkbox('Boxplots', value=False)
                    self.include_distribution = ui.checkbox('Distributions', value=False)

                    # Generate Button
                    ui.button('Generate Report', on_click=self._generate_report).classes(
                        'w-full mt-4 bg-green-600 text-white'
                    )

                # Right: Preview
                with ui.card().classes('flex-1'):
                    ui.label('Preview').classes('font-bold')
                    self.preview_area = ui.column().classes('w-full h-96 bg-gray-100 p-2')
                    with self.preview_area:
                        ui.label('Report preview will appear here...').classes('text-gray-500')

    def _generate_report(self) -> None:
        """PowerPoint Report generieren."""
        wafers = self.wafer_select.value or []

        if not wafers:
            ui.notify('Please select at least one wafer', type='warning')
            return

        ui.notify('Generating report...', type='info')

        # TODO: Implement PPTX generation
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt

            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)

            # Title Slide
            if self.include_title.value:
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                # Add title text box

            # Save to bytes
            buffer = io.BytesIO()
            prs.save(buffer)
            buffer.seek(0)

            # Trigger download
            filename = f"Report_{datetime.now().strftime('%Y%m%d_%H%M')}.pptx"
            ui.download(buffer.getvalue(), filename)
            ui.notify('Report generated successfully!', type='positive')

        except ImportError:
            ui.notify('python-pptx not installed. Run: pip install python-pptx', type='negative')
        except Exception as e:
            ui.notify(f'Error generating report: {e}', type='negative')

    def update_wafers(self, wafers: List[str]) -> None:
        """Wafer-Liste aktualisieren."""
        self.wafer_select.options = wafers
        self.wafer_select.value = wafers
        self.wafer_select.update()
