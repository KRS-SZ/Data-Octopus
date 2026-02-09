#!/usr/bin/env python3
# Python
# from Semi_ATE.STDF.STDFFile import STDFFile

import sys

print(sys.executable)

import os
import threading
import queue
from concurrent.futures import ThreadPoolExecutor, as_completed
from dataclasses import dataclass, field
from typing import Any, Optional, Callable
import time

# Try to set MATPLOTLIBDATA for Meta environments
try:
    from libfb import parutil

    os.environ["MATPLOTLIBDATA"] = parutil.get_dir_path("matplotlib/mpl-data")
except ImportError:
    pass

import tkinter as tk
from tkinter import filedialog, ttk

import matplotlib
import matplotlib.pyplot as plt
import numpy as np
from matplotlib.backends.backend_tkagg import FigureCanvasTkAgg, NavigationToolbar2Tk
from PIL import Image, ImageTk

# Robust STDF import - try to import the module
STDF_MODULE = None
STDF_TYPE = None

try:
    import Semi_ATE.STDF

    STDF_MODULE = Semi_ATE.STDF
    STDF_TYPE = "Semi_ATE"
    print(f"Successfully imported Semi_ATE.STDF module")
    print(
        f"Available functions: {[x for x in dir(STDF_MODULE) if not x.startswith('_')][:15]}"
    )
except ImportError as e:
    print(f"Failed to import Semi_ATE.STDF: {e}")
    try:
        import pystdf

        STDF_MODULE = pystdf
        STDF_TYPE = "pystdf"
        print("Successfully imported pystdf")
    except ImportError:
        print("Warning: No STDF library found. STDF file features will not work.")
        STDF_MODULE = None
        STDF_TYPE = None

import pandas as pd

# Try to import python-pptx for PowerPoint creation
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from io import BytesIO
    PPTX_AVAILABLE = True
    print("Successfully imported python-pptx module")
except ImportError as e:
    PPTX_AVAILABLE = False
    print(f"Warning: python-pptx import failed: {e}")
    print("Install with: pip install python-pptx")

# Try to import Google API libraries for Google Drive/Slides integration
GOOGLE_API_AVAILABLE = False
GEMINI_AVAILABLE = False
try:
    from google.oauth2.credentials import Credentials
    from google_auth_oauthlib.flow import InstalledAppFlow
    from google.auth.transport.requests import Request
    from googleapiclient.discovery import build
    from googleapiclient.http import MediaFileUpload, MediaIoBaseUpload
    import pickle
    GOOGLE_API_AVAILABLE = True
    print("Successfully imported Google API libraries")
except ImportError as e:
    print(f"Warning: Google API libraries not available: {e}")
    print("Install with: pip install google-api-python-client google-auth-httplib2 google-auth-oauthlib")

# Try to import Google Generative AI (Gemini)
try:
    import google.generativeai as genai
    GEMINI_AVAILABLE = True
    print("Successfully imported Google Generative AI (Gemini)")
except ImportError as e:
    print(f"Warning: Google Generative AI not available: {e}")
    print("Install with: pip install google-generativeai")

# Google API Scopes
GOOGLE_SCOPES = [
    'https://www.googleapis.com/auth/drive.file',
    'https://www.googleapis.com/auth/presentations',
]

print("Starting application...")

# ============================================================================
# Multithreading Support with Priority Queue
# ============================================================================

# Thread pool for background loading
MAX_WORKERS = 4  # Number of parallel loading threads
thread_pool = ThreadPoolExecutor(max_workers=MAX_WORKERS)

# Priority levels (lower number = higher priority)
PRIORITY_HIGH = 1
PRIORITY_NORMAL = 2
PRIORITY_LOW = 3


@dataclass(order=True)
class PrioritizedTask:
    """Task with priority for the loading queue"""
    priority: int
    task_id: int = field(compare=False)
    func: Callable = field(compare=False)
    args: tuple = field(compare=False, default=())
    kwargs: dict = field(compare=False, default_factory=dict)
    callback: Optional[Callable] = field(compare=False, default=None)
    error_callback: Optional[Callable] = field(compare=False, default=None)


# Global task counter for unique task IDs
_task_counter = 0
_task_lock = threading.Lock()

# Loading state tracking
_loading_tasks = {}
_loading_lock = threading.Lock()


def get_next_task_id():
    """Get a unique task ID"""
    global _task_counter
    with _task_lock:
        _task_counter += 1
        return _task_counter


def submit_loading_task(func, args=(), kwargs=None, callback=None,
                        error_callback=None, priority=PRIORITY_NORMAL):
    """Submit a loading task to the thread pool with priority"""
    if kwargs is None:
        kwargs = {}

    task_id = get_next_task_id()

    def wrapped_task():
        try:
            result = func(*args, **kwargs)
            if callback:
                # Schedule callback on main thread
                main_win.after(0, lambda: callback(result))
            return result
        except Exception as e:
            print(f"Error in loading task {task_id}: {e}")
            if error_callback:
                main_win.after(0, lambda: error_callback(e))
            raise

    future = thread_pool.submit(wrapped_task)

    with _loading_lock:
        _loading_tasks[task_id] = {
            'future': future,
            'priority': priority,
            'start_time': time.time()
        }

    return task_id, future


def cancel_loading_task(task_id):
    """Cancel a loading task if possible"""
    with _loading_lock:
        if task_id in _loading_tasks:
            future = _loading_tasks[task_id]['future']
            cancelled = future.cancel()
            if cancelled:
                del _loading_tasks[task_id]
            return cancelled
    return False


def get_loading_progress():
    """Get current loading progress"""
    with _loading_lock:
        total = len(_loading_tasks)
        completed = sum(1 for t in _loading_tasks.values() if t['future'].done())
        return completed, total


# Initialize empty data storage
data_arrays = []
file_names = []

# ============================================================================
# Google Drive / Google Slides Integration
# ============================================================================

# Store Google credentials globally
google_credentials = None

def get_google_credentials():
    """Get or refresh Google API credentials"""
    global google_credentials

    if not GOOGLE_API_AVAILABLE:
        return None

    creds = None
    token_file = os.path.join(os.path.dirname(__file__), 'google_token.pickle')
    credentials_file = os.path.join(os.path.dirname(__file__), 'credentials.json')

    # Load existing token
    if os.path.exists(token_file):
        try:
            with open(token_file, 'rb') as token:
                creds = pickle.load(token)
        except Exception as e:
            print(f"Error loading token: {e}")

    # Refresh or get new credentials
    if not creds or not creds.valid:
        if creds and creds.expired and creds.refresh_token:
            try:
                creds.refresh(Request())
            except Exception as e:
                print(f"Error refreshing credentials: {e}")
                creds = None

        if not creds:
            if os.path.exists(credentials_file):
                try:
                    flow = InstalledAppFlow.from_client_secrets_file(credentials_file, GOOGLE_SCOPES)
                    creds = flow.run_local_server(port=0)
                except Exception as e:
                    print(f"Error getting credentials: {e}")
                    return None
            else:
                print(f"credentials.json not found at {credentials_file}")
                return None

        # Save credentials
        try:
            with open(token_file, 'wb') as token:
                pickle.dump(creds, token)
        except Exception as e:
            print(f"Error saving token: {e}")

    google_credentials = creds
    return creds


def upload_to_google_drive(file_path, file_name=None, convert_to_slides=True):
    """Upload a file to Google Drive and optionally convert to Google Slides"""
    creds = get_google_credentials()
    if not creds:
        return None, "Google credentials not available. Please set up credentials.json"

    try:
        service = build('drive', 'v3', credentials=creds)

        if file_name is None:
            file_name = os.path.basename(file_path)

        # File metadata
        file_metadata = {'name': file_name}

        # If converting to Google Slides
        if convert_to_slides and file_path.endswith('.pptx'):
            file_metadata['mimeType'] = 'application/vnd.google-apps.presentation'

        media = MediaFileUpload(file_path, resumable=True)

        file = service.files().create(
            body=file_metadata,
            media_body=media,
            fields='id, webViewLink'
        ).execute()

        file_id = file.get('id')
        web_link = file.get('webViewLink')

        print(f"File uploaded successfully. ID: {file_id}")
        return file_id, web_link

    except Exception as e:
        print(f"Error uploading to Google Drive: {e}")
        return None, str(e)


def upload_pptx_bytes_to_drive(pptx_bytes, file_name, convert_to_slides=True):
    """Upload PowerPoint bytes directly to Google Drive"""
    creds = get_google_credentials()
    if not creds:
        return None, "Google credentials not available"

    try:
        service = build('drive', 'v3', credentials=creds)

        file_metadata = {'name': file_name}
        if convert_to_slides:
            file_metadata['mimeType'] = 'application/vnd.google-apps.presentation'

        media = MediaIoBaseUpload(
            BytesIO(pptx_bytes),
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
            resumable=True
        )

        file = service.files().create(
            body=file_metadata,
            media_body=media,
            fields='id, webViewLink'
        ).execute()

        return file.get('id'), file.get('webViewLink')

    except Exception as e:
        return None, str(e)


def beautify_slides_with_gemini(presentation_id, api_key=None):
    """Use Gemini AI to suggest improvements for slides"""
    if not GEMINI_AVAILABLE:
        return None, "Gemini AI not available. Install with: pip install google-generativeai"

    if not api_key:
        api_key = os.environ.get('GEMINI_API_KEY') or os.environ.get('GOOGLE_API_KEY')

    if not api_key:
        return None, "Gemini API key not found. Set GEMINI_API_KEY environment variable."

    try:
        genai.configure(api_key=api_key)
        model = genai.GenerativeModel('gemini-pro')

        # Get slides content
        creds = get_google_credentials()
        if not creds:
            return None, "Google credentials not available"

        slides_service = build('slides', 'v1', credentials=creds)
        presentation = slides_service.presentations().get(presentationId=presentation_id).execute()

        # Extract slide content for analysis
        slides_content = []
        for slide in presentation.get('slides', []):
            slide_text = []
            for element in slide.get('pageElements', []):
                if 'shape' in element and 'text' in element.get('shape', {}):
                    for text_elem in element['shape']['text'].get('textElements', []):
                        if 'textRun' in text_elem:
                            slide_text.append(text_elem['textRun'].get('content', ''))
            slides_content.append(' '.join(slide_text))

        # Generate beautification suggestions
        prompt = f"""Analyze this presentation content and provide specific suggestions to improve the visual design and clarity.

Slide contents:
{chr(10).join([f"Slide {i+1}: {content[:500]}" for i, content in enumerate(slides_content)])}

Provide suggestions for:
1. Color scheme improvements
2. Layout optimization
3. Text clarity and formatting
4. Visual hierarchy
5. Chart/graph improvements

Format as actionable bullet points."""

        response = model.generate_content(prompt)
        suggestions = response.text

        return suggestions, None

    except Exception as e:
        return None, str(e)


def apply_gemini_theme_to_slides(presentation_id, theme_style="professional"):
    """Apply a theme style to Google Slides using Gemini suggestions"""
    if not GOOGLE_API_AVAILABLE:
        return False, "Google API not available"

    creds = get_google_credentials()
    if not creds:
        return False, "Google credentials not available"

    try:
        slides_service = build('slides', 'v1', credentials=creds)

        # Define color schemes
        themes = {
            "professional": {
                "primary": {"red": 0.2, "green": 0.4, "blue": 0.6},
                "secondary": {"red": 0.9, "green": 0.9, "blue": 0.9},
                "accent": {"red": 0.2, "green": 0.6, "blue": 0.4}
            },
            "modern": {
                "primary": {"red": 0.1, "green": 0.1, "blue": 0.1},
                "secondary": {"red": 1.0, "green": 1.0, "blue": 1.0},
                "accent": {"red": 0.0, "green": 0.7, "blue": 0.9}
            },
            "vibrant": {
                "primary": {"red": 0.9, "green": 0.3, "blue": 0.2},
                "secondary": {"red": 1.0, "green": 0.95, "blue": 0.9},
                "accent": {"red": 0.2, "green": 0.7, "blue": 0.3}
            }
        }

        selected_theme = themes.get(theme_style, themes["professional"])

        # Get presentation
        presentation = slides_service.presentations().get(presentationId=presentation_id).execute()

        requests = []

        # Update slide backgrounds
        for slide in presentation.get('slides', []):
            slide_id = slide.get('objectId')
            requests.append({
                'updatePageProperties': {
                    'objectId': slide_id,
                    'pageProperties': {
                        'pageBackgroundFill': {
                            'solidFill': {
                                'color': {
                                    'rgbColor': selected_theme["secondary"]
                                }
                            }
                        }
                    },
                    'fields': 'pageBackgroundFill'
                }
            })

        if requests:
            slides_service.presentations().batchUpdate(
                presentationId=presentation_id,
                body={'requests': requests}
            ).execute()

        return True, f"Applied {theme_style} theme successfully"

    except Exception as e:
        return False, str(e)


# Language translations dictionary
TRANSLATIONS = {
    "en": {
        "window_title": "Measurement Data Visualization",
        "tab_plot": "Plot",
        "tab_boxplot": "Boxplot",
        "tab_probability": "Probability",
        "tab_images": "Images",
        "tab_wafermap": "Wafermap",
        "tab_stdf_heatmap": "STDF Heatmap",
        "select_measurement_files": "Select Measurement Files",
        "no_data_loaded": "No data loaded. Click 'Select Measurement Files' to begin.",
        "select_stdf_file": "Select STDF File",
        "file": "File:",
        "parameter": "Parameter:",
        "refresh": "Refresh",
        "load_multiple_stdf": "Load Multiple STDF Files",
        "no_stdf_loaded": "No STDF files loaded",
        "show_grid": "Show Grid",
        "select_images": "Select Image(s)",
        "zoom_in": "Zoom In",
        "zoom_out": "Zoom Out",
        "language": "Language:",
        "heatmap_title": "Heatmap:",
        "wafer": "Wafer",
        "x_coordinate": "X Coordinate",
        "y_coordinate": "Y Coordinate",
        "column_index": "Column Index",
        "row_index": "Row Index",
        "min": "Min",
        "max": "Max",
        "mean": "Mean",
        "median": "Median",
        "dies": "Dies",
        "die_coordinates": "Die Coordinates",
        "additional_info": "Additional Information",
        "bin": "Bin",
        "boxplot_title": "Boxplot of Each Dataset (with Mean and Median)",
        "dataset": "Dataset",
        "luminance_value": "Luminance Value",
        "probability_title": "Probability Plot of Measurement Values",
        "value": "Value",
        "probability": "Probability",
        "histogram_title": "Combined Histogram of Measurement Values",
        "frequency": "Frequency",
        "zoomed": "Zoomed",
        "die_info": "Die Information",
        "loaded_files": "Loaded {count} file(s)",
        "loaded_params": "Loaded: {count} parameters available",
        "tab_presentation": "Create Presentation",
        "create_pptx": "Create PowerPoint",
        "select_sources": "Select Data Sources",
        "presentation_title": "Presentation Title:",
        "include_statistics": "Include Statistics",
        "include_wafermap": "Include Wafermap",
        "include_multiple_wafermaps": "Include Multiple Wafermaps",
        "include_diffmap": "Include Diffmap",
        "select_wafermaps": "Select Wafermaps to Include:",
        "generating_pptx": "Generating PowerPoint...",
        "pptx_created": "PowerPoint created successfully!",
        "pptx_error": "Error creating PowerPoint",
        "no_data_for_pptx": "No data available for presentation",
    },
    "de": {
        "window_title": "Messdaten-Visualisierung",
        "tab_plot": "Diagramm",
        "tab_boxplot": "Boxplot",
        "tab_probability": "Wahrscheinlichkeit",
        "tab_images": "Bilder",
        "tab_wafermap": "Wafermap",
        "tab_stdf_heatmap": "STDF Heatmap",
        "select_measurement_files": "Messdateien auswählen",
        "no_data_loaded": "Keine Daten geladen. Klicken Sie auf 'Messdateien auswählen'.",
        "select_stdf_file": "STDF-Datei auswählen",
        "file": "Datei:",
        "parameter": "Parameter:",
        "refresh": "Aktualisieren",
        "load_multiple_stdf": "Mehrere STDF-Dateien laden",
        "no_stdf_loaded": "Keine STDF-Dateien geladen",
        "show_grid": "Raster anzeigen",
        "select_images": "Bild(er) auswählen",
        "zoom_in": "Vergrößern",
        "zoom_out": "Verkleinern",
        "language": "Sprache:",
        "heatmap_title": "Heatmap:",
        "wafer": "Wafer",
        "x_coordinate": "X-Koordinate",
        "y_coordinate": "Y-Koordinate",
        "column_index": "Spaltenindex",
        "row_index": "Zeilenindex",
        "min": "Min",
        "max": "Max",
        "mean": "Mittelwert",
        "median": "Median",
        "dies": "Dies",
        "die_coordinates": "Die-Koordinaten",
        "additional_info": "Zusätzliche Informationen",
        "bin": "Bin",
        "boxplot_title": "Boxplot der Datensätze (mit Mittelwert und Median)",
        "dataset": "Datensatz",
        "luminance_value": "Luminanzwert",
        "probability_title": "Wahrscheinlichkeitsdiagramm der Messwerte",
        "value": "Wert",
        "probability": "Wahrscheinlichkeit",
        "histogram_title": "Kombiniertes Histogramm der Messwerte",
        "frequency": "Häufigkeit",
        "zoomed": "Vergrößert",
        "die_info": "Die-Information",
        "loaded_files": "{count} Datei(en) geladen",
        "loaded_params": "Geladen: {count} Parameter verfügbar",
        "tab_presentation": "Präsentation erstellen",
        "create_pptx": "PowerPoint erstellen",
        "select_sources": "Datenquellen auswählen",
        "presentation_title": "Präsentationstitel:",
        "include_statistics": "Statistiken einschließen",
        "include_wafermap": "Wafermap einschließen",
        "include_multiple_wafermaps": "Mehrere Wafermaps einschließen",
        "include_diffmap": "Diffmap einschließen",
        "select_wafermaps": "Wafermaps auswählen:",
        "generating_pptx": "PowerPoint wird erstellt...",
        "pptx_created": "PowerPoint erfolgreich erstellt!",
        "pptx_error": "Fehler beim Erstellen der PowerPoint",
        "no_data_for_pptx": "Keine Daten für Präsentation verfügbar",
    },
}

# Current language
current_language = "en"


def get_text(key):
    """Get translated text for the current language"""
    return TRANSLATIONS.get(current_language, TRANSLATIONS["en"]).get(key, key)

# Create main Tkinter window with tabs and large tab font
print("Creating main window...")
main_win = tk.Tk()
main_win.title("Measurement Data Visualization")
main_win.geometry("1200x800")
print("Main window created successfully")

style = ttk.Style(main_win)
style.configure("TNotebook.Tab", font=("Helvetica", 20, "bold"))

# Top bar for language selection
top_bar = tk.Frame(main_win)
top_bar.pack(side=tk.TOP, fill=tk.X, padx=10, pady=5)

language_label = tk.Label(top_bar, text="Language:", font=("Helvetica", 10))
language_label.pack(side=tk.RIGHT, padx=5)

language_var = tk.StringVar(value="English")
language_combobox = ttk.Combobox(
    top_bar,
    textvariable=language_var,
    values=["English", "Deutsch"],
    state="readonly",
    width=10,
    font=("Helvetica", 10),
)
language_combobox.pack(side=tk.RIGHT, padx=5)


def change_language(event=None):
    """Change the application language and update all UI elements"""
    global current_language

    selected_lang = language_var.get()
    current_language = "de" if selected_lang == "Deutsch" else "en"

    # Update window title
    main_win.title(get_text("window_title"))

    # Update tab names
    notebook.tab(tab7, text="Configuration")
    notebook.tab(tab6, text="Wafermap")
    notebook.tab(tab_multi_wafer, text="Multiple Wafermaps")

    # Update buttons and labels
    select_multiple_stdf_button.config(text=get_text("load_multiple_stdf"))
    heatmap_param_label.config(text=get_text("parameter"))
    heatmap_refresh_button.config(text=get_text("refresh"))
    show_grid_checkbox.config(text=get_text("show_grid"))
    language_label.config(text=get_text("language"))

    print(f"Language changed to: {selected_lang}")


language_combobox.bind("<<ComboboxSelected>>", change_language)

notebook = ttk.Notebook(main_win)
notebook.pack(fill="both", expand=True)

# Tab 1: Configuration (moved to first position)
tab7 = ttk.Frame(notebook)
notebook.add(tab7, text="⚙ Config")

# Tab 2: Wafermap (renamed from STDF Heatmap)
tab6 = ttk.Frame(notebook)
notebook.add(tab6, text="🗺 Wafer")

# Tab 3: Multiple Wafermaps
tab_multi_wafer = ttk.Frame(notebook)
notebook.add(tab_multi_wafer, text="📊 Multi-Wafer")

# Tab 4: Diffmap - Difference between reference and comparison wafermaps
tab_diffmap = ttk.Frame(notebook)
notebook.add(tab_diffmap, text="🔄 Diffmap")

# Tab 5: Gage R&R Analysis
tab_grr = ttk.Frame(notebook)
notebook.add(tab_grr, text="📏 Gage R&R")

# Tab 6: Create Presentation
tab_presentation = ttk.Frame(notebook)
notebook.add(tab_presentation, text="📑 Create PPT")

# Global variables for plot canvases
canvas1 = None
canvas1_hist = None
fig1 = None
fig1_hist = None
axes1 = None


def load_measurement_files():
    """Load measurement data files and update all visualizations"""
    global data_arrays, file_names, canvas1, canvas1_hist, fig1, fig1_hist, axes1

    file_paths = filedialog.askopenfilenames(
        title="Select one or more measurement data files",
        filetypes=(("Text files", "*.txt"), ("All files", "*.*")),
    )

    if not file_paths:
        print("No files selected.")
        return

    # Clear existing data
    data_arrays = []
    file_names = []

    # Load files
    for file_path in file_paths:
        try:
            with open(file_path, "r") as f:
                lines = f.readlines()

            data_start = 0
            for i, line in enumerate(lines):
                if line.strip() and line[0].isdigit():
                    data_start = i
                    break

            raw_data = lines[data_start:]
            matrix = [
                list(map(float, row.strip().split(",")))
                for row in raw_data
                if row.strip()
            ]

            arr = np.array(matrix)
            data_arrays.append(arr)
            file_names.append(file_path.split("/")[-1])
            print(f"Loaded: {file_path.split('/')[-1]} - Shape: {arr.shape}")

        except Exception as e:
            print(f"Error loading {file_path}: {e}")

    if not data_arrays:
        print("No valid data files loaded.")
        return

    # Clear existing canvases
    if canvas1:
        canvas1.get_tk_widget().destroy()
    if canvas1_hist:
        canvas1_hist.get_tk_widget().destroy()
    if fig1:
        plt.close(fig1)
    if fig1_hist:
        plt.close(fig1_hist)

    # Create heatmap plots
    fig1, axes1 = plt.subplots(
        1, len(data_arrays), figsize=(8 * len(data_arrays), 8), constrained_layout=True
    )
    if len(data_arrays) == 1:
        axes1 = [axes1]

    for idx, arr in enumerate(data_arrays):
        ax = axes1[idx]
        im = ax.imshow(arr, cmap="viridis", aspect="equal")
        ax.set_title(f"Heatmap: {file_names[idx]}")
        ax.set_xlabel("Column Index")
        ax.set_ylabel("Row Index")
        fig1.colorbar(im, ax=ax, fraction=0.046, pad=0.04)

    # Create histogram
    fig1_hist, ax_hist = plt.subplots(figsize=(8, 4))
    for idx, arr in enumerate(data_arrays):
        ax_hist.hist(arr.flatten(), bins=100, alpha=0.5, label=file_names[idx])
    ax_hist.set_title("Combined Histogram of Measurement Values")
    ax_hist.set_xlabel("Luminance Value")
    ax_hist.set_ylabel("Frequency")
    ax_hist.legend(loc="upper right")
    fig1_hist.tight_layout()

    # Create canvases
    canvas1 = FigureCanvasTkAgg(fig1, master=plot_frame)
    canvas1.get_tk_widget().pack(side=tk.TOP, fill=tk.BOTH, expand=True)

    canvas1_hist = FigureCanvasTkAgg(fig1_hist, master=plot_frame)
    canvas1_hist.get_tk_widget().pack(side=tk.BOTTOM, fill=tk.BOTH, expand=True)

    # Connect double-click event
    canvas1.mpl_connect("button_press_event", on_subplot_doubleclick)

    # Update other tabs
    update_boxplot_tab()
    update_probability_tab()

    print(f"Successfully loaded {len(data_arrays)} file(s)")


def on_subplot_doubleclick(event):
    """Handle double-click on subplot to show zoomed view"""
    if event.dblclick and axes1 is not None:
        for idx, ax in enumerate(axes1):
            if event.inaxes == ax:
                arr = data_arrays[idx]
                title = file_names[idx]

                zoom_win = tk.Toplevel(main_win)
                zoom_win.title(f"Zoomed: {title}")

                fig_zoom, ax_zoom = plt.subplots(figsize=(16, 12))
                im_zoom = ax_zoom.imshow(arr, cmap="viridis", aspect="equal")
                ax_zoom.set_title(f"Zoomed: {title}")
                ax_zoom.set_xlabel("Column Index")
                ax_zoom.set_ylabel("Row Index")
                fig_zoom.colorbar(im_zoom, ax=ax_zoom, fraction=0.046, pad=0.04)

                canvas_zoom = FigureCanvasTkAgg(fig_zoom, master=zoom_win)
                canvas_zoom.get_tk_widget().pack(fill=tk.BOTH, expand=True)

                def on_close():
                    plt.close(fig_zoom)
                    zoom_win.destroy()

                zoom_win.protocol("WM_DELETE_WINDOW", on_close)
                break


def update_boxplot_tab():
    """Update boxplot visualization with mean and median values - NOT USED (tab deleted)"""
    pass


def update_probability_tab():
    """Update probability plot visualization - NOT USED (tab deleted)"""
    pass

wafermap_canvas = None
current_stdf_data = None
current_wafer_id = None
test_parameters = {}
grouped_parameters = {}  # Dictionary: group_name -> list of (test_num, short_name, full_name)
test_limits = {}  # Dictionary: test_num -> {'lo_limit': value, 'hi_limit': value, 'units': str}

# Wafer Configuration (from WCR record)
current_wafer_config = {
    'notch_orientation': None,  # 'U', 'D', 'L', 'R' or None
    'wafer_size': None,         # Wafer diameter in mm
    'die_width': None,
    'die_height': None,
    'pos_x': None,              # Positive X direction
    'pos_y': None               # Positive Y direction
}


def select_stdf_file():
    stdf_path = filedialog.askopenfilename(
        title="Select an STDF file",
        filetypes=[("STDF files", "*.stdf"), ("All files", "*.*")],
    )

    if stdf_path:
        load_stdf_data(stdf_path)


def update_source_buttons():
    """Update button visibility based on selected file source (STDF or CSV)"""
    source = file_source_var.get()

    if source == "STDF":
        # Show STDF button, hide CSV button
        select_multiple_stdf_button.pack(side=tk.LEFT, padx=3)
        select_csv_button.pack_forget()
    else:  # CSV
        # Hide STDF button, show CSV button
        select_multiple_stdf_button.pack_forget()
        select_csv_button.pack(side=tk.LEFT, padx=3)


def load_csv_wafermap_file():
    """Load wafermap data from a CSV file"""
    global current_stdf_data, current_wafer_id, test_parameters, grouped_parameters, test_limits
    global multiple_stdf_data, multiple_wafer_ids, die_image_directory

    csv_path = filedialog.askopenfilename(
        title="Select a CSV file with wafermap data",
        filetypes=[("CSV files", "*.csv"), ("All files", "*.*")],
    )

    if not csv_path:
        return

    try:
        print(f"Loading CSV file: {csv_path}")

        # Read CSV file
        df = pd.read_csv(csv_path)

        # Check for required columns (x, y coordinates)
        # Common column names for coordinates
        x_col_candidates = ['x', 'X', 'x_coord', 'X_COORD', 'X_Coordinate', 'x_coordinate', 'DIE_X', 'die_x', 'col', 'COL', 'Column']
        y_col_candidates = ['y', 'Y', 'y_coord', 'Y_COORD', 'Y_Coordinate', 'y_coordinate', 'DIE_Y', 'die_y', 'row', 'ROW', 'Row']

        x_col = None
        y_col = None

        for candidate in x_col_candidates:
            if candidate in df.columns:
                x_col = candidate
                break

        for candidate in y_col_candidates:
            if candidate in df.columns:
                y_col = candidate
                break

        if x_col is None or y_col is None:
            # If standard names not found, try to guess from first few columns
            print(f"Warning: Could not find standard x/y columns. Available columns: {df.columns.tolist()}")
            # Try to prompt user or use first two numeric columns
            numeric_cols = df.select_dtypes(include=[np.number]).columns.tolist()
            if len(numeric_cols) >= 2:
                x_col = numeric_cols[0]
                y_col = numeric_cols[1]
                print(f"Using '{x_col}' as X coordinate and '{y_col}' as Y coordinate")
            else:
                print("ERROR: Could not identify coordinate columns in CSV")
                return

        # Standardize column names
        df = df.rename(columns={x_col: 'x', y_col: 'y'})

        # Look for bin column
        bin_col_candidates = ['bin', 'BIN', 'Bin', 'HARD_BIN', 'hard_bin', 'SOFT_BIN', 'soft_bin', 'HB', 'SB']
        bin_col = None
        for candidate in bin_col_candidates:
            if candidate in df.columns:
                bin_col = candidate
                df = df.rename(columns={bin_col: 'bin'})
                break

        if 'bin' not in df.columns:
            # Create a default bin column (all pass = bin 1)
            df['bin'] = 1
            print("No bin column found, created default bin column (all bins = 1)")

        # Extract wafer ID from filename
        wafer_id = os.path.basename(csv_path).replace('.csv', '').replace('.CSV', '')

        # Build test parameters from remaining numeric columns
        test_params = {}
        grouped_params = {}  # Will be populated with detected groups
        test_limits_dict = {}

        numeric_columns = df.select_dtypes(include=[np.number]).columns.tolist()
        # Exclude x, y, bin from parameter list
        exclude_cols = ['x', 'y', 'bin']
        test_columns = [col for col in numeric_columns if col not in exclude_cols]

        # Detect groups from column name prefixes (e.g., OPTIC_, DC_, ELECTRICAL_)
        def extract_group_from_column(col_name):
            """Extract group name from column prefix"""
            col_str = str(col_name).upper()
            # Common prefixes to look for
            prefixes = ['OPTIC', 'OPTICAL', 'DC', 'ELECTRICAL', 'ANALOG', 'DIGITAL',
                       'POWER', 'SIGNAL', 'TEST', 'MEAS', 'PARAM']

            # Check for underscore-separated prefix
            if '_' in col_str:
                prefix = col_str.split('_')[0]
                if len(prefix) >= 2:  # At least 2 characters
                    return prefix.title()  # Return with title case

            # Check for known prefixes at start
            for prefix in prefixes:
                if col_str.startswith(prefix):
                    return prefix.title()

            return "Other"  # Default group

        for idx, col in enumerate(test_columns):
            test_num = idx + 1
            test_key = f"test_{test_num}"

            # Detect group from column name
            group_name = extract_group_from_column(col)

            # Rename column to match expected format
            df = df.rename(columns={col: test_num})

            test_params[test_key] = col  # Use original column name as display name

            # Add to appropriate group
            if group_name not in grouped_params:
                grouped_params[group_name] = []
            grouped_params[group_name].append((test_num, col, col))

            # Calculate limits from data (min/max)
            col_data = df[test_num].dropna()
            if len(col_data) > 0:
                test_limits_dict[test_num] = {
                    'lo_limit': col_data.min(),
                    'hi_limit': col_data.max(),
                    'units': ''
                }

        # Print detected groups
        print(f"Detected groups from CSV: {list(grouped_params.keys())}")
        for grp, params in grouped_params.items():
            print(f"  {grp}: {len(params)} parameters")

        # Update global variables
        current_stdf_data = df
        current_wafer_id = wafer_id
        test_parameters = test_params
        grouped_parameters = grouped_params
        test_limits = test_limits_dict

        # Also update multiple_stdf_data for compatibility with heatmap tab
        multiple_stdf_data = [df]
        multiple_wafer_ids = [wafer_id]

        # Try to detect notch orientation from CSV columns or data
        global current_wafer_config
        notch_orientation = None

        # Check for notch/flat column in CSV
        notch_col_candidates = ['notch', 'Notch', 'NOTCH', 'flat', 'Flat', 'FLAT',
                                'WF_FLAT', 'wf_flat', 'orientation', 'Orientation',
                                'ORIENTATION', 'wafer_flat', 'WAFER_FLAT']

        for candidate in notch_col_candidates:
            if candidate in df.columns:
                # Get the first non-null value
                notch_values = df[candidate].dropna()
                if len(notch_values) > 0:
                    notch_val = str(notch_values.iloc[0]).strip().upper()
                    # Map common values to standard U/D/L/R
                    notch_mapping = {
                        'U': 'U', 'UP': 'U', '0': 'U', 'TOP': 'U', 'NORTH': 'U',
                        'D': 'D', 'DOWN': 'D', '180': 'D', 'BOTTOM': 'D', 'SOUTH': 'D',
                        'L': 'L', 'LEFT': 'L', '270': 'L', 'WEST': 'L',
                        'R': 'R', 'RIGHT': 'R', '90': 'R', 'EAST': 'R'
                    }
                    notch_orientation = notch_mapping.get(notch_val, notch_val[0] if notch_val else None)
                    print(f"Detected notch orientation from CSV column '{candidate}': {notch_orientation}")
                break

        # If not found in columns, check in the original CSV for header comments or metadata
        if notch_orientation is None:
            try:
                with open(csv_path, 'r') as f:
                    # Read first 20 lines to look for notch info in comments
                    for i, line in enumerate(f):
                        if i > 20:
                            break
                        line_upper = line.upper()
                        if 'NOTCH' in line_upper or 'FLAT' in line_upper:
                            # Try to extract orientation
                            for orient in ['UP', 'DOWN', 'LEFT', 'RIGHT', 'NORTH', 'SOUTH', 'EAST', 'WEST']:
                                if orient in line_upper:
                                    notch_mapping = {'UP': 'U', 'DOWN': 'D', 'LEFT': 'L', 'RIGHT': 'R',
                                                    'NORTH': 'U', 'SOUTH': 'D', 'WEST': 'L', 'EAST': 'R'}
                                    notch_orientation = notch_mapping.get(orient, 'D')
                                    print(f"Detected notch orientation from CSV header: {notch_orientation}")
                                    break
                            if notch_orientation:
                                break
            except:
                pass

        # Default to 'D' (down) if not found - common default
        if notch_orientation is None:
            notch_orientation = 'D'
            print(f"No notch orientation found in CSV, using default: {notch_orientation}")

        # Update wafer config
        current_wafer_config = {
            'notch_orientation': notch_orientation,
            'wafer_size': None,
            'die_width': None,
            'die_height': None,
            'pos_x': None,
            'pos_y': None
        }
        print(f"Wafer config set: Notch={current_wafer_config.get('notch_orientation')}")

        # Try to auto-detect image folder from CSV path
        # Look for ImageCaptures folder in parent or sibling directories
        csv_dir = os.path.dirname(csv_path)
        parent_dir = os.path.dirname(csv_dir)

        possible_image_dirs = [
            os.path.join(parent_dir, "ImageCaptures"),
            os.path.join(parent_dir, "Images"),
            os.path.join(csv_dir, "ImageCaptures"),
            os.path.join(csv_dir, "Images"),
        ]

        for img_dir in possible_image_dirs:
            if os.path.isdir(img_dir):
                die_image_directory = img_dir
                short_path = os.path.basename(img_dir)
                folder_status_label.config(text=f"Img: {short_path}", fg="green")
                print(f"Auto-detected image folder: {img_dir}")
                break
        else:
            print("No image folder auto-detected. Use 'Image Folder' button to select manually.")

        # Auto-detect PLM folder from CSV path
        global plm_file_directory
        possible_plm_dirs = [
            os.path.join(parent_dir, "PLMFiles"),
            os.path.join(parent_dir, "PLM"),
            os.path.join(csv_dir, "PLMFiles"),
            os.path.join(csv_dir, "PLM"),
        ]

        for plm_dir in possible_plm_dirs:
            if os.path.isdir(plm_dir):
                plm_file_directory = plm_dir
                print(f"Auto-detected PLM folder: {plm_dir}")
                # Update folder status
                current_status = folder_status_label.cget("text")
                if current_status and "Img:" in current_status:
                    folder_status_label.config(text=f"{current_status} | PLM:✓", fg="green")
                else:
                    folder_status_label.config(text=f"PLM:✓", fg="green")
                break
        else:
            print("No PLM folder auto-detected.")

        print(f"Loaded CSV with {len(df)} dies and {len(test_params)} parameters")
        print(f"Wafer ID: {wafer_id}")
        print(f"Parameters: {list(test_params.values())[:10]}...")  # Show first 10

        # Update UI
        update_group_combobox()

        param_options = ["BIN (Bin Number)"]
        for test_key, test_name in sorted(test_parameters.items()):
            param_options.append(f"{test_key}: {test_name}")

        heatmap_param_combobox["values"] = param_options
        if param_options:
            heatmap_param_combobox.current(0)

        # Update heatmap display
        refresh_heatmap_display()
        update_heatmap_parameter_list()

        # Update info label
        heatmap_info_label.config(text=f"Loaded: {wafer_id} ({len(df)} dies)")

    except Exception as e:
        print(f"Error loading CSV file: {e}")
        import traceback
        traceback.print_exc()


def read_wafermap_from_stdf(stdf_path):
    """
    Read wafermap data from STDF file.
    OPTIMIZED version with debug output for troubleshooting.
    Now also reads WCR (Wafer Configuration Record) for notch orientation.
    """
    import re

    # Pre-compile regex patterns (performance optimization)
    PATTERN_SECTION = re.compile(r'^>{3,}\s*(.+?)\s*<{3,}$')
    PATTERN_SUBGROUP = re.compile(r'^<([A-Za-z][A-Za-z0-9_]*)>$')
    PATTERN_BRACKET = re.compile(r'^\[([A-Za-z][A-Za-z0-9_]*)\]$')
    EXCLUDED_GROUPS = {"Definitions", "Initialization"}

    wafermap = []
    wafer_id = None
    test_params = {}
    test_groups = {}
    grouped_params = {}
    test_limits_dict = {}

    # Wafer configuration from WCR record
    wafer_config = {
        'notch_orientation': None,
        'wafer_size': None,
        'die_width': None,
        'die_height': None,
        'pos_x': None,
        'pos_y': None
    }

    if STDF_MODULE is None:
        print("ERROR: STDF library not installed! Install with: pip install Semi-ATE-STDF")
        return pd.DataFrame(), None, {}, {}, {}, wafer_config

    try:
        if STDF_TYPE == "pystdf":
            wafermap, wafer_id, test_params, grouped_params, test_limits_dict, wafer_config = read_with_pystdf(stdf_path)
        else:
            file_size = os.path.getsize(stdf_path)
            print(f"Loading: {os.path.basename(stdf_path)} ({file_size / 1024 / 1024:.1f} MB)...")
            load_start = time.time()

            test_info = {}
            current_die_tests = {}
            current_group = "Ungrouped"
            wafermap = []
            debug_prr_printed = False
            debug_wcr_printed = False

            with open(stdf_path, "rb") as f:
                records_gen = STDF_MODULE.records_from_file(f)
                record_count = 0
                prr_count = 0
                last_update = time.time()

                for record in records_gen:
                    record_count += 1

                    # Progress update every 2 seconds
                    if record_count % 100000 == 0:
                        now = time.time()
                        if now - last_update > 2.0:
                            elapsed = now - load_start
                            print(f"  {record_count:,} records, {len(wafermap):,} dies ({elapsed:.0f}s)...")
                            last_update = now

                    rec_type = type(record).__name__

                    # WCR - Wafer Configuration Record
                    if rec_type == "WCR":
                        try:
                            # Debug: Print all available fields in WCR record
                            if not debug_wcr_printed:
                                print(f"  DEBUG WCR Record found!")
                                print(f"    Type: {type(record)}")
                                print(f"    Dir: {[a for a in dir(record) if not a.startswith('_')]}")
                                if hasattr(record, "get_value"):
                                    # Try to get all possible field names
                                    for field_name in ['WF_FLAT', 'WAFR_SIZ', 'DIE_WID', 'DIE_HT', 'POS_X', 'POS_Y', 'CENTER_X', 'CENTER_Y']:
                                        try:
                                            val = record.get_value(field_name)
                                            print(f"    {field_name}: {val} (type: {type(val).__name__})")
                                        except:
                                            print(f"    {field_name}: <not available>")

                            if hasattr(record, "get_value"):
                                wafer_config['notch_orientation'] = record.get_value("WF_FLAT")
                                wafer_config['wafer_size'] = record.get_value("WAFR_SIZ")
                                wafer_config['die_width'] = record.get_value("DIE_WID")
                                wafer_config['die_height'] = record.get_value("DIE_HT")
                                wafer_config['pos_x'] = record.get_value("POS_X")
                                wafer_config['pos_y'] = record.get_value("POS_Y")
                            elif hasattr(record, "WF_FLAT"):
                                wafer_config['notch_orientation'] = record.WF_FLAT
                                wafer_config['wafer_size'] = getattr(record, 'WAFR_SIZ', None)
                                wafer_config['die_width'] = getattr(record, 'DIE_WID', None)
                                wafer_config['die_height'] = getattr(record, 'DIE_HT', None)
                                wafer_config['pos_x'] = getattr(record, 'POS_X', None)
                                wafer_config['pos_y'] = getattr(record, 'POS_Y', None)

                            if not debug_wcr_printed:
                                print(f"  WCR Result: Notch='{wafer_config['notch_orientation']}', Size={wafer_config['wafer_size']}mm")
                                debug_wcr_printed = True
                        except Exception as e:
                            if not debug_wcr_printed:
                                print(f"  DEBUG WCR error: {e}")
                                import traceback
                                traceback.print_exc()
                                debug_wcr_printed = True

                    elif rec_type == "PRR":
                        prr_count += 1
                        try:
                            # Try different methods to access record fields
                            x_coord = None
                            y_coord = None
                            hard_bin = None
                            soft_bin = None

                            # Method 1: get_value method (Semi_ATE.STDF uses this)
                            if hasattr(record, "get_value"):
                                x_coord = record.get_value("X_COORD")
                                y_coord = record.get_value("Y_COORD")
                                hard_bin = record.get_value("HARD_BIN")
                                soft_bin = record.get_value("SOFT_BIN")
                            # Method 2: fields dict
                            elif hasattr(record, "fields"):
                                fields = record.fields
                                x_coord = fields.get("X_COORD")
                                y_coord = fields.get("Y_COORD")
                                hard_bin = fields.get("HARD_BIN")
                                soft_bin = fields.get("SOFT_BIN")
                            # Method 3: Direct attribute access
                            elif hasattr(record, "X_COORD"):
                                x_coord = record.X_COORD
                                y_coord = record.Y_COORD if hasattr(record, "Y_COORD") else None
                                hard_bin = record.HARD_BIN if hasattr(record, "HARD_BIN") else None
                                soft_bin = record.SOFT_BIN if hasattr(record, "SOFT_BIN") else None

                            # Debug: Print first PRR record structure
                            if not debug_prr_printed:
                                print(f"  DEBUG PRR: type={type(record)}, dir={[a for a in dir(record) if not a.startswith('_')][:15]}")
                                if hasattr(record, "__dict__"):
                                    print(f"  DEBUG PRR __dict__: {list(record.__dict__.keys())[:10]}")
                                print(f"  DEBUG PRR values: x={x_coord}, y={y_coord}, hbin={hard_bin}, sbin={soft_bin}")
                                debug_prr_printed = True

                            if x_coord is not None and y_coord is not None:
                                die_data = {
                                    "x": x_coord,
                                    "y": y_coord,
                                    "bin": hard_bin if hard_bin is not None else soft_bin,
                                }
                                if current_die_tests:
                                    die_data.update(current_die_tests)
                                wafermap.append(die_data)
                        except Exception as e:
                            if not debug_prr_printed:
                                print(f"  DEBUG PRR error: {e}")
                                debug_prr_printed = True

                    elif rec_type == "PTR":
                        try:
                            # Try different access methods
                            test_num = None
                            if hasattr(record, "TEST_NUM"):
                                test_num = record.TEST_NUM
                            elif hasattr(record, "get_value"):
                                test_num = record.get_value("TEST_NUM")

                            if test_num is not None:
                                result = None
                                if hasattr(record, "RESULT"):
                                    result = record.RESULT
                                elif hasattr(record, "get_value"):
                                    result = record.get_value("RESULT")

                                if test_num not in test_info:
                                    test_name = None
                                    if hasattr(record, "TEST_TXT"):
                                        test_name = record.TEST_TXT
                                    elif hasattr(record, "get_value"):
                                        test_name = record.get_value("TEST_TXT")
                                    if test_name:
                                        test_info[test_num] = test_name
                                        test_groups[test_num] = current_group

                                if test_num not in test_limits_dict:
                                    lo_limit = getattr(record, "LO_LIMIT", None) or (record.get_value("LO_LIMIT") if hasattr(record, "get_value") else None)
                                    hi_limit = getattr(record, "HI_LIMIT", None) or (record.get_value("HI_LIMIT") if hasattr(record, "get_value") else None)
                                    units = getattr(record, "UNITS", "") or (record.get_value("UNITS") if hasattr(record, "get_value") else "") or ""
                                    test_limits_dict[test_num] = {
                                        'lo_limit': lo_limit,
                                        'hi_limit': hi_limit,
                                        'units': units
                                    }

                                if result is not None:
                                    current_die_tests[test_num] = result
                        except:
                            pass

                    elif rec_type == "PIR":
                        current_die_tests = {}

                    elif rec_type == "WIR":
                        try:
                            if hasattr(record, "WAFER_ID"):
                                wafer_id = record.WAFER_ID
                            elif hasattr(record, "get_value"):
                                wafer_id = record.get_value("WAFER_ID")
                        except:
                            pass

                    elif rec_type == "DTR":
                        try:
                            text_data = None
                            if hasattr(record, "TEXT_DAT"):
                                text_data = record.TEXT_DAT
                            elif hasattr(record, "get_value"):
                                text_data = record.get_value("TEXT_DAT")

                            if text_data:
                                text_stripped = text_data.strip()
                                match = PATTERN_SECTION.match(text_stripped)
                                if match:
                                    group_name = match.group(1).strip()
                                    if group_name and group_name not in EXCLUDED_GROUPS:
                                        current_group = group_name
                                else:
                                    match = PATTERN_SUBGROUP.match(text_stripped)
                                    if match:
                                        current_group = match.group(1)
                                    else:
                                        match = PATTERN_BRACKET.match(text_stripped)
                                        if match:
                                            current_group = match.group(1)
                        except:
                            pass

                # Build test_params and grouped_params
                for test_num, test_name in test_info.items():
                    short_name = shorten_test_name(test_name) if test_name else f"Test {test_num}"
                    test_params[f"test_{test_num}"] = short_name
                    group = test_groups.get(test_num, "Ungrouped")
                    if group not in grouped_params:
                        grouped_params[group] = []
                    grouped_params[group].append((test_num, short_name, test_name))

                elapsed = time.time() - load_start
                print(f"  Done: {record_count:,} records | {len(wafermap):,} dies | {len(test_params)} params | {elapsed:.1f}s")

        df = pd.DataFrame(wafermap)
        if len(df) > 0:
            print(f"  Loaded: X[{df['x'].min()}-{df['x'].max()}] Y[{df['y'].min()}-{df['y'].max()}] Bins:{len(df['bin'].unique())}")

        return df, wafer_id, test_params, grouped_params, test_limits_dict, wafer_config

    except Exception as e:
        print(f"Error reading STDF: {e}")
        import traceback
        traceback.print_exc()
        return pd.DataFrame(), None, {}, {}, {}, wafer_config


def shorten_test_name(test_name):
    """Shorten test name by truncating at _X_X_X pattern"""
    if not test_name:
        return test_name

    import re
    # Find pattern like _X_X_X (where X can be any character or empty)
    # This matches _X_X_X, _X_X_X_, _X_X_X_something
    match = re.search(r'_X_X_X', test_name)
    if match:
        # Return everything before _X_X_X
        return test_name[:match.start()]

    # If no _X_X_X pattern, return original name (truncated if too long)
    if len(test_name) > 40:
        return test_name[:37] + "..."
    return test_name


def read_with_pystdf(stdf_path):
    import pystdf

    wafermap = []
    wafer_id = None
    test_params = {}
    test_info = {}
    current_die_tests = {}
    test_limits_dict = {}  # Dictionary to store test limits

    # Wafer configuration
    wafer_config = {
        'notch_orientation': None,
        'wafer_size': None,
        'die_width': None,
        'die_height': None,
        'pos_x': None,
        'pos_y': None
    }

    with open(stdf_path, "rb") as f:
        parser = pystdf.Parser(inp=f)

        for record in parser:
            if record.id == "WIR":
                wafer_id = record.WAFER_ID

            elif record.id == "WCR":
                try:
                    wafer_config['notch_orientation'] = getattr(record, 'WF_FLAT', None)
                    wafer_config['wafer_size'] = getattr(record, 'WAFR_SIZ', None)
                    wafer_config['die_width'] = getattr(record, 'DIE_WID', None)
                    wafer_config['die_height'] = getattr(record, 'DIE_HT', None)
                    wafer_config['pos_x'] = getattr(record, 'POS_X', None)
                    wafer_config['pos_y'] = getattr(record, 'POS_Y', None)
                    print(f"  WCR: Notch={wafer_config['notch_orientation']}, Size={wafer_config['wafer_size']}mm")
                except:
                    pass

            elif record.id == "PIR":
                current_die_tests = {}

            elif record.id == "PTR":
                try:
                    test_num = record.TEST_NUM
                    test_name = record.TEST_TXT
                    result = record.RESULT

                    if test_num is not None and test_name:
                        test_info[test_num] = test_name

                    # Extract test limits
                    if test_num is not None:
                        lo_limit = getattr(record, 'LO_LIMIT', None)
                        hi_limit = getattr(record, 'HI_LIMIT', None)
                        units = getattr(record, 'UNITS', "") or ""

                        if test_num not in test_limits_dict:
                            test_limits_dict[test_num] = {
                                'lo_limit': lo_limit,
                                'hi_limit': hi_limit,
                                'units': units
                            }
                        elif lo_limit is not None or hi_limit is not None:
                            if lo_limit is not None:
                                test_limits_dict[test_num]['lo_limit'] = lo_limit
                            if hi_limit is not None:
                                test_limits_dict[test_num]['hi_limit'] = hi_limit
                            if units:
                                test_limits_dict[test_num]['units'] = units

                    if test_num is not None and result is not None:
                        current_die_tests[test_num] = result

                except (ValueError, KeyError, AttributeError) as e:
                    print(f"Warning: Could not read PTR record: {e}")

            elif record.id == "PRR":
                try:
                    if record.X_COORD is not None and record.Y_COORD is not None:
                        die_data = {
                            "x": record.X_COORD,
                            "y": record.Y_COORD,
                            "bin": record.HARD_BIN
                            if record.HARD_BIN is not None
                            else record.SOFT_BIN,
                        }

                        for test_num in current_die_tests:
                            die_data[f"test_{test_num}"] = current_die_tests[test_num]

                        wafermap.append(die_data)

                except (ValueError, KeyError, AttributeError) as e:
                    print(f"Warning: Could not read PRR record: {e}")

    test_params = {f"test_{num}": shorten_test_name(name) for num, name in test_info.items()}

    # Create grouped_params (all in "Ungrouped" for pystdf since we don't parse DTR)
    grouped_params = {"Ungrouped": [(num, shorten_test_name(name), name) for num, name in test_info.items()]}

    return wafermap, wafer_id, test_params, grouped_params, test_limits_dict, wafer_config


def load_stdf_data(stdf_path):
    """Load STDF file and populate parameter selection"""
    global current_stdf_data, current_wafer_id, test_parameters, grouped_parameters, test_limits
    global current_wafer_config

    df, wafer_id, test_params, grouped_params, limits_dict, wafer_cfg = read_wafermap_from_stdf(stdf_path)

    current_stdf_data = df
    current_wafer_id = wafer_id
    test_parameters = test_params
    grouped_parameters = grouped_params
    test_limits = limits_dict
    current_wafer_config = wafer_cfg

    current_stdf_data = df
    current_wafer_id = wafer_id
    test_parameters = test_params
    grouped_parameters = grouped_params
    test_limits = limits_dict

    # Update group combobox
    update_group_combobox()

    param_options = ["BIN (Bin Number)"]
    for test_key, test_name in sorted(test_parameters.items()):
        param_options.append(f"{test_key}: {test_name}")

    param_combobox["values"] = param_options
    if param_options:
        param_combobox.current(0)

    print(f"Available parameters: {len(param_options)}")

    update_wafermap_with_parameter()
    update_heatmap_parameter_list()


def update_wafermap_with_parameter():
    """Update wafermap display based on selected parameter"""
    global current_stdf_data, current_wafer_id

    if current_stdf_data is None or current_stdf_data.empty:
        print("No STDF data loaded")
        return

    selected = param_combobox.get()
    print(f"\n=== Parameter Selection Debug ===")
    print(f"Selected from combobox: '{selected}'")

    if selected.startswith("BIN"):
        param_column = "bin"
        param_label = "Bin"
    else:
        # Extract test number from "test_XXXXX: Test Name" format
        test_key = selected.split(":")[0].strip()
        print(f"Extracted test_key: '{test_key}'")
        # Remove 'test_' prefix if present to get the actual column name
        if test_key.startswith("test_"):
            param_column = int(test_key.replace("test_", ""))
        else:
            param_column = int(test_key)
        param_label = selected
        print(f"Looking for column: {param_column} (type: {type(param_column).__name__})")

    print(f"Available columns (first 10): {list(current_stdf_data.columns[:10])}")
    print(f"Column types: {[type(col).__name__ for col in current_stdf_data.columns[:10]]}")

    update_wafermap_display(param_column, param_label)


def update_wafermap_display(parameter="bin", param_label="Bin"):
    """Display wafermap colored by specified parameter"""
    global wafermap_canvas, current_stdf_data, current_wafer_id

    fig, ax = plt.subplots(figsize=(10, 10))

    if current_stdf_data is not None and not current_stdf_data.empty:
        if parameter not in current_stdf_data.columns:
            ax.set_title(f"Parameter '{parameter}' not found in data")
            print(f"Available columns: {current_stdf_data.columns.tolist()}")
        else:
            plot_data = current_stdf_data.dropna(subset=[parameter])

            if len(plot_data) > 0:
                sc = ax.scatter(
                    plot_data["x"],
                    plot_data["y"],
                    c=plot_data[parameter],
                    cmap="viridis" if parameter != "bin" else "tab20",
                    s=80,
                    edgecolors="black",
                    linewidth=0.5,
                )

                ax.set_xlabel("X Coordinate", fontsize=12)
                ax.set_ylabel("Y Coordinate", fontsize=12)
                ax.set_title(
                    f"Wafermap for Wafer {current_wafer_id}\n{param_label}", fontsize=14
                )
                ax.grid(True, alpha=0.3)
                ax.set_aspect("equal")

                cbar = plt.colorbar(sc, ax=ax, label=param_label)
                cbar.ax.tick_params(labelsize=10)

                print(f"Plotted {len(plot_data)} dies with {parameter}")
                print(
                    f"Value range: {plot_data[parameter].min()} to {plot_data[parameter].max()}"
                )

            else:
                ax.set_title(f"No valid data for parameter '{param_label}'")
    else:
        ax.set_title("No wafermap data available")

    fig.tight_layout()

    if wafermap_canvas:
        wafermap_canvas.get_tk_widget().destroy()

    wafermap_canvas = FigureCanvasTkAgg(fig, master=wafermap_frame)
    wafermap_canvas.get_tk_widget().pack(fill=tk.BOTH, expand=True, pady=10)
    wafermap_canvas.draw()
# Tab 6: STDF Heatmap with Parameter Selection - Already created above, just add components

heatmap_canvas = None
current_heatmap_data = None
show_grid_var = tk.BooleanVar(value=False)
selected_die_coords = None
multiple_stdf_data = []
multiple_wafer_ids = []

# Performance optimization: Cache for computed grids
_grid_cache = {}
_current_fig = None
_current_ax = None
_die_info_window = None  # Track current die info popup window

# Canvases for statistics plots
stats_boxplot_canvas = None
stats_prob_canvas = None
stats_bin_fail_canvas = None

# Image directory for die images
die_image_directory = None
die_image_refs = []  # Keep references to prevent garbage collection

# Image type filter
image_type_var = None  # Will be initialized later
available_image_types = ["All"]  # Will be populated when images are found
current_selected_die = None  # Store current die coordinates for refresh

# Project folder structure
project_folder = None
project_subfolders = {
    "stdf": "STDDatalog",
    "images": "ImageCaptures",
    "csv": "CSVFiles",
    "txt": "TXTDatalog",
    "plm": "PLMFiles"
}

# ============================================================================
# Custom Test Calculator - Define custom tests using mathematical operations
# ============================================================================
custom_tests = {}  # Dictionary to store custom test definitions: name -> formula_dict

class CustomTestCalculatorDialog:
    """Dialog for creating custom tests using mathematical formulas with existing test parameters"""

    def __init__(self, parent):
        self.parent = parent
        self.dialog = tk.Toplevel(parent)
        self.dialog.title("🧮 Custom Test Calculator")
        self.dialog.geometry("700x550")
        self.dialog.transient(parent)
        self.dialog.grab_set()

        # Center the dialog
        self.dialog.update_idletasks()
        x = (self.dialog.winfo_screenwidth() - 700) // 2
        y = (self.dialog.winfo_screenheight() - 550) // 2
        self.dialog.geometry(f"+{x}+{y}")

        self.formula_parts = []  # List of formula components: (type, value)
        self.create_widgets()

    def create_widgets(self):
        # Main container
        main_frame = tk.Frame(self.dialog, padx=10, pady=10)
        main_frame.pack(fill=tk.BOTH, expand=True)

        # === TEST NAME SECTION ===
        name_frame = tk.LabelFrame(main_frame, text="Test Definition", font=("Helvetica", 10, "bold"), padx=10, pady=5)
        name_frame.pack(fill=tk.X, pady=(0, 10))

        tk.Label(name_frame, text="New Test Name:", font=("Helvetica", 9)).grid(row=0, column=0, sticky="w", padx=5, pady=5)
        self.test_name_var = tk.StringVar(value="")
        self.test_name_entry = tk.Entry(name_frame, textvariable=self.test_name_var, width=30, font=("Helvetica", 10))
        self.test_name_entry.grid(row=0, column=1, sticky="w", padx=5, pady=5)

        tk.Label(name_frame, text="Example: Vf10, Ratio_A_B, Delta_Power", font=("Helvetica", 8), fg="gray").grid(row=0, column=2, sticky="w", padx=10)

        # === FORMULA BUILDER SECTION ===
        formula_frame = tk.LabelFrame(main_frame, text="Formula Builder", font=("Helvetica", 10, "bold"), padx=10, pady=5)
        formula_frame.pack(fill=tk.BOTH, expand=True, pady=(0, 10))

        # Current formula display
        formula_display_frame = tk.Frame(formula_frame)
        formula_display_frame.pack(fill=tk.X, pady=5)

        tk.Label(formula_display_frame, text="Formula:", font=("Helvetica", 9, "bold")).pack(side=tk.LEFT, padx=5)

        self.formula_display = tk.Text(formula_display_frame, height=2, width=60, font=("Consolas", 10), bg="#f5f5f5", wrap=tk.WORD)
        self.formula_display.pack(side=tk.LEFT, fill=tk.X, expand=True, padx=5)
        self.formula_display.config(state='disabled')

        # Parameter selection
        param_frame = tk.Frame(formula_frame)
        param_frame.pack(fill=tk.X, pady=5)

        tk.Label(param_frame, text="Select Test Parameter:", font=("Helvetica", 9)).pack(side=tk.LEFT, padx=5)

        self.param_var = tk.StringVar()
        self.param_combobox = ttk.Combobox(param_frame, textvariable=self.param_var, state="readonly", width=40, font=("Helvetica", 9))
        self.param_combobox.pack(side=tk.LEFT, padx=5)
        self.populate_parameters()

        add_param_btn = tk.Button(param_frame, text="Add Parameter", command=self.add_parameter,
                                  font=("Helvetica", 9), bg="#4CAF50", fg="white")
        add_param_btn.pack(side=tk.LEFT, padx=5)

        # Operators frame
        operators_frame = tk.Frame(formula_frame)
        operators_frame.pack(fill=tk.X, pady=10)

        tk.Label(operators_frame, text="Operators:", font=("Helvetica", 9, "bold")).pack(side=tk.LEFT, padx=5)

        operators = [
            ("+", "Add"),
            ("-", "Subtract"),
            ("*", "Multiply"),
            ("/", "Divide"),
            ("(", "Open ("),
            (")", "Close )"),
        ]

        for op, tooltip in operators:
            btn = tk.Button(operators_frame, text=op, width=3, font=("Helvetica", 11, "bold"),
                           command=lambda o=op: self.add_operator(o))
            btn.pack(side=tk.LEFT, padx=3)

        # Number input
        number_frame = tk.Frame(formula_frame)
        number_frame.pack(fill=tk.X, pady=5)

        tk.Label(number_frame, text="Add Number:", font=("Helvetica", 9)).pack(side=tk.LEFT, padx=5)

        self.number_var = tk.StringVar(value="")
        self.number_entry = tk.Entry(number_frame, textvariable=self.number_var, width=15, font=("Helvetica", 10))
        self.number_entry.pack(side=tk.LEFT, padx=5)

        add_number_btn = tk.Button(number_frame, text="Add Number", command=self.add_number,
                                   font=("Helvetica", 9), bg="#2196F3", fg="white")
        add_number_btn.pack(side=tk.LEFT, padx=5)

        # Clear and undo buttons
        edit_frame = tk.Frame(formula_frame)
        edit_frame.pack(fill=tk.X, pady=5)

        undo_btn = tk.Button(edit_frame, text="↩ Undo Last", command=self.undo_last,
                             font=("Helvetica", 9), bg="#FF9800", fg="white")
        undo_btn.pack(side=tk.LEFT, padx=5)

        clear_btn = tk.Button(edit_frame, text="🗑 Clear All", command=self.clear_formula,
                              font=("Helvetica", 9), bg="#f44336", fg="white")
        clear_btn.pack(side=tk.LEFT, padx=5)

        # === EXISTING CUSTOM TESTS SECTION ===
        existing_frame = tk.LabelFrame(main_frame, text="Existing Custom Tests", font=("Helvetica", 10, "bold"), padx=10, pady=5)
        existing_frame.pack(fill=tk.X, pady=(0, 10))

        # Listbox with existing custom tests
        list_frame = tk.Frame(existing_frame)
        list_frame.pack(fill=tk.X, pady=5)

        self.custom_tests_listbox = tk.Listbox(list_frame, height=4, font=("Consolas", 9), selectmode=tk.SINGLE)
        self.custom_tests_listbox.pack(side=tk.LEFT, fill=tk.X, expand=True, padx=5)

        # Bind double-click to load test for editing
        self.custom_tests_listbox.bind("<Double-Button-1>", self.load_custom_test_for_edit)

        scrollbar = tk.Scrollbar(list_frame, orient="vertical", command=self.custom_tests_listbox.yview)
        scrollbar.pack(side=tk.RIGHT, fill=tk.Y)
        self.custom_tests_listbox.config(yscrollcommand=scrollbar.set)

        # Populate with existing custom tests
        self.refresh_custom_tests_list()

        # Hint label
        hint_label = tk.Label(existing_frame, text="💡 Double-click a test to edit it", font=("Helvetica", 8), fg="gray")
        hint_label.pack(side=tk.TOP, anchor="w", padx=5)

        # Delete button for selected custom test
        delete_btn = tk.Button(existing_frame, text="🗑 Delete Selected", command=self.delete_selected_custom_test,
                               font=("Helvetica", 9), bg="#f44336", fg="white")
        delete_btn.pack(side=tk.LEFT, padx=5, pady=5)

        # === BOTTOM BUTTONS ===
        button_frame = tk.Frame(main_frame)
        button_frame.pack(fill=tk.X, pady=10)

        save_btn = tk.Button(button_frame, text="💾 Save Custom Test", command=self.save_custom_test,
                             font=("Helvetica", 10, "bold"), bg="#4CAF50", fg="white", padx=20, pady=5)
        save_btn.pack(side=tk.LEFT, padx=10)

        close_btn = tk.Button(button_frame, text="Close", command=self.dialog.destroy,
                              font=("Helvetica", 10), padx=20, pady=5)
        close_btn.pack(side=tk.RIGHT, padx=10)

    def populate_parameters(self):
        """Populate parameter dropdown with available test parameters"""
        global test_parameters, grouped_parameters

        param_options = []

        # Add grouped parameters if available
        if grouped_parameters:
            for group_name, params in sorted(grouped_parameters.items()):
                for test_num, short_name, full_name in params:
                    param_options.append(f"test_{test_num}: {short_name}")
        elif test_parameters:
            for test_key, test_name in sorted(test_parameters.items()):
                param_options.append(f"{test_key}: {test_name}")

        if not param_options:
            param_options = ["(No parameters loaded - load STDF/CSV first)"]

        self.param_combobox["values"] = param_options
        if param_options:
            self.param_combobox.current(0)

    def add_parameter(self):
        """Add selected parameter to the formula"""
        param = self.param_var.get()
        if param and not param.startswith("(No parameters"):
            # Extract test_key from the selection (e.g., "test_1234: TestName" -> "test_1234")
            test_key = param.split(":")[0].strip()
            self.formula_parts.append(("param", test_key))
            self.update_formula_display()

    def add_operator(self, operator):
        """Add operator to the formula"""
        self.formula_parts.append(("op", operator))
        self.update_formula_display()

    def add_number(self):
        """Add number to the formula"""
        try:
            number = float(self.number_var.get())
            self.formula_parts.append(("num", number))
            self.number_var.set("")
            self.update_formula_display()
        except ValueError:
            tk.messagebox.showwarning("Invalid Number", "Please enter a valid number.")

    def undo_last(self):
        """Remove the last added element"""
        if self.formula_parts:
            self.formula_parts.pop()
            self.update_formula_display()

    def clear_formula(self):
        """Clear the entire formula"""
        self.formula_parts = []
        self.update_formula_display()

    def update_formula_display(self):
        """Update the formula display text"""
        formula_str = self.get_formula_string()
        self.formula_display.config(state='normal')
        self.formula_display.delete('1.0', tk.END)
        self.formula_display.insert('1.0', formula_str)
        self.formula_display.config(state='disabled')

    def get_formula_string(self):
        """Convert formula parts to a readable string"""
        parts = []
        for part_type, value in self.formula_parts:
            if part_type == "param":
                # Show shortened parameter name
                param_name = test_parameters.get(value, value)
                parts.append(f"[{param_name}]")
            elif part_type == "op":
                parts.append(f" {value} ")
            elif part_type == "num":
                parts.append(str(value))
        return "".join(parts)

    def refresh_custom_tests_list(self):
        """Refresh the list of existing custom tests"""
        global custom_tests
        self.custom_tests_listbox.delete(0, tk.END)
        for name, formula_dict in custom_tests.items():
            formula_str = formula_dict.get('display', '')
            self.custom_tests_listbox.insert(tk.END, f"{name}: {formula_str}")

    def load_custom_test_for_edit(self, event=None):
        """Load a custom test from the list into the editor for re-editing"""
        global custom_tests
        selection = self.custom_tests_listbox.curselection()
        if not selection:
            return

        item = self.custom_tests_listbox.get(selection[0])
        test_name = item.split(":")[0].strip()

        if test_name not in custom_tests:
            return

        # Load the test name
        self.test_name_var.set(test_name)

        # Load the formula parts
        self.formula_parts = custom_tests[test_name]['parts'].copy()

        # Update the display
        self.update_formula_display()

        print(f"Loaded custom test '{test_name}' for editing")

    def delete_selected_custom_test(self):
        """Delete the selected custom test"""
        global custom_tests
        selection = self.custom_tests_listbox.curselection()
        if selection:
            item = self.custom_tests_listbox.get(selection[0])
            test_name = item.split(":")[0].strip()
            if test_name in custom_tests:
                del custom_tests[test_name]
                self.refresh_custom_tests_list()
                update_heatmap_parameter_list()  # Update the dropdown in main UI
                tk.messagebox.showinfo("Deleted", f"Custom test '{test_name}' has been deleted.")

    def save_custom_test(self):
        """Save the custom test definition"""
        global custom_tests

        test_name = self.test_name_var.get().strip()

        if not test_name:
            tk.messagebox.showwarning("Missing Name", "Please enter a name for the custom test.")
            return

        if not self.formula_parts:
            tk.messagebox.showwarning("Empty Formula", "Please build a formula before saving.")
            return

        # Validate the formula has proper structure (basic check)
        if not self.validate_formula():
            tk.messagebox.showwarning("Invalid Formula", "The formula structure is invalid. Please check operators and parameters.")
            return

        # Save the custom test
        custom_tests[test_name] = {
            'parts': self.formula_parts.copy(),
            'display': self.get_formula_string()
        }

        print(f"Saved custom test '{test_name}' with formula: {self.get_formula_string()}")

        # Clear the form
        self.test_name_var.set("")
        self.clear_formula()

        # Refresh lists
        self.refresh_custom_tests_list()
        update_heatmap_parameter_list()  # Update the parameter dropdown in main UI
        update_group_combobox()  # Update the group dropdown to include Custom Tests group

        tk.messagebox.showinfo("Saved", f"Custom test '{test_name}' has been saved.\n\nYou can now select it from the parameter dropdown\nor select 'Custom Tests' from the Group dropdown.")

    def validate_formula(self):
        """Basic validation of formula structure"""
        if not self.formula_parts:
            return False

        # Check parentheses balance
        paren_count = 0
        for part_type, value in self.formula_parts:
            if part_type == "op":
                if value == "(":
                    paren_count += 1
                elif value == ")":
                    paren_count -= 1
                    if paren_count < 0:
                        return False

        if paren_count != 0:
            return False

        return True


def evaluate_custom_test(test_name, die_data):
    """
    Evaluate a custom test formula for a given die's data.

    Args:
        test_name: Name of the custom test
        die_data: Dictionary or Series containing the die's test values

    Returns:
        Calculated value or None if calculation fails
    """
    global custom_tests

    if test_name not in custom_tests:
        return None

    formula = custom_tests[test_name]['parts']

    # Build expression string
    expr_parts = []
    for part_type, value in formula:
        if part_type == "param":
            # Get value from die_data
            param_value = None

            # Extract the numeric part from test_key (e.g., "test_1234" -> 1234)
            test_num = None
            if isinstance(value, str) and value.startswith("test_"):
                try:
                    test_num = int(value.replace("test_", ""))
                except ValueError:
                    pass

            # Try different ways to access the value
            if isinstance(die_data, dict):
                param_value = die_data.get(test_num) if test_num else die_data.get(value)
                if param_value is None:
                    param_value = die_data.get(value)
            else:
                # pandas Series - try multiple approaches to handle different column types
                try:
                    if test_num is not None:
                        # Try integer key first
                        if test_num in die_data.index:
                            param_value = die_data[test_num]
                        # Try string version of the integer
                        elif str(test_num) in die_data.index:
                            param_value = die_data[str(test_num)]
                        # Try with test_ prefix
                        elif f"test_{test_num}" in die_data.index:
                            param_value = die_data[f"test_{test_num}"]

                    # If still None, try the original value
                    if param_value is None and value in die_data.index:
                        param_value = die_data[value]

                    # Last resort - try using .get() method
                    if param_value is None:
                        param_value = die_data.get(test_num) if test_num else die_data.get(value)
                except Exception as e:
                    print(f"Error accessing param {value} (test_num={test_num}): {e}")
                    param_value = None

            if param_value is None or (isinstance(param_value, float) and np.isnan(param_value)):
                return None

            expr_parts.append(str(float(param_value)))
        elif part_type == "op":
            expr_parts.append(value)
        elif part_type == "num":
            expr_parts.append(str(float(value)))

    expr_string = "".join(expr_parts)

    try:
        # Safe evaluation using only basic math operations
        # Only allow numbers, operators, parentheses
        allowed_chars = set("0123456789.+-*/() eE")
        if not all(c in allowed_chars for c in expr_string):
            print(f"Invalid characters in expression: {expr_string}")
            return None

        result = eval(expr_string)
        return float(result)
    except (ZeroDivisionError, ValueError, SyntaxError, TypeError) as e:
        print(f"Error evaluating custom test '{test_name}': {e}, expression: {expr_string}")
        return None


def open_custom_test_calculator():
    """Open the Custom Test Calculator dialog"""
    CustomTestCalculatorDialog(main_win)


# ============================================================================
# Save Modified Data Dialog - Export data with custom tests to CSV or STDF
# ============================================================================

class SaveDataDialog:
    """Dialog for saving wafer data with custom tests to CSV or STDF format"""

    def __init__(self, parent):
        self.parent = parent
        self.dialog = tk.Toplevel(parent)
        self.dialog.title("💾 Save Modified Data")
        self.dialog.geometry("550x520")
        self.dialog.minsize(550, 520)  # Set minimum size
        self.dialog.transient(parent)
        self.dialog.grab_set()

        # Center the dialog
        self.dialog.update_idletasks()
        x = (self.dialog.winfo_screenwidth() - 550) // 2
        y = (self.dialog.winfo_screenheight() - 520) // 2
        self.dialog.geometry(f"+{x}+{y}")

        self.create_widgets()

    def create_widgets(self):
        # Main container with scrollable frame
        main_frame = tk.Frame(self.dialog, padx=15, pady=15)
        main_frame.pack(fill=tk.BOTH, expand=True)

        # === DATA INFO SECTION ===
        info_frame = tk.LabelFrame(main_frame, text="Data Information", font=("Helvetica", 10, "bold"), padx=10, pady=5)
        info_frame.pack(fill=tk.X, pady=(0, 10))

        # Show loaded data info
        global multiple_stdf_data, multiple_wafer_ids, custom_tests

        num_files = len(multiple_stdf_data) if multiple_stdf_data else 0
        num_custom_tests = len(custom_tests) if custom_tests else 0
        total_dies = sum(len(df) for df in multiple_stdf_data) if multiple_stdf_data else 0

        info_text = f"Loaded files: {num_files} | Total dies: {total_dies} | Custom tests: {num_custom_tests}"

        if custom_tests and len(custom_tests) <= 3:
            info_text += f" ({', '.join(custom_tests.keys())})"

        info_label = tk.Label(info_frame, text=info_text, font=("Helvetica", 9), justify=tk.LEFT, anchor="w")
        info_label.pack(fill=tk.X, pady=5)

        # === FORMAT SELECTION ===
        format_frame = tk.LabelFrame(main_frame, text="Export Format", font=("Helvetica", 10, "bold"), padx=10, pady=5)
        format_frame.pack(fill=tk.X, pady=(0, 10))

        self.format_var = tk.StringVar(value="CSV")

        csv_radio = tk.Radiobutton(format_frame, text="CSV (Comma-Separated Values)", variable=self.format_var,
                                   value="CSV", font=("Helvetica", 9), command=self.update_options)
        csv_radio.pack(anchor="w", pady=2)

        stdf_radio = tk.Radiobutton(format_frame, text="STDF (Standard Test Data Format)", variable=self.format_var,
                                    value="STDF", font=("Helvetica", 9), command=self.update_options)
        stdf_radio.pack(anchor="w", pady=2)

        # STDF availability note
        stdf_note = "(Note: STDF export requires Semi_ATE library with write support)"
        stdf_note_label = tk.Label(format_frame, text=stdf_note, font=("Helvetica", 8), fg="gray")
        stdf_note_label.pack(anchor="w", padx=20)

        # === OPTIONS SECTION ===
        options_frame = tk.LabelFrame(main_frame, text="Export Options", font=("Helvetica", 10, "bold"), padx=10, pady=5)
        options_frame.pack(fill=tk.X, pady=(0, 10))

        # Include custom tests checkbox
        self.include_custom_var = tk.BooleanVar(value=True)
        custom_check = tk.Checkbutton(options_frame, text="Include computed custom test values",
                                       variable=self.include_custom_var, font=("Helvetica", 9))
        custom_check.pack(anchor="w", pady=2)

        # Include original test data
        self.include_original_var = tk.BooleanVar(value=True)
        original_check = tk.Checkbutton(options_frame, text="Include original test parameters",
                                         variable=self.include_original_var, font=("Helvetica", 9))
        original_check.pack(anchor="w", pady=2)

        # Include limits in CSV
        self.include_limits_var = tk.BooleanVar(value=False)
        self.limits_check = tk.Checkbutton(options_frame, text="Include test limits (CSV only - adds LoLimit/HiLimit rows)",
                                            variable=self.include_limits_var, font=("Helvetica", 9))
        self.limits_check.pack(anchor="w", pady=2)

        # Single file or multiple files
        self.file_mode_var = tk.StringVar(value="combined")

        file_mode_frame = tk.Frame(options_frame)
        file_mode_frame.pack(fill=tk.X, pady=5)

        tk.Label(file_mode_frame, text="File output:", font=("Helvetica", 9)).pack(side=tk.LEFT)

        combined_radio = tk.Radiobutton(file_mode_frame, text="Combined (all wafers in one file)",
                                        variable=self.file_mode_var, value="combined", font=("Helvetica", 9))
        combined_radio.pack(side=tk.LEFT, padx=10)

        separate_radio = tk.Radiobutton(file_mode_frame, text="Separate files per wafer",
                                        variable=self.file_mode_var, value="separate", font=("Helvetica", 9))
        separate_radio.pack(side=tk.LEFT, padx=10)

        # === BUTTONS AT BOTTOM (pack first to ensure visibility) ===
        button_frame = tk.Frame(self.dialog, pady=10, padx=15)
        button_frame.pack(side=tk.BOTTOM, fill=tk.X)

        # Separator line
        separator = tk.Frame(self.dialog, height=2, bg="gray")
        separator.pack(side=tk.BOTTOM, fill=tk.X, padx=15)

        save_btn = tk.Button(button_frame, text="💾 Save", command=self.save_data,
                             font=("Helvetica", 11, "bold"), bg="#4CAF50", fg="white", padx=30, pady=8)
        save_btn.pack(side=tk.LEFT, padx=10)

        cancel_btn = tk.Button(button_frame, text="Cancel", command=self.dialog.destroy,
                               font=("Helvetica", 11), padx=30, pady=8)
        cancel_btn.pack(side=tk.RIGHT, padx=10)

        # === STATUS LABEL ===
        self.status_var = tk.StringVar(value="Click 'Save' to choose file location")
        status_label = tk.Label(main_frame, textvariable=self.status_var, font=("Helvetica", 9), fg="blue")
        status_label.pack(fill=tk.X, pady=5)

    def update_options(self):
        """Update available options based on selected format"""
        if self.format_var.get() == "STDF":
            self.limits_check.config(state="disabled")
        else:
            self.limits_check.config(state="normal")

    def save_data(self):
        """Save the data based on selected options"""
        global multiple_stdf_data, multiple_wafer_ids, custom_tests, test_parameters, test_limits

        if not multiple_stdf_data:
            self.status_var.set("Error: No data loaded to save")
            return

        format_type = self.format_var.get()
        file_mode = self.file_mode_var.get()
        include_custom = self.include_custom_var.get()
        include_original = self.include_original_var.get()
        include_limits = self.include_limits_var.get()

        try:
            if format_type == "CSV":
                self.save_as_csv(file_mode, include_custom, include_original, include_limits)
            else:
                self.save_as_stdf(file_mode, include_custom, include_original)
        except Exception as e:
            self.status_var.set(f"Error: {str(e)}")
            print(f"Save error: {e}")
            import traceback
            traceback.print_exc()

    def prepare_dataframe_with_custom_tests(self, df, include_custom, include_original):
        """Prepare a dataframe with custom test columns computed"""
        global custom_tests, test_parameters

        # Create a copy to avoid modifying original
        export_df = df.copy()

        # Rename numeric columns to their test names for readability
        if include_original:
            rename_map = {}
            for test_key, test_name in test_parameters.items():
                if test_key.startswith("test_"):
                    try:
                        test_num = int(test_key.replace("test_", ""))
                        if test_num in export_df.columns:
                            rename_map[test_num] = test_name
                    except ValueError:
                        pass
            if rename_map:
                export_df = export_df.rename(columns=rename_map)
        else:
            # Remove original test columns (keep x, y, bin)
            cols_to_keep = ['x', 'y', 'bin']
            export_df = export_df[[c for c in cols_to_keep if c in export_df.columns]]

        # Add custom test columns
        if include_custom and custom_tests:
            for test_name in custom_tests.keys():
                custom_values = []
                for idx, row in df.iterrows():
                    value = evaluate_custom_test(test_name, row)
                    custom_values.append(value)
                export_df[f"CUSTOM_{test_name}"] = custom_values
                print(f"Added custom test column: CUSTOM_{test_name}")

        return export_df

    def save_as_csv(self, file_mode, include_custom, include_original, include_limits):
        """Save data as CSV file(s)"""
        global multiple_stdf_data, multiple_wafer_ids, test_limits

        # Release grab to allow file dialog to appear
        self.dialog.grab_release()

        if file_mode == "combined":
            # Save all wafers in one file
            save_path = filedialog.asksaveasfilename(
                parent=self.dialog,
                title="Save Combined CSV",
                defaultextension=".csv",
                filetypes=[("CSV files", "*.csv"), ("All files", "*.*")],
                initialfile="combined_wafer_data.csv"
            )

            if not save_path:
                self.dialog.grab_set()  # Restore grab
                return

            # Combine all dataframes
            all_dfs = []
            for df, wafer_id in zip(multiple_stdf_data, multiple_wafer_ids):
                export_df = self.prepare_dataframe_with_custom_tests(df, include_custom, include_original)
                export_df.insert(0, 'wafer_id', wafer_id)
                all_dfs.append(export_df)

            combined_df = pd.concat(all_dfs, ignore_index=True)

            # Save with optional limits header
            if include_limits:
                self.save_csv_with_limits(combined_df, save_path)
            else:
                combined_df.to_csv(save_path, index=False)

            self.status_var.set(f"Saved: {os.path.basename(save_path)}")
            print(f"Saved combined CSV to: {save_path}")

        else:
            # Save separate files per wafer
            save_dir = filedialog.askdirectory(parent=self.dialog, title="Select folder to save CSV files")

            if not save_dir:
                self.dialog.grab_set()  # Restore grab
                return

            saved_count = 0
            for df, wafer_id in zip(multiple_stdf_data, multiple_wafer_ids):
                export_df = self.prepare_dataframe_with_custom_tests(df, include_custom, include_original)

                # Clean wafer_id for filename
                clean_wafer_id = "".join(c if c.isalnum() or c in '-_' else '_' for c in str(wafer_id))
                file_path = os.path.join(save_dir, f"{clean_wafer_id}_modified.csv")

                if include_limits:
                    self.save_csv_with_limits(export_df, file_path)
                else:
                    export_df.to_csv(file_path, index=False)

                saved_count += 1

            self.status_var.set(f"Saved {saved_count} CSV files to folder")
            print(f"Saved {saved_count} CSV files to: {save_dir}")

        self.dialog.grab_set()  # Restore grab after save
        tk.messagebox.showinfo("Save Complete", f"Data saved successfully!\n\n{self.status_var.get()}")

    def save_csv_with_limits(self, df, file_path):
        """Save CSV with test limits as header rows"""
        global test_limits

        with open(file_path, 'w', newline='', encoding='utf-8') as f:
            import csv
            writer = csv.writer(f)

            # Write header
            columns = list(df.columns)
            writer.writerow(columns)

            # Write LoLimit row
            lo_limits = []
            for col in columns:
                limit_value = ""
                if col in test_limits:
                    limit_value = test_limits[col].get('lo_limit', '')
                lo_limits.append(f"LoLimit: {limit_value}" if limit_value != "" else "")
            writer.writerow(lo_limits)

            # Write HiLimit row
            hi_limits = []
            for col in columns:
                limit_value = ""
                if col in test_limits:
                    limit_value = test_limits[col].get('hi_limit', '')
                hi_limits.append(f"HiLimit: {limit_value}" if limit_value != "" else "")
            writer.writerow(hi_limits)

            # Write data rows
            for _, row in df.iterrows():
                writer.writerow(row.tolist())

    def save_as_stdf(self, file_mode, include_custom, include_original):
        """Save data as STDF file(s)"""
        global STDF_MODULE, STDF_TYPE

        # Check if STDF writing is supported
        if STDF_MODULE is None:
            tk.messagebox.showwarning("STDF Not Available",
                "STDF library is not installed.\n\n"
                "To save as STDF, install the Semi-ATE-STDF library:\n"
                "  pip install Semi-ATE-STDF\n\n"
                "Alternatively, you can save as CSV format.")
            return

        # Note: STDF writing is complex and requires proper record structure
        # For now, we'll provide a simplified version that may need expansion

        # Release grab to allow file dialog to appear
        self.dialog.grab_release()

        try:
            if file_mode == "combined":
                save_path = filedialog.asksaveasfilename(
                    parent=self.dialog,
                    title="Save STDF File",
                    defaultextension=".stdf",
                    filetypes=[("STDF files", "*.stdf"), ("All files", "*.*")],
                    initialfile="combined_wafer_data.stdf"
                )

                if not save_path:
                    self.dialog.grab_set()  # Restore grab
                    return

                self.write_stdf_file(save_path, multiple_stdf_data, multiple_wafer_ids,
                                     include_custom, include_original)

                self.status_var.set(f"Saved: {os.path.basename(save_path)}")
                self.dialog.grab_set()  # Restore grab
                tk.messagebox.showinfo("Save Complete", f"STDF saved to:\n{save_path}")

            else:
                save_dir = filedialog.askdirectory(parent=self.dialog, title="Select folder to save STDF files")

                if not save_dir:
                    self.dialog.grab_set()  # Restore grab
                    return

                saved_count = 0
                for df, wafer_id in zip(multiple_stdf_data, multiple_wafer_ids):
                    clean_wafer_id = "".join(c if c.isalnum() or c in '-_' else '_' for c in str(wafer_id))
                    file_path = os.path.join(save_dir, f"{clean_wafer_id}_modified.stdf")

                    self.write_stdf_file(file_path, [df], [wafer_id], include_custom, include_original)
                    saved_count += 1

                self.status_var.set(f"Saved {saved_count} STDF files")
                self.dialog.grab_set()  # Restore grab
                tk.messagebox.showinfo("Save Complete", f"Saved {saved_count} STDF files to:\n{save_dir}")

        except Exception as e:
            self.dialog.grab_set()  # Restore grab
            tk.messagebox.showerror("STDF Save Error",
                f"Error saving STDF file:\n{str(e)}\n\n"
                "STDF writing may not be fully supported.\n"
                "Consider saving as CSV instead.")

    def write_stdf_file(self, file_path, data_frames, wafer_ids, include_custom, include_original):
        """Write data to STDF file format using binary structure"""
        import struct

        # STDF V4 format - manual binary writing
        # Reference: STDF V4 specification

        try:
            with open(file_path, 'wb') as f:
                # Helper functions for STDF binary format
                def write_record(rec_type, rec_sub, data):
                    """Write a single STDF record"""
                    rec_len = len(data)
                    # Header: REC_LEN (2 bytes), REC_TYP (1 byte), REC_SUB (1 byte)
                    header = struct.pack('<HBB', rec_len, rec_type, rec_sub)
                    f.write(header)
                    f.write(data)

                def pack_string(s, max_len=255):
                    """Pack a string with length prefix (Cn format)"""
                    s = str(s)[:max_len]
                    return struct.pack('B', len(s)) + s.encode('ascii', errors='replace')

                def pack_time(timestamp):
                    """Pack timestamp as U*4"""
                    return struct.pack('<I', int(timestamp))

                # FAR - File Attributes Record (type 0, sub 10)
                # CPU_TYPE=2 (x86/little-endian), STDF_VER=4
                far_data = struct.pack('BB', 2, 4)
                write_record(0, 10, far_data)

                # MIR - Master Information Record (type 1, sub 10)
                current_time = int(time.time())
                mir_data = bytearray()
                mir_data += pack_time(current_time)  # SETUP_T
                mir_data += pack_time(current_time)  # START_T
                mir_data += struct.pack('B', 1)      # STAT_NUM
                mir_data += struct.pack('c', b'P')   # MODE_COD (Production)
                mir_data += struct.pack('c', b' ')   # RTST_COD
                mir_data += struct.pack('c', b' ')   # PROT_COD
                mir_data += struct.pack('<H', 0)     # BURN_TIM
                mir_data += struct.pack('c', b' ')   # CMOD_COD
                mir_data += pack_string('MODIFIED')  # LOT_ID
                mir_data += pack_string('WAFER')     # PART_TYP
                mir_data += pack_string('EXPORT')    # NODE_NAM
                mir_data += pack_string('DEVMATE')   # TSTR_TYP
                mir_data += pack_string('CUSTOM')    # JOB_NAM
                mir_data += pack_string('')          # JOB_REV
                mir_data += pack_string('')          # SBLOT_ID
                mir_data += pack_string('')          # OPER_NAM
                mir_data += pack_string('')          # EXEC_TYP
                mir_data += pack_string('')          # EXEC_VER
                mir_data += pack_string('')          # TEST_COD
                mir_data += pack_string('')          # TST_TEMP
                mir_data += pack_string('')          # USER_TXT
                mir_data += pack_string('')          # AUX_FILE
                mir_data += pack_string('')          # PKG_TYP
                mir_data += pack_string('')          # FAMLY_ID
                mir_data += pack_string('')          # DATE_COD
                mir_data += pack_string('')          # FACIL_ID
                mir_data += pack_string('')          # FLOOR_ID
                mir_data += pack_string('')          # PROC_ID
                mir_data += pack_string('')          # OPER_FRQ
                mir_data += pack_string('')          # SPEC_NAM
                mir_data += pack_string('')          # SPEC_VER
                mir_data += pack_string('')          # FLOW_ID
                mir_data += pack_string('')          # SETUP_ID
                mir_data += pack_string('')          # DSGN_REV
                mir_data += pack_string('')          # ENG_ID
                mir_data += pack_string('')          # ROM_COD
                mir_data += pack_string('')          # SERL_NUM
                mir_data += pack_string('')          # SUPR_NAM
                write_record(1, 10, bytes(mir_data))

                part_count = 0
                test_num = 0

                for df_idx, (df, wafer_id) in enumerate(zip(data_frames, wafer_ids)):
                    # WIR - Wafer Information Record (type 2, sub 10)
                    wir_data = bytearray()
                    wir_data += struct.pack('B', 1)           # HEAD_NUM
                    wir_data += struct.pack('B', 255)         # SITE_GRP (255 = all sites)
                    wir_data += pack_time(current_time)       # START_T
                    wir_data += pack_string(str(wafer_id)[:20])  # WAFER_ID
                    write_record(2, 10, bytes(wir_data))

                    # Prepare export data with custom tests
                    export_df = self.prepare_dataframe_with_custom_tests(df, include_custom, include_original)

                    # Get test columns (exclude x, y, bin, wafer_id)
                    exclude_cols = {'x', 'y', 'bin', 'wafer_id'}
                    test_cols = [c for c in export_df.columns if c not in exclude_cols]

                    # Process each die
                    for row_idx, (_, row) in enumerate(df.iterrows()):
                        x_coord = int(row.get('x', 0)) if pd.notna(row.get('x')) else 0
                        y_coord = int(row.get('y', 0)) if pd.notna(row.get('y')) else 0
                        bin_num = int(row.get('bin', 1)) if pd.notna(row.get('bin')) else 1

                        # Clamp values to valid ranges
                        bin_num = max(0, min(65535, bin_num))  # U*2 range
                        x_coord = max(-32768, min(32767, x_coord))  # I*2 range
                        y_coord = max(-32768, min(32767, y_coord))  # I*2 range

                        # PIR - Part Information Record (type 5, sub 10)
                        pir_data = struct.pack('BB', 1, 1)  # HEAD_NUM, SITE_NUM
                        write_record(5, 10, pir_data)

                        # Track tests per part
                        tests_this_part = 0

                        # PTR - Parametric Test Record (type 15, sub 10) for each test
                        for col_idx, col in enumerate(test_cols):
                            try:
                                # Get value from export_df (which has custom tests computed)
                                value = export_df.loc[row_idx, col] if col in export_df.columns else None
                                if value is None or (isinstance(value, float) and pd.isna(value)):
                                    continue

                                value = float(value)
                                test_num += 1
                                tests_this_part += 1

                                # Use column index + 1 as test number (more stable)
                                test_number = col_idx + 1

                                ptr_data = bytearray()
                                ptr_data += struct.pack('<I', test_number)  # TEST_NUM (U*4)
                                ptr_data += struct.pack('B', 1)             # HEAD_NUM
                                ptr_data += struct.pack('B', 1)             # SITE_NUM
                                ptr_data += struct.pack('B', 0)             # TEST_FLG (pass)
                                ptr_data += struct.pack('B', 0)             # PARM_FLG
                                ptr_data += struct.pack('<f', value)        # RESULT
                                ptr_data += pack_string(str(col)[:50])      # TEST_TXT
                                ptr_data += pack_string('')                 # ALARM_ID
                                ptr_data += struct.pack('B', 0)             # OPT_FLAG
                                ptr_data += struct.pack('b', 0)             # RES_SCAL
                                ptr_data += struct.pack('b', 0)             # LLM_SCAL
                                ptr_data += struct.pack('b', 0)             # HLM_SCAL
                                ptr_data += struct.pack('<f', 0.0)          # LO_LIMIT
                                ptr_data += struct.pack('<f', 0.0)          # HI_LIMIT
                                ptr_data += pack_string('')                 # UNITS
                                write_record(15, 10, bytes(ptr_data))
                            except (ValueError, TypeError) as e:
                                continue

                        # PRR - Part Results Record (type 5, sub 20)
                        # Clamp tests_this_part to U*2 range
                        tests_this_part = min(65535, tests_this_part)

                        prr_data = bytearray()
                        prr_data += struct.pack('B', 1)                    # HEAD_NUM
                        prr_data += struct.pack('B', 1)                    # SITE_NUM
                        prr_data += struct.pack('B', 0 if bin_num == 1 else 8)  # PART_FLG
                        prr_data += struct.pack('<H', tests_this_part)     # NUM_TEST (U*2)
                        prr_data += struct.pack('<H', bin_num)             # HARD_BIN (U*2)
                        prr_data += struct.pack('<H', bin_num)             # SOFT_BIN (U*2)
                        prr_data += struct.pack('<h', x_coord)             # X_COORD (I*2)
                        prr_data += struct.pack('<h', y_coord)             # Y_COORD (I*2)
                        prr_data += struct.pack('<I', 0)                   # TEST_T (U*4)
                        prr_data += pack_string('')                        # PART_ID
                        prr_data += pack_string('')                        # PART_TXT
                        prr_data += struct.pack('B', 0)                    # PART_FIX
                        write_record(5, 20, bytes(prr_data))

                        part_count += 1

                    # WRR - Wafer Results Record (type 2, sub 20)
                    wrr_data = bytearray()
                    wrr_data += struct.pack('B', 1)           # HEAD_NUM
                    wrr_data += struct.pack('B', 255)         # SITE_GRP
                    wrr_data += pack_time(current_time)       # FINISH_T
                    wrr_data += struct.pack('<I', len(df))    # PART_CNT
                    wrr_data += struct.pack('<I', 0)          # RTST_CNT
                    wrr_data += struct.pack('<I', 0)          # ABRT_CNT
                    wrr_data += struct.pack('<I', len(df))    # GOOD_CNT
                    wrr_data += struct.pack('<I', 0)          # FUNC_CNT
                    wrr_data += pack_string(str(wafer_id)[:20])  # WAFER_ID
                    write_record(2, 20, bytes(wrr_data))

                # MRR - Master Results Record (type 1, sub 20)
                mrr_data = bytearray()
                mrr_data += pack_time(current_time)  # FINISH_T
                mrr_data += struct.pack('c', b' ')   # DISP_COD
                mrr_data += pack_string('')          # USR_DESC
                mrr_data += pack_string('')          # EXC_DESC
                write_record(1, 20, bytes(mrr_data))

            print(f"Successfully wrote STDF file: {file_path}")
            print(f"  Total parts: {part_count}, Total tests: {test_num}")

        except Exception as e:
            print(f"Error writing STDF: {e}")
            import traceback
            traceback.print_exc()
            raise Exception(f"STDF write error: {e}")


def open_save_data_dialog():
    """Open the Save Data dialog"""
    global multiple_stdf_data

    if not multiple_stdf_data:
        tk.messagebox.showwarning("No Data", "No data loaded to save.\n\nPlease load STDF or CSV files first.")
        return

    SaveDataDialog(main_win)


def update_heatmap_parameter_list():
    """Update the parameter dropdown in the heatmap tab with custom tests included"""
    global test_parameters, custom_tests

    param_options = ["BIN (Bin Number)"]

    # Add regular test parameters
    for test_key, test_name in sorted(test_parameters.items()):
        param_options.append(f"{test_key}: {test_name}")

    # Add custom tests with a special prefix
    if custom_tests:
        param_options.append("─── Custom Tests ───")  # Separator
        for test_name in sorted(custom_tests.keys()):
            param_options.append(f"CUSTOM: {test_name}")
        print(f"Added {len(custom_tests)} custom tests to dropdown: {list(custom_tests.keys())}")

    heatmap_param_combobox["values"] = param_options
    print(f"Updated parameter dropdown with {len(param_options)} options")


# Control frame for STDF heatmap - Using TWO rows for better fit on smaller screens
control_frame_heatmap = tk.Frame(tab6)
control_frame_heatmap.pack(side=tk.TOP, fill=tk.X, padx=5, pady=2)

# Row 1: File loading controls
control_row1 = tk.Frame(control_frame_heatmap)
control_row1.pack(side=tk.TOP, fill=tk.X, pady=2)

# File source selector (Dropdown/Combobox)
file_source_var = tk.StringVar(value="STDF")

source_label = tk.Label(control_row1, text="Format:", font=("Helvetica", 9))
source_label.pack(side=tk.LEFT, padx=(0, 3))

file_source_combobox = ttk.Combobox(
    control_row1,
    textvariable=file_source_var,
    values=["STDF", "CSV"],
    state="readonly",
    width=6,
    font=("Helvetica", 9)
)
file_source_combobox.pack(side=tk.LEFT, padx=2)
file_source_combobox.bind("<<ComboboxSelected>>", lambda e: update_source_buttons())

# Frame to hold the file-type specific load buttons (STDF or CSV)
load_buttons_frame = tk.Frame(control_row1)
load_buttons_frame.pack(side=tk.LEFT, padx=2)

select_multiple_stdf_button = tk.Button(
    load_buttons_frame,
    text="Load STDF",
    command=lambda: load_multiple_stdf_files(),
    font=("Helvetica", 9),
)
select_multiple_stdf_button.pack(side=tk.LEFT, padx=2)

# CSV Load button (initially hidden)
select_csv_button = tk.Button(
    load_buttons_frame,
    text="Load CSV",
    command=lambda: load_csv_wafermap_file(),
    font=("Helvetica", 9),
    bg="#4CAF50",
    fg="white",
)
# Don't pack initially - will be managed by update_source_buttons()

# Button to load entire project folder (always visible)
load_project_button = tk.Button(
    control_row1,
    text="Project Folder",
    command=lambda: load_project_folder(),
    font=("Helvetica", 9),
    bg="#FF9800",
    fg="white",
)
load_project_button.pack(side=tk.LEFT, padx=2)

# Separator
tk.Label(control_row1, text="|", font=("Helvetica", 10), fg="gray").pack(side=tk.LEFT, padx=3)

# Group selection dropdown
tk.Label(control_row1, text="Group:", font=("Helvetica", 9)).pack(side=tk.LEFT, padx=2)

heatmap_group_combobox = ttk.Combobox(
    control_row1, state="readonly", width=15, font=("Helvetica", 9)
)
heatmap_group_combobox.pack(side=tk.LEFT, padx=2)
heatmap_group_combobox.bind("<<ComboboxSelected>>", lambda e: on_group_selected())

tk.Label(control_row1, text="Param:", font=("Helvetica", 9)).pack(side=tk.LEFT, padx=2)

heatmap_param_combobox = ttk.Combobox(
    control_row1, state="readonly", width=30, font=("Helvetica", 9)
)
heatmap_param_combobox.pack(side=tk.LEFT, padx=2)
heatmap_param_combobox.bind("<<ComboboxSelected>>", lambda e: refresh_heatmap_display())

heatmap_refresh_button = tk.Button(
    control_row1,
    text="Refresh",
    command=lambda: refresh_heatmap_display(),
    font=("Helvetica", 9),
)
heatmap_refresh_button.pack(side=tk.LEFT, padx=3)

# Custom Test Calculator button
custom_test_button = tk.Button(
    control_row1,
    text="🧮 Custom Test",
    command=open_custom_test_calculator,
    font=("Helvetica", 9),
    bg="#9C27B0",
    fg="white",
)
custom_test_button.pack(side=tk.LEFT, padx=3)

# Save Modified Data button
save_data_button = tk.Button(
    control_row1,
    text="💾 Save Data",
    command=lambda: open_save_data_dialog(),
    font=("Helvetica", 9),
    bg="#2196F3",
    fg="white",
)
save_data_button.pack(side=tk.LEFT, padx=3)

# Info label at the end of row 1
heatmap_info_label = tk.Label(
    control_row1,
    text="No files loaded",
    font=("Helvetica", 8),
    fg="gray"
)
heatmap_info_label.pack(side=tk.RIGHT, padx=5)

# Row 2: View controls and image options
control_row2 = tk.Frame(control_frame_heatmap)
control_row2.pack(side=tk.TOP, fill=tk.X, pady=2)

show_grid_var = tk.BooleanVar(value=False)
show_grid_checkbox = tk.Checkbutton(
    control_row2,
    text="Grid",
    variable=show_grid_var,
    command=lambda: refresh_heatmap_display(),
    font=("Helvetica", 9),
)
show_grid_checkbox.pack(side=tk.LEFT, padx=3)

# Zoom buttons
zoom_in_button = tk.Button(
    control_row2,
    text="Zoom+",
    command=lambda: zoom_heatmap(True),
    font=("Helvetica", 9),
)
zoom_in_button.pack(side=tk.LEFT, padx=1)

zoom_out_button = tk.Button(
    control_row2,
    text="Zoom-",
    command=lambda: zoom_heatmap(False),
    font=("Helvetica", 9),
)
zoom_out_button.pack(side=tk.LEFT, padx=1)

reset_zoom_button = tk.Button(
    control_row2,
    text="Reset",
    command=lambda: refresh_heatmap_display(),
    font=("Helvetica", 9),
)
reset_zoom_button.pack(side=tk.LEFT, padx=1)

clear_selection_button = tk.Button(
    control_row2,
    text="Clear Sel",
    command=lambda: clear_die_selection(),
    font=("Helvetica", 9),
)
clear_selection_button.pack(side=tk.LEFT, padx=2)

# Separator
tk.Label(control_row2, text="|", font=("Helvetica", 10), fg="gray").pack(side=tk.LEFT, padx=3)

# View Type dropdown (Data, Images, PLM Files)
tk.Label(control_row2, text="View:", font=("Helvetica", 9, "bold")).pack(side=tk.LEFT, padx=2)

view_type_var = tk.StringVar(value="Data")
view_type_combobox = ttk.Combobox(
    control_row2,
    textvariable=view_type_var,
    values=["Data", "Images", "PLM Files"],
    state="readonly",
    width=9,
    font=("Helvetica", 9)
)
view_type_combobox.pack(side=tk.LEFT, padx=2)

# Sub-options dropdown (changes based on View Type selection)
tk.Label(control_row2, text="Type:", font=("Helvetica", 9)).pack(side=tk.LEFT, padx=2)

view_subtype_var = tk.StringVar(value="Heatmap")
view_subtype_combobox = ttk.Combobox(
    control_row2,
    textvariable=view_subtype_var,
    values=["Heatmap"],
    state="readonly",
    width=10,
    font=("Helvetica", 9)
)
view_subtype_combobox.pack(side=tk.LEFT, padx=2)

def update_view_subtype_options(*args):
    """Update the subtype dropdown options based on selected view type"""
    view_type = view_type_var.get()

    if view_type == "Data":
        view_subtype_combobox["values"] = ["Heatmap"]
        view_subtype_var.set("Heatmap")
    elif view_type == "Images":
        view_subtype_combobox["values"] = ["All", "ALLON", "XHAIR", "CHKBRD", "ALLOFF"]
        view_subtype_var.set("All")
    elif view_type == "PLM Files":
        # Detect available PLM types from PLM directory
        plm_types = detect_plm_types()
        view_subtype_combobox["values"] = plm_types
        view_subtype_var.set("All")

    refresh_heatmap_display()

def detect_plm_types():
    """Detect available PLM file types from PLM directory"""
    global plm_file_directory

    plm_types = ["All"]

    # Get PLM directory
    plm_dir = plm_file_directory
    if not plm_dir and "plm" in config_dirs:
        try:
            plm_dir = config_dirs["plm"].get()
        except:
            pass

    if not plm_dir or not os.path.exists(plm_dir):
        return plm_types

    found_types = set()

    try:
        import re
        for filename in os.listdir(plm_dir):
            if filename.lower().endswith(('.plm', '.txt', '.csv', '.dat')):
                # Extract type from end of filename before extension
                # Pattern: ..._TypeName_timestamp.txt
                # Examples: CheckerSyn, InvCheckerSyn, UniformitySyn, Bridged-Pixels, PLM-Stitched-Image
                match = re.search(r'_([A-Za-z][A-Za-z0-9\-]+)_\d{14}\.', filename)
                if match:
                    plm_type = match.group(1)
                    found_types.add(plm_type)
    except Exception as e:
        print(f"Error detecting PLM types: {e}")

    if found_types:
        # Sort and add to list
        plm_types.extend(sorted(found_types))
        print(f"Detected PLM types: {plm_types}")

    return plm_types

view_type_combobox.bind("<<ComboboxSelected>>", update_view_subtype_options)
view_subtype_combobox.bind("<<ComboboxSelected>>", lambda e: refresh_heatmap_display())

# Keep backward compatibility variables
image_view_var = tk.BooleanVar(value=False)
plm_view_var = tk.BooleanVar(value=False)
image_type_view_var = tk.StringVar(value="All")

# Folder status label (shows what folders are loaded from project)
folder_status_label = tk.Label(
    control_row2,
    text="",
    font=("Helvetica", 8),
    fg="gray",
)
folder_status_label.pack(side=tk.LEFT, padx=5)

# PLM folder directory
plm_file_directory = None

# Frame for heatmap display - now with left stats panel
heatmap_main_container = tk.Frame(tab6)
heatmap_main_container.pack(side=tk.BOTTOM, fill=tk.BOTH, expand=True)

# Left panel for statistics (boxplot and probability distribution)
stats_panel = tk.Frame(heatmap_main_container, width=320, bg="#f0f0f0")
stats_panel.pack(side=tk.LEFT, fill=tk.Y, padx=5, pady=5)
stats_panel.pack_propagate(False)  # Keep fixed width

# Stats panel title
stats_title_label = tk.Label(
    stats_panel,
    text="Statistics",
    font=("Helvetica", 12, "bold"),
    bg="#f0f0f0"
)
stats_title_label.pack(side=tk.TOP, pady=5)

# Frame for boxplot / bin distribution (with selector)
boxplot_frame = tk.Frame(stats_panel, bg="#f0f0f0")
boxplot_frame.pack(side=tk.TOP, fill=tk.BOTH, expand=True, padx=5, pady=5)

# Header with selector for Boxplot vs Bin Distribution
boxplot_header_frame = tk.Frame(boxplot_frame, bg="#f0f0f0")
boxplot_header_frame.pack(side=tk.TOP, fill=tk.X, pady=2)

boxplot_type_var = tk.StringVar(value="Boxplot")
boxplot_type_combobox = ttk.Combobox(
    boxplot_header_frame,
    textvariable=boxplot_type_var,
    values=["Boxplot", "Bin Distribution"],
    state="readonly",
    width=14,
    font=("Helvetica", 9)
)
boxplot_type_combobox.pack(side=tk.LEFT, padx=2)
boxplot_type_combobox.bind("<<ComboboxSelected>>", lambda e: update_stats_plots())

# Container frame for boxplot/bin plot (below header)
boxplot_plot_frame = tk.Frame(boxplot_frame, bg="#f0f0f0")
boxplot_plot_frame.pack(side=tk.TOP, fill=tk.BOTH, expand=True)

# Frame for probability distribution
prob_frame = tk.Frame(stats_panel, bg="#f0f0f0")
prob_frame.pack(side=tk.TOP, fill=tk.BOTH, expand=True, padx=5, pady=5)

# Probability plot header with option selector (always visible at top)
prob_header_frame = tk.Frame(prob_frame, bg="#f0f0f0")
prob_header_frame.pack(side=tk.TOP, fill=tk.X, pady=2)

prob_label = tk.Label(
    prob_header_frame,
    text="Distribution:",
    font=("Helvetica", 10, "bold"),
    bg="#f0f0f0"
)
prob_label.pack(side=tk.LEFT, padx=2)

# Dropdown to choose between CDF and PDF
prob_type_var = tk.StringVar(value="CDF")
prob_type_combobox = ttk.Combobox(
    prob_header_frame,
    textvariable=prob_type_var,
    values=["CDF", "PDF"],
    state="readonly",
    width=6,
    font=("Helvetica", 9)
)
prob_type_combobox.pack(side=tk.LEFT, padx=5)
prob_type_combobox.bind("<<ComboboxSelected>>", lambda e: update_stats_plots())

# Container frame for probability plot (below header)
prob_plot_frame = tk.Frame(prob_frame, bg="#f0f0f0")
prob_plot_frame.pack(side=tk.TOP, fill=tk.BOTH, expand=True)

# Die Images & PLM panel - MUST be packed BEFORE heatmap_display_frame
# This panel shows die images and PLM files for the selected die
die_image_panel = tk.Frame(heatmap_main_container, width=280, bg="#e8e8e8")
die_image_panel.pack(side=tk.LEFT, fill=tk.Y, padx=5, pady=5)
die_image_panel.pack_propagate(False)  # Keep fixed width

# Die image panel title
die_image_title = tk.Label(
    die_image_panel,
    text="Die Images & PLM",
    font=("Helvetica", 10, "bold"),
    bg="#e8e8e8"
)
die_image_title.pack(side=tk.TOP, pady=3)

# ============== DIE IMAGES SECTION ==============
# Image section header with filter
image_section_frame = tk.Frame(die_image_panel, bg="#d0d0d0", relief="groove", bd=1)
image_section_frame.pack(side=tk.TOP, fill=tk.BOTH, expand=True, padx=3, pady=3)

image_header_frame = tk.Frame(image_section_frame, bg="#d0d0d0")
image_header_frame.pack(side=tk.TOP, fill=tk.X, padx=2, pady=2)

image_section_label = tk.Label(
    image_header_frame,
    text="📷 Images",
    font=("Helvetica", 9, "bold"),
    bg="#d0d0d0"
)
image_section_label.pack(side=tk.LEFT, padx=2)

image_type_var = tk.StringVar(value="All")
image_type_combobox = ttk.Combobox(
    image_header_frame,
    textvariable=image_type_var,
    values=["All", "ALLON", "XHAIR", "CHKBRD", "ALLOFF"],
    state="readonly",
    width=8,
    font=("Helvetica", 8)
)
image_type_combobox.pack(side=tk.RIGHT, padx=2)

# Scrollable frame for die images
die_image_canvas = tk.Canvas(image_section_frame, bg="#e8e8e8")
die_image_scrollbar_v = tk.Scrollbar(image_section_frame, orient="vertical", command=die_image_canvas.yview)
die_image_scrollbar_h = tk.Scrollbar(image_section_frame, orient="horizontal", command=die_image_canvas.xview)
die_image_scrollable_frame = tk.Frame(die_image_canvas, bg="#e8e8e8")

die_image_scrollable_frame.bind(
    "<Configure>",
    lambda e: die_image_canvas.configure(scrollregion=die_image_canvas.bbox("all"))
)

die_image_canvas.create_window((0, 0), window=die_image_scrollable_frame, anchor="nw")
die_image_canvas.configure(xscrollcommand=die_image_scrollbar_h.set, yscrollcommand=die_image_scrollbar_v.set)

die_image_scrollbar_v.pack(side=tk.RIGHT, fill=tk.Y)
die_image_scrollbar_h.pack(side=tk.BOTTOM, fill=tk.X)
die_image_canvas.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)

# ============== PLM FILES SECTION ==============
# PLM section with same structure as images
plm_section_frame = tk.Frame(die_image_panel, bg="#d0d0d0", relief="groove", bd=1)
plm_section_frame.pack(side=tk.TOP, fill=tk.BOTH, expand=True, padx=3, pady=3)

plm_header_frame = tk.Frame(plm_section_frame, bg="#d0d0d0")
plm_header_frame.pack(side=tk.TOP, fill=tk.X, padx=2, pady=2)

plm_section_label = tk.Label(
    plm_header_frame,
    text="📊 PLM Files",
    font=("Helvetica", 9, "bold"),
    bg="#d0d0d0"
)
plm_section_label.pack(side=tk.LEFT, padx=2)

plm_type_var = tk.StringVar(value="All")
plm_type_combobox = ttk.Combobox(
    plm_header_frame,
    textvariable=plm_type_var,
    values=["All"],
    state="readonly",
    width=10,
    font=("Helvetica", 8)
)
plm_type_combobox.pack(side=tk.RIGHT, padx=2)

# Scrollable frame for PLM files
plm_canvas = tk.Canvas(plm_section_frame, bg="#e8e8e8")
plm_scrollbar = tk.Scrollbar(plm_section_frame, orient="vertical", command=plm_canvas.yview)
plm_scrollable_frame = tk.Frame(plm_canvas, bg="#e8e8e8")

plm_scrollable_frame.bind(
    "<Configure>",
    lambda e: plm_canvas.configure(scrollregion=plm_canvas.bbox("all"))
)

plm_canvas.create_window((0, 0), window=plm_scrollable_frame, anchor="nw")
plm_canvas.configure(yscrollcommand=plm_scrollbar.set)

plm_scrollbar.pack(side=tk.RIGHT, fill=tk.Y)
plm_canvas.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)

# PLM file references
plm_file_refs = []
available_plm_types = ["All"]
current_selected_die_plm = None
_plm_highlight_rect = None  # Highlight rectangle for PLM wafermap
_image_highlight_rect = None  # Highlight rectangle for Image wafermap
die_image_refs = []
current_selected_die = None
available_image_types = ["All"]

# Right panel for heatmap display
heatmap_display_frame = tk.Frame(heatmap_main_container)
heatmap_display_frame.pack(side=tk.RIGHT, fill=tk.BOTH, expand=True)

# Limits panel - overlaid on heatmap in top-left corner
# Variables to store current limit values for editing
limit_lo_var = tk.StringVar(value="")
limit_hi_var = tk.StringVar(value="")
limit_units_var = tk.StringVar(value="")

# Create a floating limits frame that will be placed over the heatmap
limits_overlay_frame = tk.Frame(heatmap_display_frame, bg="white", relief="solid", bd=1)

# Title for limits
limits_title_label = tk.Label(
    limits_overlay_frame,
    text="Test Limits",
    font=("Helvetica", 10, "bold"),
    bg="white"
)
limits_title_label.grid(row=0, column=0, columnspan=3, pady=3, padx=5)

# Low limit row
lo_label = tk.Label(limits_overlay_frame, text="Lo:", font=("Helvetica", 9), bg="white")
lo_label.grid(row=1, column=0, padx=3, pady=2, sticky="e")
lo_entry = tk.Entry(limits_overlay_frame, textvariable=limit_lo_var, width=12, font=("Helvetica", 9))
lo_entry.grid(row=1, column=1, padx=3, pady=2)

# High limit row
hi_label = tk.Label(limits_overlay_frame, text="Hi:", font=("Helvetica", 9), bg="white")
hi_label.grid(row=2, column=0, padx=3, pady=2, sticky="e")
hi_entry = tk.Entry(limits_overlay_frame, textvariable=limit_hi_var, width=12, font=("Helvetica", 9))
hi_entry.grid(row=2, column=1, padx=3, pady=2)

# Units label
units_display_label = tk.Label(limits_overlay_frame, textvariable=limit_units_var, font=("Helvetica", 9), bg="white", fg="gray")
units_display_label.grid(row=1, column=2, rowspan=2, padx=3, pady=2, sticky="w")

# Apply button
def apply_limits():
    """Apply the edited limits to the colorbar and refresh display"""
    global _colorbar_range_slider

    try:
        lo_val = float(limit_lo_var.get()) if limit_lo_var.get() else None
        hi_val = float(limit_hi_var.get()) if limit_hi_var.get() else None

        if lo_val is not None and hi_val is not None and lo_val < hi_val:
            # Update the colorbar range slider if available
            if _colorbar_range_slider is not None:
                _colorbar_range_slider['low_val'] = lo_val
                _colorbar_range_slider['high_val'] = hi_val

                # Update markers and lines
                if 'low_marker' in _colorbar_range_slider:
                    _colorbar_range_slider['low_marker'].set_offsets([[0.5, lo_val]])
                if 'high_marker' in _colorbar_range_slider:
                    _colorbar_range_slider['high_marker'].set_offsets([[0.5, hi_val]])
                if 'low_line' in _colorbar_range_slider:
                    _colorbar_range_slider['low_line'].set_ydata([lo_val, lo_val])
                if 'high_line' in _colorbar_range_slider:
                    _colorbar_range_slider['high_line'].set_ydata([hi_val, hi_val])
                if 'low_text' in _colorbar_range_slider:
                    _colorbar_range_slider['low_text'].set_position((1.8, lo_val))
                    _colorbar_range_slider['low_text'].set_text(f'{lo_val:.2f}')
                if 'high_text' in _colorbar_range_slider:
                    _colorbar_range_slider['high_text'].set_position((1.8, hi_val))
                    _colorbar_range_slider['high_text'].set_text(f'{hi_val:.2f}')

                # Update all heatmap color limits
                if heatmap_canvas:
                    fig = heatmap_canvas.figure
                    for ax in fig.axes:
                        for child in ax.get_children():
                            if hasattr(child, 'set_clim'):
                                child.set_clim(vmin=lo_val, vmax=hi_val)
                    heatmap_canvas.draw_idle()

            print(f"Applied limits: Lo={lo_val}, Hi={hi_val}")
        else:
            print("Invalid limits: Lo must be less than Hi")
    except ValueError as e:
        print(f"Invalid limit values: {e}")

apply_limits_button = tk.Button(
    limits_overlay_frame,
    text="Apply",
    command=apply_limits,
    font=("Helvetica", 8),
    bg="#4CAF50",
    fg="white",
    width=6
)
apply_limits_button.grid(row=3, column=0, columnspan=3, pady=3, padx=5)

# Function to update limits display when parameter changes
def update_limits_display(test_num=None):
    """Update the limits display based on selected parameter"""
    global test_limits

    selected = heatmap_param_combobox.get()

    if selected.startswith("BIN"):
        # BIN doesn't have limits
        limit_lo_var.set("")
        limit_hi_var.set("")
        limit_units_var.set("")
        limits_title_label.config(text="Test Limits (N/A)")
        return

    # Extract test number from selection
    test_key = selected.split(":")[0].strip()
    if test_key.startswith("test_"):
        test_num = int(test_key.replace("test_", ""))
    else:
        try:
            test_num = int(test_key)
        except ValueError:
            limit_lo_var.set("")
            limit_hi_var.set("")
            limit_units_var.set("")
            return

    # Get limits for this test
    if test_num in test_limits:
        limits = test_limits[test_num]
        lo = limits.get('lo_limit')
        hi = limits.get('hi_limit')
        units = limits.get('units', '')

        limit_lo_var.set(f"{lo:.4g}" if lo is not None else "")
        limit_hi_var.set(f"{hi:.4g}" if hi is not None else "")
        limit_units_var.set(f"[{units}]" if units else "")
        limits_title_label.config(text="Test Limits")
    else:
        limit_lo_var.set("")
        limit_hi_var.set("")
        limit_units_var.set("")
        limits_title_label.config(text="Test Limits (N/A)")

# Initially hide the limits overlay (will be shown when heatmap is displayed)
# limits_overlay_frame.place_forget()  # Start hidden


def update_group_combobox():
    """Update the group selection combobox with available groups"""
    global grouped_parameters, custom_tests

    group_names = ["All Groups"]

    if grouped_parameters:
        group_names += sorted(grouped_parameters.keys())

    # Add Custom Tests group if there are custom tests
    if custom_tests:
        group_names.append("─── Custom Tests ───")

    heatmap_group_combobox["values"] = group_names
    heatmap_group_combobox.current(0)
    print(f"Updated group combobox with {len(group_names)} groups")


def on_group_selected():
    """Handle group selection - update parameter dropdown with tests from selected group"""
    global grouped_parameters, test_parameters, custom_tests

    selected_group = heatmap_group_combobox.get()

    param_options = ["BIN (Bin Number)"]

    if selected_group == "─── Custom Tests ───":
        # Show only custom tests
        for test_name in sorted(custom_tests.keys()):
            param_options.append(f"CUSTOM: {test_name}")
    elif selected_group == "All Groups" or not grouped_parameters:
        # Show all parameters
        for test_key, test_name in sorted(test_parameters.items()):
            param_options.append(f"{test_key}: {test_name}")

        # Also add custom tests at the end
        if custom_tests:
            param_options.append("─── Custom Tests ───")
            for test_name in sorted(custom_tests.keys()):
                param_options.append(f"CUSTOM: {test_name}")
    else:
        # Show only parameters from selected group
        if selected_group in grouped_parameters:
            for test_num, short_name, full_name in sorted(grouped_parameters[selected_group], key=lambda x: x[0]):
                param_options.append(f"test_{test_num}: {short_name}")

    heatmap_param_combobox["values"] = param_options
    if param_options:
        heatmap_param_combobox.current(0)

    print(f"Group '{selected_group}' selected: {len(param_options)-1} parameters")
    refresh_heatmap_display()


def open_plot_zoom_window(fig_func, title):
    """Open a zoomed view of a plot in a separate window with zoom functionality"""
    zoom_win = tk.Toplevel(main_win)
    zoom_win.title(title)
    zoom_win.geometry("800x600")

    # Create larger figure
    fig, ax = fig_func(figsize=(10, 8))

    # Create canvas with navigation toolbar for zoom
    canvas = FigureCanvasTkAgg(fig, master=zoom_win)
    canvas.get_tk_widget().pack(side=tk.TOP, fill=tk.BOTH, expand=True)

    # Add navigation toolbar for zoom/pan
    toolbar = NavigationToolbar2Tk(canvas, zoom_win)
    toolbar.update()
    canvas.get_tk_widget().pack(side=tk.TOP, fill=tk.BOTH, expand=True)

    canvas.draw()

    def on_close():
        plt.close(fig)
        zoom_win.destroy()

    zoom_win.protocol("WM_DELETE_WINDOW", on_close)
    return zoom_win


def update_stats_plots():
    """Update boxplot/bin distribution and probability distribution plots based on selected parameter"""
    global stats_boxplot_canvas, stats_prob_canvas, multiple_stdf_data, current_stdf_data

    # Clear existing canvases in boxplot plot frame (but keep header)
    for widget in boxplot_plot_frame.winfo_children():
        widget.destroy()

    # Clear existing canvases in probability plot frame (but keep header)
    for widget in prob_plot_frame.winfo_children():
        widget.destroy()

    # Determine data source
    if multiple_stdf_data and len(multiple_stdf_data) > 0:
        data_sources = multiple_stdf_data
        wafer_labels = multiple_wafer_ids
    elif current_stdf_data is not None and not current_stdf_data.empty:
        data_sources = [current_stdf_data]
        wafer_labels = [current_wafer_id if current_wafer_id else "Wafer"]
    else:
        return

    # Get selected parameter
    selected = heatmap_param_combobox.get()
    if not selected:
        return

    # Check if this is a custom test or separator
    if selected.startswith("───"):
        return  # Skip separator

    if selected.startswith("BIN"):
        param_column = "bin"
        param_label_text = "Bin"
    elif selected.startswith("CUSTOM:"):
        # Handle custom test
        custom_test_name = selected.replace("CUSTOM:", "").strip()
        param_column = f"_custom_{custom_test_name}"
        param_label_text = f"Custom: {custom_test_name}"

        # Compute custom test values for all data sources
        for df in data_sources:
            if param_column not in df.columns:
                custom_values = []
                for idx, row in df.iterrows():
                    value = evaluate_custom_test(custom_test_name, row)
                    custom_values.append(value)
                df[param_column] = custom_values
                print(f"Computed custom test '{custom_test_name}' for stats plots: {len(df)} dies")
    else:
        test_key = selected.split(":")[0].strip()
        if test_key.startswith("test_"):
            param_column = int(test_key.replace("test_", ""))
        else:
            try:
                param_column = int(test_key)
            except ValueError:
                return  # Invalid parameter
        param_label_text = selected.split(":")[-1].strip() if ":" in selected else selected

    # Check which plot type is selected (Boxplot or Bin Distribution)
    plot_selection = boxplot_type_var.get()

    if plot_selection == "Bin Distribution":
        # Create Bin Distribution bar chart
        all_bins = []
        for df in data_sources:
            if "bin" in df.columns:
                bins = df["bin"].dropna().values
                all_bins.extend(bins)

        if len(all_bins) > 0:
            all_bins = np.array(all_bins)
            unique_bins, bin_counts = np.unique(all_bins, return_counts=True)
            total_dies = len(all_bins)

            # Calculate percentages
            bin_percentages = (bin_counts / total_dies) * 100

            # Bin 0 is pass bin
            pass_bin = 0

            # Create figure for bin distribution chart
            fig_bin, ax_bin = plt.subplots(figsize=(2.8, 2.5))
            fig_bin.patch.set_facecolor('#f0f0f0')

            # Create color list (green for pass bin 0, red for fail bins)
            colors = []
            for b in unique_bins:
                if b == pass_bin:
                    colors.append('#4CAF50')  # Green for pass
                else:
                    colors.append('#F44336')  # Red for fail

            # Create bar chart
            x_pos = np.arange(len(unique_bins))
            bars = ax_bin.bar(x_pos, bin_percentages, color=colors, edgecolor='black', linewidth=0.5)

            # Add percentage labels on bars
            for bar, pct in zip(bars, bin_percentages):
                height = bar.get_height()
                ax_bin.annotate(f'{pct:.1f}%',
                    xy=(bar.get_x() + bar.get_width() / 2, height),
                    xytext=(0, 2),
                    textcoords="offset points",
                    ha='center', va='bottom',
                    fontsize=6, fontweight='bold')

            # Customize chart
            ax_bin.set_xticks(x_pos)
            ax_bin.set_xticklabels([f'Bin {int(b)}' for b in unique_bins], fontsize=6, rotation=45, ha='right')
            ax_bin.set_ylabel("Percentage (%)", fontsize=7)
            ax_bin.set_title("Bin Distribution", fontsize=8, fontweight="bold")
            ax_bin.tick_params(axis='both', which='major', labelsize=6)
            ax_bin.grid(True, alpha=0.3, linestyle='--', linewidth=0.5, axis='y')
            ax_bin.set_ylim(0, max(bin_percentages) * 1.15)  # Add space for labels

            # Add total yield info (Bin 0 = pass)
            pass_count = bin_counts[unique_bins == pass_bin].sum() if pass_bin in unique_bins else 0
            yield_pct = (pass_count / total_dies) * 100
            ax_bin.text(0.98, 0.98, f'Yield: {yield_pct:.1f}%\nTotal: {total_dies}',
                transform=ax_bin.transAxes, fontsize=6, fontweight='bold',
                va='top', ha='right',
                bbox=dict(boxstyle='round', facecolor='white', alpha=0.8))

            fig_bin.tight_layout()

            stats_boxplot_canvas = FigureCanvasTkAgg(fig_bin, master=boxplot_plot_frame)
            stats_boxplot_canvas.get_tk_widget().pack(fill=tk.BOTH, expand=True)
            stats_boxplot_canvas.draw()

    else:
        # Create Boxplot (default)
        # Collect data from all sources
        all_data = []
        labels = []

        for df, label in zip(data_sources, wafer_labels):
            if param_column in df.columns:
                values = df[param_column].dropna().values
                if len(values) > 0:
                    all_data.append(values)
                    # Truncate long labels
                    short_label = label[:15] + "..." if len(str(label)) > 15 else str(label)
                    labels.append(short_label)

        if not all_data:
            return

        fig_box, ax_box = plt.subplots(figsize=(2.8, 2.5))
        fig_box.patch.set_facecolor('white')

        # Professional boxplot styling
        bp = ax_box.boxplot(
            all_data,
            tick_labels=labels if len(labels) <= 3 else [f"W{i+1}" for i in range(len(labels))],
            vert=True,
            patch_artist=True,
            showmeans=True,
            widths=0.6,
            meanprops=dict(marker="D", markerfacecolor="#E74C3C", markeredgecolor="white", markersize=5, markeredgewidth=1),
            medianprops=dict(color="#2C3E50", linewidth=2),
            whiskerprops=dict(color="#2C3E50", linewidth=1.5, linestyle='-'),
            capprops=dict(color="#2C3E50", linewidth=1.5),
            flierprops=dict(marker='o', markerfacecolor='#95A5A6', markeredgecolor='#7F8C8D', markersize=3, alpha=0.6),
            boxprops=dict(linewidth=1.5)
        )

        # Professional color palette
        professional_colors = ['#3498DB', '#2ECC71', '#9B59B6', '#E67E22', '#1ABC9C', '#E74C3C']
        for idx, patch in enumerate(bp["boxes"]):
            color = professional_colors[idx % len(professional_colors)]
            patch.set_facecolor(color)
            patch.set_alpha(0.75)
            patch.set_edgecolor('#2C3E50')

        # Calculate statistics for all combined data
        combined_data = np.concatenate(all_data)
        stats_max = np.max(combined_data)
        stats_min = np.min(combined_data)
        stats_q1 = np.percentile(combined_data, 25)
        stats_q3 = np.percentile(combined_data, 75)
        stats_mean = np.mean(combined_data)
        stats_median = np.median(combined_data)

        # Add statistics box in top left corner
        stats_text = (
            f"Max: {stats_max:.3g}\n"
            f"Min: {stats_min:.3g}\n"
            f"Q1: {stats_q1:.3g}\n"
            f"Q3: {stats_q3:.3g}\n"
            f"Mean: {stats_mean:.3g}\n"
            f"Median: {stats_median:.3g}"
        )
        ax_box.text(0.02, 0.98, stats_text,
            transform=ax_box.transAxes,
            fontsize=5,
            fontweight='normal',
            fontfamily='monospace',
            verticalalignment='top',
            horizontalalignment='left',
            bbox=dict(boxstyle='round,pad=0.3', facecolor='white', edgecolor='#BDC3C7', alpha=0.9, linewidth=1)
        )

        ax_box.set_title(param_label_text[:20], fontsize=8, fontweight="bold", color='#2C3E50')
        ax_box.tick_params(axis='both', which='major', labelsize=6, colors='#2C3E50')
        ax_box.set_ylabel("Value", fontsize=7, color='#2C3E50')
        ax_box.set_facecolor('#FAFAFA')
        ax_box.grid(True, alpha=0.4, linestyle='-', linewidth=0.5, color='#BDC3C7')
        ax_box.spines['top'].set_visible(False)
        ax_box.spines['right'].set_visible(False)
        ax_box.spines['left'].set_color('#BDC3C7')
        ax_box.spines['bottom'].set_color('#BDC3C7')

        fig_box.tight_layout()

        stats_boxplot_canvas = FigureCanvasTkAgg(fig_box, master=boxplot_plot_frame)
        stats_boxplot_canvas.get_tk_widget().pack(fill=tk.BOTH, expand=True)
        stats_boxplot_canvas.draw()

    # Create probability distribution plot (CDF or PDF based on selection)
    # Collect data for probability plot
    all_data = []
    labels = []
    for df, label in zip(data_sources, wafer_labels):
        if param_column in df.columns:
            values = df[param_column].dropna().values
            if len(values) > 0:
                all_data.append(values)
                short_label = label[:15] + "..." if len(str(label)) > 15 else str(label)
                labels.append(short_label)

    if not all_data:
        return

    fig_prob, ax_prob = plt.subplots(figsize=(2.8, 2.5))
    fig_prob.patch.set_facecolor('white')

    plot_type = prob_type_var.get()  # Get CDF or PDF selection

    # Professional color palette
    professional_colors = ['#3498DB', '#2ECC71', '#9B59B6', '#E67E22', '#1ABC9C', '#E74C3C']

    if plot_type == "PDF":
        # Probability Density Function (histogram-based) - professional style
        for idx, (data, label) in enumerate(zip(all_data, labels)):
            color = professional_colors[idx % len(professional_colors)]
            ax_prob.hist(data, bins=30, density=True, alpha=0.6, label=label,
                        color=color, edgecolor='white', linewidth=0.5)
        ax_prob.set_title("PDF", fontsize=8, fontweight="bold", color='#2C3E50')
        ax_prob.set_ylabel("Density", fontsize=7, color='#2C3E50')
    else:
        # Cumulative Distribution Function (default) - professional style
        for idx, (data, label) in enumerate(zip(all_data, labels)):
            color = professional_colors[idx % len(professional_colors)]
            sorted_data = np.sort(data)
            prob = np.linspace(0, 1, len(sorted_data), endpoint=False)
            ax_prob.plot(sorted_data, prob, label=label, linewidth=2, alpha=0.85, color=color)
            ax_prob.fill_between(sorted_data, prob, alpha=0.15, color=color)
        ax_prob.set_title("CDF", fontsize=8, fontweight="bold", color='#2C3E50')
        ax_prob.set_ylabel("Probability", fontsize=7, color='#2C3E50')

    ax_prob.set_xlabel("Value", fontsize=7, color='#2C3E50')
    ax_prob.tick_params(axis='both', which='major', labelsize=6, colors='#2C3E50')
    ax_prob.set_facecolor('#FAFAFA')
    ax_prob.grid(True, alpha=0.4, linestyle='-', linewidth=0.5, color='#BDC3C7')
    ax_prob.spines['top'].set_visible(False)
    ax_prob.spines['right'].set_visible(False)
    ax_prob.spines['left'].set_color('#BDC3C7')
    ax_prob.spines['bottom'].set_color('#BDC3C7')

    if len(labels) <= 3:
        ax_prob.legend(fontsize=5, loc='lower right' if plot_type == "CDF" else 'upper right',
                      framealpha=0.9, edgecolor='#BDC3C7')

    fig_prob.tight_layout()

    stats_prob_canvas = FigureCanvasTkAgg(fig_prob, master=prob_plot_frame)
    stats_prob_canvas.get_tk_widget().pack(fill=tk.BOTH, expand=True)
    stats_prob_canvas.draw()


def load_multiple_stdf_files():
    """Load multiple STDF files for multi-plot heatmap display - THREADED VERSION"""
    global multiple_stdf_data, multiple_wafer_ids, test_parameters

    stdf_paths = filedialog.askopenfilenames(
        title="Select multiple STDF files",
        filetypes=[("STDF files", "*.stdf"), ("All files", "*.*")],
    )

    if not stdf_paths:
        print("No files selected.")
        return

    # Start threaded loading with progress dialog
    load_stdf_files_threaded(list(stdf_paths), "Loading STDF Files")


def load_stdf_files_threaded(stdf_paths, title="Loading"):
    """Load multiple STDF files using threads with progress dialog"""
    global multiple_stdf_data, multiple_wafer_ids, test_parameters

    num_files = len(stdf_paths)

    # Create progress dialog
    progress_win = tk.Toplevel(main_win)
    progress_win.title(title)
    progress_win.geometry("500x200")
    progress_win.transient(main_win)
    progress_win.grab_set()

    # Center the dialog
    progress_win.update_idletasks()
    x = main_win.winfo_x() + (main_win.winfo_width() - 500) // 2
    y = main_win.winfo_y() + (main_win.winfo_height() - 200) // 2
    progress_win.geometry(f"+{x}+{y}")

    # Progress UI elements
    title_label = tk.Label(
        progress_win,
        text=f"Loading {num_files} STDF file(s)...",
        font=("Helvetica", 12, "bold")
    )
    title_label.pack(pady=10)

    progress_var = tk.DoubleVar(value=0)
    progress_bar = ttk.Progressbar(
        progress_win,
        variable=progress_var,
        maximum=num_files,
        length=400,
        mode='determinate'
    )
    progress_bar.pack(pady=10)

    status_label = tk.Label(
        progress_win,
        text="Initializing...",
        font=("Helvetica", 10)
    )
    status_label.pack(pady=5)

    file_label = tk.Label(
        progress_win,
        text="",
        font=("Helvetica", 9),
        fg="gray"
    )
    file_label.pack(pady=5)

    time_label = tk.Label(
        progress_win,
        text="",
        font=("Helvetica", 9),
        fg="blue"
    )
    time_label.pack(pady=5)

    # Storage for results
    results = []
    results_lock = threading.Lock()
    completed_count = [0]
    start_time = [time.time()]
    cancelled = [False]

    def load_single_file(stdf_path):
        """Load a single STDF file (runs in thread)"""
        if cancelled[0]:
            return None

        filename = os.path.basename(stdf_path)
        print(f"[Thread] Loading: {filename}")

        try:
            df, wafer_id, test_params, grouped_params, limits_dict, wafer_cfg = read_wafermap_from_stdf(stdf_path)

            if not df.empty:
                return {
                    'df': df,
                    'wafer_id': wafer_id if wafer_id else filename,
                    'test_params': test_params,
                    'grouped_params': grouped_params,
                    'test_limits': limits_dict,
                    'wafer_config': wafer_cfg,
                    'filename': filename
                }
        except Exception as e:
            print(f"[Thread] Error loading {filename}: {e}")

        return None

    def on_file_complete(future, stdf_path):
        """Called when a file finishes loading"""
        if cancelled[0]:
            return

        try:
            result = future.result()

            with results_lock:
                if result:
                    results.append(result)
                completed_count[0] += 1
                count = completed_count[0]

            # Update UI (must be on main thread)
            def update_ui():
                if cancelled[0] or not progress_win.winfo_exists():
                    return

                elapsed = time.time() - start_time[0]
                files_per_sec = count / elapsed if elapsed > 0 else 0
                remaining = num_files - count
                eta = remaining / files_per_sec if files_per_sec > 0 else 0

                progress_var.set(count)
                status_label.config(text=f"Loaded {count}/{num_files} files")
                file_label.config(text=f"Last: {os.path.basename(stdf_path)}")
                time_label.config(text=f"Speed: {files_per_sec:.1f} files/sec | ETA: {eta:.1f}s")

                # Check if all done
                if count >= num_files:
                    finish_loading()

            main_win.after(0, update_ui)

        except Exception as e:
            print(f"Error in callback: {e}")

    def finish_loading():
        """Called when all files are loaded"""
        global multiple_stdf_data, multiple_wafer_ids, test_parameters, grouped_parameters, test_limits
        global current_wafer_config

        elapsed = time.time() - start_time[0]

        # Process results
        multiple_stdf_data = []
        multiple_wafer_ids = []

        for result in results:
            multiple_stdf_data.append(result['df'])
            multiple_wafer_ids.append(result['wafer_id'])

            if not test_parameters:
                test_parameters = result['test_params']
            else:
                test_parameters.update(result['test_params'])

            # Handle grouped parameters
            if 'grouped_params' in result:
                if not grouped_parameters:
                    grouped_parameters = result['grouped_params']
                else:
                    for group, params in result['grouped_params'].items():
                        if group not in grouped_parameters:
                            grouped_parameters[group] = []
                        grouped_parameters[group].extend(params)

            # Handle test limits
            if 'test_limits' in result:
                if not test_limits:
                    test_limits = result['test_limits']
                else:
                    test_limits.update(result['test_limits'])

        # Set wafer config from first loaded file (for notch display)
        if results and results[0].get('wafer_config'):
            current_wafer_config = results[0]['wafer_config']
            print(f"Wafer Config set: Notch={current_wafer_config.get('notch_orientation')}")

        # Update group combobox
        update_group_combobox()

        # Update parameter combobox
        if multiple_stdf_data:
            param_options = ["BIN (Bin Number)"]
            for test_key, test_name in sorted(test_parameters.items()):
                param_options.append(f"{test_key}: {test_name}")

            heatmap_param_combobox["values"] = param_options
            if param_options:
                heatmap_param_combobox.current(0)

        # Update info label
        heatmap_info_label.config(
            text=f"Loaded {len(multiple_stdf_data)} STDF files in {elapsed:.1f}s"
        )

        print(f"\n{'='*60}")
        print(f"Successfully loaded {len(multiple_stdf_data)} STDF files in {elapsed:.1f}s")
        print(f"Average: {elapsed/len(results):.2f}s per file" if results else "")
        print(f"{'='*60}\n")

        # Close progress dialog
        progress_win.destroy()

        # Update heatmap display
        if multiple_stdf_data:
            update_multi_stdf_heatmap()

    def on_cancel():
        """Cancel loading"""
        cancelled[0] = True
        progress_win.destroy()
        print("Loading cancelled by user")

    # Cancel button
    cancel_btn = tk.Button(
        progress_win,
        text="Cancel",
        command=on_cancel,
        font=("Helvetica", 10)
    )
    cancel_btn.pack(pady=10)

    # Submit all files to thread pool
    futures = []
    for stdf_path in stdf_paths:
        future = thread_pool.submit(load_single_file, stdf_path)
        future.add_done_callback(lambda f, p=stdf_path: on_file_complete(f, p))
        futures.append(future)

    print(f"Submitted {len(futures)} files to thread pool ({MAX_WORKERS} workers)")


def load_project_folder():
    """Load entire project folder with standard subfolder structure - THREADED VERSION"""
    global project_folder, die_image_directory, plm_file_directory, multiple_stdf_data, multiple_wafer_ids, test_parameters

    folder_path = filedialog.askdirectory(
        title="Select Project Folder (with STDDatalog, ImageCaptures, etc.)"
    )

    if not folder_path:
        print("No folder selected.")
        return

    project_folder = folder_path
    print(f"\n{'='*60}")
    print(f"Loading project folder: {folder_path}")
    print(f"{'='*60}")

    # Check for expected subfolders
    found_folders = {}
    for key, subfolder in project_subfolders.items():
        subfolder_path = os.path.join(folder_path, subfolder)
        if os.path.exists(subfolder_path):
            found_folders[key] = subfolder_path
            print(f"  Found {subfolder}: {subfolder_path}")
        else:
            print(f"  Missing {subfolder}")

    # Build folder status text
    folder_status_parts = []

    # Set image folder from ImageCaptures (quick, no threading needed)
    if "images" in found_folders:
        die_image_directory = found_folders["images"]
        folder_status_parts.append("Img:✓")
        print(f"Set image folder: {die_image_directory}")

        # Count images in background
        def count_images():
            count = len([f for f in os.listdir(die_image_directory)
                        if f.lower().endswith(('.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tif', '.tiff'))])
            main_win.after(0, lambda: print(f"  Image folder has {count} images"))

        thread_pool.submit(count_images)

    # Set PLM folder from PLMFiles
    if "plm" in found_folders:
        plm_file_directory = found_folders["plm"]
        folder_status_parts.append("PLM:✓")
        print(f"Set PLM folder: {plm_file_directory}")

        # Count PLM files in background
        def count_plm():
            count = len([f for f in os.listdir(plm_file_directory)
                        if f.lower().endswith(('.plm', '.txt', '.csv', '.dat'))])
            main_win.after(0, lambda: print(f"  PLM folder has {count} files"))

        thread_pool.submit(count_plm)

    # Update folder status label
    if folder_status_parts:
        folder_status_label.config(text=" | ".join(folder_status_parts), fg="green")
    else:
        folder_status_label.config(text="No folders", fg="gray")

    # Load STDF files from STDDatalog using threads
    if "stdf" in found_folders:
        stdf_folder = found_folders["stdf"]
        stdf_files = [f for f in os.listdir(stdf_folder)
                      if f.lower().endswith(('.stdf', '.std'))]

        if stdf_files:
            print(f"\nFound {len(stdf_files)} STDF file(s):")
            stdf_paths = [os.path.join(stdf_folder, f) for f in stdf_files]

            # Store project info for later update
            project_name = folder_path.split("/")[-1] if "/" in folder_path else folder_path.split("\\")[-1]

            # Use threaded loading
            load_stdf_files_threaded(stdf_paths, f"Loading Project: {project_name}")
        else:
            print("  No STDF files found in STDDatalog folder")
            # Update info label without STDF
            project_name = folder_path.split("/")[-1] if "/" in folder_path else folder_path.split("\\")[-1]
            info_parts = []
            if die_image_directory:
                info_parts.append("Images")
            if plm_file_directory:
                info_parts.append("PLM")
            if "csv" in found_folders:
                info_parts.append("CSV")
            if "txt" in found_folders:
                info_parts.append("TXT")
            heatmap_info_label.config(text=f"Project: {project_name} ({', '.join(info_parts)})")
    else:
        # No STDF folder found
        project_name = folder_path.split("/")[-1] if "/" in folder_path else folder_path.split("\\")[-1]
        info_parts = []
        if die_image_directory:
            info_parts.append("Images")
        if plm_file_directory:
            info_parts.append("PLM")
        if "csv" in found_folders:
            info_parts.append("CSV")
        if "txt" in found_folders:
            info_parts.append("TXT")
        heatmap_info_label.config(text=f"Project: {project_name} ({', '.join(info_parts)})")

    print(f"\n{'='*60}")
    print(f"Project loading initiated!")
    print(f"{'='*60}\n")


def update_multi_stdf_heatmap():
    """Update multi-plot STDF heatmap display - OPTIMIZED VERSION with integrated stats"""
    global heatmap_canvas, multiple_stdf_data, multiple_wafer_ids

    if not multiple_stdf_data:
        print("No STDF data loaded for multi-plot")
        return

    selected = heatmap_param_combobox.get()

    if not selected:
        print("No parameter selected")
        return

    # Check if this is a custom test
    is_custom_test = selected.startswith("CUSTOM:")
    custom_test_name = None

    if selected.startswith("BIN"):
        param_column = "bin"
        param_label = "Bin"
    elif is_custom_test:
        # Handle custom test - we'll compute values for each die in each file
        custom_test_name = selected.replace("CUSTOM:", "").strip()
        param_column = f"_custom_{custom_test_name}"
        param_label = f"Custom: {custom_test_name}"

        # Compute custom test values for all dies in all files
        for df in multiple_stdf_data:
            custom_values = []
            for idx, row in df.iterrows():
                value = evaluate_custom_test(custom_test_name, row)
                custom_values.append(value)
            df[param_column] = custom_values
        print(f"Computed custom test '{custom_test_name}' for {len(multiple_stdf_data)} files")
    elif selected.startswith("───"):
        # This is a separator, ignore it
        print("Separator selected - please select a valid parameter")
        return
    else:
        test_key = selected.split(":")[0].strip()
        if test_key.startswith("test_"):
            param_column = int(test_key.replace("test_", ""))
        else:
            param_column = int(test_key)
        param_label = selected

    num_files = len(multiple_stdf_data)
    cols = min(3, num_files)
    rows = (num_files + cols - 1) // cols

    # Adjusted figure size - taller to fill more vertical space
    fig, axes = plt.subplots(
        rows, cols, figsize=(8 * cols, 8 * rows), constrained_layout=True
    )

    if num_files == 1:
        axes = [axes]
    else:
        axes = axes.flatten() if num_files > 1 else [axes]

    cmap = "tab20" if param_column == "bin" else "viridis"

    for idx, (df, wafer_id) in enumerate(zip(multiple_stdf_data, multiple_wafer_ids)):
        ax = axes[idx]

        if param_column not in df.columns:
            ax.set_title(f"Parameter '{param_column}' not found")
            continue

        # Use mask for faster filtering
        mask = df[param_column].notna()
        plot_data = df[mask]

        if len(plot_data) == 0:
            ax.set_title(f"No data for {wafer_id}")
            continue

        # Estimate full wafer extent for proper display
        full_extent, center_info = _estimate_full_wafer_extent(plot_data)

        # Use vectorized grid computation with full wafer extent
        grid, x_min, y_min, x_max, y_max, grid_width, grid_height = _compute_grid_fast(
            plot_data, param_column, full_wafer_extent=full_extent
        )

        # Set background color for empty dies (light gray)
        ax.set_facecolor('#e8e8e8')

        im = ax.imshow(
            grid,
            cmap=cmap,
            aspect="equal",
            interpolation="nearest",
            origin="upper",
        )

        # Wafer circle outline removed per user request
        # _draw_wafer_circle(ax, grid_width, grid_height, x_min, y_min, center_info)

        # Truncate title to fit width, reduce font size by 50%
        short_wafer_id = wafer_id[:25] + "..." if len(str(wafer_id)) > 25 else wafer_id
        short_param = param_label[:30] + "..." if len(str(param_label)) > 30 else param_label

        # Move X-axis to top
        ax.xaxis.set_label_position('top')
        ax.xaxis.tick_top()
        ax.set_xlabel("X", fontsize=10)
        ax.set_ylabel("Y", fontsize=10)

        # Place wafer name at bottom
        ax.set_title("")  # Remove top title
        ax.text(0.5, -0.08, f"{short_wafer_id}\n{short_param}", transform=ax.transAxes,
                ha='center', fontsize=6, fontweight='bold')

        if show_grid_var.get():
            ax.set_xticks(np.arange(-0.5, grid_width, 1), minor=True)
            ax.set_yticks(np.arange(-0.5, grid_height, 1), minor=True)
            ax.grid(which="minor", color="black", linewidth=0.3)
            ax.tick_params(which="minor", size=0)

        # Add die selection highlight with brackets (only for first wafer plot for now)
        if selected_die_coords is not None and idx == 0:
            sel_x, sel_y = selected_die_coords
            x_idx = sel_x - x_min
            y_idx = sel_y - y_min

            # Check if the selected die is within this grid
            if 0 <= x_idx < grid_width and 0 <= y_idx < grid_height:
                from matplotlib.patches import Rectangle

                # Draw a prominent selection rectangle
                rect = Rectangle(
                    (x_idx - 0.5, y_idx - 0.5),
                    1,
                    1,
                    linewidth=2.0,
                    edgecolor="red",
                    facecolor="none",
                    linestyle="-",
                )
                ax.add_patch(rect)

                # Add corner brackets for better visibility
                bracket_size = 0.3
                bracket_color = "black"
                bracket_lw = 1.5

                # Top-left bracket
                ax.plot([x_idx - 0.5, x_idx - 0.5 + bracket_size], [y_idx + 0.5, y_idx + 0.5],
                        color=bracket_color, linewidth=bracket_lw)
                ax.plot([x_idx - 0.5, x_idx - 0.5], [y_idx + 0.5, y_idx + 0.5 - bracket_size],
                        color=bracket_color, linewidth=bracket_lw)

                # Top-right bracket
                ax.plot([x_idx + 0.5, x_idx + 0.5 - bracket_size], [y_idx + 0.5, y_idx + 0.5],
                        color=bracket_color, linewidth=bracket_lw)
                ax.plot([x_idx + 0.5, x_idx + 0.5], [y_idx + 0.5, y_idx + 0.5 - bracket_size],
                        color=bracket_color, linewidth=bracket_lw)

                # Bottom-left bracket
                ax.plot([x_idx - 0.5, x_idx - 0.5 + bracket_size], [y_idx - 0.5, y_idx - 0.5],
                        color=bracket_color, linewidth=bracket_lw)
                ax.plot([x_idx - 0.5, x_idx - 0.5], [y_idx - 0.5, y_idx - 0.5 + bracket_size],
                        color=bracket_color, linewidth=bracket_lw)

                # Bottom-right bracket
                ax.plot([x_idx + 0.5, x_idx + 0.5 - bracket_size], [y_idx - 0.5, y_idx - 0.5],
                        color=bracket_color, linewidth=bracket_lw)
                ax.plot([x_idx + 0.5, x_idx + 0.5], [y_idx - 0.5, y_idx - 0.5 + bracket_size],
                        color=bracket_color, linewidth=bracket_lw)

        # Draw notch marker if orientation is available
        if current_wafer_config and current_wafer_config.get('notch_orientation'):
            draw_notch_marker(ax, grid_width, grid_height, current_wafer_config['notch_orientation'])

        # Create colorbar - double width
        cbar = fig.colorbar(im, ax=ax, fraction=0.024, pad=0.01)

        # Store imshow reference for this axis
        if idx == 0:
            first_im = im
            first_cbar = cbar

    # Add professional interactive range sliders on the colorbar
    if 'first_cbar' in dir() and first_cbar is not None:
        # Get data range from first valid dataset
        for df in multiple_stdf_data:
            if param_column in df.columns:
                mask = df[param_column].notna()
                values = df.loc[mask, param_column].values
                if len(values) > 0:
                    data_min = np.nanmin(values)
                    data_max = np.nanmax(values)
                    data_range = data_max - data_min

                    # Get colorbar axes limits for proper coordinate mapping
                    cbar_ylim = first_cbar.ax.get_ylim()
                    cbar_ymin, cbar_ymax = min(cbar_ylim), max(cbar_ylim)

                    # Professional slider design using patches
                    from matplotlib.patches import FancyBboxPatch, Rectangle

                    # Create triangular arrow markers pointing inward
                    # Low limit slider (bottom) - triangle pointing up
                    low_marker = first_cbar.ax.scatter(
                        [0.5], [data_min],
                        marker='^',  # Triangle pointing up
                        s=200,
                        c='#D32F2F',  # Material Red
                        edgecolors='white',
                        linewidths=1.5,
                        zorder=15,
                        clip_on=False
                    )

                    # High limit slider (top) - triangle pointing down
                    high_marker = first_cbar.ax.scatter(
                        [0.5], [data_max],
                        marker='v',  # Triangle pointing down
                        s=200,
                        c='#1976D2',  # Material Blue
                        edgecolors='white',
                        linewidths=1.5,
                        zorder=15,
                        clip_on=False
                    )

                    # Add horizontal limit lines extending beyond colorbar
                    low_line, = first_cbar.ax.plot(
                        [-0.6, 1.6], [data_min, data_min],
                        color='#D32F2F', linewidth=2.5,
                        solid_capstyle='round', zorder=12,
                        transform=first_cbar.ax.get_yaxis_transform()
                    )

                    high_line, = first_cbar.ax.plot(
                        [-0.6, 1.6], [data_max, data_max],
                        color='#1976D2', linewidth=2.5,
                        solid_capstyle='round', zorder=12,
                        transform=first_cbar.ax.get_yaxis_transform()
                    )

                    # Add value labels next to sliders
                    low_text = first_cbar.ax.text(
                        1.8, data_min, f'{data_min:.2f}',
                        fontsize=7, fontweight='bold',
                        color='#D32F2F', va='center', ha='left',
                        transform=first_cbar.ax.get_yaxis_transform(),
                        zorder=20
                    )

                    high_text = first_cbar.ax.text(
                        1.8, data_max, f'{data_max:.2f}',
                        fontsize=7, fontweight='bold',
                        color='#1976D2', va='center', ha='left',
                        transform=first_cbar.ax.get_yaxis_transform(),
                        zorder=20
                    )

                    # Store references and state
                    slider_state = {
                        'low_val': data_min,
                        'high_val': data_max,
                        'data_min': data_min,
                        'data_max': data_max,
                        'cbar_ymin': cbar_ymin,
                        'cbar_ymax': cbar_ymax,
                        'dragging': None,
                        'low_marker': low_marker,
                        'high_marker': high_marker,
                        'low_line': low_line,
                        'high_line': high_line,
                        'low_text': low_text,
                        'high_text': high_text,
                        'cbar_ax': first_cbar.ax
                    }

                    # Get all imshow objects
                    imshow_list = []
                    for ax_item in fig.axes:
                        for child in ax_item.get_children():
                            if hasattr(child, 'set_clim'):
                                imshow_list.append(child)

                    def on_press(event):
                        if event.inaxes == slider_state['cbar_ax']:
                            y_click = event.ydata
                            if y_click is not None:
                                # Calculate distance to each slider in data coordinates
                                dist_to_low = abs(y_click - slider_state['low_val'])
                                dist_to_high = abs(y_click - slider_state['high_val'])

                                # Threshold for grabbing (10% of data range)
                                threshold = 0.1 * (slider_state['data_max'] - slider_state['data_min'])

                                # Select the closest slider if within threshold
                                if dist_to_low < dist_to_high and dist_to_low < threshold:
                                    slider_state['dragging'] = 'low'
                                    low_marker.set_sizes([300])  # Visual feedback - enlarge
                                elif dist_to_high < threshold:
                                    slider_state['dragging'] = 'high'
                                    high_marker.set_sizes([300])  # Visual feedback - enlarge

                                fig.canvas.draw_idle()

                    def on_motion(event):
                        if slider_state['dragging']:
                            # Allow motion even outside colorbar axes for smooth dragging
                            if event.ydata is not None:
                                y_click = event.ydata
                            elif event.y is not None:
                                # Convert pixel to data coordinates if outside axes
                                inv = slider_state['cbar_ax'].transData.inverted()
                                _, y_click = inv.transform((event.x, event.y))
                            else:
                                return

                            # Clamp to data range
                            y_click = max(slider_state['data_min'], min(slider_state['data_max'], y_click))

                            min_gap = 0.02 * (slider_state['data_max'] - slider_state['data_min'])

                            if slider_state['dragging'] == 'low':
                                # Low slider can't go above high slider
                                new_val = min(y_click, slider_state['high_val'] - min_gap)
                                new_val = max(new_val, slider_state['data_min'])
                                slider_state['low_val'] = new_val

                                # Update marker position
                                low_marker.set_offsets([[0.5, new_val]])
                                # Update line position
                                low_line.set_ydata([new_val, new_val])
                                # Update text
                                low_text.set_position((1.8, new_val))
                                low_text.set_text(f'{new_val:.2f}')

                            elif slider_state['dragging'] == 'high':
                                # High slider can't go below low slider
                                new_val = max(y_click, slider_state['low_val'] + min_gap)
                                new_val = min(new_val, slider_state['data_max'])
                                slider_state['high_val'] = new_val

                                # Update marker position
                                high_marker.set_offsets([[0.5, new_val]])
                                # Update line position
                                high_line.set_ydata([new_val, new_val])
                                # Update text
                                high_text.set_position((1.8, new_val))
                                high_text.set_text(f'{new_val:.2f}')

                            # Update colormap limits for all heatmaps
                            for img in imshow_list:
                                img.set_clim(vmin=slider_state['low_val'], vmax=slider_state['high_val'])

                            fig.canvas.draw_idle()

                    def on_release(event):
                        if slider_state['dragging'] == 'low':
                            low_marker.set_sizes([200])  # Reset size
                        elif slider_state['dragging'] == 'high':
                            high_marker.set_sizes([200])  # Reset size
                        slider_state['dragging'] = None
                        fig.canvas.draw_idle()

                    # Connect events
                    fig.canvas.mpl_connect('button_press_event', on_press)
                    fig.canvas.mpl_connect('motion_notify_event', on_motion)
                    fig.canvas.mpl_connect('button_release_event', on_release)

                    # Store reference
                    global _colorbar_range_slider
                    _colorbar_range_slider = slider_state

                    break

    for idx in range(num_files, len(axes)):
        fig.delaxes(axes[idx])

    if heatmap_canvas:
        heatmap_canvas.get_tk_widget().destroy()

    heatmap_canvas = FigureCanvasTkAgg(fig, master=heatmap_display_frame)
    heatmap_canvas.get_tk_widget().pack(fill=tk.BOTH, expand=True, pady=10)

    # Connect mouse events for rectangle zoom, pan, scroll, and click
    heatmap_canvas.mpl_connect("button_press_event", on_heatmap_press)
    heatmap_canvas.mpl_connect("motion_notify_event", on_heatmap_motion)
    heatmap_canvas.mpl_connect("button_release_event", on_heatmap_release)
    heatmap_canvas.mpl_connect("scroll_event", on_heatmap_scroll)

    heatmap_canvas.draw()

    # Update limits display and show overlay in top-left corner
    update_limits_display()
    limits_overlay_frame.place(x=15, y=15)
    limits_overlay_frame.lift()  # Raise above canvas

    print(f"Multi-plot heatmap updated with {num_files} wafers")


def update_heatmap_parameter_list():
    """Update the heatmap parameter combobox when STDF data is loaded"""
    global current_stdf_data, test_parameters

    if current_stdf_data is None or current_stdf_data.empty:
        heatmap_param_combobox["values"] = []
        heatmap_info_label.config(
            text="No STDF data loaded. Load file in Wafermap tab first."
        )
        return

    param_options = ["BIN (Bin Number)"]
    for test_key, test_name in sorted(test_parameters.items()):
        param_options.append(f"{test_key}: {test_name}")

    heatmap_param_combobox["values"] = param_options
    if param_options:
        heatmap_param_combobox.current(0)

    heatmap_info_label.config(text=f"Loaded: {len(param_options)} parameters available")

    update_stdf_heatmap()


def _compute_grid_fast(plot_data, param_column, full_wafer_extent=None):
    """
    Compute heatmap grid using vectorized numpy operations.
    This is 10-100x faster than iterrows() for large datasets.

    If full_wafer_extent is provided as (x_min, y_min, x_max, y_max),
    it will use those bounds instead of data-derived bounds for full wafer display.
    """
    x_vals = plot_data["x"].values
    y_vals = plot_data["y"].values
    param_vals = plot_data[param_column].values

    data_x_min, data_x_max = x_vals.min(), x_vals.max()
    data_y_min, data_y_max = y_vals.min(), y_vals.max()

    # Use full wafer extent if provided, otherwise use data extent
    if full_wafer_extent is not None:
        x_min, y_min, x_max, y_max = full_wafer_extent
    else:
        x_min, x_max = data_x_min, data_x_max
        y_min, y_max = data_y_min, data_y_max

    grid_width = int(x_max - x_min + 1)
    grid_height = int(y_max - y_min + 1)

    # Create grid filled with NaN
    grid = np.full((grid_height, grid_width), np.nan)

    # Vectorized index computation
    x_indices = (x_vals - x_min).astype(int)
    y_indices = (y_vals - y_min).astype(int)

    # Use advanced indexing for fast assignment
    grid[y_indices, x_indices] = param_vals

    return grid, x_min, y_min, x_max, y_max, grid_width, grid_height


def draw_notch_marker(ax, grid_width, grid_height, notch_orientation, marker_size=None):
    """
    Draw a notch marker on the wafermap to indicate wafer orientation.

    Parameters:
    - ax: matplotlib axes
    - grid_width: width of the grid
    - grid_height: height of the grid
    - notch_orientation: 'U' (Up/0°), 'D' (Down/180°), 'L' (Left/270°), 'R' (Right/90°)
    - marker_size: size of the notch marker (default: auto-calculated)
    """
    if notch_orientation is None or notch_orientation == ' ':
        return  # No orientation info available

    # Calculate marker size based on grid dimensions (halved again)
    if marker_size is None:
        marker_size = min(grid_width, grid_height) * 0.02

    # Calculate center of the grid
    center_x = grid_width / 2 - 0.5
    center_y = grid_height / 2 - 0.5

    # Determine notch position based on orientation
    # Using a triangle marker pointing inward
    margin = 1.5  # Distance from edge

    if notch_orientation == 'U':  # Up (0°) - Notch at top
        notch_x = center_x
        notch_y = grid_height - margin
        # Triangle pointing down (into wafer)
        triangle = plt.Polygon([
            [notch_x, notch_y],
            [notch_x - marker_size, notch_y + marker_size],
            [notch_x + marker_size, notch_y + marker_size]
        ], closed=True, facecolor='red', edgecolor='darkred', linewidth=1.5, zorder=100)
        text_y = notch_y + marker_size + 1

    elif notch_orientation == 'D':  # Down (180°) - Notch at bottom
        notch_x = center_x
        notch_y = margin - 1
        # Triangle pointing up (into wafer)
        triangle = plt.Polygon([
            [notch_x, notch_y],
            [notch_x - marker_size, notch_y - marker_size],
            [notch_x + marker_size, notch_y - marker_size]
        ], closed=True, facecolor='red', edgecolor='darkred', linewidth=1.5, zorder=100)
        text_y = notch_y - marker_size - 1.5

    elif notch_orientation == 'L':  # Left (270°) - Notch at left
        notch_x = margin - 1
        notch_y = center_y
        # Triangle pointing right (into wafer)
        triangle = plt.Polygon([
            [notch_x, notch_y],
            [notch_x - marker_size, notch_y - marker_size],
            [notch_x - marker_size, notch_y + marker_size]
        ], closed=True, facecolor='red', edgecolor='darkred', linewidth=1.5, zorder=100)
        text_y = notch_y

    elif notch_orientation == 'R':  # Right (90°) - Notch at right
        notch_x = grid_width - margin
        notch_y = center_y
        # Triangle pointing left (into wafer)
        triangle = plt.Polygon([
            [notch_x, notch_y],
            [notch_x + marker_size, notch_y - marker_size],
            [notch_x + marker_size, notch_y + marker_size]
        ], closed=True, facecolor='red', edgecolor='darkred', linewidth=1.5, zorder=100)
        text_y = notch_y
    else:
        return  # Unknown orientation

    ax.add_patch(triangle)

    # Add "NOTCH" label
    if notch_orientation in ['U', 'D']:
        ax.text(notch_x, text_y, 'NOTCH', ha='center', va='center',
                fontsize=7, fontweight='bold', color='darkred', zorder=101)
    else:
        # For left/right, position text differently
        if notch_orientation == 'L':
            ax.text(notch_x - marker_size - 1, text_y, 'N', ha='center', va='center',
                    fontsize=7, fontweight='bold', color='darkred', zorder=101, rotation=90)
        else:
            ax.text(notch_x + marker_size + 1, text_y, 'N', ha='center', va='center',
                    fontsize=7, fontweight='bold', color='darkred', zorder=101, rotation=90)


def _estimate_full_wafer_extent(plot_data):
    """
    Estimate full wafer extent based on measured die positions.
    Assumes a circular wafer and estimates the full extent from partial data.
    """
    x_vals = plot_data["x"].values
    y_vals = plot_data["y"].values

    data_x_min, data_x_max = x_vals.min(), x_vals.max()
    data_y_min, data_y_max = y_vals.min(), y_vals.max()

    # Calculate the center and radius from measured data
    x_center = (data_x_min + data_x_max) / 2
    y_center = (data_y_min + data_y_max) / 2

    # Estimate radius based on the furthest point from center
    dx = np.abs(x_vals - x_center)
    dy = np.abs(y_vals - y_center)
    measured_radius = np.sqrt(dx**2 + dy**2).max()

    # For partial wafers, estimate full wafer dimensions
    # Typical wafer might be larger than what was measured
    # We'll use the maximum extent from center as the radius
    data_width = data_x_max - data_x_min
    data_height = data_y_max - data_y_min

    # Estimate full wafer radius - add some margin for unmeasured areas
    # Use the maximum of width/height as base estimate
    estimated_radius = max(data_width, data_height) / 2 * 1.05

    # If measured radius is significantly smaller than data extent radius,
    # the wafer center might not be in the data - use data extent + margin
    if measured_radius < estimated_radius * 0.5:
        # Data might be off-center, extend from data bounds
        margin = max(3, int(max(data_width, data_height) * 0.15))
        x_min = data_x_min - margin
        y_min = data_y_min - margin
        x_max = data_x_max + margin
        y_max = data_y_max + margin
    else:
        # Extend from estimated center
        x_min = int(x_center - estimated_radius)
        y_min = int(y_center - estimated_radius)
        x_max = int(x_center + estimated_radius)
        y_max = int(y_center + estimated_radius)

    return (x_min, y_min, x_max, y_max), (x_center, y_center, estimated_radius)


def _draw_wafer_circle(ax, grid_width, grid_height, x_min, y_min, center_info=None):
    """
    Draw wafer outline circle on the axes.

    Args:
        ax: matplotlib axes
        grid_width, grid_height: grid dimensions
        x_min, y_min: grid offset
        center_info: tuple of (x_center, y_center, radius) in die coordinates
    """
    from matplotlib.patches import Circle, Wedge

    # Calculate center in grid coordinates
    if center_info is not None:
        x_center, y_center, radius = center_info
        # Convert to grid coordinates
        grid_x_center = x_center - x_min
        grid_y_center = y_center - y_min
        grid_radius = radius
    else:
        # Use grid center and estimate radius
        grid_x_center = grid_width / 2
        grid_y_center = grid_height / 2
        grid_radius = min(grid_width, grid_height) / 2 * 0.98

    # Draw wafer outline circle
    wafer_circle = Circle(
        (grid_x_center, grid_y_center),
        grid_radius,
        fill=False,
        edgecolor='#333333',
        linewidth=2.5,
        linestyle='-',
        zorder=10
    )
    ax.add_patch(wafer_circle)

    # Draw wafer flat/notch at bottom (standard orientation)
    # Flat is typically at 6 o'clock position, about 10-15% of diameter
    flat_angle_deg = 15  # degrees from bottom center

    # For a more realistic look, draw a small notch instead of a flat for 200mm+ wafers
    notch_size = grid_radius * 0.05
    notch_x = grid_x_center
    notch_y = grid_y_center - grid_radius

    # Draw a small triangular notch
    notch_points = [
        (notch_x, notch_y),
        (notch_x - notch_size, notch_y + notch_size * 2),
        (notch_x + notch_size, notch_y + notch_size * 2),
    ]
    from matplotlib.patches import Polygon
    notch = Polygon(notch_points, fill=True, facecolor='white',
                   edgecolor='#333333', linewidth=2, zorder=11)
    ax.add_patch(notch)

    return (grid_x_center, grid_y_center, grid_radius)


def _get_cached_grid(data_id, param_column, plot_data):
    """
    Get grid from cache or compute and cache it.
    Cache key is based on data identity and parameter.
    """
    global _grid_cache

    cache_key = (data_id, param_column)

    if cache_key in _grid_cache:
        return _grid_cache[cache_key]

    # Compute grid
    result = _compute_grid_fast(plot_data, param_column)

    # Cache with size limit (keep last 10 grids)
    if len(_grid_cache) > 10:
        # Remove oldest entry
        oldest_key = next(iter(_grid_cache))
        del _grid_cache[oldest_key]

    _grid_cache[cache_key] = result
    return result


def update_stdf_heatmap():
    """Update STDF heatmap display based on selected parameter - OPTIMIZED VERSION"""
    global heatmap_canvas, current_stdf_data, current_wafer_id, _current_fig, _current_ax

    if current_stdf_data is None or current_stdf_data.empty:
        print("No STDF data loaded for heatmap")
        return

    selected = heatmap_param_combobox.get()

    if not selected:
        print("No parameter selected")
        return

    # Check if this is a custom test
    is_custom_test = selected.startswith("CUSTOM:")
    custom_test_name = None

    if selected.startswith("BIN"):
        param_column = "bin"
        param_label = "Bin"
    elif is_custom_test:
        # Handle custom test - we'll compute values for each die
        custom_test_name = selected.replace("CUSTOM:", "").strip()
        param_column = f"_custom_{custom_test_name}"
        param_label = f"Custom: {custom_test_name}"

        # Debug: print the custom test formula
        if custom_test_name in custom_tests:
            print(f"Custom test '{custom_test_name}' formula: {custom_tests[custom_test_name]}")
        else:
            print(f"ERROR: Custom test '{custom_test_name}' not found in custom_tests dict!")

        # Debug: print column names in the data
        print(f"Available columns in data: {list(current_stdf_data.columns)[:20]}...")

        # Compute custom test values for all dies
        custom_values = []
        none_count = 0
        for idx, row in current_stdf_data.iterrows():
            value = evaluate_custom_test(custom_test_name, row)
            custom_values.append(value)
            if value is None:
                none_count += 1

        # Add custom values as a temporary column
        current_stdf_data[param_column] = custom_values
        valid_count = len(custom_values) - none_count
        print(f"Computed custom test '{custom_test_name}' for {len(current_stdf_data)} dies: {valid_count} valid, {none_count} None values")

        # Debug: show sample calculated values
        sample_values = [v for v in custom_values[:5] if v is not None]
        print(f"Sample calculated values: {sample_values}")
    elif selected.startswith("───"):
        # This is a separator, ignore it
        print("Separator selected - please select a valid parameter")
        return
    else:
        test_key = selected.split(":")[0].strip()
        if test_key.startswith("test_"):
            param_column = int(test_key.replace("test_", ""))
        else:
            param_column = int(test_key)
        param_label = selected

    if param_column not in current_stdf_data.columns:
        print(f"ERROR: Parameter '{param_column}' not found in data")
        return

    # Use mask instead of dropna for better performance
    mask = current_stdf_data[param_column].notna()
    plot_data = current_stdf_data[mask]

    if len(plot_data) == 0:
        fig, ax = plt.subplots(figsize=(12, 10))
        ax.set_title(f"No valid data for parameter '{param_label}'")
        fig.tight_layout()

        if heatmap_canvas:
            heatmap_canvas.get_tk_widget().destroy()

        heatmap_canvas = FigureCanvasTkAgg(fig, master=heatmap_display_frame)
        heatmap_canvas.get_tk_widget().pack(fill=tk.BOTH, expand=True, pady=10)
        heatmap_canvas.draw()
        return

    # Estimate full wafer extent for proper display
    full_extent, center_info = _estimate_full_wafer_extent(plot_data)

    # Use grid computation with full wafer extent
    grid, x_min, y_min, x_max, y_max, grid_width, grid_height = _compute_grid_fast(
        plot_data, param_column, full_wafer_extent=full_extent
    )

    # Create figure with optimized settings
    fig, ax = plt.subplots(figsize=(12, 10))

    # Set background color for empty dies (light gray)
    ax.set_facecolor('#e8e8e8')

    # Use appropriate colormap
    cmap = "tab20" if param_column == "bin" else "viridis"

    # imshow is already fast, no changes needed here
    im = ax.imshow(
        grid,
        cmap=cmap,
        aspect="equal",
        interpolation="nearest",
        origin="upper",
    )

    # Wafer circle outline removed per user request
    # _draw_wafer_circle(ax, grid_width, grid_height, x_min, y_min, center_info)

    cbar = fig.colorbar(im, ax=ax, label=param_label)
    cbar.ax.tick_params(labelsize=10)

    # Move X-axis to top
    ax.xaxis.set_label_position('top')
    ax.xaxis.tick_top()
    ax.set_xlabel("X Coordinate", fontsize=12)
    ax.set_ylabel("Y Coordinate", fontsize=12)

    # Truncate title to fit width, reduce font size by 50%
    short_param = param_label[:40] + "..." if len(str(param_label)) > 40 else param_label
    short_wafer_id = current_wafer_id[:30] + "..." if len(str(current_wafer_id)) > 30 else current_wafer_id

    # Place title at the bottom using figtext or xlabel area
    ax.set_title("")  # Remove top title
    fig.text(0.5, 0.02, f"Heatmap: {short_param} | Wafer: {short_wafer_id}",
             ha='center', fontsize=7, fontweight='bold')

    # Optimized tick computation
    num_x_ticks = min(10, grid_width)
    num_y_ticks = min(10, grid_height)
    x_tick_positions = np.linspace(0, grid_width - 1, num_x_ticks)
    y_tick_positions = np.linspace(0, grid_height - 1, num_y_ticks)

    ax.set_xticks(x_tick_positions)
    ax.set_yticks(y_tick_positions)
    ax.set_xticklabels([f"{int(x_min + pos)}" for pos in x_tick_positions])
    # With origin="upper", position 0 is at top, so y_min should be at top
    ax.set_yticklabels([f"{int(y_min + pos)}" for pos in y_tick_positions])

    # Grid lines - optimized: only draw if enabled
    if show_grid_var.get():
        # Use LineCollection for faster grid line rendering
        ax.set_xticks(np.arange(-0.5, grid_width, 1), minor=True)
        ax.set_yticks(np.arange(-0.5, grid_height, 1), minor=True)
        ax.grid(which="minor", color="black", linewidth=0.5)
        ax.tick_params(which="minor", size=0)
    else:
        ax.grid(True, alpha=0.3, linestyle="--", linewidth=0.5)

    # Draw notch marker if orientation is available
    if current_wafer_config and current_wafer_config.get('notch_orientation'):
        draw_notch_marker(ax, grid_width, grid_height, current_wafer_config['notch_orientation'])

    # Selected die highlight - make it more visible
    if selected_die_coords is not None:
        sel_x, sel_y = selected_die_coords
        x_idx = sel_x - x_min
        y_idx = sel_y - y_min

        from matplotlib.patches import Rectangle

        # Draw a prominent selection rectangle
        rect = Rectangle(
            (x_idx - 0.5, y_idx - 0.5),
            1,
            1,
            linewidth=2.0,
            edgecolor="red",
            facecolor="none",
            linestyle="-",
        )
        ax.add_patch(rect)

        # Add corner brackets for better visibility
        bracket_size = 0.3
        bracket_color = "black"
        bracket_lw = 1.5

        # Top-left bracket
        ax.plot([x_idx - 0.5, x_idx - 0.5 + bracket_size], [y_idx + 0.5, y_idx + 0.5],
                color=bracket_color, linewidth=bracket_lw)
        ax.plot([x_idx - 0.5, x_idx - 0.5], [y_idx + 0.5, y_idx + 0.5 - bracket_size],
                color=bracket_color, linewidth=bracket_lw)

        # Top-right bracket
        ax.plot([x_idx + 0.5, x_idx + 0.5 - bracket_size], [y_idx + 0.5, y_idx + 0.5],
                color=bracket_color, linewidth=bracket_lw)
        ax.plot([x_idx + 0.5, x_idx + 0.5], [y_idx + 0.5, y_idx + 0.5 - bracket_size],
                color=bracket_color, linewidth=bracket_lw)

        # Bottom-left bracket
        ax.plot([x_idx - 0.5, x_idx - 0.5 + bracket_size], [y_idx - 0.5, y_idx - 0.5],
                color=bracket_color, linewidth=bracket_lw)
        ax.plot([x_idx - 0.5, x_idx - 0.5], [y_idx - 0.5, y_idx - 0.5 + bracket_size],
                color=bracket_color, linewidth=bracket_lw)

        # Bottom-right bracket
        ax.plot([x_idx + 0.5, x_idx + 0.5 - bracket_size], [y_idx - 0.5, y_idx - 0.5],
                color=bracket_color, linewidth=bracket_lw)
        ax.plot([x_idx + 0.5, x_idx + 0.5], [y_idx - 0.5, y_idx - 0.5 + bracket_size],
                color=bracket_color, linewidth=bracket_lw)

    # Pre-compute statistics using numpy for speed
    param_vals = plot_data[param_column].values
    stats_text = (
        f"Min: {np.nanmin(param_vals):.2f}\n"
        f"Max: {np.nanmax(param_vals):.2f}\n"
        f"Mean: {np.nanmean(param_vals):.2f}\n"
        f"Median: {np.nanmedian(param_vals):.2f}\n"
        f"Dies: {len(plot_data)}"
    )

    ax.text(
        0.02,
        0.98,
        stats_text,
        transform=ax.transAxes,
        fontsize=10,
        verticalalignment="top",
        bbox=dict(boxstyle="round", facecolor="wheat", alpha=0.8),
    )

    fig.tight_layout()

    # Destroy old canvas before creating new one
    if heatmap_canvas:
        heatmap_canvas.get_tk_widget().destroy()

    heatmap_canvas = FigureCanvasTkAgg(fig, master=heatmap_display_frame)
    canvas_widget = heatmap_canvas.get_tk_widget()
    canvas_widget.pack(fill=tk.BOTH, expand=True, pady=10)

    # Connect mouse events for rectangle zoom, pan, scroll, and click
    heatmap_canvas.mpl_connect("button_press_event", on_heatmap_press)
    heatmap_canvas.mpl_connect("motion_notify_event", on_heatmap_motion)
    heatmap_canvas.mpl_connect("button_release_event", on_heatmap_release)
    heatmap_canvas.mpl_connect("scroll_event", on_heatmap_scroll)

    # Also bind tkinter native scroll events for Windows compatibility
    canvas_widget.bind("<MouseWheel>", on_tk_mousewheel)
    canvas_widget.bind("<Button-4>", lambda e: on_tk_mousewheel_linux(e, True))  # Linux scroll up
    canvas_widget.bind("<Button-5>", lambda e: on_tk_mousewheel_linux(e, False))  # Linux scroll down

    heatmap_canvas.draw()

    # Update limits display and show overlay in top-left corner
    update_limits_display()
    limits_overlay_frame.place(x=15, y=15)
    limits_overlay_frame.lift()  # Raise above canvas


def refresh_heatmap_display():
    """Refresh heatmap display - calls appropriate function based on view type selection"""
    global selected_die_coords

    # Get the current view type from dropdown
    try:
        view_type = view_type_var.get()
        view_subtype = view_subtype_var.get()
    except:
        view_type = "Data"
        view_subtype = "Heatmap"

    print(f"Refresh display - View: {view_type}, Subtype: {view_subtype}")

    # Check view type and call appropriate function
    if view_type == "PLM Files":
        update_plm_wafermap()
        return

    if view_type == "Images":
        # Update the image type filter variable for compatibility
        image_type_view_var.set(view_subtype)
        update_image_wafermap()
        return

    # Default: Data view (heatmap)
    if multiple_stdf_data:
        update_multi_stdf_heatmap()
    elif current_stdf_data is not None and not current_stdf_data.empty:
        update_stdf_heatmap()
    else:
        print("No STDF data loaded")

    # Update statistics plots (boxplot and probability distribution)
    update_stats_plots()

    # Update limits display
    update_limits_display()

    # Show limits overlay in top-left corner of heatmap
    limits_overlay_frame.place(x=15, y=15)
    limits_overlay_frame.lift()  # Raise above canvas


def update_image_wafermap():
    """Display wafermap with die images at their coordinates - images scale with zoom"""
    global heatmap_canvas, multiple_stdf_data, multiple_wafer_ids, die_image_directory

    print(f"update_image_wafermap called")
    print(f"  die_image_directory: {die_image_directory}")
    print(f"  multiple_stdf_data: {len(multiple_stdf_data) if multiple_stdf_data else 'None'}")

    # Close any existing figure to prevent memory warnings
    plt.close('all')

    if not multiple_stdf_data and (current_stdf_data is None or current_stdf_data.empty):
        print("No wafermap data loaded - showing message")
        # Show message on the heatmap
        fig, ax = plt.subplots(figsize=(10, 10))
        ax.text(0.5, 0.5, "No wafermap data loaded.\n\nPlease load STDF or CSV data first.",
                ha='center', va='center', fontsize=14, transform=ax.transAxes)
        ax.set_title("Image Wafermap - No Data", fontsize=14)
        ax.axis('off')

        if heatmap_canvas:
            heatmap_canvas.get_tk_widget().destroy()

        heatmap_canvas = FigureCanvasTkAgg(fig, master=heatmap_display_frame)
        heatmap_canvas.get_tk_widget().pack(fill=tk.BOTH, expand=True)
        heatmap_canvas.draw()
        return

    if not die_image_directory:
        print("No image folder selected. Please select an image folder first.")
        # Show message on the heatmap
        fig, ax = plt.subplots(figsize=(10, 10))
        ax.text(0.5, 0.5, "No image folder selected.\n\nPlease use 'Img Folder' button\nto select a folder with die images.",
                ha='center', va='center', fontsize=14, transform=ax.transAxes)
        ax.set_title("Image Wafermap - No Images", fontsize=14)
        ax.axis('off')

        if heatmap_canvas:
            heatmap_canvas.get_tk_widget().destroy()

        heatmap_canvas = FigureCanvasTkAgg(fig, master=heatmap_display_frame)
        heatmap_canvas.get_tk_widget().pack(fill=tk.BOTH, expand=True)
        heatmap_canvas.draw()
        return

    # Get wafermap data
    if multiple_stdf_data:
        df = multiple_stdf_data[0]
        wafer_id = multiple_wafer_ids[0] if multiple_wafer_ids else "Wafer"
    else:
        df = current_stdf_data
        wafer_id = current_wafer_id if current_wafer_id else "Wafer"

    print(f"  Using dataframe with {len(df)} rows")

    # Estimate full wafer extent for proper display
    full_extent, center_info = _estimate_full_wafer_extent(df)
    x_min, y_min, x_max, y_max = full_extent

    grid_width = x_max - x_min + 1
    grid_height = y_max - y_min + 1

    print(f"Creating image wafermap: {grid_width}x{grid_height} grid")
    print(f"X range: {x_min} to {x_max}, Y range: {y_min} to {y_max}")

    # Show loading indicator
    print("Loading and resizing images for optimal performance...")

    # Create figure
    fig_size = max(12, min(20, grid_width * 0.5))
    fig, ax = plt.subplots(figsize=(fig_size, fig_size))

    # Set background color for the entire wafer area
    ax.set_facecolor('#e8e8e8')

    from matplotlib.patches import Rectangle
    import re

    # Draw light gray background for all die positions first (vectorized for speed)
    die_positions = list(zip(df['x'].values, df['y'].values))
    for x, y in die_positions:
        x_idx = x - x_min
        y_idx = y - y_min
        rect = Rectangle((x_idx - 0.5, y_idx - 0.5), 1, 1,
                         facecolor='#f0f0f0', edgecolor='#cccccc', linewidth=0.3)
        ax.add_patch(rect)

    # Wafer circle outline removed per user request
    # _draw_wafer_circle(ax, grid_width, grid_height, x_min, y_min, center_info)

    images_found = 0
    images_not_found = 0

    # Get all image files in directory and map to coordinates
    image_extensions = ('.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tif', '.tiff')
    all_images = {}

    # Get the selected image type filter from the view subtype dropdown
    try:
        selected_image_type = view_subtype_var.get()
    except:
        selected_image_type = "All"
    print(f"Image type filter: {selected_image_type}")

    try:
        for filename in os.listdir(die_image_directory):
            if filename.lower().endswith(image_extensions):
                # Filter by image type if not "All"
                if selected_image_type != "All":
                    # Check if filename contains the selected type
                    # Common patterns: ALLON, XHAIR, CHKBRD, ALLOFF
                    filename_upper = filename.upper()
                    if selected_image_type.upper() not in filename_upper:
                        continue

                # Extract coordinates from filename
                # Patterns: X19_Y46, X19-Y46, _X19_Y46_
                match = re.search(r'[_\-]?X(\d+)[_\-]Y(\d+)', filename, re.IGNORECASE)
                if match:
                    img_x = int(match.group(1))
                    img_y = int(match.group(2))
                    coord_key = (img_x, img_y)
                    if coord_key not in all_images:
                        all_images[coord_key] = []
                    all_images[coord_key].append(os.path.join(die_image_directory, filename))
    except Exception as e:
        print(f"Error scanning image directory: {e}")

    print(f"Found images for {len(all_images)} die positions")

    # Target resolution for resized images (good balance of quality vs performance)
    # 256x256 gives good detail when zoomed while keeping memory low
    TARGET_SIZE = 256

    # Image cache to store resized images
    resized_image_cache = {}

    # Pre-load and resize all images first for better performance
    total_to_load = len([k for k in all_images.keys() if k in [(int(row['x']), int(row['y'])) for _, row in df.iterrows()]])
    loaded_count = 0

    # Place images at die coordinates using imshow with extent (scales with zoom!)
    for _, row in df.iterrows():
        x = int(row['x'])
        y = int(row['y'])
        x_idx = x - x_min
        y_idx = y - y_min

        coord_key = (x, y)

        if coord_key in all_images and all_images[coord_key]:
            # Use first image for this die
            img_path = all_images[coord_key][0]

            try:
                # Check cache first
                if img_path in resized_image_cache:
                    img_array = resized_image_cache[img_path]
                else:
                    # Load and resize image for optimal performance
                    img = Image.open(img_path)

                    # Convert to RGB if necessary (handles RGBA, grayscale, etc.)
                    if img.mode != 'RGB':
                        img = img.convert('RGB')

                    # Resize to target size using fast resampling
                    img = img.resize((TARGET_SIZE, TARGET_SIZE), Image.Resampling.BILINEAR)

                    # Convert to numpy array
                    img_array = np.array(img, dtype=np.uint8)

                    # Cache the resized image
                    resized_image_cache[img_path] = img_array

                # Calculate extent for this die cell
                # extent = [left, right, bottom, top]
                extent = [x_idx - 0.48, x_idx + 0.48, y_idx - 0.48, y_idx + 0.48]

                # Plot image using imshow - it will scale automatically with zoom!
                # Use 'nearest' interpolation for faster rendering
                ax.imshow(img_array, extent=extent, aspect='auto',
                         interpolation='nearest', zorder=2)

                images_found += 1
                loaded_count += 1

                # Progress update every 50 images
                if loaded_count % 50 == 0:
                    print(f"  Loaded {loaded_count} images...")

            except Exception as e:
                print(f"Error loading image {img_path}: {e}")
                images_not_found += 1
        else:
            images_not_found += 1

    # Set axis properties
    ax.set_xlim(-0.5, grid_width - 0.5)
    ax.set_ylim(-0.5, grid_height - 0.5)
    ax.set_aspect('equal')
    ax.set_xlabel('X Coordinate', fontsize=10)
    ax.set_ylabel('Y Coordinate', fontsize=10)
    # Truncate title to fit width, reduce font size by 50%
    short_wafer_id = wafer_id[:20] + "..." if len(str(wafer_id)) > 20 else wafer_id
    ax.set_title(f"Image: {short_wafer_id} [{selected_image_type}]\n{images_found} imgs, {images_not_found} empty",
                fontsize=6, fontweight='bold')

    # Add coordinate labels
    tick_step = max(1, grid_width // 10)
    ax.set_xticks(range(0, grid_width, tick_step))
    ax.set_xticklabels([str(x_min + i) for i in range(0, grid_width, tick_step)])
    ax.set_yticks(range(0, grid_height, tick_step))
    ax.set_yticklabels([str(y_min + i) for i in range(0, grid_height, tick_step)])

    fig.tight_layout()

    # Display on canvas
    if heatmap_canvas:
        heatmap_canvas.get_tk_widget().destroy()

    heatmap_canvas = FigureCanvasTkAgg(fig, master=heatmap_display_frame)
    canvas_widget = heatmap_canvas.get_tk_widget()
    canvas_widget.pack(fill=tk.BOTH, expand=True)

    # Add navigation toolbar for better zoom/pan control
    from matplotlib.backends.backend_tkagg import NavigationToolbar2Tk

    # Connect scroll events for zooming
    heatmap_canvas.mpl_connect("scroll_event", on_heatmap_scroll)
    canvas_widget.bind("<MouseWheel>", on_tk_mousewheel)

    # Connect click event for die selection
    heatmap_canvas.mpl_connect("button_press_event", lambda e: on_image_wafermap_click(e, x_min, y_min, df, all_images))

    heatmap_canvas.draw()

    print(f"Image wafermap displayed: {images_found} images, {images_not_found} empty positions")
    print("Zoom in to see image details!")


def on_image_wafermap_click(event, x_min, y_min, df, all_images):
    """Handle click on image wafermap to show full-size image with die highlighting"""
    global current_selected_die, heatmap_canvas, _image_highlight_rect

    if event.inaxes is None or event.button != 1:
        return

    # Get clicked grid position
    x_idx = int(round(event.xdata))
    y_idx = int(round(event.ydata))

    # Convert to actual coordinates
    actual_x = x_idx + x_min
    actual_y = y_idx + y_min

    coord_key = (actual_x, actual_y)

    # Remove previous highlight rectangle if it exists
    if '_image_highlight_rect' in globals() and _image_highlight_rect is not None:
        try:
            _image_highlight_rect.remove()
        except:
            pass
        _image_highlight_rect = None

    # Add black border highlight to the clicked die (like in Data view)
    from matplotlib.patches import Rectangle
    ax = event.inaxes
    _image_highlight_rect = Rectangle(
        (x_idx - 0.5, y_idx - 0.5), 1, 1,
        fill=False, edgecolor='black', linewidth=3, zorder=100
    )
    ax.add_patch(_image_highlight_rect)

    # Update the canvas to show the highlight
    if heatmap_canvas:
        heatmap_canvas.draw_idle()

    # Store the selected die coordinates
    current_selected_die = (actual_x, actual_y)

    # Update the Image and PLM Files boxes on the left side
    display_die_images(actual_x, actual_y)
    display_plm_files(actual_x, actual_y)

    print(f"Clicked die ({actual_x}, {actual_y}) - Images and PLM Files updated")


def update_plm_wafermap():
    """Display wafermap with PLM data visualizations at their coordinates - like PLM Files panel"""
    global heatmap_canvas, multiple_stdf_data, multiple_wafer_ids, plm_file_directory

    print(f"update_plm_wafermap called")
    print(f"  plm_file_directory: {plm_file_directory}")

    # Close any existing figure to prevent memory warnings
    plt.close('all')

    if not multiple_stdf_data and (current_stdf_data is None or current_stdf_data.empty):
        print("No wafermap data loaded - showing message")
        fig, ax = plt.subplots(figsize=(10, 10))
        ax.text(0.5, 0.5, "No wafermap data loaded.\n\nPlease load STDF or CSV data first.",
                ha='center', va='center', fontsize=14, transform=ax.transAxes)
        ax.set_title("PLM Wafermap - No Data", fontsize=14)
        ax.axis('off')

        if heatmap_canvas:
            heatmap_canvas.get_tk_widget().destroy()

        heatmap_canvas = FigureCanvasTkAgg(fig, master=heatmap_display_frame)
        heatmap_canvas.get_tk_widget().pack(fill=tk.BOTH, expand=True)
        heatmap_canvas.draw()
        return

    # Get PLM directory - first try plm_file_directory, then config_dirs
    plm_dir = plm_file_directory
    if not plm_dir and "plm" in config_dirs:
        plm_dir = config_dirs["plm"].get()

    if not plm_dir or not os.path.exists(plm_dir):
        print("No PLM folder available.")
        # Show message on the heatmap
        fig, ax = plt.subplots(figsize=(10, 10))
        ax.text(0.5, 0.5, "No PLM folder available.\n\nPlease load a Project Folder\nthat contains PLMFiles subfolder.",
                ha='center', va='center', fontsize=14, transform=ax.transAxes)
        ax.set_title("PLM Wafermap - No PLM Files", fontsize=14)
        ax.axis('off')

        if heatmap_canvas:
            heatmap_canvas.get_tk_widget().destroy()

        heatmap_canvas = FigureCanvasTkAgg(fig, master=heatmap_display_frame)
        heatmap_canvas.get_tk_widget().pack(fill=tk.BOTH, expand=True)
        heatmap_canvas.draw()
        return

    # Get wafermap data
    if multiple_stdf_data:
        df = multiple_stdf_data[0]
        wafer_id = multiple_wafer_ids[0] if multiple_wafer_ids else "Wafer"
    else:
        df = current_stdf_data
        wafer_id = current_wafer_id if current_wafer_id else "Wafer"

    print(f"  Using dataframe with {len(df)} rows")

    # Estimate full wafer extent for proper display
    full_extent, center_info = _estimate_full_wafer_extent(df)
    x_min, y_min, x_max, y_max = full_extent

    grid_width = x_max - x_min + 1
    grid_height = y_max - y_min + 1

    print(f"Creating PLM wafermap: {grid_width}x{grid_height} grid")
    print(f"X range: {x_min} to {x_max}, Y range: {y_min} to {y_max}")
    print(f"PLM directory: {plm_dir}")

    # Show loading indicator
    print("Loading PLM data files...")

    # Create figure
    fig_size = max(12, min(20, grid_width * 0.5))
    fig, ax = plt.subplots(figsize=(fig_size, fig_size))

    # Set background color for the entire wafer area
    ax.set_facecolor('#e8e8e8')

    from matplotlib.patches import Rectangle
    import re

    # Draw light gray background for all die positions first
    die_positions = list(zip(df['x'].values, df['y'].values))
    for x, y in die_positions:
        x_idx = x - x_min
        y_idx = y - y_min
        rect = Rectangle((x_idx - 0.5, y_idx - 0.5), 1, 1,
                         facecolor='#f0f0f0', edgecolor='#cccccc', linewidth=0.3)
        ax.add_patch(rect)

    # Wafer circle outline removed per user request
    # _draw_wafer_circle(ax, grid_width, grid_height, x_min, y_min, center_info)

    plm_found = 0
    plm_not_found = 0

    # Get the selected PLM type filter
    try:
        selected_plm_type = view_subtype_var.get()
    except:
        selected_plm_type = "All"
    print(f"PLM type filter: {selected_plm_type}")

    # Get all PLM files in directory and map to coordinates
    # PLM files: .plm, .txt, .csv, .dat extensions
    all_plm_files = {}
    plm_extensions = ('.plm', '.txt', '.csv', '.dat')

    try:
        for filename in os.listdir(plm_dir):
            if filename.lower().endswith(plm_extensions):
                # Filter by PLM type if not "All"
                if selected_plm_type != "All":
                    # Check if filename contains the selected PLM type
                    # Pattern: ..._TypeName_timestamp.txt
                    if selected_plm_type not in filename:
                        continue

                # Extract coordinates from filename using multiple patterns
                # Pattern: Die_X19_Y46, X19_Y46, X19-Y46, -X19-Y46-
                match = re.search(r'[_\-]X(\d+)[_\-]Y(\d+)', filename, re.IGNORECASE)
                if match:
                    plm_x = int(match.group(1))
                    plm_y = int(match.group(2))
                    coord_key = (plm_x, plm_y)
                    if coord_key not in all_plm_files:
                        all_plm_files[coord_key] = []
                    all_plm_files[coord_key].append(os.path.join(plm_dir, filename))
    except Exception as e:
        print(f"Error scanning PLM directory: {e}")

    print(f"Found PLM files for {len(all_plm_files)} die positions")

    # Cache for PLM data converted to images
    plm_image_cache = {}
    loaded_count = 0

    # Place PLM visualizations at die coordinates
    for _, row in df.iterrows():
        x = int(row['x'])
        y = int(row['y'])
        x_idx = x - x_min
        y_idx = y - y_min

        coord_key = (x, y)

        if coord_key in all_plm_files and all_plm_files[coord_key]:
            # Use first PLM file for this die
            plm_path = all_plm_files[coord_key][0]

            try:
                if plm_path in plm_image_cache:
                    img_array = plm_image_cache[plm_path]
                else:
                    # Use the existing load_plm_as_matrix function
                    plm_data = load_plm_as_matrix(plm_path)

                    if plm_data is not None and plm_data.size > 0:
                        print(f"    Loaded PLM data: shape={plm_data.shape}, min={plm_data.min():.1f}, max={plm_data.max():.1f}")
                        # Convert PLM data to color-mapped image
                        # Normalize data to 0-1 range
                        data_min = np.nanmin(plm_data)
                        data_max = np.nanmax(plm_data)
                        if data_max > data_min:
                            normalized = (plm_data - data_min) / (data_max - data_min)
                        else:
                            normalized = np.zeros_like(plm_data)

                        # Apply viridis colormap
                        from matplotlib import cm
                        colormap = cm.get_cmap('viridis')
                        colored = (colormap(normalized)[:, :, :3] * 255).astype(np.uint8)

                        # Resize to 64x64 for display
                        img = Image.fromarray(colored)
                        img = img.resize((64, 64), Image.Resampling.BILINEAR)
                        img_array = np.array(img, dtype=np.uint8)
                        plm_image_cache[plm_path] = img_array
                    else:
                        img_array = None
                        plm_image_cache[plm_path] = None

                if img_array is not None:
                    extent = [x_idx - 0.48, x_idx + 0.48, y_idx - 0.48, y_idx + 0.48]
                    ax.imshow(img_array, extent=extent, aspect='auto',
                             interpolation='nearest', zorder=2)

                    plm_found += 1
                    loaded_count += 1

                    if loaded_count % 20 == 0:
                        print(f"  Loaded {loaded_count} PLM files...")
                else:
                    plm_not_found += 1

            except Exception as e:
                print(f"Error processing PLM file {plm_path}: {e}")
                plm_not_found += 1
        else:
            plm_not_found += 1

    # Set axis properties
    ax.set_xlim(-0.5, grid_width - 0.5)
    ax.set_ylim(-0.5, grid_height - 0.5)
    ax.set_aspect('equal')
    ax.set_xlabel('X', fontsize=10, fontweight='bold')
    ax.set_ylabel('Y', fontsize=10, fontweight='bold')
    # Truncate title to fit width, reduce font size by 50%
    short_wafer_id = wafer_id[:20] + "..." if len(str(wafer_id)) > 20 else wafer_id
    ax.set_title(f"PLM: {short_wafer_id} [{selected_plm_type}]\n{plm_found} files, {plm_not_found} empty",
                fontsize=6, fontweight='bold')

    # Add coordinate labels - show actual X/Y coordinates
    tick_step = max(1, grid_width // 10)
    x_ticks = range(0, grid_width, tick_step)
    ax.set_xticks(x_ticks)
    ax.set_xticklabels([str(x_min + i) for i in x_ticks], fontsize=8)

    tick_step_y = max(1, grid_height // 10)
    y_ticks = range(0, grid_height, tick_step_y)
    ax.set_yticks(y_ticks)
    ax.set_yticklabels([str(y_min + i) for i in y_ticks], fontsize=8)

    fig.tight_layout()

    # Display on canvas
    if heatmap_canvas:
        heatmap_canvas.get_tk_widget().destroy()

    heatmap_canvas = FigureCanvasTkAgg(fig, master=heatmap_display_frame)
    canvas_widget = heatmap_canvas.get_tk_widget()
    canvas_widget.pack(fill=tk.BOTH, expand=True)

    # Connect scroll events for zooming
    heatmap_canvas.mpl_connect("scroll_event", on_heatmap_scroll)
    canvas_widget.bind("<MouseWheel>", on_tk_mousewheel)

    # Connect click event for die selection
    heatmap_canvas.mpl_connect("button_press_event", lambda e: on_plm_wafermap_click(e, x_min, y_min, df, all_plm_files))

    heatmap_canvas.draw()

    print(f"PLM wafermap displayed: {plm_found} files, {plm_not_found} empty positions")
    print("Zoom in to see PLM data details!")


def on_plm_wafermap_click(event, x_min, y_min, df, all_plm_files):
    """Handle click on PLM wafermap to show detailed PLM data with die highlighting"""
    global current_selected_die_plm, current_selected_die, heatmap_canvas, _plm_highlight_rect

    if event.inaxes is None or event.button != 1:
        return

    x_idx = int(round(event.xdata))
    y_idx = int(round(event.ydata))

    actual_x = x_idx + x_min
    actual_y = y_idx + y_min

    coord_key = (actual_x, actual_y)

    # Remove previous highlight rectangle if it exists
    if '_plm_highlight_rect' in globals() and _plm_highlight_rect is not None:
        try:
            _plm_highlight_rect.remove()
        except:
            pass
        _plm_highlight_rect = None

    # Add black border highlight to the clicked die
    from matplotlib.patches import Rectangle
    ax = event.inaxes
    _plm_highlight_rect = Rectangle(
        (x_idx - 0.5, y_idx - 0.5), 1, 1,
        fill=False, edgecolor='black', linewidth=3, zorder=100
    )
    ax.add_patch(_plm_highlight_rect)

    # Update the canvas to show the highlight
    if heatmap_canvas:
        heatmap_canvas.draw_idle()

    # Store the selected die coordinates
    current_selected_die_plm = (actual_x, actual_y)
    current_selected_die = (actual_x, actual_y)

    # Update the Image and PLM Files boxes on the left side
    display_die_images(actual_x, actual_y)
    display_plm_files(actual_x, actual_y)

    print(f"Clicked die ({actual_x}, {actual_y}) - Images and PLM Files updated")


def show_plm_data_popup(plm_path, x_coord, y_coord):
    """Show PLM data in a detailed popup window with heatmap visualization"""
    try:
        with open(plm_path, 'r') as f:
            lines = f.readlines()

        # Parse header and data
        cols = 18
        rows = 24
        description = ""
        export_quantity = ""
        data_lines = []
        header_ended = False

        for line in lines:
            line_stripped = line.strip()
            if line_stripped.startswith('Description:'):
                description = line_stripped.split(':', 1)[1].strip()
            elif line_stripped.startswith('Columns:'):
                cols = int(line_stripped.split(':')[1].strip())
            elif line_stripped.startswith('Rows:'):
                rows = int(line_stripped.split(':')[1].strip())
            elif line_stripped.startswith('Export Quantity:'):
                export_quantity = line_stripped.split(':', 1)[1].strip()
                header_ended = True
            elif header_ended and line_stripped:
                values = [float(v) for v in line_stripped.split(',') if v.strip()]
                data_lines.extend(values)

        # Determine data format: pairs (luminance, other) or single values
        total_expected = rows * cols

        if len(data_lines) >= total_expected * 2:
            # Format with pairs (take every other value for luminance)
            luminance_values = data_lines[::2]
        elif len(data_lines) >= total_expected:
            # Format with single values (all values are luminance)
            luminance_values = data_lines
        else:
            luminance_values = data_lines

        if len(luminance_values) >= total_expected:
            data_array = np.array(luminance_values[:total_expected]).reshape(rows, cols)

            # Create popup window
            popup = tk.Toplevel()
            popup.title(f"PLM Data - Die ({x_coord}, {y_coord})")
            popup.geometry("700x600")

            # Info label
            info_text = f"File: {os.path.basename(plm_path)}\n"
            info_text += f"Description: {description}\n"
            info_text += f"Export: {export_quantity} | Size: {cols}x{rows}"
            info_label = tk.Label(popup, text=info_text, font=("Helvetica", 9), justify=tk.LEFT)
            info_label.pack(pady=5)

            # Create matplotlib figure for the heatmap
            fig, ax = plt.subplots(figsize=(8, 6))
            im = ax.imshow(data_array, cmap='viridis', aspect='auto')
            ax.set_title(f"PLM Luminance Data - Die ({x_coord}, {y_coord})", fontsize=12)
            ax.set_xlabel('Column')
            ax.set_ylabel('Row')
            fig.colorbar(im, ax=ax, label=export_quantity)
            fig.tight_layout()

            # Embed in tkinter
            canvas = FigureCanvasTkAgg(fig, master=popup)
            canvas.get_tk_widget().pack(fill=tk.BOTH, expand=True, padx=10, pady=5)
            canvas.draw()

            # Stats label
            stats_text = f"Min: {data_array.min():.1f} | Max: {data_array.max():.1f} | "
            stats_text += f"Mean: {data_array.mean():.1f} | Std: {data_array.std():.1f}"
            stats_label = tk.Label(popup, text=stats_text, font=("Helvetica", 9))
            stats_label.pack(pady=5)

        else:
            print(f"Not enough data in PLM file: {len(luminance_values)} values, expected {total_expected}")

    except Exception as e:
        print(f"Error showing PLM data: {e}")


def zoom_heatmap(zoom_in):
    """Zoom in or out on the heatmap using buttons"""
    global heatmap_canvas

    if heatmap_canvas is None:
        print("No heatmap canvas available")
        return

    # Get the figure and axes from the canvas
    fig = heatmap_canvas.figure
    if not fig.axes:
        print("No axes found in figure")
        return

    ax = fig.axes[0]

    # Get current axis limits
    cur_xlim = ax.get_xlim()
    cur_ylim = ax.get_ylim()

    # Zoom factor
    zoom_factor = 0.8 if zoom_in else 1.25

    # Calculate center of current view
    x_center = (cur_xlim[0] + cur_xlim[1]) / 2
    y_center = (cur_ylim[0] + cur_ylim[1]) / 2

    # Calculate current ranges
    x_range = cur_xlim[1] - cur_xlim[0]
    y_range = cur_ylim[1] - cur_ylim[0]

    # New ranges after zoom
    new_x_range = x_range * zoom_factor
    new_y_range = y_range * zoom_factor

    # New limits centered on current view center
    new_xlim = [x_center - new_x_range / 2, x_center + new_x_range / 2]
    new_ylim = [y_center - new_y_range / 2, y_center + new_y_range / 2]

    ax.set_xlim(new_xlim)
    ax.set_ylim(new_ylim)

    # Force redraw
    heatmap_canvas.draw()
    heatmap_canvas.flush_events()

    print(f"Zoom {'in' if zoom_in else 'out'}: xlim {cur_xlim} -> {new_xlim}")


# ============== RECTANGLE ZOOM AND PAN STATE ==============
_heatmap_interaction_state = {
    'rect_zoom_active': False,
    'rect_zoom_start': None,
    'rect_zoom_rect': None,
    'pan_active': False,
    'pan_start': None,
    'pan_xlim': None,
    'pan_ylim': None
}

# Store current die popup window reference
_wafer_die_popup = None


def show_die_data_popup(x_coord, y_coord, data_source, selected_param_column):
    """Show a popup window with all measurement parameters for the clicked die
    Design consistent with show_die_data_popup_all_wafers from Multiwafer tab"""
    global _wafer_die_popup, test_parameters, multiple_stdf_data, multiple_wafer_ids

    # Close previous popup if exists
    if _wafer_die_popup is not None:
        try:
            _wafer_die_popup.destroy()
        except:
            pass
        _wafer_die_popup = None

    # Collect data from all loaded wafers (like Multiwafer tab)
    all_wafer_data = []

    if multiple_stdf_data and len(multiple_stdf_data) > 0:
        # Multiple files loaded - show all wafers
        for i, df in enumerate(multiple_stdf_data):
            wafer_id = multiple_wafer_ids[i] if i < len(multiple_wafer_ids) else f"Wafer {i+1}"
            die_data = df[(df['x'] == x_coord) & (df['y'] == y_coord)]
            if len(die_data) > 0:
                all_wafer_data.append({
                    'wafer_id': wafer_id,
                    'die_row': die_data.iloc[0]
                })
    else:
        # Single file loaded
        die_data = data_source[(data_source['x'] == x_coord) & (data_source['y'] == y_coord)]
        if len(die_data) > 0:
            all_wafer_data.append({
                'wafer_id': "Wafer",
                'die_row': die_data.iloc[0]
            })

    if not all_wafer_data:
        print(f"No data found at position X={x_coord}, Y={y_coord}")
        return

    # Create popup window (75% of original size: 450x375 instead of 600x500)
    popup = tk.Toplevel(main_win)
    popup.title(f"Die Data - Position: X={x_coord}, Y={y_coord} - {len(all_wafer_data)} Wafer(s)")
    popup.geometry("450x375")
    popup.transient(main_win)

    # Store reference to current popup
    _wafer_die_popup = popup

    def on_popup_close():
        global _wafer_die_popup
        _wafer_die_popup = None
        popup.destroy()

    popup.protocol("WM_DELETE_WINDOW", on_popup_close)

    # Center the popup
    popup.update_idletasks()
    x = main_win.winfo_x() + (main_win.winfo_width() - 450) // 2
    y = main_win.winfo_y() + (main_win.winfo_height() - 375) // 2
    popup.geometry(f"+{x}+{y}")

    # Header frame
    header_frame = tk.Frame(popup, bg="#2196F3", pady=5)
    header_frame.pack(fill=tk.X)

    tk.Label(
        header_frame,
        text=f"Position: X = {x_coord}, Y = {y_coord}",
        font=("Helvetica", 10, "bold"),
        bg="#2196F3",
        fg="white"
    ).pack()

    tk.Label(
        header_frame,
        text=f"Data from {len(all_wafer_data)} Wafer(s)",
        font=("Helvetica", 9),
        bg="#2196F3",
        fg="white"
    ).pack()

    # Legend frame
    legend_frame = tk.Frame(popup, pady=3)
    legend_frame.pack(fill=tk.X)
    tk.Label(legend_frame, text="■ Selected Parameter", fg="red", font=("Helvetica", 8, "bold")).pack(side=tk.LEFT, padx=5)

    # Create Treeview with scrollbars - dynamic columns based on number of wafers
    tree_frame = tk.Frame(popup)
    tree_frame.pack(fill=tk.BOTH, expand=True, padx=5, pady=3)

    tree_scroll_y = ttk.Scrollbar(tree_frame, orient="vertical")
    tree_scroll_x = ttk.Scrollbar(tree_frame, orient="horizontal")

    # Build column list: Parameter + one column per wafer
    columns = ["parameter"]
    for wd in all_wafer_data:
        short_id = str(wd['wafer_id'])[:15] + "..." if len(str(wd['wafer_id'])) > 15 else str(wd['wafer_id'])
        columns.append(short_id)

    tree = ttk.Treeview(
        tree_frame,
        columns=columns,
        show="headings",
        yscrollcommand=tree_scroll_y.set,
        xscrollcommand=tree_scroll_x.set
    )

    # Set headings
    tree.heading("parameter", text="Parameter")
    tree.column("parameter", width=150, anchor="w")

    col_width = max(60, 250 // len(all_wafer_data)) if all_wafer_data else 80
    for i, wd in enumerate(all_wafer_data):
        short_id = str(wd['wafer_id'])[:15] + "..." if len(str(wd['wafer_id'])) > 15 else str(wd['wafer_id'])
        tree.heading(columns[i+1], text=short_id)
        tree.column(columns[i+1], width=col_width, anchor="center")

    tree_scroll_y.config(command=tree.yview)
    tree_scroll_x.config(command=tree.xview)

    tree_scroll_y.pack(side=tk.RIGHT, fill=tk.Y)
    tree_scroll_x.pack(side=tk.BOTTOM, fill=tk.X)
    tree.pack(fill=tk.BOTH, expand=True)

    # Configure tags for highlighting
    tree.tag_configure("selected", background="#FFCDD2", foreground="red")
    tree.tag_configure("normal", background="white")
    tree.tag_configure("coordinate", background="#E3F2FD")

    # Get all unique columns from all wafers
    all_columns = set()
    for wd in all_wafer_data:
        all_columns.update(wd['die_row'].index)

    # Remove x, y from display columns
    display_columns = [c for c in all_columns if c not in ['x', 'y']]

    # Sort columns: bin first, then numeric test numbers
    def sort_key(col):
        if col == 'bin':
            return (0, 0)
        elif isinstance(col, int):
            return (1, col)
        else:
            return (2, str(col))

    display_columns = sorted(display_columns, key=sort_key)

    # Add data rows
    for col in display_columns:
        # Determine parameter name
        if col == 'bin':
            param_name = "BIN"
        else:
            test_key = f"test_{col}"
            if test_parameters and test_key in test_parameters:
                param_name = f"{col}: {test_parameters[test_key]}"
            else:
                param_name = f"Test {col}"

        # Determine if this is the selected parameter
        is_selected = False
        if col == selected_param_column:
            is_selected = True
        elif isinstance(col, int) and col == selected_param_column:
            is_selected = True

        # Collect values from all wafers
        row_values = [param_name]
        for wd in all_wafer_data:
            die_row = wd['die_row']
            if col in die_row.index:
                value = die_row[col]
                if pd.notna(value):
                    if isinstance(value, float):
                        row_values.append(f"{value:.4g}")
                    else:
                        row_values.append(str(value))
                else:
                    row_values.append("-")
            else:
                row_values.append("-")

        tag = "selected" if is_selected else "normal"
        tree.insert("", "end", values=row_values, tags=(tag,))

    # Button frame
    btn_frame = tk.Frame(popup, pady=5)
    btn_frame.pack(fill=tk.X)

    def copy_to_clipboard():
        """Copy die data to clipboard"""
        data_str = f"Position: X={x_coord}, Y={y_coord}\n\n"
        # Header row
        header = "Parameter"
        for wd in all_wafer_data:
            header += f"\t{wd['wafer_id']}"
        data_str += header + "\n"

        for item in tree.get_children():
            values = tree.item(item, "values")
            data_str += "\t".join(str(v) for v in values) + "\n"

        popup.clipboard_clear()
        popup.clipboard_append(data_str)
        print("Die data copied to clipboard")

    tk.Button(
        btn_frame,
        text="Copy",
        command=copy_to_clipboard,
        font=("Helvetica", 8),
        bg="#4CAF50",
        fg="white"
    ).pack(side=tk.LEFT, padx=5)

    tk.Button(
        btn_frame,
        text="Close",
        command=on_popup_close,
        font=("Helvetica", 8)
    ).pack(side=tk.RIGHT, padx=5)


def on_heatmap_press(event):
    """Handle mouse button press for rectangle zoom (left) and pan (middle)"""
    global heatmap_canvas, _heatmap_interaction_state

    if event.inaxes is None or heatmap_canvas is None:
        return

    ax = event.inaxes

    # Left button (button 1) - start rectangle zoom
    if event.button == 1:
        _heatmap_interaction_state['rect_zoom_active'] = True
        _heatmap_interaction_state['rect_zoom_start'] = (event.xdata, event.ydata)

        # Create initial rectangle (invisible until drag)
        from matplotlib.patches import Rectangle
        rect = Rectangle(
            (event.xdata, event.ydata), 0, 0,
            linewidth=2, edgecolor='red', facecolor='yellow',
            alpha=0.3, linestyle='--', zorder=100
        )
        ax.add_patch(rect)
        _heatmap_interaction_state['rect_zoom_rect'] = rect
        _heatmap_interaction_state['rect_zoom_ax'] = ax
        heatmap_canvas.draw_idle()

    # Middle button (button 2) - start pan
    elif event.button == 2:
        _heatmap_interaction_state['pan_active'] = True
        _heatmap_interaction_state['pan_start'] = (event.xdata, event.ydata)
        _heatmap_interaction_state['pan_xlim'] = ax.get_xlim()
        _heatmap_interaction_state['pan_ylim'] = ax.get_ylim()
        _heatmap_interaction_state['pan_ax'] = ax
        # Change cursor to indicate panning
        heatmap_canvas.get_tk_widget().config(cursor="fleur")


def on_heatmap_motion(event):
    """Handle mouse motion for rectangle zoom and pan"""
    global heatmap_canvas, _heatmap_interaction_state

    if heatmap_canvas is None:
        return

    # Rectangle zoom drag
    if _heatmap_interaction_state['rect_zoom_active'] and _heatmap_interaction_state['rect_zoom_start']:
        if event.xdata is None or event.ydata is None:
            return

        x0, y0 = _heatmap_interaction_state['rect_zoom_start']
        x1, y1 = event.xdata, event.ydata

        # Update rectangle
        rect = _heatmap_interaction_state['rect_zoom_rect']
        if rect:
            rect.set_x(min(x0, x1))
            rect.set_y(min(y0, y1))
            rect.set_width(abs(x1 - x0))
            rect.set_height(abs(y1 - y0))
            heatmap_canvas.draw_idle()

    # Pan drag
    elif _heatmap_interaction_state['pan_active'] and _heatmap_interaction_state['pan_start']:
        if event.xdata is None or event.ydata is None:
            return

        ax = _heatmap_interaction_state['pan_ax']
        x0, y0 = _heatmap_interaction_state['pan_start']
        dx = x0 - event.xdata
        dy = y0 - event.ydata

        xlim = _heatmap_interaction_state['pan_xlim']
        ylim = _heatmap_interaction_state['pan_ylim']

        ax.set_xlim(xlim[0] + dx, xlim[1] + dx)
        ax.set_ylim(ylim[0] + dy, ylim[1] + dy)

        heatmap_canvas.draw_idle()


def on_heatmap_release(event):
    """Handle mouse button release for rectangle zoom and pan"""
    global heatmap_canvas, _heatmap_interaction_state

    if heatmap_canvas is None:
        return

    # Rectangle zoom release (left button)
    if event.button == 1 and _heatmap_interaction_state['rect_zoom_active']:
        _heatmap_interaction_state['rect_zoom_active'] = False

        # Remove the rectangle
        rect = _heatmap_interaction_state['rect_zoom_rect']
        if rect:
            rect.remove()
            _heatmap_interaction_state['rect_zoom_rect'] = None

        # Apply zoom if we have valid start and current positions
        start = _heatmap_interaction_state['rect_zoom_start']
        if start and event.xdata is not None and event.ydata is not None:
            x0, y0 = start
            x1, y1 = event.xdata, event.ydata

            # Only zoom if the rectangle is large enough (not just a click)
            if abs(x1 - x0) > 0.5 and abs(y1 - y0) > 0.5:
                ax = _heatmap_interaction_state.get('rect_zoom_ax')
                if ax:
                    # Set new limits to the rectangle bounds
                    ax.set_xlim(min(x0, x1), max(x0, x1))
                    ax.set_ylim(min(y0, y1), max(y0, y1))
                    print(f"Rectangle zoom: x=[{min(x0,x1):.1f}, {max(x0,x1):.1f}], y=[{min(y0,y1):.1f}, {max(y0,y1):.1f}]")
            else:
                # Small movement = regular click, trigger die selection
                on_heatmap_die_click(event)

        _heatmap_interaction_state['rect_zoom_start'] = None
        heatmap_canvas.draw_idle()

    # Pan release (middle button)
    elif event.button == 2 and _heatmap_interaction_state['pan_active']:
        _heatmap_interaction_state['pan_active'] = False
        _heatmap_interaction_state['pan_start'] = None
        # Restore cursor
        heatmap_canvas.get_tk_widget().config(cursor="")
        print("Pan complete")


def on_heatmap_die_click(event):
    """Handle click on heatmap to display die value (called after rectangle zoom check)"""
    global current_stdf_data, selected_die_coords, heatmap_canvas, multiple_stdf_data, _die_info_window

    if event.inaxes is None:
        return

    # Use multiple_stdf_data if available, otherwise use current_stdf_data
    if multiple_stdf_data and len(multiple_stdf_data) > 0:
        data_source = multiple_stdf_data[0]
    elif current_stdf_data is not None and not current_stdf_data.empty:
        data_source = current_stdf_data
    else:
        return

    selected = heatmap_param_combobox.get()
    if not selected:
        return

    # Check if this is a custom test or separator
    if selected.startswith("───"):
        return  # Skip separator

    if selected.startswith("BIN"):
        param_column = "bin"
        param_label = "Bin"
    elif selected.startswith("CUSTOM:"):
        # Handle custom test
        custom_test_name = selected.replace("CUSTOM:", "").strip()
        param_column = f"_custom_{custom_test_name}"
        param_label = f"Custom: {custom_test_name}"
    else:
        test_key = selected.split(":")[0].strip()
        if test_key.startswith("test_"):
            param_column = int(test_key.replace("test_", ""))
        else:
            try:
                param_column = int(test_key)
            except ValueError:
                return  # Invalid parameter
        param_label = selected

    x_click = event.xdata
    y_click = event.ydata

    if x_click is None or y_click is None:
        return

    if param_column not in data_source.columns:
        return

    mask = data_source[param_column].notna()
    plot_data = data_source[mask]

    if len(plot_data) == 0:
        return

    # Use the same extent calculation as the display function
    full_extent, center_info = _estimate_full_wafer_extent(plot_data)
    x_min, y_min, x_max, y_max = full_extent

    actual_x = int(round(x_click + x_min))
    actual_y = int(round(y_click + y_min))

    # Search in full data source to find the die
    die_data = data_source[(data_source["x"] == actual_x) & (data_source["y"] == actual_y)]

    if len(die_data) > 0:
        die_row = die_data.iloc[0]
        value = die_row[param_column]

        # Set selected coordinates
        selected_die_coords = (actual_x, actual_y)

        # Close previous die info window if it exists
        if _die_info_window is not None:
            try:
                _die_info_window.destroy()
            except:
                pass

        print(f"Clicked die at ({actual_x}, {actual_y}): {param_label} = {value}")

        # Show popup with all measurement parameters
        show_die_data_popup(actual_x, actual_y, data_source, param_column)

        # Display images and PLM files for the selected die
        display_die_images(actual_x, actual_y)
        display_plm_files(actual_x, actual_y)

        # Refresh heatmap with selection highlight
        refresh_heatmap_display()


def clear_die_selection():
    """Clear the selected die and refresh the heatmap"""
    global selected_die_coords
    selected_die_coords = None
    print("Die selection cleared")
    refresh_heatmap_display()


def on_heatmap_click(event):
    """Handle left-click on heatmap to display die value"""
    global current_stdf_data, selected_die_coords, heatmap_canvas, multiple_stdf_data, _die_info_window

    print(f"Click event: button={event.button}, inaxes={event.inaxes is not None}, dblclick={event.dblclick}")

    # Middle button handled by pan, left button handled by rectangle zoom/click
    if event.button != 1:
        return

    if event.inaxes is None:
        print("Click outside axes")
        return

    # Use multiple_stdf_data if available, otherwise use current_stdf_data
    if multiple_stdf_data and len(multiple_stdf_data) > 0:
        data_source = multiple_stdf_data[0]  # Use first file for now
    elif current_stdf_data is not None and not current_stdf_data.empty:
        data_source = current_stdf_data
    else:
        print("No STDF data available")
        return

    selected = heatmap_param_combobox.get()

    if not selected:
        print("No parameter selected")
        return

    if selected.startswith("BIN"):
        param_column = "bin"
        param_label = "Bin"
    else:
        test_key = selected.split(":")[0].strip()
        if test_key.startswith("test_"):
            param_column = int(test_key.replace("test_", ""))
        else:
            param_column = int(test_key)
        param_label = selected

    x_click = event.xdata
    y_click = event.ydata

    print(f"Click at grid position: ({x_click}, {y_click})")

    if x_click is None or y_click is None:
        return

    if param_column not in data_source.columns:
        print(f"Parameter {param_column} not in data")
        return

    mask = data_source[param_column].notna()
    plot_data = data_source[mask]

    if len(plot_data) == 0:
        return

    # Use the same extent calculation as the display function
    full_extent, center_info = _estimate_full_wafer_extent(plot_data)
    x_min, y_min, x_max, y_max = full_extent

    actual_x = int(round(x_click + x_min))
    actual_y = int(round(y_click + y_min))

    print(f"Actual die coordinates: ({actual_x}, {actual_y})")

    # Search in full data source (not just filtered plot_data) to find the die
    die_data = data_source[(data_source["x"] == actual_x) & (data_source["y"] == actual_y)]

    if len(die_data) > 0:
        die_row = die_data.iloc[0]
        value = die_row[param_column]

        print(f"Found die with value: {value}")

        # Set selected coordinates
        selected_die_coords = (actual_x, actual_y)
        print(f"Set selected_die_coords to: {selected_die_coords}")

        # Close previous die info window if it exists
        if _die_info_window is not None:
            try:
                _die_info_window.destroy()
            except:
                pass  # Window may already be destroyed

        # Show info popup FIRST (before redraw which may take time)
        info_win = tk.Toplevel(main_win)
        info_win.title(f"Die Information - ({actual_x}, {actual_y})")
        info_win.geometry("400x300")

        # Store reference to current window
        _die_info_window = info_win

        text_frame = tk.Frame(info_win)
        text_frame.pack(fill=tk.BOTH, expand=True, padx=10, pady=10)

        scrollbar = tk.Scrollbar(text_frame)
        scrollbar.pack(side=tk.RIGHT, fill=tk.Y)

        text_widget = tk.Text(
            text_frame, wrap=tk.WORD, yscrollcommand=scrollbar.set, font=("Courier", 10)
        )
        text_widget.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)

        scrollbar.config(command=text_widget.yview)

        info_text = f"Die Coordinates: ({actual_x}, {actual_y})\n"
        info_text += f"{'='*40}\n\n"

        text_widget.insert("end", info_text)

        try:
            text_widget.insert("end", f"{param_label}: {value:.4f}\n\n", "selected")
        except (TypeError, ValueError):
            text_widget.insert("end", f"{param_label}: {value}\n\n", "selected")

        text_widget.insert("end", f"Additional Information:\n")
        text_widget.insert("end", f"{'-'*40}\n")

        if "bin" in die_row:
            if param_column == "bin":
                text_widget.insert("end", f"Bin: {die_row['bin']}\n", "selected")
            else:
                text_widget.insert("end", f"Bin: {die_row['bin']}\n")

        for col in die_row.index:
            if isinstance(col, int) and col != param_column:
                test_name = test_parameters.get(f"test_{col}", f"Test {col}")
                try:
                    text_widget.insert("end", f"{test_name}: {die_row[col]:.4f}\n")
                except (TypeError, ValueError):
                    text_widget.insert("end", f"{test_name}: {die_row[col]}\n")

        text_widget.tag_config(
            "selected", foreground="red", font=("Courier", 10, "bold")
        )

        text_widget.config(state=tk.DISABLED)

        # Clear reference when window is closed
        def on_info_close():
            global _die_info_window
            _die_info_window = None
            info_win.destroy()

        info_win.protocol("WM_DELETE_WINDOW", on_info_close)

        print(f"Clicked die at ({actual_x}, {actual_y}): {param_label} = {value}")

        # Display images for the selected die
        display_die_images(actual_x, actual_y)

        # Display PLM files for the selected die
        display_plm_files(actual_x, actual_y)

        # Now redraw with selection highlight
        refresh_heatmap_display()

    else:
        print(f"No die found at position ({actual_x}, {actual_y})")


def on_heatmap_scroll(event):
    """Handle mouse wheel scroll to zoom in/out on heatmap"""
    global heatmap_canvas

    print(f"Scroll event: step={getattr(event, 'step', 'N/A')}, button={getattr(event, 'button', 'N/A')}")

    if event.inaxes is None:
        print("Scroll outside axes")
        return

    ax = event.inaxes

    # Get current axis limits
    cur_xlim = ax.get_xlim()
    cur_ylim = ax.get_ylim()

    x_data = event.xdata
    y_data = event.ydata

    if x_data is None or y_data is None:
        print("No valid scroll position")
        return

    # Determine zoom direction
    # On Windows TkAgg: event.step > 0 means scroll up (zoom in)
    # event.button can be 'up' or 'down'
    zoom_in = False

    if hasattr(event, 'step') and event.step is not None:
        zoom_in = event.step > 0
        print(f"Using step: {event.step}, zoom_in={zoom_in}")
    elif hasattr(event, 'button') and event.button in ['up', 'down']:
        zoom_in = (event.button == 'up')
        print(f"Using button: {event.button}, zoom_in={zoom_in}")
    else:
        print("Could not determine scroll direction")
        return

    # Zoom factor: 0.8 = zoom in (smaller view), 1.25 = zoom out (larger view)
    zoom_factor = 0.8 if zoom_in else 1.25

    # Calculate current ranges
    x_range = cur_xlim[1] - cur_xlim[0]
    y_range = cur_ylim[1] - cur_ylim[0]

    # New ranges after zoom
    new_x_range = x_range * zoom_factor
    new_y_range = y_range * zoom_factor

    # Calculate relative position of mouse in current view (0 to 1)
    rel_x = (x_data - cur_xlim[0]) / x_range
    rel_y = (y_data - cur_ylim[0]) / y_range

    # New limits keeping mouse position fixed
    new_xlim = [x_data - rel_x * new_x_range, x_data + (1 - rel_x) * new_x_range]
    new_ylim = [y_data - rel_y * new_y_range, y_data + (1 - rel_y) * new_y_range]

    print(f"Zoom {'in' if zoom_in else 'out'}: xlim {cur_xlim} -> {new_xlim}")

    ax.set_xlim(new_xlim)
    ax.set_ylim(new_ylim)

    # Force redraw
    if heatmap_canvas:
        heatmap_canvas.draw()
        heatmap_canvas.flush_events()


def on_tk_mousewheel(event):
    """Handle Windows native mouse wheel events for zoom"""
    global heatmap_canvas

    print(f"TK MouseWheel event: delta={event.delta}")

    if heatmap_canvas is None:
        return

    fig = heatmap_canvas.figure
    if not fig.axes:
        return

    ax = fig.axes[0]

    # Get current axis limits
    cur_xlim = ax.get_xlim()
    cur_ylim = ax.get_ylim()

    # On Windows, event.delta is positive for scroll up, negative for scroll down
    zoom_in = event.delta > 0
    zoom_factor = 0.8 if zoom_in else 1.25

    # Calculate center of current view
    x_center = (cur_xlim[0] + cur_xlim[1]) / 2
    y_center = (cur_ylim[0] + cur_ylim[1]) / 2

    # Calculate current ranges
    x_range = cur_xlim[1] - cur_xlim[0]
    y_range = cur_ylim[1] - cur_ylim[0]

    # New ranges after zoom
    new_x_range = x_range * zoom_factor
    new_y_range = y_range * zoom_factor

    # New limits centered on current view center
    new_xlim = [x_center - new_x_range / 2, x_center + new_x_range / 2]
    new_ylim = [y_center - new_y_range / 2, y_center + new_y_range / 2]

    ax.set_xlim(new_xlim)
    ax.set_ylim(new_ylim)

    # Force redraw
    heatmap_canvas.draw()
    heatmap_canvas.flush_events()

    print(f"TK Zoom {'in' if zoom_in else 'out'}: xlim {cur_xlim} -> {new_xlim}")


def on_tk_mousewheel_linux(event, zoom_in):
    """Handle Linux native mouse wheel events for zoom"""
    global heatmap_canvas

    print(f"TK Linux MouseWheel event: zoom_in={zoom_in}")

    if heatmap_canvas is None:
        return

    fig = heatmap_canvas.figure
    if not fig.axes:
        return

    ax = fig.axes[0]

    # Get current axis limits
    cur_xlim = ax.get_xlim()
    cur_ylim = ax.get_ylim()

    zoom_factor = 0.8 if zoom_in else 1.25

    # Calculate center of current view
    x_center = (cur_xlim[0] + cur_xlim[1]) / 2
    y_center = (cur_ylim[0] + cur_ylim[1]) / 2

    # Calculate current ranges
    x_range = cur_xlim[1] - cur_xlim[0]
    y_range = cur_ylim[1] - cur_ylim[0]

    # New ranges after zoom
    new_x_range = x_range * zoom_factor
    new_y_range = y_range * zoom_factor

    # New limits centered on current view center
    new_xlim = [x_center - new_x_range / 2, x_center + new_x_range / 2]
    new_ylim = [y_center - new_y_range / 2, y_center + new_y_range / 2]

    ax.set_xlim(new_xlim)
    ax.set_ylim(new_ylim)

    # Force redraw
    heatmap_canvas.draw()
    heatmap_canvas.flush_events()

    print(f"TK Linux Zoom {'in' if zoom_in else 'out'}: xlim {cur_xlim} -> {new_xlim}")


def select_die_image_folder():
    """Select the folder containing die images"""
    global die_image_directory

    folder_path = filedialog.askdirectory(
        title="Select Image Folder for Die Images"
    )

    if folder_path:
        die_image_directory = folder_path
        # Show short version of path
        short_path = folder_path.split("/")[-1] if "/" in folder_path else folder_path.split("\\")[-1]
        folder_status_label.config(text=f"Img: {short_path}", fg="green")
        print(f"Die image folder set to: {folder_path}")


def show_image_gallery():
    """Show all images from the image folder in a gallery/grid view"""
    global die_image_directory

    if not die_image_directory:
        # Ask user to select folder first
        folder_path = filedialog.askdirectory(
            title="Select Image Folder to Display"
        )
        if not folder_path:
            return
        die_image_directory = folder_path
        short_path = folder_path.split("/")[-1] if "/" in folder_path else folder_path.split("\\")[-1]
        folder_status_label.config(text=f"Img: {short_path}", fg="green")

    # Get all image files
    image_extensions = ('.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tif', '.tiff')
    image_files = []

    try:
        for filename in os.listdir(die_image_directory):
            if filename.lower().endswith(image_extensions):
                image_files.append(os.path.join(die_image_directory, filename))
    except Exception as e:
        print(f"Error reading image folder: {e}")
        return

    if not image_files:
        tk.messagebox.showinfo("No Images", f"No images found in:\n{die_image_directory}")
        return

    # Sort files by name
    image_files.sort()

    # Create gallery window
    gallery_win = tk.Toplevel(main_win)
    gallery_win.title(f"Image Gallery - {len(image_files)} images")
    gallery_win.geometry("1400x900")

    # Control frame at top
    control_frame = tk.Frame(gallery_win)
    control_frame.pack(side=tk.TOP, fill=tk.X, padx=10, pady=5)

    # Thumbnail size selector
    tk.Label(control_frame, text="Thumbnail Size:", font=("Helvetica", 10)).pack(side=tk.LEFT, padx=5)

    thumb_size_var = tk.StringVar(value="Medium")
    thumb_sizes = {"Small": 100, "Medium": 150, "Large": 200, "XLarge": 300}

    thumb_size_combo = ttk.Combobox(
        control_frame,
        textvariable=thumb_size_var,
        values=list(thumb_sizes.keys()),
        state="readonly",
        width=10
    )
    thumb_size_combo.pack(side=tk.LEFT, padx=5)

    # Columns selector
    tk.Label(control_frame, text="Columns:", font=("Helvetica", 10)).pack(side=tk.LEFT, padx=15)

    cols_var = tk.StringVar(value="Auto")
    cols_combo = ttk.Combobox(
        control_frame,
        textvariable=cols_var,
        values=["Auto", "2", "3", "4", "5", "6", "8", "10"],
        state="readonly",
        width=8
    )
    cols_combo.pack(side=tk.LEFT, padx=5)

    # Filter by die coordinates
    tk.Label(control_frame, text="Filter:", font=("Helvetica", 10)).pack(side=tk.LEFT, padx=15)

    filter_var = tk.StringVar(value="")
    filter_entry = tk.Entry(control_frame, textvariable=filter_var, width=20, font=("Helvetica", 10))
    filter_entry.pack(side=tk.LEFT, padx=5)

    # Info label
    info_label = tk.Label(control_frame, text=f"Total: {len(image_files)} images", font=("Helvetica", 10))
    info_label.pack(side=tk.RIGHT, padx=10)

    # Create scrollable canvas for thumbnails
    canvas_frame = tk.Frame(gallery_win)
    canvas_frame.pack(side=tk.TOP, fill=tk.BOTH, expand=True, padx=10, pady=5)

    gallery_canvas = tk.Canvas(canvas_frame, bg="#e0e0e0")
    scrollbar_y = tk.Scrollbar(canvas_frame, orient=tk.VERTICAL, command=gallery_canvas.yview)
    scrollbar_x = tk.Scrollbar(canvas_frame, orient=tk.HORIZONTAL, command=gallery_canvas.xview)

    scrollable_frame = tk.Frame(gallery_canvas, bg="#e0e0e0")

    scrollable_frame.bind(
        "<Configure>",
        lambda e: gallery_canvas.configure(scrollregion=gallery_canvas.bbox("all"))
    )

    gallery_canvas.create_window((0, 0), window=scrollable_frame, anchor="nw")
    gallery_canvas.configure(yscrollcommand=scrollbar_y.set, xscrollcommand=scrollbar_x.set)

    scrollbar_y.pack(side=tk.RIGHT, fill=tk.Y)
    scrollbar_x.pack(side=tk.BOTTOM, fill=tk.X)
    gallery_canvas.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)

    # Enable mousewheel scrolling
    def on_mousewheel(event):
        gallery_canvas.yview_scroll(int(-1 * (event.delta / 120)), "units")

    gallery_canvas.bind_all("<MouseWheel>", on_mousewheel)

    # Store image references to prevent garbage collection
    gallery_image_refs = []

    def load_thumbnails():
        """Load and display thumbnails"""
        nonlocal gallery_image_refs

        # Clear existing
        for widget in scrollable_frame.winfo_children():
            widget.destroy()
        gallery_image_refs.clear()

        # Get settings
        thumb_size = thumb_sizes.get(thumb_size_var.get(), 150)
        filter_text = filter_var.get().strip().lower()

        # Determine columns
        cols_setting = cols_var.get()
        if cols_setting == "Auto":
            # Calculate based on window width and thumbnail size
            try:
                win_width = gallery_win.winfo_width() - 50
                num_cols = max(2, win_width // (thumb_size + 20))
            except:
                num_cols = 5
        else:
            num_cols = int(cols_setting)

        # Filter images
        filtered_images = image_files
        if filter_text:
            filtered_images = [f for f in image_files if filter_text in os.path.basename(f).lower()]

        info_label.config(text=f"Showing: {len(filtered_images)} of {len(image_files)} images")

        # Create thumbnails
        row = 0
        col = 0

        for img_path in filtered_images:
            try:
                # Load and resize image
                img = Image.open(img_path)
                img.thumbnail((thumb_size, thumb_size), Image.Resampling.LANCZOS)
                photo = ImageTk.PhotoImage(img)
                gallery_image_refs.append(photo)

                # Create frame for each thumbnail
                thumb_frame = tk.Frame(scrollable_frame, bg="#e0e0e0", padx=5, pady=5)
                thumb_frame.grid(row=row, column=col, padx=5, pady=5, sticky="n")

                # Image label
                img_label = tk.Label(thumb_frame, image=photo, bg="white", relief="solid", bd=1)
                img_label.pack()

                # Filename label (shortened)
                filename = os.path.basename(img_path)
                if len(filename) > 25:
                    display_name = filename[:22] + "..."
                else:
                    display_name = filename

                name_label = tk.Label(
                    thumb_frame,
                    text=display_name,
                    font=("Helvetica", 8),
                    bg="#e0e0e0",
                    wraplength=thumb_size
                )
                name_label.pack()

                # Bind click to show full image
                def show_full_image(path=img_path):
                    show_zoomed_die_image(path)

                img_label.bind("<Button-1>", lambda e, p=img_path: show_zoomed_die_image(p))
                img_label.bind("<Enter>", lambda e, lbl=img_label: lbl.config(relief="raised", bd=2))
                img_label.bind("<Leave>", lambda e, lbl=img_label: lbl.config(relief="solid", bd=1))

                # Move to next position
                col += 1
                if col >= num_cols:
                    col = 0
                    row += 1

            except Exception as e:
                print(f"Error loading thumbnail for {img_path}: {e}")
                continue

        # Update scroll region
        scrollable_frame.update_idletasks()
        gallery_canvas.configure(scrollregion=gallery_canvas.bbox("all"))

    # Bind refresh on settings change
    thumb_size_combo.bind("<<ComboboxSelected>>", lambda e: load_thumbnails())
    cols_combo.bind("<<ComboboxSelected>>", lambda e: load_thumbnails())
    filter_entry.bind("<Return>", lambda e: load_thumbnails())
    filter_entry.bind("<KeyRelease>", lambda e: gallery_win.after(500, load_thumbnails))

    # Refresh button
    refresh_btn = tk.Button(
        control_frame,
        text="Refresh",
        command=load_thumbnails,
        font=("Helvetica", 10)
    )
    refresh_btn.pack(side=tk.LEFT, padx=10)

    # Load initial thumbnails after window is displayed
    gallery_win.after(100, load_thumbnails)

    # Cleanup on close
    def on_gallery_close():
        gallery_canvas.unbind_all("<MouseWheel>")
        gallery_image_refs.clear()
        gallery_win.destroy()

    gallery_win.protocol("WM_DELETE_WINDOW", on_gallery_close)


def find_die_images(x_coord, y_coord, image_type_filter=None):
    """Find images matching the given die coordinates and optional type filter"""
    global die_image_directory, available_image_types

    if not die_image_directory:
        return []

    matching_images = []
    found_types = set(["All"])  # Always include "All" option

    try:
        import re

        # List all files in the directory
        for filename in os.listdir(die_image_directory):
            # Check for image file extensions
            if not filename.lower().endswith(('.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tif', '.tiff')):
                continue

            # Look for coordinate patterns in filename
            # Patterns: X19_Y46, X19-Y46, _X19_Y46_, Die_X19_Y46
            # Pattern 1: X##_Y## or X##-Y##
            pattern1 = rf'[_\-]X{x_coord}[_\-]Y{y_coord}[_\-\.]'
            # Pattern 2: _X##_Y##_ (with underscores)
            pattern2 = rf'_X{x_coord}_Y{y_coord}_'
            # Pattern 3: -X##-Y##- (with dashes)
            pattern3 = rf'-X{x_coord}-Y{y_coord}-'

            if (re.search(pattern1, filename, re.IGNORECASE) or
                re.search(pattern2, filename, re.IGNORECASE) or
                re.search(pattern3, filename, re.IGNORECASE)):

                # Extract image type from filename
                # Pattern: VPG-XHAIR-12FRAME or VPG-ALLON-12FRAME
                type_match = re.search(r'VPG-([A-Z0-9]+)-', filename, re.IGNORECASE)
                if type_match:
                    img_type = type_match.group(1).upper()
                    found_types.add(img_type)

                    # Apply filter if specified
                    if image_type_filter and image_type_filter != "All":
                        if img_type != image_type_filter.upper():
                            continue

                full_path = os.path.join(die_image_directory, filename)
                matching_images.append(full_path)
                print(f"Found matching image: {filename}")

        # Update available types in combobox
        available_image_types = sorted(list(found_types))
        if "All" in available_image_types:
            available_image_types.remove("All")
            available_image_types.insert(0, "All")

        # Update combobox values
        current_selection = image_type_var.get()
        image_type_combobox["values"] = available_image_types
        if current_selection not in available_image_types:
            image_type_var.set("All")

    except Exception as e:
        print(f"Error searching for die images: {e}")

    return matching_images


def refresh_die_images():
    """Refresh die images based on current filter selection"""
    global current_selected_die

    if current_selected_die is not None:
        x_coord, y_coord = current_selected_die
        display_die_images(x_coord, y_coord)



# Note: die_image_panel, plm panels, and related widgets are defined earlier
# (around line 1926) before heatmap_display_frame to ensure proper pack order

# Add bindings for refresh functions
image_type_combobox.bind("<<ComboboxSelected>>", lambda e: refresh_die_images())
plm_type_combobox.bind("<<ComboboxSelected>>", lambda e: refresh_plm_files())


def find_plm_files(x_coord, y_coord, plm_type_filter=None):
    """Find PLM files matching the given die coordinates and optional type filter"""
    global available_plm_types, plm_file_directory

    # First try the global plm_file_directory (auto-detected), then fall back to config_dirs
    plm_dir = plm_file_directory
    if not plm_dir or not os.path.exists(plm_dir):
        plm_dir = config_dirs["plm"].get() if "plm" in config_dirs else None
    print(f"  find_plm_files: plm_dir = {plm_dir}")

    if not plm_dir or not os.path.exists(plm_dir):
        print(f"  find_plm_files: PLM dir not set or doesn't exist")
        return []

    matching_files = []
    found_types = set(["All"])

    try:
        import re

        files_in_dir = os.listdir(plm_dir)
        print(f"  find_plm_files: {len(files_in_dir)} files in PLM directory")

        for filename in files_in_dir:
            # Check for PLM file extensions (common ones)
            if not filename.lower().endswith(('.plm', '.txt', '.csv', '.dat', '.xml', '.json')):
                continue

            # Look for coordinate patterns in filename
            # Pattern for Die_X##_Y## format (e.g., Die_X31_Y28)
            pattern1 = rf'Die_X{x_coord}_Y{y_coord}[_\.]'
            # Pattern for W##-X##-Y## format (e.g., W21-X31-Y28)
            pattern2 = rf'-X{x_coord}-Y{y_coord}[_\-]'
            # Pattern for generic underscore/dash separated coordinates
            pattern3 = rf'[_\-]X{x_coord}[_\-]Y{y_coord}[_\-\.]'

            match1 = re.search(pattern1, filename, re.IGNORECASE)
            match2 = re.search(pattern2, filename, re.IGNORECASE)
            match3 = re.search(pattern3, filename, re.IGNORECASE)

            if match1 or match2 or match3:

                # Extract PLM type from filename if present
                type_match = re.search(r'PLM-([A-Z0-9]+)-', filename, re.IGNORECASE)
                if type_match:
                    plm_type = type_match.group(1).upper()
                    found_types.add(plm_type)

                    if plm_type_filter and plm_type_filter != "All":
                        if plm_type != plm_type_filter.upper():
                            continue
                else:
                    # Try to extract type from file extension or other patterns
                    ext = os.path.splitext(filename)[1].upper().replace('.', '')
                    found_types.add(ext)

                    if plm_type_filter and plm_type_filter != "All":
                        if ext != plm_type_filter.upper():
                            continue

                full_path = os.path.join(plm_dir, filename)
                matching_files.append(full_path)
                print(f"Found matching PLM file: {filename}")

        # Update available types
        available_plm_types = sorted(list(found_types))
        if "All" in available_plm_types:
            available_plm_types.remove("All")
            available_plm_types.insert(0, "All")

        current_selection = plm_type_var.get()
        plm_type_combobox["values"] = available_plm_types
        if current_selection not in available_plm_types:
            plm_type_var.set("All")

    except Exception as e:
        print(f"Error searching for PLM files: {e}")

    return matching_files


def refresh_plm_files():
    """Refresh PLM files based on current filter selection"""
    global current_selected_die_plm

    if current_selected_die_plm is not None:
        x_coord, y_coord = current_selected_die_plm
        display_plm_files(x_coord, y_coord)


def display_plm_files(x_coord, y_coord):
    """Display PLM files for the selected die coordinates as heatmaps like Plot tab"""
    global plm_file_refs, current_selected_die_plm

    print(f"display_plm_files called for die ({x_coord}, {y_coord})")

    current_selected_die_plm = (x_coord, y_coord)

    # Clear existing content
    for widget in plm_scrollable_frame.winfo_children():
        widget.destroy()
    plm_file_refs.clear()

    # Get current filter
    current_filter = plm_type_var.get()
    print(f"  PLM filter: {current_filter}")

    # Find matching PLM files
    matching_files = find_plm_files(x_coord, y_coord, current_filter)
    print(f"  Found {len(matching_files)} matching PLM files")

    if not matching_files:
        filter_text = f" (filter: {current_filter})" if current_filter != "All" else ""
        no_plm_label = tk.Label(
            plm_scrollable_frame,
            text=f"No PLM files\nfor die ({x_coord}, {y_coord}){filter_text}",
            font=("Helvetica", 9),
            bg="#e8e8e8",
            fg="gray"
        )
        no_plm_label.pack(pady=10)
        return

    # Display file count
    count_text = f"{len(matching_files)} PLM file(s)"
    if current_filter != "All":
        count_text += f" [{current_filter}]"
    count_label = tk.Label(
        plm_scrollable_frame,
        text=count_text,
        font=("Helvetica", 8),
        bg="#e8e8e8",
        fg="gray"
    )
    count_label.pack(pady=2)

    # Display each PLM file as a heatmap (like Plot tab)
    for plm_path in matching_files:
        try:
            file_frame = tk.Frame(plm_scrollable_frame, bg="#e8e8e8")
            file_frame.pack(pady=3, padx=3, fill=tk.X)

            # Try to load PLM file as measurement data
            plm_data = load_plm_as_matrix(plm_path)

            if plm_data is not None and plm_data.size > 0:
                # Create small heatmap figure
                fig_plm, ax_plm = plt.subplots(figsize=(2.6, 2.0))
                fig_plm.patch.set_facecolor('#e8e8e8')

                im = ax_plm.imshow(plm_data, cmap='viridis', aspect='auto')

                # Extract filename for title
                filename = os.path.basename(plm_path)
                short_name = filename[:20] + "..." if len(filename) > 20 else filename
                ax_plm.set_title(short_name, fontsize=7, fontweight="bold")
                ax_plm.tick_params(axis='both', which='major', labelsize=5)

                # Add small colorbar
                cbar = fig_plm.colorbar(im, ax=ax_plm, fraction=0.046, pad=0.04)
                cbar.ax.tick_params(labelsize=5)

                fig_plm.tight_layout()

                # Embed in tkinter
                plm_canvas = FigureCanvasTkAgg(fig_plm, master=file_frame)
                plm_canvas.get_tk_widget().pack(fill=tk.BOTH, expand=True)
                plm_canvas.draw()

                # Keep reference to prevent garbage collection
                plm_file_refs.append(plm_canvas)

                # Make clickable to open zoomed view
                plm_canvas.get_tk_widget().bind("<Button-1>",
                    lambda e, path=plm_path, data=plm_data: show_zoomed_plm(path, data))
            else:
                # Fallback: show as text file link if can't parse as matrix
                ext = os.path.splitext(plm_path)[1].upper().replace('.', '')
                type_label = tk.Label(
                    file_frame,
                    text=f"[{ext}]",
                    font=("Helvetica", 8, "bold"),
                    bg="#e8e8e8",
                    fg="#9C27B0"
                )
                type_label.pack(side=tk.LEFT, padx=2)

                filename = os.path.basename(plm_path)
                if len(filename) > 25:
                    filename = filename[:22] + "..."

                file_label = tk.Label(
                    file_frame,
                    text=filename,
                    font=("Helvetica", 8),
                    bg="#e8e8e8",
                    fg="blue",
                    cursor="hand2"
                )
                file_label.pack(side=tk.LEFT, padx=2)
                file_label.bind("<Button-1>", lambda e, path=plm_path: open_plm_file(path))

        except Exception as e:
            print(f"Error displaying PLM file {plm_path}: {e}")

    print(f"Displayed {len(matching_files)} PLM files for die ({x_coord}, {y_coord})")


def load_plm_as_matrix(file_path):
    """Load PLM file as a numpy matrix for heatmap display"""
    try:
        ext = os.path.splitext(file_path)[1].lower()

        if ext in ('.csv',):
            try:
                data = np.genfromtxt(file_path, delimiter=',', skip_header=0)
                if np.isnan(data).all():
                    data = np.genfromtxt(file_path, delimiter=',', skip_header=1)
                return data
            except:
                pass

        elif ext in ('.txt', '.dat'):
            try:
                with open(file_path, 'r') as f:
                    lines = f.readlines()

                cols = None
                rows = None
                data_start = None

                for i, line in enumerate(lines):
                    stripped = line.strip()
                    if stripped.startswith('Columns:'):
                        try:
                            cols = int(stripped.split(':')[1].strip())
                        except:
                            pass
                    elif stripped.startswith('Rows:'):
                        try:
                            rows = int(stripped.split(':')[1].strip())
                        except:
                            pass
                    elif stripped.startswith('Export Quantity:'):
                        data_start = i + 1
                        break

                if data_start is None:
                    # Try alternative: look for first line with only numbers and commas
                    for i, line in enumerate(lines):
                        stripped = line.strip()
                        # Check if line looks like data (numbers separated by commas)
                        if stripped and ',' in stripped:
                            parts = stripped.split(',')
                            try:
                                # Try to parse first few values as floats
                                float(parts[0].strip())
                                float(parts[1].strip())
                                data_start = i
                                break
                            except:
                                pass

                    if data_start is None:
                        return None

                # Read all data lines
                raw_data = lines[data_start:]

                if not raw_data:
                    return None

                # Collect all values from all lines
                all_values = []
                for line in raw_data:
                    stripped = line.strip()
                    if stripped:
                        parts = stripped.split(',')
                        for p in parts:
                            p = p.strip()
                            if p:
                                try:
                                    all_values.append(float(p))
                                except ValueError:
                                    pass

                if not all_values:
                    return None

                # If we have cols and rows from header, use them
                if cols and rows:
                    expected_single = rows * cols
                    expected_paired = rows * cols * 2

                    if len(all_values) >= expected_paired:
                        # Paired data - take every other value
                        luminance_values = all_values[::2]
                        if len(luminance_values) >= expected_single:
                            return np.array(luminance_values[:expected_single]).reshape(rows, cols)
                    elif len(all_values) >= expected_single:
                        # Single values
                        return np.array(all_values[:expected_single]).reshape(rows, cols)
                else:
                    # Try to infer shape from data
                    # Count values per line to guess columns
                    first_data_line = raw_data[0].strip()
                    if first_data_line:
                        inferred_cols = len([p for p in first_data_line.split(',') if p.strip()])
                        if inferred_cols > 0 and len(all_values) % inferred_cols == 0:
                            inferred_rows = len(all_values) // inferred_cols
                            return np.array(all_values).reshape(inferred_rows, inferred_cols)

                return None

            except Exception as e:
                print(f"Error parsing PLM file {os.path.basename(file_path)}: {e}")
                return None

        return None
    except Exception as e:
        print(f"Error loading PLM file as matrix: {e}")
        return None


def show_zoomed_plm(file_path, plm_data):
    """Show zoomed view of PLM heatmap with mouse wheel zoom and statistics plots"""
    zoom_win = tk.Toplevel(main_win)
    filename = os.path.basename(file_path)
    zoom_win.title(f"PLM Heatmap: {filename}")
    zoom_win.geometry("1200x700")

    # Main container frame
    main_container = tk.Frame(zoom_win)
    main_container.pack(fill=tk.BOTH, expand=True)

    # Left frame for PLM heatmap
    left_frame = tk.Frame(main_container)
    left_frame.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)

    # Right frame for statistics (boxplot and probability plot)
    right_frame = tk.Frame(main_container, bg="#f0f0f0", width=280)
    right_frame.pack(side=tk.RIGHT, fill=tk.Y, padx=5, pady=5)
    right_frame.pack_propagate(False)

    # --- Normalized Profile Plot Section ---
    norm_profile_label = tk.Label(
        right_frame,
        text="Normalized Profile",
        font=("Helvetica", 10, "bold"),
        bg="#f0f0f0"
    )
    norm_profile_label.pack(pady=(5, 2))

    # Direction selection frame
    norm_frame = tk.Frame(right_frame, bg="#f0f0f0")
    norm_frame.pack(fill=tk.X, pady=(2, 5), padx=5)

    norm_label = tk.Label(
        norm_frame,
        text="Direction:",
        font=("Helvetica", 9),
        bg="#f0f0f0"
    )
    norm_label.pack(side=tk.LEFT, padx=2)

    plm_norm_var = tk.StringVar(value="X (Columns)")
    norm_combo = ttk.Combobox(
        norm_frame,
        textvariable=plm_norm_var,
        values=["X (Columns)", "Y (Rows)"],
        state="readonly",
        width=12,
        font=("Helvetica", 9)
    )
    norm_combo.pack(side=tk.LEFT, padx=5)

    # Container frame for normalized profile plot
    norm_plot_frame = tk.Frame(right_frame, bg="#f0f0f0")
    norm_plot_frame.pack(fill=tk.BOTH, expand=False, padx=5, pady=2, ipady=0)

    # Store figure reference for cleanup
    norm_profile_fig = [None]
    norm_profile_canvas = [None]

    def update_normalized_profile_plot():
        """Update the normalized profile line plot based on selected direction"""
        # Clear previous plot
        for widget in norm_plot_frame.winfo_children():
            widget.destroy()
        if norm_profile_fig[0] is not None:
            plt.close(norm_profile_fig[0])

        direction = plm_norm_var.get()

        # Create figure for normalized profile
        fig_norm, ax_norm = plt.subplots(figsize=(2.6, 2.0))
        fig_norm.patch.set_facecolor('#f0f0f0')
        norm_profile_fig[0] = fig_norm

        if direction == "X (Columns)":
            # Calculate mean of each column, then normalize so overall mean = 1
            col_means = np.nanmean(plm_data, axis=0)  # Mean along rows for each column
            overall_mean = np.nanmean(col_means)
            if overall_mean != 0:
                normalized_values = col_means / overall_mean
            else:
                normalized_values = col_means

            x_indices = np.arange(len(normalized_values))
            ax_norm.plot(x_indices, normalized_values, color='#2E86AB', linewidth=1.5, marker='o', markersize=2)
            ax_norm.axhline(y=1.0, color='#D32F2F', linestyle='--', linewidth=1, label='Mean (1.0)')
            ax_norm.set_xlabel("Column Index", fontsize=7)
            ax_norm.set_ylabel("Normalized Value", fontsize=7)
            ax_norm.set_title("Column Profile (norm. to mean)", fontsize=8, fontweight="bold")

            # Add min/max annotations
            min_idx = np.nanargmin(normalized_values)
            max_idx = np.nanargmax(normalized_values)
            ax_norm.annotate(f'Min: {normalized_values[min_idx]:.3f}',
                xy=(min_idx, normalized_values[min_idx]),
                xytext=(5, -10), textcoords='offset points',
                fontsize=5, color='blue')
            ax_norm.annotate(f'Max: {normalized_values[max_idx]:.3f}',
                xy=(max_idx, normalized_values[max_idx]),
                xytext=(5, 5), textcoords='offset points',
                fontsize=5, color='red')

        else:  # Y (Rows)
            # Calculate mean of each row, then normalize so overall mean = 1
            row_means = np.nanmean(plm_data, axis=1)  # Mean along columns for each row
            overall_mean = np.nanmean(row_means)
            if overall_mean != 0:
                normalized_values = row_means / overall_mean
            else:
                normalized_values = row_means

            y_indices = np.arange(len(normalized_values))
            ax_norm.plot(y_indices, normalized_values, color='#4CAF50', linewidth=1.5, marker='o', markersize=2)
            ax_norm.axhline(y=1.0, color='#D32F2F', linestyle='--', linewidth=1, label='Mean (1.0)')
            ax_norm.set_xlabel("Row Index", fontsize=7)
            ax_norm.set_ylabel("Normalized Value", fontsize=7)
            ax_norm.set_title("Row Profile (norm. to mean)", fontsize=8, fontweight="bold")

            # Add min/max annotations
            min_idx = np.nanargmin(normalized_values)
            max_idx = np.nanargmax(normalized_values)
            ax_norm.annotate(f'Min: {normalized_values[min_idx]:.3f}',
                xy=(min_idx, normalized_values[min_idx]),
                xytext=(5, -10), textcoords='offset points',
                fontsize=5, color='blue')
            ax_norm.annotate(f'Max: {normalized_values[max_idx]:.3f}',
                xy=(max_idx, normalized_values[max_idx]),
                xytext=(5, 5), textcoords='offset points',
                fontsize=5, color='red')

        ax_norm.tick_params(axis='both', which='major', labelsize=6)
        ax_norm.grid(True, alpha=0.3, linestyle='--', linewidth=0.5)
        ax_norm.legend(fontsize=5, loc='upper right')

        # Set y-axis to show variation around 1.0
        y_min = np.nanmin(normalized_values)
        y_max = np.nanmax(normalized_values)
        margin = (y_max - y_min) * 0.1 if y_max != y_min else 0.1
        ax_norm.set_ylim(y_min - margin, y_max + margin)

        fig_norm.tight_layout()

        canvas_norm = FigureCanvasTkAgg(fig_norm, master=norm_plot_frame)
        canvas_norm.get_tk_widget().pack(fill=tk.BOTH, expand=True)
        canvas_norm.draw()
        norm_profile_canvas[0] = canvas_norm

    # Bind dropdown selection
    norm_combo.bind("<<ComboboxSelected>>", lambda e: update_normalized_profile_plot())

    # Initial normalized profile plot
    update_normalized_profile_plot()

    # Separator after normalized profile
    ttk.Separator(right_frame, orient='horizontal').pack(fill=tk.X, padx=5, pady=5)

    # Create larger heatmap
    fig_zoom, ax_zoom = plt.subplots(figsize=(9, 7))

    im = ax_zoom.imshow(plm_data, cmap='viridis', aspect='auto')
    ax_zoom.set_title(f"PLM: {filename}", fontsize=12, fontweight="bold")
    ax_zoom.set_xlabel("Column Index", fontsize=10)
    ax_zoom.set_ylabel("Row Index", fontsize=10)

    cbar = fig_zoom.colorbar(im, ax=ax_zoom, fraction=0.024, pad=0.01)
    cbar.ax.tick_params(labelsize=9)

    # Add professional interactive range sliders on the colorbar
    data_min = np.nanmin(plm_data)
    data_max = np.nanmax(plm_data)
    data_range = data_max - data_min

    # Create triangular arrow markers pointing inward
    # Low limit slider (bottom) - triangle pointing up
    low_marker = cbar.ax.scatter(
        [0.5], [data_min],
        marker='^',  # Triangle pointing up
        s=200,
        c='#D32F2F',  # Material Red
        edgecolors='white',
        linewidths=1.5,
        zorder=15,
        clip_on=False
    )

    # High limit slider (top) - triangle pointing down
    high_marker = cbar.ax.scatter(
        [0.5], [data_max],
        marker='v',  # Triangle pointing down
        s=200,
        c='#1976D2',  # Material Blue
        edgecolors='white',
        linewidths=1.5,
        zorder=15,
        clip_on=False
    )

    # Add horizontal limit lines extending beyond colorbar
    low_line, = cbar.ax.plot(
        [-0.6, 1.6], [data_min, data_min],
        color='#D32F2F', linewidth=2.5,
        solid_capstyle='round', zorder=12,
        transform=cbar.ax.get_yaxis_transform()
    )

    high_line, = cbar.ax.plot(
        [-0.6, 1.6], [data_max, data_max],
        color='#1976D2', linewidth=2.5,
        solid_capstyle='round', zorder=12,
        transform=cbar.ax.get_yaxis_transform()
    )

    # Add value labels next to sliders
    low_text = cbar.ax.text(
        1.8, data_min, f'{data_min:.2f}',
        fontsize=7, fontweight='bold',
        color='#D32F2F', va='center', ha='left',
        transform=cbar.ax.get_yaxis_transform(),
        zorder=20
    )

    high_text = cbar.ax.text(
        1.8, data_max, f'{data_max:.2f}',
        fontsize=7, fontweight='bold',
        color='#1976D2', va='center', ha='left',
        transform=cbar.ax.get_yaxis_transform(),
        zorder=20
    )

    # Store references and state for PLM sliders
    plm_slider_state = {
        'low_val': data_min,
        'high_val': data_max,
        'data_min': data_min,
        'data_max': data_max,
        'dragging': None,
        'low_marker': low_marker,
        'high_marker': high_marker,
        'low_line': low_line,
        'high_line': high_line,
        'low_text': low_text,
        'high_text': high_text,
        'cbar_ax': cbar.ax,
        'im': im
    }

    def on_plm_slider_press(event):
        if event.inaxes == plm_slider_state['cbar_ax']:
            y_click = event.ydata
            if y_click is not None:
                # Calculate distance to each slider in data coordinates
                dist_to_low = abs(y_click - plm_slider_state['low_val'])
                dist_to_high = abs(y_click - plm_slider_state['high_val'])

                # Threshold for grabbing (10% of data range)
                threshold = 0.1 * (plm_slider_state['data_max'] - plm_slider_state['data_min'])

                # Select the closest slider if within threshold
                if dist_to_low < dist_to_high and dist_to_low < threshold:
                    plm_slider_state['dragging'] = 'low'
                    low_marker.set_sizes([300])  # Visual feedback - enlarge
                elif dist_to_high < threshold:
                    plm_slider_state['dragging'] = 'high'
                    high_marker.set_sizes([300])  # Visual feedback - enlarge

                fig_zoom.canvas.draw_idle()

    def on_plm_slider_motion(event):
        if plm_slider_state['dragging']:
            # Allow motion even outside colorbar axes for smooth dragging
            if event.ydata is not None:
                y_click = event.ydata
            elif event.y is not None:
                # Convert pixel to data coordinates if outside axes
                inv = plm_slider_state['cbar_ax'].transData.inverted()
                _, y_click = inv.transform((event.x, event.y))
            else:
                return

            # Clamp to data range
            y_click = max(plm_slider_state['data_min'], min(plm_slider_state['data_max'], y_click))

            min_gap = 0.02 * (plm_slider_state['data_max'] - plm_slider_state['data_min'])

            if plm_slider_state['dragging'] == 'low':
                # Low slider can't go above high slider
                new_val = min(y_click, plm_slider_state['high_val'] - min_gap)
                new_val = max(new_val, plm_slider_state['data_min'])
                plm_slider_state['low_val'] = new_val

                # Update marker position
                low_marker.set_offsets([[0.5, new_val]])
                # Update line position
                low_line.set_ydata([new_val, new_val])
                # Update text
                low_text.set_position((1.8, new_val))
                low_text.set_text(f'{new_val:.2f}')

            elif plm_slider_state['dragging'] == 'high':
                # High slider can't go below low slider
                new_val = max(y_click, plm_slider_state['low_val'] + min_gap)
                new_val = min(new_val, plm_slider_state['data_max'])
                plm_slider_state['high_val'] = new_val

                # Update marker position
                high_marker.set_offsets([[0.5, new_val]])
                # Update line position
                high_line.set_ydata([new_val, new_val])
                # Update text
                high_text.set_position((1.8, new_val))
                high_text.set_text(f'{new_val:.2f}')

            # Update colormap limits
            plm_slider_state['im'].set_clim(vmin=plm_slider_state['low_val'], vmax=plm_slider_state['high_val'])

            fig_zoom.canvas.draw_idle()

    def on_plm_slider_release(event):
        if plm_slider_state['dragging'] == 'low':
            low_marker.set_sizes([200])  # Reset size
        elif plm_slider_state['dragging'] == 'high':
            high_marker.set_sizes([200])  # Reset size
        plm_slider_state['dragging'] = None
        fig_zoom.canvas.draw_idle()

    # Add statistics text overlay
    stats_text = (
        f"Min: {np.nanmin(plm_data):.4f}\n"
        f"Max: {np.nanmax(plm_data):.4f}\n"
        f"Mean: {np.nanmean(plm_data):.4f}\n"
        f"Shape: {plm_data.shape}\n"
        f"Use mouse wheel to zoom\n"
        f"Drag sliders to adjust range"
    )
    ax_zoom.text(
        0.02, 0.98, stats_text,
        transform=ax_zoom.transAxes,
        fontsize=9,
        verticalalignment='top',
        bbox=dict(boxstyle='round', facecolor='wheat', alpha=0.8)
    )

    fig_zoom.tight_layout()

    # Embed heatmap in left frame
    canvas_zoom = FigureCanvasTkAgg(fig_zoom, master=left_frame)
    canvas_zoom.get_tk_widget().pack(fill=tk.BOTH, expand=True)

    # Connect slider events
    fig_zoom.canvas.mpl_connect('button_press_event', on_plm_slider_press)
    fig_zoom.canvas.mpl_connect('motion_notify_event', on_plm_slider_motion)
    fig_zoom.canvas.mpl_connect('button_release_event', on_plm_slider_release)

    # Add toolbar
    toolbar = NavigationToolbar2Tk(canvas_zoom, left_frame)
    toolbar.update()

    # --- Right side: Statistics plots ---

    # Flatten PLM data for statistics (remove NaN values)
    flat_data = plm_data.flatten()
    valid_data = flat_data[~np.isnan(flat_data)]

    # Boxplot section
    boxplot_label = tk.Label(
        right_frame,
        text="Boxplot",
        font=("Helvetica", 10, "bold"),
        bg="#f0f0f0"
    )
    boxplot_label.pack(pady=(5, 2))

    # Create boxplot figure - professional styling
    fig_box, ax_box = plt.subplots(figsize=(2.6, 2.4))
    fig_box.patch.set_facecolor('white')

    bp = ax_box.boxplot(
        [valid_data],
        tick_labels=["PLM"],
        vert=True,
        patch_artist=True,
        showmeans=True,
        widths=0.6,
        meanprops=dict(marker="D", markerfacecolor="#E74C3C", markeredgecolor="white", markersize=5, markeredgewidth=1),
        medianprops=dict(color="#2C3E50", linewidth=2),
        whiskerprops=dict(color="#2C3E50", linewidth=1.5, linestyle='-'),
        capprops=dict(color="#2C3E50", linewidth=1.5),
        flierprops=dict(marker='o', markerfacecolor='#95A5A6', markeredgecolor='#7F8C8D', markersize=3, alpha=0.6),
        boxprops=dict(linewidth=1.5)
    )

    # Professional color for box
    for patch in bp['boxes']:
        patch.set_facecolor('#3498DB')
        patch.set_alpha(0.75)
        patch.set_edgecolor('#2C3E50')

    # Calculate statistics
    stats_max = np.max(valid_data)
    stats_min = np.min(valid_data)
    stats_q1 = np.percentile(valid_data, 25)
    stats_q3 = np.percentile(valid_data, 75)
    stats_mean = np.mean(valid_data)
    stats_median = np.median(valid_data)

    # Add statistics box in top left corner
    stats_text = (
        f"Max: {stats_max:.3g}\n"
        f"Min: {stats_min:.3g}\n"
        f"Q1: {stats_q1:.3g}\n"
        f"Q3: {stats_q3:.3g}\n"
        f"Mean: {stats_mean:.3g}\n"
        f"Median: {stats_median:.3g}"
    )
    ax_box.text(0.02, 0.98, stats_text,
        transform=ax_box.transAxes,
        fontsize=5,
        fontweight='normal',
        fontfamily='monospace',
        verticalalignment='top',
        horizontalalignment='left',
        bbox=dict(boxstyle='round,pad=0.3', facecolor='white', edgecolor='#BDC3C7', alpha=0.9, linewidth=1)
    )

    ax_box.set_title("PLM Values", fontsize=8, fontweight="bold", color='#2C3E50')
    ax_box.tick_params(axis='both', which='major', labelsize=6, colors='#2C3E50')
    ax_box.set_ylabel("Value", fontsize=7, color='#2C3E50')
    ax_box.set_facecolor('#FAFAFA')
    ax_box.grid(True, alpha=0.4, linestyle='-', linewidth=0.5, color='#BDC3C7')
    ax_box.spines['top'].set_visible(False)
    ax_box.spines['right'].set_visible(False)
    ax_box.spines['left'].set_color('#BDC3C7')
    ax_box.spines['bottom'].set_color('#BDC3C7')

    fig_box.tight_layout()

    canvas_box = FigureCanvasTkAgg(fig_box, master=right_frame)
    canvas_box.get_tk_widget().pack(fill=tk.X, padx=5, pady=2)
    canvas_box.draw()

    # Probability plot section with dropdown
    prob_header_frame = tk.Frame(right_frame, bg="#f0f0f0")
    prob_header_frame.pack(fill=tk.X, pady=(10, 2))

    prob_label = tk.Label(
        prob_header_frame,
        text="Distribution:",
        font=("Helvetica", 10, "bold"),
        bg="#f0f0f0"
    )
    prob_label.pack(side=tk.LEFT, padx=2)

    # Dropdown to choose between CDF and PDF
    plm_prob_type_var = tk.StringVar(value="CDF")
    prob_type_combo = ttk.Combobox(
        prob_header_frame,
        textvariable=plm_prob_type_var,
        values=["CDF", "PDF"],
        state="readonly",
        width=6,
        font=("Helvetica", 9)
    )
    prob_type_combo.pack(side=tk.LEFT, padx=5)

    # Container for probability plot
    prob_plot_frame = tk.Frame(right_frame, bg="#f0f0f0")
    prob_plot_frame.pack(fill=tk.BOTH, expand=True, padx=5, pady=2)

    # Store figure reference for cleanup
    plm_prob_fig = [None]
    plm_prob_canvas = [None]

    def update_plm_prob_plot():
        # Clear previous plot
        for widget in prob_plot_frame.winfo_children():
            widget.destroy()
        if plm_prob_fig[0] is not None:
            plt.close(plm_prob_fig[0])

        # Create new probability figure - professional styling
        fig_prob, ax_prob = plt.subplots(figsize=(2.6, 2.4))
        fig_prob.patch.set_facecolor('white')
        plm_prob_fig[0] = fig_prob

        plot_type = plm_prob_type_var.get()

        if plot_type == "CDF":
            # CDF plot - professional style
            sorted_data = np.sort(valid_data)
            cumulative = np.arange(1, len(sorted_data) + 1) / len(sorted_data)
            ax_prob.plot(sorted_data, cumulative, color='#3498DB', linewidth=2)
            ax_prob.fill_between(sorted_data, cumulative, alpha=0.15, color='#3498DB')
            ax_prob.set_ylabel("Cumulative Prob.", fontsize=7, color='#2C3E50')
            ax_prob.set_title("CDF", fontsize=8, fontweight="bold", color='#2C3E50')
        else:
            # PDF plot (histogram) - professional style
            ax_prob.hist(valid_data, bins=30, density=True, color='#3498DB', alpha=0.6, edgecolor='white', linewidth=0.5)
            ax_prob.set_ylabel("Density", fontsize=7, color='#2C3E50')
            ax_prob.set_title("PDF", fontsize=8, fontweight="bold", color='#2C3E50')

        ax_prob.set_xlabel("Value", fontsize=7, color='#2C3E50')
        ax_prob.tick_params(axis='both', which='major', labelsize=6, colors='#2C3E50')
        ax_prob.set_facecolor('#FAFAFA')
        ax_prob.grid(True, alpha=0.4, linestyle='-', linewidth=0.5, color='#BDC3C7')
        ax_prob.spines['top'].set_visible(False)
        ax_prob.spines['right'].set_visible(False)
        ax_prob.spines['left'].set_color('#BDC3C7')
        ax_prob.spines['bottom'].set_color('#BDC3C7')

        fig_prob.tight_layout()

        canvas_prob = FigureCanvasTkAgg(fig_prob, master=prob_plot_frame)
        canvas_prob.get_tk_widget().pack(fill=tk.BOTH, expand=True)
        canvas_prob.draw()
        plm_prob_canvas[0] = canvas_prob

    # Bind dropdown selection
    prob_type_combo.bind("<<ComboboxSelected>>", lambda e: update_plm_prob_plot())

    # Initial probability plot
    update_plm_prob_plot()

    # Mouse wheel zoom function for heatmap
    def on_scroll_zoom(event):
        if event.inaxes != ax_zoom:
            return

        cur_xlim = ax_zoom.get_xlim()
        cur_ylim = ax_zoom.get_ylim()

        x_data = event.xdata
        y_data = event.ydata

        if x_data is None or y_data is None:
            return

        if event.button == 'up':
            zoom_factor = 0.8
        elif event.button == 'down':
            zoom_factor = 1.25
        else:
            return

        x_range = cur_xlim[1] - cur_xlim[0]
        y_range = cur_ylim[1] - cur_ylim[0]

        new_x_range = x_range * zoom_factor
        new_y_range = y_range * zoom_factor

        rel_x = (x_data - cur_xlim[0]) / x_range
        rel_y = (y_data - cur_ylim[0]) / y_range

        new_xlim = [x_data - rel_x * new_x_range, x_data + (1 - rel_x) * new_x_range]
        new_ylim = [y_data - rel_y * new_y_range, y_data + (1 - rel_y) * new_y_range]

        ax_zoom.set_xlim(new_xlim)
        ax_zoom.set_ylim(new_ylim)

        canvas_zoom.draw()

    def on_tk_scroll(event):
        class FakeEvent:
            pass
        fake_event = FakeEvent()
        fake_event.inaxes = ax_zoom
        fake_event.xdata = (ax_zoom.get_xlim()[0] + ax_zoom.get_xlim()[1]) / 2
        fake_event.ydata = (ax_zoom.get_ylim()[0] + ax_zoom.get_ylim()[1]) / 2
        fake_event.button = 'up' if event.delta > 0 else 'down'
        on_scroll_zoom(fake_event)

    canvas_zoom.mpl_connect('scroll_event', on_scroll_zoom)
    canvas_zoom.get_tk_widget().bind('<MouseWheel>', on_tk_scroll)

    canvas_zoom.draw()

    # Button frame
    btn_frame = tk.Frame(zoom_win)
    btn_frame.pack(side=tk.BOTTOM, pady=5)

    def reset_zoom():
        rows, cols = plm_data.shape
        ax_zoom.set_xlim(-0.5, cols - 0.5)
        ax_zoom.set_ylim(rows - 0.5, -0.5)
        canvas_zoom.draw()

    tk.Button(
        btn_frame,
        text="Reset Zoom",
        command=reset_zoom,
        font=("Helvetica", 10)
    ).pack(side=tk.LEFT, padx=5)

    tk.Button(
        btn_frame,
        text="View Raw Text",
        command=lambda: open_plm_file(file_path),
        font=("Helvetica", 10)
    ).pack(side=tk.LEFT, padx=5)

    # Grid on/off checkbox
    plm_grid_var = tk.BooleanVar(value=False)

    def toggle_plm_grid():
        if plm_grid_var.get():
            # Turn grid on
            ax_zoom.grid(True, color='white', linestyle='-', linewidth=0.5, alpha=0.7)
        else:
            # Turn grid off
            ax_zoom.grid(False)
        canvas_zoom.draw()

    grid_checkbox = tk.Checkbutton(
        btn_frame,
        text="Show Grid",
        variable=plm_grid_var,
        command=toggle_plm_grid,
        font=("Helvetica", 10)
    )
    grid_checkbox.pack(side=tk.LEFT, padx=10)

    def on_close():
        plt.close(fig_zoom)
        plt.close(fig_box)
        if plm_prob_fig[0] is not None:
            plt.close(plm_prob_fig[0])
        if norm_profile_fig[0] is not None:
            plt.close(norm_profile_fig[0])
        zoom_win.destroy()

    zoom_win.protocol("WM_DELETE_WINDOW", on_close)


def open_plm_file(file_path):
    """Open and display PLM file contents"""
    plm_win = tk.Toplevel(main_win)
    filename = os.path.basename(file_path)
    plm_win.title(f"PLM File: {filename}")
    plm_win.geometry("700x500")

    # Text widget with scrollbars
    text_frame = tk.Frame(plm_win)
    text_frame.pack(fill=tk.BOTH, expand=True, padx=10, pady=10)

    v_scrollbar = tk.Scrollbar(text_frame, orient=tk.VERTICAL)
    h_scrollbar = tk.Scrollbar(text_frame, orient=tk.HORIZONTAL)

    text_widget = tk.Text(
        text_frame,
        wrap=tk.NONE,
        xscrollcommand=h_scrollbar.set,
        yscrollcommand=v_scrollbar.set,
        font=("Courier", 9)
    )

    v_scrollbar.config(command=text_widget.yview)
    h_scrollbar.config(command=text_widget.xview)

    v_scrollbar.pack(side=tk.RIGHT, fill=tk.Y)
    h_scrollbar.pack(side=tk.BOTTOM, fill=tk.X)
    text_widget.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)

    try:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            content = f.read()
        text_widget.insert("1.0", content)
    except Exception as e:
        text_widget.insert("1.0", f"Error reading file: {e}")

    text_widget.config(state=tk.DISABLED)

    # Info bar
    info_bar = tk.Label(
        plm_win,
        text=f"File: {file_path}",
        font=("Helvetica", 8),
        fg="gray",
        anchor="w"
    )
    info_bar.pack(side=tk.BOTTOM, fill=tk.X, padx=10, pady=5)


def display_die_images(x_coord, y_coord):
    """Display images for the selected die coordinates"""
    global die_image_refs, current_selected_die

    # Store current die coordinates for filter refresh
    current_selected_die = (x_coord, y_coord)

    # Clear existing images
    for widget in die_image_scrollable_frame.winfo_children():
        widget.destroy()
    die_image_refs.clear()

    # Get current filter selection
    current_filter = image_type_var.get()

    # Find matching images with filter
    matching_images = find_die_images(x_coord, y_coord, current_filter)

    if not matching_images:
        filter_text = f" (filter: {current_filter})" if current_filter != "All" else ""
        no_image_label = tk.Label(
            die_image_scrollable_frame,
            text=f"No images found\nfor die ({x_coord}, {y_coord}){filter_text}",
            font=("Helvetica", 10),
            bg="#e8e8e8",
            fg="gray"
        )
        no_image_label.pack(pady=20)
        return

    # Display coordinate label
    coord_label = tk.Label(
        die_image_scrollable_frame,
        text=f"Die ({x_coord}, {y_coord})",
        font=("Helvetica", 11, "bold"),
        bg="#e8e8e8"
    )
    coord_label.pack(pady=5)

    # Display image count
    count_text = f"{len(matching_images)} image(s)"
    if current_filter != "All":
        count_text += f" [{current_filter}]"
    count_label = tk.Label(
        die_image_scrollable_frame,
        text=count_text,
        font=("Helvetica", 9),
        bg="#e8e8e8",
        fg="gray"
    )
    count_label.pack(pady=2)

    # Display each image
    for img_path in matching_images:
        try:
            # Create frame for image and label
            img_frame = tk.Frame(die_image_scrollable_frame, bg="#e8e8e8")
            img_frame.pack(pady=5, padx=5, fill=tk.X)

            # Extract image type for label
            import re
            type_match = re.search(r'VPG-([A-Z0-9]+)-', os.path.basename(img_path), re.IGNORECASE)
            img_type = type_match.group(1).upper() if type_match else "Unknown"

            # Type label
            type_label = tk.Label(
                img_frame,
                text=f"[{img_type}]",
                font=("Helvetica", 9, "bold"),
                bg="#e8e8e8",
                fg="#2196F3"
            )
            type_label.pack()

            # Load and resize image - handle TIF files specially
            img = Image.open(img_path)

            # Convert TIF/TIFF to RGB if needed (handles various color modes)
            if img.mode in ('I;16', 'I;16B', 'I;16L', 'I;16N', 'I', 'F'):
                # 16-bit or float images - normalize to 8-bit
                img_array = np.array(img, dtype=np.float32)
                img_array = (img_array - img_array.min()) / (img_array.max() - img_array.min() + 1e-10) * 255
                img = Image.fromarray(img_array.astype(np.uint8))

            if img.mode not in ('RGB', 'RGBA'):
                img = img.convert('RGB')

            # Calculate thumbnail size maintaining aspect ratio
            max_width = 280
            max_height = 180
            img.thumbnail((max_width, max_height), Image.Resampling.LANCZOS)

            tk_img = ImageTk.PhotoImage(img, master=main_win)
            die_image_refs.append(tk_img)

            # Image label (clickable)
            img_label = tk.Label(img_frame, image=tk_img, bg="#e8e8e8", cursor="hand2")
            img_label.pack()
            img_label.bind("<Button-1>", lambda e, path=img_path: show_zoomed_die_image(path))

            # Filename label (truncated if too long)
            filename = os.path.basename(img_path)
            if len(filename) > 40:
                filename = filename[:37] + "..."
            name_label = tk.Label(
                img_frame,
                text=filename,
                font=("Helvetica", 8),
                bg="#e8e8e8",
                fg="gray"
            )
            name_label.pack()

        except Exception as e:
            print(f"Error loading image {img_path}: {e}")
            error_label = tk.Label(
                die_image_scrollable_frame,
                text=f"Error loading: {os.path.basename(img_path)}",
                font=("Helvetica", 8),
                bg="#e8e8e8",
                fg="red"
            )
            error_label.pack()

    print(f"Displayed {len(matching_images)} images for die ({x_coord}, {y_coord})")


def show_zoomed_die_image(image_path):
    """Show zoomed view of die image in a new window with mouse wheel zoom"""
    zoom_win = tk.Toplevel(main_win)
    filename = os.path.basename(image_path)
    zoom_win.title(f"Die Image: {filename}")
    zoom_win.geometry("900x700")

    # Load image
    img = Image.open(image_path)

    # Convert TIF/TIFF to RGB if needed (handles various color modes)
    if img.mode in ('I;16', 'I;16B', 'I;16L', 'I;16N', 'I', 'F'):
        # 16-bit or float images - normalize to 8-bit
        img_array = np.array(img, dtype=np.float32)
        img_array = (img_array - img_array.min()) / (img_array.max() - img_array.min() + 1e-10) * 255
        img = Image.fromarray(img_array.astype(np.uint8))

    if img.mode not in ('RGB', 'RGBA'):
        img = img.convert('RGB')

    orig_width, orig_height = img.size
    img_array = np.array(img)

    # Create figure for the image - clean display without axis
    fig_zoom, ax_zoom = plt.subplots(figsize=(10, 8))

    # Display image without axis
    ax_zoom.imshow(img_array)
    ax_zoom.set_title(f"{filename}", fontsize=11)

    # Hide axis for clean image display
    ax_zoom.axis('off')

    fig_zoom.tight_layout()

    # Embed in window
    canvas_zoom = FigureCanvasTkAgg(fig_zoom, master=zoom_win)
    canvas_zoom.get_tk_widget().pack(fill=tk.BOTH, expand=True)

    # Mouse wheel zoom function
    def on_scroll_zoom(event):
        if event.inaxes != ax_zoom:
            return

        # Get current limits
        cur_xlim = ax_zoom.get_xlim()
        cur_ylim = ax_zoom.get_ylim()

        x_data = event.xdata
        y_data = event.ydata

        if x_data is None or y_data is None:
            return

        # Determine zoom direction
        if event.button == 'up':
            zoom_factor = 0.8  # Zoom in
        elif event.button == 'down':
            zoom_factor = 1.25  # Zoom out
        else:
            return

        # Calculate new limits
        x_range = cur_xlim[1] - cur_xlim[0]
        y_range = cur_ylim[1] - cur_ylim[0]

        new_x_range = x_range * zoom_factor
        new_y_range = y_range * zoom_factor

        # Keep mouse position fixed
        rel_x = (x_data - cur_xlim[0]) / x_range
        rel_y = (y_data - cur_ylim[0]) / y_range

        new_xlim = [x_data - rel_x * new_x_range, x_data + (1 - rel_x) * new_x_range]
        new_ylim = [y_data - rel_y * new_y_range, y_data + (1 - rel_y) * new_y_range]

        ax_zoom.set_xlim(new_xlim)
        ax_zoom.set_ylim(new_ylim)

        canvas_zoom.draw()

    # Windows mouse wheel handler
    def on_tk_scroll(event):
        # Create a fake matplotlib event
        class FakeEvent:
            pass
        fake_event = FakeEvent()
        fake_event.inaxes = ax_zoom
        fake_event.xdata = (ax_zoom.get_xlim()[0] + ax_zoom.get_xlim()[1]) / 2
        fake_event.ydata = (ax_zoom.get_ylim()[0] + ax_zoom.get_ylim()[1]) / 2
        fake_event.button = 'up' if event.delta > 0 else 'down'
        on_scroll_zoom(fake_event)

    # Connect scroll events
    canvas_zoom.mpl_connect('scroll_event', on_scroll_zoom)
    canvas_zoom.get_tk_widget().bind('<MouseWheel>', on_tk_scroll)  # Windows

    canvas_zoom.draw()

    # Button frame
    btn_frame = tk.Frame(zoom_win)
    btn_frame.pack(side=tk.BOTTOM, pady=5)

    # Reset zoom button
    def reset_zoom():
        ax_zoom.set_xlim(-0.5, orig_width - 0.5)
        ax_zoom.set_ylim(orig_height - 0.5, -0.5)
        canvas_zoom.draw()

    tk.Button(
        btn_frame,
        text="Reset Zoom",
        command=reset_zoom,
        font=("Helvetica", 10)
    ).pack(side=tk.LEFT, padx=5)

    # Open in default viewer button
    def open_external():
        import subprocess
        try:
            os.startfile(image_path)  # Windows
        except AttributeError:
            subprocess.call(['xdg-open', image_path])  # Linux
        except Exception as e:
            print(f"Error opening file externally: {e}")

    tk.Button(
        btn_frame,
        text="Open in External Viewer",
        command=open_external,
        font=("Helvetica", 10)
    ).pack(side=tk.LEFT, padx=5)

    # Grid on/off checkbox for image
    img_grid_var = tk.BooleanVar(value=False)

    def toggle_image_grid():
        if img_grid_var.get():
            # Turn grid on - show axis and grid
            ax_zoom.axis('on')
            ax_zoom.grid(True, color='white', linestyle='-', linewidth=0.5, alpha=0.7)
            ax_zoom.tick_params(axis='both', which='major', labelsize=8)
        else:
            # Turn grid off - hide axis
            ax_zoom.grid(False)
            ax_zoom.axis('off')
        canvas_zoom.draw()

    grid_checkbox = tk.Checkbutton(
        btn_frame,
        text="Show Grid",
        variable=img_grid_var,
        command=toggle_image_grid,
        font=("Helvetica", 10)
    )
    grid_checkbox.pack(side=tk.LEFT, padx=10)

    # Image info label
    info_label = tk.Label(
        zoom_win,
        text=f"Size: {orig_width} x {orig_height} | Use mouse wheel to zoom",
        font=("Helvetica", 9),
        fg="gray"
    )
    info_label.pack(side=tk.BOTTOM, pady=2)

    def on_close():
        plt.close(fig_zoom)
        zoom_win.destroy()

    zoom_win.protocol("WM_DELETE_WINDOW", on_close)


# Configuration tab layout (already created above as tab7)

# Configuration variables
config_dirs = {
    "stdf": tk.StringVar(value=""),
    "images": tk.StringVar(value=""),
    "csv": tk.StringVar(value=""),
    "txt": tk.StringVar(value=""),
    "plm": tk.StringVar(value="")
}

# Configuration file path
import json
CONFIG_FILE = os.path.join(os.path.expanduser("~"), ".wafermap_config.json")


def load_config():
    """Load configuration from file"""
    global config_dirs
    try:
        if os.path.exists(CONFIG_FILE):
            with open(CONFIG_FILE, 'r') as f:
                saved_config = json.load(f)

            for key, value in saved_config.items():
                if key in config_dirs:
                    config_dirs[key].set(value)

            print(f"Configuration loaded from {CONFIG_FILE}")
            update_config_status()
    except Exception as e:
        print(f"Error loading configuration: {e}")


def save_config():
    """Save configuration to file"""
    global config_dirs
    try:
        config_data = {key: var.get() for key, var in config_dirs.items()}

        with open(CONFIG_FILE, 'w') as f:
            json.dump(config_data, f, indent=2)

        print(f"Configuration saved to {CONFIG_FILE}")
        config_status_label.config(text="Configuration saved!", fg="green")
    except Exception as e:
        print(f"Error saving configuration: {e}")
        config_status_label.config(text=f"Error: {e}", fg="red")


def browse_directory(dir_type):
    """Browse for a directory and update the corresponding config"""
    folder_path = filedialog.askdirectory(
        title=f"Select {dir_type.upper()} Directory"
    )

    if folder_path:
        config_dirs[dir_type].set(folder_path)
        update_config_status()
        print(f"Set {dir_type} directory to: {folder_path}")


def update_config_status():
    """Update the status of configured directories"""
    configured = sum(1 for var in config_dirs.values() if var.get())
    config_status_label.config(
        text=f"{configured}/5 directories configured",
        fg="green" if configured == 5 else "orange" if configured > 0 else "gray"
    )


def apply_config_and_load():
    """Apply configuration and load data from configured directories"""
    global die_image_directory, multiple_stdf_data, multiple_wafer_ids, test_parameters

    print(f"\n{'='*60}")
    print("Applying configuration and loading data...")
    print(f"{'='*60}")

    # Set image directory
    images_dir = config_dirs["images"].get()
    if images_dir and os.path.exists(images_dir):
        die_image_directory = images_dir
        folder_status_label.config(text=f"Img: {os.path.basename(images_dir)}", fg="green")
        print(f"Set image directory: {images_dir}")

    # Load STDF files
    stdf_dir = config_dirs["stdf"].get()
    if stdf_dir and os.path.exists(stdf_dir):
        stdf_files = [f for f in os.listdir(stdf_dir)
                      if f.lower().endswith(('.stdf', '.std'))]

        if stdf_files:
            print(f"Found {len(stdf_files)} STDF file(s)")
            multiple_stdf_data = []
            multiple_wafer_ids = []

            for stdf_file in stdf_files:
                stdf_path = os.path.join(stdf_dir, stdf_file)
                print(f"  Loading: {stdf_file}")

                df, wafer_id, test_params, grouped_params = read_wafermap_from_stdf(stdf_path)

                if not df.empty:
                    multiple_stdf_data.append(df)
                    multiple_wafer_ids.append(wafer_id if wafer_id else stdf_file)

                    if not test_parameters:
                        test_parameters = test_params
                    if not grouped_parameters:
                        grouped_parameters = grouped_params

            # Update parameter combobox
            if multiple_stdf_data:
                param_options = ["BIN (Bin Number)"]
                for test_key, test_name in sorted(test_parameters.items()):
                    param_options.append(f"{test_key}: {test_name}")

                heatmap_param_combobox["values"] = param_options
                if param_options:
                    heatmap_param_combobox.current(0)

                heatmap_info_label.config(text=f"Loaded {len(multiple_stdf_data)} STDF files")

                # Update heatmap display
                update_multi_stdf_heatmap()

    config_status_label.config(text="Configuration applied!", fg="green")
    print(f"\n{'='*60}")
    print("Configuration applied successfully!")
    print(f"{'='*60}\n")


# Configuration tab header
config_header = tk.Label(
    tab7,
    text="Directory Configuration",
    font=("Helvetica", 16, "bold")
)
config_header.pack(pady=20)

config_description = tk.Label(
    tab7,
    text="Configure directories for different file types. These can be local or network paths.",
    font=("Helvetica", 10),
    fg="gray"
)
config_description.pack(pady=5)

# Frame for directory settings
config_frame = tk.Frame(tab7)
config_frame.pack(fill=tk.BOTH, expand=True, padx=40, pady=20)

# Directory configuration entries
dir_configs = [
    ("stdf", "STDF Files (STDDatalog)", "Binary STDF test data files (.stdf, .std)"),
    ("images", "Image Files (ImageCaptures)", "Die images (.jpg, .png, .tif, etc.)"),
    ("csv", "CSV Files", "Comma-separated value files (.csv)"),
    ("txt", "TXT Files (TXTDatalog)", "Text datalog files (.txt, .atdf)"),
    ("plm", "PLM Files", "PLM data files")
]

for idx, (key, label, description) in enumerate(dir_configs):
    # Create frame for each directory setting
    dir_frame = tk.Frame(config_frame)
    dir_frame.pack(fill=tk.X, pady=10)

    # Label
    dir_label = tk.Label(
        dir_frame,
        text=label,
        font=("Helvetica", 11, "bold"),
        width=25,
        anchor="w"
    )
    dir_label.pack(side=tk.LEFT, padx=5)

    # Entry for path
    dir_entry = tk.Entry(
        dir_frame,
        textvariable=config_dirs[key],
        width=60,
        font=("Helvetica", 10)
    )
    dir_entry.pack(side=tk.LEFT, padx=5)

    # Browse button
    browse_btn = tk.Button(
        dir_frame,
        text="Browse...",
        command=lambda k=key: browse_directory(k),
        font=("Helvetica", 9)
    )
    browse_btn.pack(side=tk.LEFT, padx=5)

    # Status indicator
    status_indicator = tk.Label(
        dir_frame,
        text="●",
        font=("Helvetica", 12),
        fg="gray"
    )
    status_indicator.pack(side=tk.LEFT, padx=5)

    # Description label
    desc_label = tk.Label(
        config_frame,
        text=f"    {description}",
        font=("Helvetica", 9),
        fg="gray",
        anchor="w"
    )
    desc_label.pack(fill=tk.X, padx=30)

# Separator
separator = ttk.Separator(tab7, orient='horizontal')
separator.pack(fill=tk.X, padx=40, pady=20)

# Buttons frame
buttons_frame = tk.Frame(tab7)
buttons_frame.pack(pady=10)

# Save Configuration button
save_config_btn = tk.Button(
    buttons_frame,
    text="Save Configuration",
    command=save_config,
    font=("Helvetica", 11),
    bg="#4CAF50",
    fg="white",
    padx=20,
    pady=5
)
save_config_btn.pack(side=tk.LEFT, padx=10)

# Load Configuration button
load_config_btn = tk.Button(
    buttons_frame,
    text="Load Configuration",
    command=load_config,
    font=("Helvetica", 11),
    bg="#2196F3",
    fg="white",
    padx=20,
    pady=5
)
load_config_btn.pack(side=tk.LEFT, padx=10)

# Apply & Load Data button
apply_config_btn = tk.Button(
    buttons_frame,
    text="Apply & Load Data",
    command=apply_config_and_load,
    font=("Helvetica", 11),
    bg="#FF9800",
    fg="white",
    padx=20,
    pady=5
)
apply_config_btn.pack(side=tk.LEFT, padx=10)

# Status label
config_status_label = tk.Label(
    tab7,
    text="0/5 directories configured",
    font=("Helvetica", 10),
    fg="gray"
)
config_status_label.pack(pady=10)

# Quick Setup section
quick_setup_frame = tk.LabelFrame(tab7, text="Quick Setup", font=("Helvetica", 11, "bold"))
quick_setup_frame.pack(fill=tk.X, padx=40, pady=20)

quick_setup_description = tk.Label(
    quick_setup_frame,
    text="Select a project folder with standard subfolder structure (STDDatalog, ImageCaptures, CSVFiles, TXTDatalog, PLMFiles)",
    font=("Helvetica", 9),
    fg="gray"
)
quick_setup_description.pack(pady=5)


def quick_setup_from_project():
    """Auto-configure from project folder with standard structure"""
    folder_path = filedialog.askdirectory(
        title="Select Project Root Folder"
    )

    if not folder_path:
        return

    print(f"Auto-configuring from project folder: {folder_path}")

    # Map standard subfolders
    subfolder_map = {
        "stdf": "STDDatalog",
        "images": "ImageCaptures",
        "csv": "CSVFiles",
        "txt": "TXTDatalog",
        "plm": "PLMFiles"
    }

    found_count = 0
    for key, subfolder in subfolder_map.items():
        subfolder_path = os.path.join(folder_path, subfolder)
        if os.path.exists(subfolder_path):
            config_dirs[key].set(subfolder_path)
            found_count += 1
            print(f"  Found {subfolder}: {subfolder_path}")
        else:
            print(f"  Missing {subfolder}")

    update_config_status()
    config_status_label.config(
        text=f"Found {found_count}/5 subfolders in project",
        fg="green" if found_count > 0 else "orange"
    )


quick_setup_btn = tk.Button(
    quick_setup_frame,
    text="Auto-Configure from Project Folder",
    command=quick_setup_from_project,
    font=("Helvetica", 10),
    bg="#9C27B0",
    fg="white",
    padx=15,
    pady=5
)
quick_setup_btn.pack(pady=10)

# Load saved configuration on startup
load_config()

# ============================================================================
# Tab 3: Multiple Wafermaps - Display multiple wafermaps for comparison
# ============================================================================

# Main container with left panel for wafer selection and right panel for content
multi_wafer_main_container = tk.Frame(tab_multi_wafer)
multi_wafer_main_container.pack(fill=tk.BOTH, expand=True, padx=5, pady=5)

# Left panel with sub-notebook for Waferselection, Statistic, Images
multi_wafer_left_panel = tk.Frame(multi_wafer_main_container, width=280, relief=tk.GROOVE, borderwidth=1)
multi_wafer_left_panel.pack(side=tk.LEFT, fill=tk.Y, padx=(5, 0), pady=5)
multi_wafer_left_panel.pack_propagate(False)

# Title for left panel
multi_wafer_left_title = tk.Label(
    multi_wafer_left_panel,
    text="Wafer & Analysis",
    font=("Helvetica", 11, "bold"),
    bg="#2196F3",
    fg="white"
)
multi_wafer_left_title.pack(fill=tk.X, pady=(0, 2))

# Sub-notebook for left panel with 3 tabs: Waferselection, Statistic, Images
# Create a custom style for subtab fonts (30% larger than default)
style.configure("LargeTab.TNotebook.Tab", font=("Helvetica", 10), padding=[6, 3])

multi_wafer_left_notebook = ttk.Notebook(multi_wafer_left_panel, style="LargeTab.TNotebook")
multi_wafer_left_notebook.pack(fill=tk.BOTH, expand=True, padx=2, pady=2)

# === Sub-Tab 1: Wafer Selection ===
multi_wafer_left_tab_selection = ttk.Frame(multi_wafer_left_notebook)
multi_wafer_left_notebook.add(multi_wafer_left_tab_selection, text="Waferselection")

# Buttons for Select All / Deselect All in Waferselection tab
wafer_select_btn_frame = tk.Frame(multi_wafer_left_tab_selection)
wafer_select_btn_frame.pack(fill=tk.X, padx=5, pady=2)

def select_all_wafers():
    for var in multi_wafer_checkbox_vars:
        var.set(True)
    on_wafer_selection_changed()

def deselect_all_wafers():
    for var in multi_wafer_checkbox_vars:
        var.set(False)
    on_wafer_selection_changed()

wafer_select_all_btn = tk.Button(
    wafer_select_btn_frame,
    text="Select All",
    command=select_all_wafers,
    font=("Helvetica", 8),
    width=10
)
wafer_select_all_btn.pack(side=tk.LEFT, padx=2)

wafer_deselect_all_btn = tk.Button(
    wafer_select_btn_frame,
    text="Deselect All",
    command=deselect_all_wafers,
    font=("Helvetica", 8),
    width=10
)
wafer_deselect_all_btn.pack(side=tk.LEFT, padx=2)

# Scrollable frame for wafer checkboxes (in Waferselection tab)
wafer_list_canvas = tk.Canvas(multi_wafer_left_tab_selection, highlightthickness=0)
wafer_list_scrollbar = ttk.Scrollbar(multi_wafer_left_tab_selection, orient="vertical", command=wafer_list_canvas.yview)
wafer_list_frame = tk.Frame(wafer_list_canvas)

wafer_list_frame.bind(
    "<Configure>",
    lambda e: wafer_list_canvas.configure(scrollregion=wafer_list_canvas.bbox("all"))
)

wafer_list_canvas.create_window((0, 0), window=wafer_list_frame, anchor="nw")
wafer_list_canvas.configure(yscrollcommand=wafer_list_scrollbar.set)

wafer_list_scrollbar.pack(side=tk.RIGHT, fill=tk.Y)
wafer_list_canvas.pack(side=tk.LEFT, fill=tk.BOTH, expand=True, padx=5, pady=5)

# Mouse wheel scrolling for wafer list
def on_wafer_list_mousewheel(event):
    wafer_list_canvas.yview_scroll(int(-1 * (event.delta / 120)), "units")

wafer_list_canvas.bind("<MouseWheel>", on_wafer_list_mousewheel)
wafer_list_frame.bind("<MouseWheel>", on_wafer_list_mousewheel)

# Storage for checkbox variables
multi_wafer_checkbox_vars = []
multi_wafer_checkbox_widgets = []

# === Sub-Tab 2: Statistic ===
multi_wafer_left_tab_statistic = ttk.Frame(multi_wafer_left_notebook)
multi_wafer_left_notebook.add(multi_wafer_left_tab_statistic, text="Statistic")

# Statistics section header
multi_wafer_stats_header = tk.Label(
    multi_wafer_left_tab_statistic,
    text="Statistics Overview",
    font=("Helvetica", 10, "bold"),
    bg="#e8e8e8"
)
multi_wafer_stats_header.pack(fill=tk.X, pady=(0, 5))

# Boxplot type selection
multi_wafer_boxplot_frame = tk.Frame(multi_wafer_left_tab_statistic)
multi_wafer_boxplot_frame.pack(fill=tk.X, padx=5, pady=2)

tk.Label(multi_wafer_boxplot_frame, text="Plot Type:", font=("Helvetica", 9)).pack(side=tk.LEFT, padx=2)
multi_wafer_boxplot_type_var = tk.StringVar(value="Boxplot")
multi_wafer_boxplot_type_combo = ttk.Combobox(
    multi_wafer_boxplot_frame,
    textvariable=multi_wafer_boxplot_type_var,
    values=["Boxplot", "Bin Distribution"],
    state="readonly",
    width=15,
    font=("Helvetica", 9)
)
multi_wafer_boxplot_type_combo.pack(side=tk.LEFT, padx=5)
multi_wafer_boxplot_type_combo.bind("<<ComboboxSelected>>", lambda e: update_multi_wafer_statistics())

# Probability plot type selection
multi_wafer_prob_frame = tk.Frame(multi_wafer_left_tab_statistic)
multi_wafer_prob_frame.pack(fill=tk.X, padx=5, pady=2)

tk.Label(multi_wafer_prob_frame, text="Prob. Plot:", font=("Helvetica", 9)).pack(side=tk.LEFT, padx=2)
multi_wafer_prob_type_var = tk.StringVar(value="CDF")
multi_wafer_prob_type_combo = ttk.Combobox(
    multi_wafer_prob_frame,
    textvariable=multi_wafer_prob_type_var,
    values=["CDF", "PDF"],
    state="readonly",
    width=15,
    font=("Helvetica", 9)
)
multi_wafer_prob_type_combo.pack(side=tk.LEFT, padx=5)
multi_wafer_prob_type_combo.bind("<<ComboboxSelected>>", lambda e: update_multi_wafer_statistics())

# Separator
ttk.Separator(multi_wafer_left_tab_statistic, orient='horizontal').pack(fill=tk.X, padx=5, pady=5)

# Frame for boxplot display in statistics tab
multi_wafer_stats_boxplot_frame = tk.Frame(multi_wafer_left_tab_statistic, bg="#e8e8e8", height=180)
multi_wafer_stats_boxplot_frame.pack(fill=tk.X, padx=5, pady=2)
multi_wafer_stats_boxplot_frame.pack_propagate(False)

# Frame for probability plot display in statistics tab
multi_wafer_stats_prob_frame = tk.Frame(multi_wafer_left_tab_statistic, bg="#e8e8e8", height=180)
multi_wafer_stats_prob_frame.pack(fill=tk.X, padx=5, pady=2)
multi_wafer_stats_prob_frame.pack_propagate(False)

# Canvas references for statistics plots
multi_wafer_stats_boxplot_canvas = None
multi_wafer_stats_prob_canvas = None

# === Sub-Tab 3: Images ===
multi_wafer_left_tab_images = ttk.Frame(multi_wafer_left_notebook)
multi_wafer_left_notebook.add(multi_wafer_left_tab_images, text="Images")

# Images section header
multi_wafer_images_header = tk.Label(
    multi_wafer_left_tab_images,
    text="Die Images & PLM Files",
    font=("Helvetica", 10, "bold"),
    bg="#e8e8e8"
)
multi_wafer_images_header.pack(fill=tk.X, pady=(0, 5))

# Folder selection buttons
multi_wafer_img_folder_frame = tk.Frame(multi_wafer_left_tab_images)
multi_wafer_img_folder_frame.pack(fill=tk.X, padx=5, pady=2)

multi_wafer_img_folder_btn = tk.Button(
    multi_wafer_img_folder_frame,
    text="📁 Img Folder",
    command=lambda: select_multi_wafer_image_folder(),
    font=("Helvetica", 8),
    bg="#2196F3",
    fg="white"
)
multi_wafer_img_folder_btn.pack(side=tk.LEFT, padx=2)

multi_wafer_plm_folder_btn = tk.Button(
    multi_wafer_img_folder_frame,
    text="📁 PLM Folder",
    command=lambda: select_multi_wafer_plm_folder(),
    font=("Helvetica", 8),
    bg="#9C27B0",
    fg="white"
)
multi_wafer_plm_folder_btn.pack(side=tk.LEFT, padx=2)

# Folder status labels
multi_wafer_img_folder_status = tk.Label(
    multi_wafer_left_tab_images,
    text="Img: Not set",
    font=("Helvetica", 8),
    fg="gray"
)
multi_wafer_img_folder_status.pack(fill=tk.X, padx=5, pady=1)

multi_wafer_plm_folder_status = tk.Label(
    multi_wafer_left_tab_images,
    text="PLM: Not set",
    font=("Helvetica", 8),
    fg="gray"
)
multi_wafer_plm_folder_status.pack(fill=tk.X, padx=5, pady=1)

# Separator
ttk.Separator(multi_wafer_left_tab_images, orient='horizontal').pack(fill=tk.X, padx=5, pady=5)

# Image Files label frame
multi_wafer_image_files_label = tk.Label(
    multi_wafer_left_tab_images,
    text="Image Files for selected Die",
    font=("Helvetica", 9, "bold"),
    bg="#e8e8e8"
)
multi_wafer_image_files_label.pack(fill=tk.X, padx=5, pady=(0, 2))

# Scrollable frame for die images
multi_wafer_image_canvas = tk.Canvas(multi_wafer_left_tab_images, height=150, bg="#f0f0f0")
multi_wafer_image_scrollbar = ttk.Scrollbar(multi_wafer_left_tab_images, orient="vertical", command=multi_wafer_image_canvas.yview)
multi_wafer_image_scrollable_frame = tk.Frame(multi_wafer_image_canvas, bg="#f0f0f0")

multi_wafer_image_scrollable_frame.bind(
    "<Configure>",
    lambda e: multi_wafer_image_canvas.configure(scrollregion=multi_wafer_image_canvas.bbox("all"))
)

multi_wafer_image_canvas.create_window((0, 0), window=multi_wafer_image_scrollable_frame, anchor="nw")
multi_wafer_image_canvas.configure(yscrollcommand=multi_wafer_image_scrollbar.set)

multi_wafer_image_scrollbar.pack(side=tk.RIGHT, fill=tk.Y)
multi_wafer_image_canvas.pack(side=tk.LEFT, fill=tk.BOTH, expand=True, padx=5, pady=2)

# Separator before PLM Files
ttk.Separator(multi_wafer_left_tab_images, orient='horizontal').pack(fill=tk.X, padx=5, pady=5)

# PLM Files label frame
multi_wafer_plm_files_label = tk.Label(
    multi_wafer_left_tab_images,
    text="PLM Files for selected Die",
    font=("Helvetica", 9, "bold"),
    bg="#e8e8e8"
)
multi_wafer_plm_files_label.pack(fill=tk.X, padx=5, pady=(0, 2))

# Scrollable frame for PLM files
multi_wafer_plm_canvas = tk.Canvas(multi_wafer_left_tab_images, height=150, bg="#f0f0f0")
multi_wafer_plm_scrollbar = ttk.Scrollbar(multi_wafer_left_tab_images, orient="vertical", command=multi_wafer_plm_canvas.yview)
multi_wafer_plm_scrollable_frame = tk.Frame(multi_wafer_plm_canvas, bg="#f0f0f0")

multi_wafer_plm_scrollable_frame.bind(
    "<Configure>",
    lambda e: multi_wafer_plm_canvas.configure(scrollregion=multi_wafer_plm_canvas.bbox("all"))
)

multi_wafer_plm_canvas.create_window((0, 0), window=multi_wafer_plm_scrollable_frame, anchor="nw")
multi_wafer_plm_canvas.configure(yscrollcommand=multi_wafer_plm_scrollbar.set)

multi_wafer_plm_scrollbar.pack(side=tk.RIGHT, fill=tk.Y)
multi_wafer_plm_canvas.pack(side=tk.LEFT, fill=tk.BOTH, expand=True, padx=5, pady=2)

# Storage for multiwafer image/PLM directories
multi_wafer_image_directory = None
multi_wafer_plm_directory = None
multi_wafer_die_image_refs = []
multi_wafer_plm_image_refs = []
multi_wafer_selected_die = None  # Currently selected die coordinates in multiwafer

def select_multi_wafer_image_folder():
    """Select parent folder that contains images and plm subfolders"""
    global multi_wafer_image_directory, multi_wafer_plm_directory
    folder_path = filedialog.askdirectory(title="Select Parent Folder (containing images/plm)")
    if folder_path:
        # Check for images and plm subfolders
        import os
        images_path = None
        plm_path = None

        # Check direct subfolders for images/plm
        for subfolder in os.listdir(folder_path):
            subfolder_lower = subfolder.lower()
            subfolder_full = os.path.join(folder_path, subfolder)
            if os.path.isdir(subfolder_full):
                if 'image' in subfolder_lower or 'img' in subfolder_lower:
                    images_path = subfolder_full
                elif 'plm' in subfolder_lower:
                    plm_path = subfolder_full

        # If not found, use the selected folder as image folder
        if images_path is None:
            images_path = folder_path

        multi_wafer_image_directory = images_path
        short_img = images_path.split("/")[-1] if "/" in images_path else images_path.split("\\")[-1]
        multi_wafer_img_folder_status.config(text=f"Img: {short_img}", fg="green")
        print(f"Multiwafer image folder set to: {images_path}")

        if plm_path:
            multi_wafer_plm_directory = plm_path
            short_plm = plm_path.split("/")[-1] if "/" in plm_path else plm_path.split("\\")[-1]
            multi_wafer_plm_folder_status.config(text=f"PLM: {short_plm}", fg="green")
            print(f"Multiwafer PLM folder set to: {plm_path}")

def select_multi_wafer_plm_folder():
    """Select PLM folder for multiwafer tab"""
    global multi_wafer_plm_directory
    folder_path = filedialog.askdirectory(title="Select PLM Folder for Multiwafer")
    if folder_path:
        multi_wafer_plm_directory = folder_path
        short_path = folder_path.split("/")[-1] if "/" in folder_path else folder_path.split("\\")[-1]
        multi_wafer_plm_folder_status.config(text=f"PLM: {short_path}", fg="green")
        print(f"Multiwafer PLM folder set to: {folder_path}")

def update_multi_wafer_statistics():
    """Update boxplot and probability plot in the Statistics subtab"""
    global multi_wafer_stats_boxplot_canvas, multi_wafer_stats_prob_canvas
    global multi_wafer_stdf_data, multi_wafer_wafer_ids

    if not multi_wafer_stdf_data:
        return

    # Get selected parameter
    param_selection = multi_wafer_param_combobox.get()
    if not param_selection:
        return

    param_key = param_selection.split(":")[0].strip()

    # Collect data from all selected wafers
    all_data = []
    labels = []

    for i, (df, wafer_id) in enumerate(zip(multi_wafer_stdf_data, multi_wafer_wafer_ids)):
        if i < len(multi_wafer_checkbox_vars) and multi_wafer_checkbox_vars[i].get():
            if param_key in df.columns:
                values = df[param_key].dropna().values
                if len(values) > 0:
                    all_data.append(values)
                    short_id = str(wafer_id)[:10] if len(str(wafer_id)) > 10 else str(wafer_id)
                    labels.append(short_id)

    if not all_data:
        return

    # Clear old plots
    for widget in multi_wafer_stats_boxplot_frame.winfo_children():
        widget.destroy()
    for widget in multi_wafer_stats_prob_frame.winfo_children():
        widget.destroy()

    # === Boxplot / Bin Distribution ===
    plot_type = multi_wafer_boxplot_type_var.get()

    fig_box, ax_box = plt.subplots(figsize=(3.5, 2.2))
    fig_box.patch.set_facecolor('#e8e8e8')

    if plot_type == "Boxplot":
        bp = ax_box.boxplot(all_data, labels=labels, patch_artist=True)
        colors = plt.cm.tab10(np.linspace(0, 1, len(all_data)))
        for patch, color in zip(bp['boxes'], colors):
            patch.set_facecolor(color)
            patch.set_alpha(0.7)
        ax_box.set_title("Boxplot", fontsize=8, fontweight='bold')
    else:  # Bin Distribution
        combined_data = np.concatenate(all_data)
        ax_box.hist(combined_data, bins=30, edgecolor='black', alpha=0.7, color='#2196F3')
        ax_box.set_title("Bin Distribution", fontsize=8, fontweight='bold')

    ax_box.tick_params(axis='both', labelsize=6)
    ax_box.set_facecolor('#f5f5f5')
    plt.setp(ax_box.get_xticklabels(), rotation=45, ha='right', fontsize=5)
    fig_box.tight_layout()

    multi_wafer_stats_boxplot_canvas = FigureCanvasTkAgg(fig_box, master=multi_wafer_stats_boxplot_frame)
    multi_wafer_stats_boxplot_canvas.draw()
    multi_wafer_stats_boxplot_canvas.get_tk_widget().pack(fill=tk.BOTH, expand=True)

    # === Probability Plot (CDF/PDF) ===
    prob_type = multi_wafer_prob_type_var.get()

    fig_prob, ax_prob = plt.subplots(figsize=(3.5, 2.2))
    fig_prob.patch.set_facecolor('#e8e8e8')

    colors = plt.cm.tab10(np.linspace(0, 1, len(all_data)))

    for i, (data, label) in enumerate(zip(all_data, labels)):
        sorted_data = np.sort(data)
        if prob_type == "CDF":
            y = np.arange(1, len(sorted_data) + 1) / len(sorted_data)
            ax_prob.plot(sorted_data, y, label=label, color=colors[i], linewidth=1.5)
            ax_prob.set_ylabel("Cumulative Prob.", fontsize=6)
            ax_prob.set_title("CDF", fontsize=8, fontweight='bold')
        else:  # PDF
            ax_prob.hist(data, bins=25, density=True, alpha=0.5, label=label, color=colors[i])
            ax_prob.set_ylabel("Density", fontsize=6)
            ax_prob.set_title("PDF", fontsize=8, fontweight='bold')

    ax_prob.tick_params(axis='both', labelsize=6)
    ax_prob.set_facecolor('#f5f5f5')
    ax_prob.legend(fontsize=5, loc='best')
    fig_prob.tight_layout()

    multi_wafer_stats_prob_canvas = FigureCanvasTkAgg(fig_prob, master=multi_wafer_stats_prob_frame)
    multi_wafer_stats_prob_canvas.draw()
    multi_wafer_stats_prob_canvas.get_tk_widget().pack(fill=tk.BOTH, expand=True)

# Right panel for content (controls + sub-notebook)
multi_wafer_right_panel = tk.Frame(multi_wafer_main_container)
multi_wafer_right_panel.pack(side=tk.LEFT, fill=tk.BOTH, expand=True, padx=5, pady=5)

# === View Type Selection Frame (Data / Images / PLM) - like in Wafermap Tab ===
multi_wafer_view_frame = tk.Frame(multi_wafer_right_panel, bg="#e0e0e0", relief=tk.RAISED, borderwidth=1)
multi_wafer_view_frame.pack(side=tk.TOP, fill=tk.X, padx=5, pady=(0, 5))

# View type label
tk.Label(
    multi_wafer_view_frame,
    text="View:",
    font=("Helvetica", 10, "bold"),
    bg="#e0e0e0"
).pack(side=tk.LEFT, padx=5, pady=3)

# View type selection
multi_wafer_view_type_var = tk.StringVar(value="Data")
multi_wafer_view_type_combo = ttk.Combobox(
    multi_wafer_view_frame,
    textvariable=multi_wafer_view_type_var,
    values=["Data", "Images", "PLM Files"],
    state="readonly",
    width=12,
    font=("Helvetica", 10)
)
multi_wafer_view_type_combo.pack(side=tk.LEFT, padx=5, pady=3)
multi_wafer_view_type_combo.bind("<<ComboboxSelected>>", lambda e: on_multi_wafer_view_changed())

# View subtype selection (for Images: image type filter)
tk.Label(
    multi_wafer_view_frame,
    text="Subtype:",
    font=("Helvetica", 10),
    bg="#e0e0e0"
).pack(side=tk.LEFT, padx=5, pady=3)

multi_wafer_view_subtype_var = tk.StringVar(value="Heatmap")
multi_wafer_view_subtype_combo = ttk.Combobox(
    multi_wafer_view_frame,
    textvariable=multi_wafer_view_subtype_var,
    values=["Heatmap", "Multi-Plot"],
    state="readonly",
    width=12,
    font=("Helvetica", 10)
)
multi_wafer_view_subtype_combo.pack(side=tk.LEFT, padx=5, pady=3)
multi_wafer_view_subtype_combo.bind("<<ComboboxSelected>>", lambda e: on_multi_wafer_view_changed())

# Separator
ttk.Separator(multi_wafer_view_frame, orient='vertical').pack(side=tk.LEFT, fill='y', padx=10, pady=3)

# Group selection
tk.Label(
    multi_wafer_view_frame,
    text="Group:",
    font=("Helvetica", 10),
    bg="#e0e0e0"
).pack(side=tk.LEFT, padx=5, pady=3)

multi_wafer_group_var = tk.StringVar(value="All Groups")
multi_wafer_group_combo = ttk.Combobox(
    multi_wafer_view_frame,
    textvariable=multi_wafer_group_var,
    values=["All Groups"],
    state="readonly",
    width=18,
    font=("Helvetica", 10)
)
multi_wafer_group_combo.pack(side=tk.LEFT, padx=5, pady=3)
multi_wafer_group_combo.bind("<<ComboboxSelected>>", lambda e: on_multi_wafer_group_changed())

def on_multi_wafer_view_changed():
    """Called when view type changes in multiwafer tab"""
    view_type = multi_wafer_view_type_var.get()

    # Update subtype options based on view type
    if view_type == "Data":
        multi_wafer_view_subtype_combo['values'] = ["Heatmap", "Multi-Plot"]
        multi_wafer_view_subtype_var.set("Heatmap")
    elif view_type == "Images":
        multi_wafer_view_subtype_combo['values'] = ["All", "BF", "DF", "OBIRCH", "TS"]
        multi_wafer_view_subtype_var.set("All")
    elif view_type == "PLM Files":
        multi_wafer_view_subtype_combo['values'] = ["All", "Heatmap"]
        multi_wafer_view_subtype_var.set("All")

    # Refresh display
    refresh_current_multi_wafer_tab()

def on_multi_wafer_group_changed():
    """Called when group selection changes in multiwafer tab"""
    # Update parameter combobox to show only parameters from selected group
    update_multi_wafer_param_by_group()
    refresh_current_multi_wafer_tab()

def update_multi_wafer_param_by_group():
    """Update multiwafer parameter combobox based on selected group"""
    global multi_wafer_test_params

    selected_group = multi_wafer_group_var.get()

    if not multi_wafer_test_params:
        return

    if selected_group == "All Groups" or selected_group not in grouped_parameters:
        # Show all parameters
        current_values = list(multi_wafer_param_combobox['values'])
        if current_values:
            return  # Keep existing values
    else:
        # Filter to selected group
        if selected_group in grouped_parameters:
            params = grouped_parameters[selected_group]
            new_values = []
            for test_num, short_name, full_name in params:
                test_key = f"test_{test_num}"
                if test_key in multi_wafer_test_params:
                    new_values.append(f"{test_key}: {short_name}")
            if new_values:
                multi_wafer_param_combobox['values'] = new_values
                multi_wafer_param_combobox.set(new_values[0])
                on_multi_wafer_param_changed()

# Control frame for multiple wafermaps - Row 1 (File operations)
multi_wafer_control_frame = tk.Frame(multi_wafer_right_panel)
multi_wafer_control_frame.pack(side=tk.TOP, fill=tk.X, padx=5, pady=(5, 2))

# Format selection dropdown
multi_wafer_format_label = tk.Label(
    multi_wafer_control_frame,
    text="Format:",
    font=("Helvetica", 10)
)
multi_wafer_format_label.pack(side=tk.LEFT, padx=5)

multi_wafer_format_var = tk.StringVar(value="CSV")
multi_wafer_format_combobox = ttk.Combobox(
    multi_wafer_control_frame,
    textvariable=multi_wafer_format_var,
    values=["CSV", "STDF"],
    state="readonly",
    width=8,
    font=("Helvetica", 10)
)
multi_wafer_format_combobox.pack(side=tk.LEFT, padx=5)

# Button to add files (appends to existing)
multi_wafer_add_btn = tk.Button(
    multi_wafer_control_frame,
    text="Add Files",
    command=lambda: add_multi_wafer_files(),
    font=("Helvetica", 10),
    bg="#4CAF50",
    fg="white"
)
multi_wafer_add_btn.pack(side=tk.LEFT, padx=5)

# Button to clear all loaded files
multi_wafer_clear_btn = tk.Button(
    multi_wafer_control_frame,
    text="Clear All",
    command=lambda: clear_multi_wafer_files(),
    font=("Helvetica", 10),
    bg="#f44336",
    fg="white"
)
multi_wafer_clear_btn.pack(side=tk.LEFT, padx=5)

# Info label
multi_wafer_info_label = tk.Label(
    multi_wafer_control_frame,
    text="No wafermaps loaded",
    font=("Helvetica", 9),
    fg="gray"
)
multi_wafer_info_label.pack(side=tk.LEFT, padx=10)

# Parameter selection
multi_wafer_param_label = tk.Label(
    multi_wafer_control_frame,
    text="Parameter:",
    font=("Helvetica", 10)
)
multi_wafer_param_label.pack(side=tk.LEFT, padx=5)

multi_wafer_param_combobox = ttk.Combobox(
    multi_wafer_control_frame,
    state="readonly",
    width=40,
    font=("Helvetica", 10)
)
multi_wafer_param_combobox.pack(side=tk.LEFT, padx=5)
multi_wafer_param_combobox.bind("<<ComboboxSelected>>", lambda e: on_multi_wafer_param_changed())

# Refresh button
multi_wafer_refresh_btn = tk.Button(
    multi_wafer_control_frame,
    text="Refresh",
    command=lambda: refresh_current_multi_wafer_tab(),
    font=("Helvetica", 10)
)
multi_wafer_refresh_btn.pack(side=tk.LEFT, padx=5)

# Reset Zoom button
def reset_multi_wafer_zoom_click():
    """Reset zoom to original view"""
    global multi_wafer_plot_data_cache
    if multi_wafer_plot_data_cache and 'reset_zoom_func' in multi_wafer_plot_data_cache:
        multi_wafer_plot_data_cache['reset_zoom_func']()
    else:
        # Fallback: refresh display
        refresh_current_multi_wafer_tab()

multi_wafer_reset_zoom_btn = tk.Button(
    multi_wafer_control_frame,
    text="🔍 Reset Zoom",
    command=reset_multi_wafer_zoom_click,
    font=("Helvetica", 10),
    bg="#FF9800",
    fg="white"
)
multi_wafer_reset_zoom_btn.pack(side=tk.LEFT, padx=5)

# Control frame for multiple wafermaps - Row 2 (Display options)
multi_wafer_control_frame_row2 = tk.Frame(multi_wafer_right_panel)
multi_wafer_control_frame_row2.pack(side=tk.TOP, fill=tk.X, padx=5, pady=(2, 5))

# Show grid checkbox
multi_wafer_grid_var = tk.BooleanVar(value=False)
multi_wafer_grid_checkbox = tk.Checkbutton(
    multi_wafer_control_frame_row2,
    text="Show Grid",
    variable=multi_wafer_grid_var,
    command=lambda: update_multi_wafer_display(),
    font=("Helvetica", 10)
)
multi_wafer_grid_checkbox.pack(side=tk.LEFT, padx=5)

# Layout options
layout_label = tk.Label(
    multi_wafer_control_frame_row2,
    text="Layout:",
    font=("Helvetica", 10)
)
layout_label.pack(side=tk.LEFT, padx=5)

multi_wafer_layout_var = tk.StringVar(value="Auto")
multi_wafer_layout_combobox = ttk.Combobox(
    multi_wafer_control_frame_row2,
    textvariable=multi_wafer_layout_var,
    values=["Auto", "1 Column", "2 Columns", "3 Columns", "4 Columns"],
    state="readonly",
    width=10,
    font=("Helvetica", 9)
)
multi_wafer_layout_combobox.pack(side=tk.LEFT, padx=5)
multi_wafer_layout_combobox.bind("<<ComboboxSelected>>", lambda e: update_multi_wafer_display())

# Compare Mode selection (Parameter Name vs Test Number)
compare_mode_label = tk.Label(
    multi_wafer_control_frame_row2,
    text="Compare by:",
    font=("Helvetica", 10)
)
compare_mode_label.pack(side=tk.LEFT, padx=5)

multi_wafer_compare_mode_var = tk.StringVar(value="Parameter Name")
multi_wafer_compare_mode_combobox = ttk.Combobox(
    multi_wafer_control_frame_row2,
    textvariable=multi_wafer_compare_mode_var,
    values=["Parameter Name", "Test Number"],
    state="readonly",
    width=15,
    font=("Helvetica", 9)
)
multi_wafer_compare_mode_combobox.pack(side=tk.LEFT, padx=5)
multi_wafer_compare_mode_combobox.bind("<<ComboboxSelected>>", lambda e: on_compare_mode_changed())

# Separator
ttk.Separator(multi_wafer_control_frame_row2, orient='vertical').pack(side=tk.LEFT, fill='y', padx=10, pady=2)

# Compare Independent checkbox - allows independent parameter selection per wafer
multi_wafer_compare_independent_var = tk.BooleanVar(value=False)
multi_wafer_compare_independent_checkbox = tk.Checkbutton(
    multi_wafer_control_frame_row2,
    text="Compare Independent",
    variable=multi_wafer_compare_independent_var,
    command=lambda: on_compare_independent_changed(),
    font=("Helvetica", 10)
)
multi_wafer_compare_independent_checkbox.pack(side=tk.LEFT, padx=5)

# Separator before alignment status
ttk.Separator(multi_wafer_control_frame_row2, orient='vertical').pack(side=tk.LEFT, fill='y', padx=10, pady=2)

# Multi-wafer alignment status frame
multi_wafer_align_frame = tk.LabelFrame(multi_wafer_control_frame_row2, text="Alignment", font=("Helvetica", 8))
multi_wafer_align_frame.pack(side=tk.LEFT, padx=5, pady=2)

# X/Y Offset status label for multi-wafer
multi_wafer_offset_label = tk.Label(
    multi_wafer_align_frame,
    text=" X/Y: -- ",
    font=("Consolas", 8),
    fg="gray",
    padx=4,
    pady=1
)
multi_wafer_offset_label.pack(side=tk.LEFT, padx=2, pady=1)

# Notch orientation status label for multi-wafer
multi_wafer_notch_label = tk.Label(
    multi_wafer_align_frame,
    text=" Notch: -- ",
    font=("Consolas", 8),
    fg="gray",
    padx=4,
    pady=1
)
multi_wafer_notch_label.pack(side=tk.LEFT, padx=2, pady=1)

# Overlay Coordinates Button for multi-wafer
multi_wafer_overlay_btn = tk.Button(
    multi_wafer_control_frame_row2,
    text="Align All",
    command=lambda: multi_wafer_overlay_coordinates(),
    font=("Helvetica", 8),
    bg="#9C27B0",
    fg="white",
    state=tk.DISABLED
)
multi_wafer_overlay_btn.pack(side=tk.LEFT, padx=5, pady=2)

# Reference wafer selection for alignment
multi_wafer_ref_label = tk.Label(
    multi_wafer_control_frame_row2,
    text="Ref:",
    font=("Helvetica", 8)
)
multi_wafer_ref_label.pack(side=tk.LEFT, padx=(5, 2))

multi_wafer_ref_var = tk.StringVar(value="First")
multi_wafer_ref_combobox_align = ttk.Combobox(
    multi_wafer_control_frame_row2,
    textvariable=multi_wafer_ref_var,
    values=["First"],
    state="readonly",
    width=15,
    font=("Helvetica", 8)
)
multi_wafer_ref_combobox_align.pack(side=tk.LEFT, padx=2)
multi_wafer_ref_combobox_align.bind("<<ComboboxSelected>>", lambda e: detect_multi_wafer_alignment())

# Store alignment info for multi-wafer
multi_wafer_alignment_info = {}

def detect_multi_wafer_notch(df):
    """Detect notch orientation for a wafer dataframe"""
    if df is None or len(df) == 0:
        return "?"

    x_min, x_max = df['x'].min(), df['x'].max()
    y_min, y_max = df['y'].min(), df['y'].max()

    x_range = x_max - x_min
    y_range = y_max - y_min

    if x_range == 0 or y_range == 0:
        return "?"

    edge_threshold_x = x_range * 0.15
    edge_threshold_y = y_range * 0.15

    top_edge = len(df[df['y'] >= (y_max - edge_threshold_y)])
    bottom_edge = len(df[df['y'] <= (y_min + edge_threshold_y)])
    left_edge = len(df[df['x'] <= (x_min + edge_threshold_x)])
    right_edge = len(df[df['x'] >= (x_max - edge_threshold_x)])

    edges = {'12h': top_edge, '6h': bottom_edge, '9h': left_edge, '3h': right_edge}
    return min(edges, key=edges.get)

def detect_multi_wafer_alignment():
    """Detect alignment status across all loaded wafers compared to reference"""
    global multi_wafer_stdf_data, multi_wafer_wafer_ids, multi_wafer_alignment_info

    if not multi_wafer_stdf_data or len(multi_wafer_stdf_data) < 2:
        multi_wafer_offset_label.config(text=" X/Y: -- ", fg="gray", bg="#f0f0f0")
        multi_wafer_notch_label.config(text=" Notch: -- ", fg="gray", bg="#f0f0f0")
        multi_wafer_overlay_btn.config(state=tk.DISABLED)
        return

    # Get reference index
    ref_selection = multi_wafer_ref_var.get()
    if ref_selection == "First":
        ref_idx = 0
    else:
        try:
            ref_idx = int(ref_selection.split(":")[0]) - 1
        except:
            ref_idx = 0

    if ref_idx >= len(multi_wafer_stdf_data):
        ref_idx = 0

    ref_df = multi_wafer_stdf_data[ref_idx]
    ref_x_min = ref_df['x'].min()
    ref_y_min = ref_df['y'].min()
    ref_notch = detect_multi_wafer_notch(ref_df)

    # Check all other wafers
    multi_wafer_alignment_info = {}
    has_offset = False
    has_notch_diff = False
    offset_count = 0
    notch_diff_count = 0

    for i, df in enumerate(multi_wafer_stdf_data):
        if i == ref_idx:
            multi_wafer_alignment_info[i] = {'x_offset': 0, 'y_offset': 0, 'notch': ref_notch, 'ref_notch': ref_notch}
            continue

        comp_x_min = df['x'].min()
        comp_y_min = df['y'].min()
        x_offset = int(comp_x_min - ref_x_min)
        y_offset = int(comp_y_min - ref_y_min)
        comp_notch = detect_multi_wafer_notch(df)

        multi_wafer_alignment_info[i] = {
            'x_offset': x_offset,
            'y_offset': y_offset,
            'notch': comp_notch,
            'ref_notch': ref_notch
        }

        if x_offset != 0 or y_offset != 0:
            has_offset = True
            offset_count += 1
        if comp_notch != ref_notch:
            has_notch_diff = True
            notch_diff_count += 1

    # Update status labels
    total_compare = len(multi_wafer_stdf_data) - 1

    if not has_offset:
        multi_wafer_offset_label.config(text=" X/Y: ✓ ", fg="white", bg="#4CAF50")
    else:
        multi_wafer_offset_label.config(text=f" X/Y: {offset_count}/{total_compare} ", fg="white", bg="#f44336")

    if not has_notch_diff:
        multi_wafer_notch_label.config(text=f" Notch: ✓ {ref_notch} ", fg="white", bg="#4CAF50")
    else:
        multi_wafer_notch_label.config(text=f" Notch: {notch_diff_count}/{total_compare} ", fg="white", bg="#f44336")

    # Enable overlay button if misalignment detected
    if has_offset or has_notch_diff:
        multi_wafer_overlay_btn.config(state=tk.NORMAL)
    else:
        multi_wafer_overlay_btn.config(state=tk.DISABLED)

def update_multi_wafer_ref_combobox():
    """Update the reference wafer combobox with loaded wafer names"""
    global multi_wafer_wafer_ids

    if not multi_wafer_wafer_ids:
        multi_wafer_ref_combobox_align['values'] = ["First"]
        multi_wafer_ref_var.set("First")
        return

    options = ["First"]
    for i, wafer_id in enumerate(multi_wafer_wafer_ids):
        short_name = str(wafer_id)[:15] if len(str(wafer_id)) > 15 else str(wafer_id)
        options.append(f"{i+1}: {short_name}")

    multi_wafer_ref_combobox_align['values'] = options
    if multi_wafer_ref_var.get() not in options:
        multi_wafer_ref_var.set("First")

def multi_wafer_rotate_data(df, from_notch, to_notch):
    """Rotate wafer data to align notch orientations"""
    notch_positions = {'12h': 0, '3h': 90, '6h': 180, '9h': 270}

    if from_notch not in notch_positions or to_notch not in notch_positions:
        return df

    from_angle = notch_positions[from_notch]
    to_angle = notch_positions[to_notch]
    rotation = (to_angle - from_angle) % 360

    if rotation == 0:
        return df

    x_center = (df['x'].min() + df['x'].max()) / 2
    y_center = (df['y'].min() + df['y'].max()) / 2

    rotated_df = df.copy()
    x_rel = df['x'] - x_center
    y_rel = df['y'] - y_center

    if rotation == 90:
        rotated_df['x'] = (y_rel + x_center).round().astype(int)
        rotated_df['y'] = (-x_rel + y_center).round().astype(int)
    elif rotation == 180:
        rotated_df['x'] = (-x_rel + x_center).round().astype(int)
        rotated_df['y'] = (-y_rel + y_center).round().astype(int)
    elif rotation == 270:
        rotated_df['x'] = (-y_rel + x_center).round().astype(int)
        rotated_df['y'] = (x_rel + y_center).round().astype(int)

    return rotated_df

def multi_wafer_overlay_coordinates():
    """Align all wafer coordinates to the reference wafer"""
    global multi_wafer_stdf_data, multi_wafer_alignment_info, multi_wafer_wafer_ids

    if not multi_wafer_stdf_data or not multi_wafer_alignment_info:
        return

    print(f"\n{'='*60}")
    print("Aligning Multi-Wafer Coordinates to Reference")
    print(f"{'='*60}")

    aligned_count = 0

    for i, info in multi_wafer_alignment_info.items():
        if i >= len(multi_wafer_stdf_data):
            continue

        x_offset = info['x_offset']
        y_offset = info['y_offset']
        notch = info['notch']
        ref_notch = info['ref_notch']

        needs_alignment = x_offset != 0 or y_offset != 0 or notch != ref_notch

        if needs_alignment:
            wafer_name = multi_wafer_wafer_ids[i] if i < len(multi_wafer_wafer_ids) else f"Wafer {i+1}"
            print(f"  Aligning: {wafer_name}")

            # Apply X/Y offset correction
            if x_offset != 0 or y_offset != 0:
                print(f"    Offset: dX={-x_offset}, dY={-y_offset}")
                multi_wafer_stdf_data[i]['x'] = multi_wafer_stdf_data[i]['x'] - x_offset
                multi_wafer_stdf_data[i]['y'] = multi_wafer_stdf_data[i]['y'] - y_offset

            # Apply rotation if notch differs
            if notch != ref_notch:
                print(f"    Notch: {notch} -> {ref_notch}")
                multi_wafer_stdf_data[i] = multi_wafer_rotate_data(
                    multi_wafer_stdf_data[i], notch, ref_notch
                )

            aligned_count += 1

    print(f"Aligned {aligned_count} wafer(s)")

    # Re-detect alignment status
    detect_multi_wafer_alignment()

    # Refresh display
    if multi_wafer_compare_independent_var.get():
        update_multi_wafer_independent_display()
    else:
        update_multi_wafer_display()

    # Show info
    from tkinter import messagebox
    messagebox.showinfo("Alignment Complete", f"Successfully aligned {aligned_count} wafer(s) to reference.")

# Storage for individual parameter selections per wafer (wafer_index -> selected_param)
multi_wafer_independent_params = {}
multi_wafer_independent_comboboxes = {}
multi_wafer_independent_selector_frame = None

def on_compare_independent_changed():
    """Called when Compare Independent checkbox changes"""
    global multi_wafer_independent_selector_frame

    if multi_wafer_compare_independent_var.get():
        # Enable independent mode - hide global param selector and show individual ones
        multi_wafer_param_label.pack_forget()
        multi_wafer_param_combobox.pack_forget()
        update_multi_wafer_independent_display()
    else:
        # Disable independent mode - remove selector frame, show global param selector
        if hasattr(create_independent_param_selectors, 'selector_frame') and create_independent_param_selectors.selector_frame:
            create_independent_param_selectors.selector_frame.destroy()
            create_independent_param_selectors.selector_frame = None
        # Re-show global parameter selection (after info label)
        multi_wafer_param_label.pack(side=tk.LEFT, padx=5, after=multi_wafer_info_label)
        multi_wafer_param_combobox.pack(side=tk.LEFT, padx=5, after=multi_wafer_param_label)
        multi_wafer_param_label.pack(side=tk.LEFT, padx=5, after=multi_wafer_info_label)
        multi_wafer_param_combobox.pack(side=tk.LEFT, padx=5, after=multi_wafer_param_label)
        update_multi_wafer_display()

# Sub-Notebook (Untertabs) for Wafermaps, Boxplot, Distribution
multi_wafer_sub_notebook = ttk.Notebook(multi_wafer_right_panel)
multi_wafer_sub_notebook.pack(fill=tk.BOTH, expand=True, padx=5, pady=5)

# Sub-Tab 0: Testheader Comparison (NEW - before Wafermap)
multi_wafer_subtab_header = ttk.Frame(multi_wafer_sub_notebook)
multi_wafer_sub_notebook.add(multi_wafer_subtab_header, text="Testheader Comparison")

# Control frame for header comparison
multi_wafer_header_control_frame = tk.Frame(multi_wafer_subtab_header)
multi_wafer_header_control_frame.pack(side=tk.TOP, fill=tk.X, padx=5, pady=5)

# Reference wafer selection
tk.Label(multi_wafer_header_control_frame, text="Reference Wafer:", font=("Helvetica", 9)).pack(side=tk.LEFT, padx=5)
multi_wafer_ref_combobox = ttk.Combobox(multi_wafer_header_control_frame, width=30, state="readonly")
multi_wafer_ref_combobox.pack(side=tk.LEFT, padx=5)
multi_wafer_ref_combobox.bind("<<ComboboxSelected>>", lambda e: update_testheader_comparison())

# Compare wafer selection
tk.Label(multi_wafer_header_control_frame, text="Compare with:", font=("Helvetica", 9)).pack(side=tk.LEFT, padx=5)
multi_wafer_compare_combobox = ttk.Combobox(multi_wafer_header_control_frame, width=30, state="readonly")
multi_wafer_compare_combobox.pack(side=tk.LEFT, padx=5)
multi_wafer_compare_combobox.bind("<<ComboboxSelected>>", lambda e: update_testheader_comparison())

# Export button
def export_testheader_comparison():
    """Export the testheader comparison table to CSV or Excel"""
    # Check if there is data to export
    items = multi_wafer_header_tree.get_children()
    if not items:
        from tkinter import messagebox
        messagebox.showwarning("Export", "No data to export. Please load wafers and select comparison first.")
        return

    # Get reference and compare wafer names for filename
    ref_name = multi_wafer_ref_combobox.get().replace(".", "_").replace(" ", "_")[:20]
    compare_name = multi_wafer_compare_combobox.get().replace(".", "_").replace(" ", "_")[:20]

    # Ask for save location
    from datetime import datetime
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    default_filename = f"Testheader_Comparison_{timestamp}"

    filepath = filedialog.asksaveasfilename(
        title="Export Testheader Comparison",
        defaultextension=".csv",
        initialfile=default_filename,
        filetypes=[
            ("CSV files", "*.csv"),
            ("Excel files", "*.xlsx"),
            ("All files", "*.*")
        ]
    )

    if not filepath:
        return

    # Collect data from treeview
    headers = ["Test #", "Reference Name", "Compare Name", "Ref Limits", "Compare Limits", "Status"]
    data = []
    for item in items:
        values = multi_wafer_header_tree.item(item, "values")
        data.append(values)

    try:
        if filepath.endswith(".xlsx"):
            # Export to Excel
            df = pd.DataFrame(data, columns=headers)
            df.to_excel(filepath, index=False, sheet_name="Testheader Comparison")
            print(f"Exported to Excel: {filepath}")
        else:
            # Export to CSV
            df = pd.DataFrame(data, columns=headers)
            df.to_csv(filepath, index=False, encoding='utf-8-sig')
            print(f"Exported to CSV: {filepath}")

        from tkinter import messagebox
        messagebox.showinfo("Export Successful", f"Testheader comparison exported to:\n{filepath}")
    except Exception as e:
        from tkinter import messagebox
        messagebox.showerror("Export Error", f"Failed to export: {str(e)}")
        print(f"Export error: {e}")

multi_wafer_export_btn = tk.Button(
    multi_wafer_header_control_frame,
    text="Export Table",
    command=export_testheader_comparison,
    font=("Helvetica", 9),
    bg="#4CAF50",
    fg="white"
)
multi_wafer_export_btn.pack(side=tk.LEFT, padx=15)

# Legend frame
multi_wafer_header_legend_frame = tk.Frame(multi_wafer_subtab_header)
multi_wafer_header_legend_frame.pack(side=tk.TOP, fill=tk.X, padx=5, pady=2)

tk.Label(multi_wafer_header_legend_frame, text="Legende:", font=("Helvetica", 9, "bold")).pack(side=tk.LEFT, padx=5)
tk.Label(multi_wafer_header_legend_frame, text="■ Match", fg="green", font=("Helvetica", 9)).pack(side=tk.LEFT, padx=5)
tk.Label(multi_wafer_header_legend_frame, text="■ Missing (not in Compare)", fg="red", font=("Helvetica", 9)).pack(side=tk.LEFT, padx=5)
tk.Label(multi_wafer_header_legend_frame, text="■ Additional (only in Compare)", fg="#DAA520", font=("Helvetica", 9)).pack(side=tk.LEFT, padx=5)
tk.Label(multi_wafer_header_legend_frame, text="■ Name differs", fg="blue", font=("Helvetica", 9)).pack(side=tk.LEFT, padx=5)

# Treeview frame for comparison table
multi_wafer_header_tree_frame = tk.Frame(multi_wafer_subtab_header)
multi_wafer_header_tree_frame.pack(fill=tk.BOTH, expand=True, padx=5, pady=5)

# Create Treeview with scrollbars
multi_wafer_header_tree_scroll_y = ttk.Scrollbar(multi_wafer_header_tree_frame, orient="vertical")
multi_wafer_header_tree_scroll_x = ttk.Scrollbar(multi_wafer_header_tree_frame, orient="horizontal")

multi_wafer_header_tree = ttk.Treeview(
    multi_wafer_header_tree_frame,
    columns=("test_num", "ref_name", "compare_name", "ref_limits", "compare_limits", "status"),
    show="headings",
    yscrollcommand=multi_wafer_header_tree_scroll_y.set,
    xscrollcommand=multi_wafer_header_tree_scroll_x.set
)

multi_wafer_header_tree.heading("test_num", text="Test #")
multi_wafer_header_tree.heading("ref_name", text="Reference Name")
multi_wafer_header_tree.heading("compare_name", text="Compare Name")
multi_wafer_header_tree.heading("ref_limits", text="Ref Limits")
multi_wafer_header_tree.heading("compare_limits", text="Compare Limits")
multi_wafer_header_tree.heading("status", text="Status")

multi_wafer_header_tree.column("test_num", width=70, anchor="center")
multi_wafer_header_tree.column("ref_name", width=250, anchor="w")
multi_wafer_header_tree.column("compare_name", width=250, anchor="w")
multi_wafer_header_tree.column("ref_limits", width=150, anchor="center")
multi_wafer_header_tree.column("compare_limits", width=150, anchor="center")
multi_wafer_header_tree.column("status", width=120, anchor="center")

multi_wafer_header_tree_scroll_y.config(command=multi_wafer_header_tree.yview)
multi_wafer_header_tree_scroll_x.config(command=multi_wafer_header_tree.xview)

multi_wafer_header_tree_scroll_y.pack(side=tk.RIGHT, fill=tk.Y)
multi_wafer_header_tree_scroll_x.pack(side=tk.BOTTOM, fill=tk.X)
multi_wafer_header_tree.pack(fill=tk.BOTH, expand=True)

# Configure Treeview tags for coloring
multi_wafer_header_tree.tag_configure("match", background="#C8E6C9")  # Light green
multi_wafer_header_tree.tag_configure("missing", background="#FFCDD2")  # Light red
multi_wafer_header_tree.tag_configure("additional", background="#FFF9C4")  # Light yellow
multi_wafer_header_tree.tag_configure("name_diff", background="#BBDEFB")  # Light blue

# Summary label for header comparison
multi_wafer_header_summary_label = tk.Label(
    multi_wafer_subtab_header,
    text="Select Reference and Compare wafers to see comparison",
    font=("Helvetica", 10)
)
multi_wafer_header_summary_label.pack(side=tk.BOTTOM, fill=tk.X, padx=5, pady=5)

# Storage for test parameters per wafer (for comparison)
multi_wafer_test_params_per_wafer = []
multi_wafer_test_limits_per_wafer = []

# Sub-Tab 1: Wafermaps
multi_wafer_subtab_maps = ttk.Frame(multi_wafer_sub_notebook)
multi_wafer_sub_notebook.add(multi_wafer_subtab_maps, text="Wafermap")

# Sub-Tab 2: Boxplot
multi_wafer_subtab_boxplot = ttk.Frame(multi_wafer_sub_notebook)
multi_wafer_sub_notebook.add(multi_wafer_subtab_boxplot, text="Boxplot")

# Sub-Tab 3: Distribution (CDF)
multi_wafer_subtab_distribution = ttk.Frame(multi_wafer_sub_notebook)
multi_wafer_sub_notebook.add(multi_wafer_subtab_distribution, text="Distribution")

# Display frame for wafermaps (in sub-tab 1)
multi_wafer_display_frame = tk.Frame(multi_wafer_subtab_maps)
multi_wafer_display_frame.pack(fill=tk.BOTH, expand=True, padx=5, pady=5)

# Control frame for boxplot options (in sub-tab 2)
multi_wafer_boxplot_control_frame = tk.Frame(multi_wafer_subtab_boxplot)
multi_wafer_boxplot_control_frame.pack(side=tk.TOP, fill=tk.X, padx=5, pady=5)

# Checkbox for showing statistics
multi_wafer_boxplot_stats_var = tk.BooleanVar(value=True)
multi_wafer_boxplot_stats_checkbox = tk.Checkbutton(
    multi_wafer_boxplot_control_frame,
    text="Show Statistics (Mean/Median)",
    variable=multi_wafer_boxplot_stats_var,
    command=lambda: update_multi_wafer_boxplot(),
    font=("Helvetica", 9)
)
multi_wafer_boxplot_stats_checkbox.pack(side=tk.LEFT, padx=5)

# Display frame for boxplot (in sub-tab 2)
multi_wafer_boxplot_frame = tk.Frame(multi_wafer_subtab_boxplot)
multi_wafer_boxplot_frame.pack(fill=tk.BOTH, expand=True, padx=5, pady=5)

# Display frame for distribution (in sub-tab 3)
multi_wafer_distribution_frame = tk.Frame(multi_wafer_subtab_distribution)
multi_wafer_distribution_frame.pack(fill=tk.BOTH, expand=True, padx=5, pady=5)

# Storage for multiple wafermap tab data
multi_wafer_stdf_data = []
multi_wafer_wafer_ids = []
multi_wafer_test_params = {}
multi_wafer_wafer_configs = []  # Wafer configurations (notch orientation) for each loaded wafer
multi_wafer_canvas = None
multi_wafer_boxplot_canvas = None
multi_wafer_distribution_canvas = None

def update_testheader_comparison():
    """Update the testheader comparison display between reference and compare wafer"""
    global multi_wafer_test_params_per_wafer, multi_wafer_test_limits_per_wafer

    # Clear existing items in treeview
    for item in multi_wafer_header_tree.get_children():
        multi_wafer_header_tree.delete(item)

    if len(multi_wafer_wafer_ids) < 2:
        multi_wafer_header_summary_label.config(
            text="Load at least 2 wafers to compare testheaders", fg="gray"
        )
        return

    ref_selection = multi_wafer_ref_combobox.current()
    compare_selection = multi_wafer_compare_combobox.current()

    if ref_selection < 0 or compare_selection < 0:
        multi_wafer_header_summary_label.config(
            text="Select both Reference and Compare wafers", fg="gray"
        )
        return

    if ref_selection >= len(multi_wafer_test_params_per_wafer) or \
       compare_selection >= len(multi_wafer_test_params_per_wafer):
        multi_wafer_header_summary_label.config(
            text="Wafer data not available", fg="red"
        )
        return

    # Get test params and limits for both wafers
    ref_params = multi_wafer_test_params_per_wafer[ref_selection]
    compare_params = multi_wafer_test_params_per_wafer[compare_selection]
    ref_limits = multi_wafer_test_limits_per_wafer[ref_selection] if ref_selection < len(multi_wafer_test_limits_per_wafer) else {}
    compare_limits = multi_wafer_test_limits_per_wafer[compare_selection] if compare_selection < len(multi_wafer_test_limits_per_wafer) else {}

    # Get all unique test numbers
    all_test_nums = set()
    for key in ref_params.keys():
        if key.startswith("test_"):
            test_num = int(key.replace("test_", ""))
            all_test_nums.add(test_num)
    for key in compare_params.keys():
        if key.startswith("test_"):
            test_num = int(key.replace("test_", ""))
            all_test_nums.add(test_num)

    # Statistics counters
    match_count = 0
    missing_count = 0
    additional_count = 0
    name_diff_count = 0

    # Compare each test number
    for test_num in sorted(all_test_nums):
        test_key = f"test_{test_num}"
        ref_name = ref_params.get(test_key, "")
        compare_name = compare_params.get(test_key, "")

        # Format limits
        ref_lim = ref_limits.get(test_num, {})
        compare_lim = compare_limits.get(test_num, {})

        ref_limits_str = format_limits(ref_lim)
        compare_limits_str = format_limits(compare_lim)

        # Determine status and tag
        if ref_name and not compare_name:
            status = "MISSING"
            tag = "missing"
            missing_count += 1
        elif not ref_name and compare_name:
            status = "ADDITIONAL"
            tag = "additional"
            additional_count += 1
        elif ref_name != compare_name:
            status = "NAME DIFFERS"
            tag = "name_diff"
            name_diff_count += 1
        else:
            status = "Match"
            tag = "match"
            match_count += 1

        # Insert row into treeview
        multi_wafer_header_tree.insert(
            "", "end",
            values=(test_num, ref_name, compare_name, ref_limits_str, compare_limits_str, status),
            tags=(tag,)
        )

    # Update summary label
    total = len(all_test_nums)
    summary_text = (
        f"Total: {total} parameters | "
        f"✓ Match: {match_count} | "
        f"✗ Missing: {missing_count} | "
        f"+ Additional: {additional_count} | "
        f"≠ Name differs: {name_diff_count}"
    )

    if missing_count > 0 or additional_count > 0 or name_diff_count > 0:
        summary_color = "orange"
    else:
        summary_color = "green"

    multi_wafer_header_summary_label.config(text=summary_text, fg=summary_color)

def format_limits(limits_dict):
    """Format limits dictionary into a readable string"""
    if not limits_dict:
        return "-"

    lo = limits_dict.get('lo_limit')
    hi = limits_dict.get('hi_limit')
    units = limits_dict.get('units', '')

    if lo is None and hi is None:
        return "-"

    parts = []
    if lo is not None:
        parts.append(f"Lo:{lo:.4g}")
    if hi is not None:
        parts.append(f"Hi:{hi:.4g}")

    result = " ".join(parts)
    if units:
        result += f" {units}"

    return result

def update_header_comparison_comboboxes():
    """Update the comboboxes for header comparison with loaded wafers"""
    wafer_options = []
    for idx, wafer_id in enumerate(multi_wafer_wafer_ids):
        short_id = str(wafer_id)[:30] + "..." if len(str(wafer_id)) > 30 else str(wafer_id)
        wafer_options.append(f"{idx+1}. {short_id}")

    multi_wafer_ref_combobox["values"] = wafer_options
    multi_wafer_compare_combobox["values"] = wafer_options

    # Auto-select first and second wafer if available
    if len(wafer_options) >= 1:
        multi_wafer_ref_combobox.current(0)
    if len(wafer_options) >= 2:
        multi_wafer_compare_combobox.current(1)

def get_selected_wafer_indices():
    """Return list of indices of selected wafers"""
    return [i for i, var in enumerate(multi_wafer_checkbox_vars) if var.get()]

def update_wafer_selection_list():
    """Update the wafer selection checkboxes based on loaded data"""
    global multi_wafer_checkbox_vars, multi_wafer_checkbox_widgets

    # Clear existing checkboxes
    for widget in multi_wafer_checkbox_widgets:
        widget.destroy()
    multi_wafer_checkbox_vars.clear()
    multi_wafer_checkbox_widgets.clear()

    # Create checkboxes for each loaded wafer
    for idx, wafer_id in enumerate(multi_wafer_wafer_ids):
        var = tk.BooleanVar(value=True)  # Default selected
        multi_wafer_checkbox_vars.append(var)

        # Truncate long wafer IDs
        short_id = str(wafer_id)[:25] + "..." if len(str(wafer_id)) > 25 else str(wafer_id)

        cb = tk.Checkbutton(
            wafer_list_frame,
            text=f"{idx+1}. {short_id}",
            variable=var,
            command=on_wafer_selection_changed,
            font=("Helvetica", 9),
            anchor="w"
        )
        cb.pack(fill=tk.X, padx=2, pady=1)
        cb.bind("<MouseWheel>", on_wafer_list_mousewheel)
        multi_wafer_checkbox_widgets.append(cb)

def on_wafer_selection_changed():
    """Called when wafer selection changes - update displays"""
    refresh_current_multi_wafer_tab()
    update_multi_wafer_statistics()

def on_multi_wafer_param_changed():
    """Called when parameter selection changes"""
    refresh_current_multi_wafer_tab()
    update_multi_wafer_statistics()

def on_compare_mode_changed():
    """Called when compare mode changes between 'Parameter Name' and 'Test Number'"""
    global multi_wafer_test_params

    compare_mode = multi_wafer_compare_mode_var.get()

    if not multi_wafer_stdf_data:
        return

    if compare_mode == "Test Number":
        # Show only test numbers (e.g., "test_1234" -> "1234")
        param_options = ["BIN (Bin Number)"]
        test_numbers = set()

        # Collect all unique test numbers from all loaded wafers
        for test_key in multi_wafer_test_params.keys():
            if test_key.startswith("test_"):
                test_num = test_key.replace("test_", "")
                test_numbers.add(int(test_num))

        for test_num in sorted(test_numbers):
            param_options.append(f"test_{test_num}: Test #{test_num}")

        multi_wafer_param_combobox["values"] = param_options
        if param_options:
            multi_wafer_param_combobox.current(0)
    else:
        # Show full parameter names (default behavior)
        param_options = ["BIN (Bin Number)"]
        for test_key, test_name in sorted(multi_wafer_test_params.items()):
            param_options.append(f"{test_key}: {test_name}")

        multi_wafer_param_combobox["values"] = param_options
        if param_options:
            multi_wafer_param_combobox.current(0)

    # Refresh the display
    refresh_current_multi_wafer_tab()

def refresh_current_multi_wafer_tab():
    """Refresh the currently visible sub-tab"""
    current_tab = multi_wafer_sub_notebook.index(multi_wafer_sub_notebook.select())
    if current_tab == 0:
        update_testheader_comparison()
    elif current_tab == 1:
        update_multi_wafer_display()
    elif current_tab == 2:
        update_multi_wafer_boxplot()
    elif current_tab == 3:
        update_multi_wafer_distribution()

# Bind tab change event
def on_multi_wafer_subtab_changed(event):
    refresh_current_multi_wafer_tab()

multi_wafer_sub_notebook.bind("<<NotebookTabChanged>>", on_multi_wafer_subtab_changed)


def add_multi_wafer_files():
    """Add files based on selected format (appends to existing loaded files)"""
    selected_format = multi_wafer_format_var.get()

    if selected_format == "CSV":
        add_multi_wafer_csv_files()
    else:
        add_multi_wafer_stdf_files()


def clear_multi_wafer_files():
    """Clear all loaded wafermap files"""
    global multi_wafer_stdf_data, multi_wafer_wafer_ids, multi_wafer_test_params
    global multi_wafer_canvas, multi_wafer_boxplot_canvas, multi_wafer_distribution_canvas
    global multi_wafer_test_params_per_wafer, multi_wafer_test_limits_per_wafer
    global multi_wafer_wafer_configs

    multi_wafer_stdf_data = []
    multi_wafer_wafer_ids = []
    multi_wafer_test_params = {}
    multi_wafer_test_params_per_wafer = []
    multi_wafer_test_limits_per_wafer = []
    multi_wafer_wafer_configs = []

    # Clear wafermap display
    if multi_wafer_canvas:
        multi_wafer_canvas.get_tk_widget().destroy()
        multi_wafer_canvas = None

    # Clear boxplot display
    if multi_wafer_boxplot_canvas:
        multi_wafer_boxplot_canvas.get_tk_widget().destroy()
        multi_wafer_boxplot_canvas = None

    # Clear distribution display
    if multi_wafer_distribution_canvas:
        multi_wafer_distribution_canvas.get_tk_widget().destroy()
        multi_wafer_distribution_canvas = None

    # Clear parameter combobox
    multi_wafer_param_combobox["values"] = []
    multi_wafer_param_combobox.set("")

    # Clear header comparison comboboxes
    multi_wafer_ref_combobox["values"] = []
    multi_wafer_ref_combobox.set("")
    multi_wafer_compare_combobox["values"] = []
    multi_wafer_compare_combobox.set("")

    # Clear header comparison treeview
    for item in multi_wafer_header_tree.get_children():
        multi_wafer_header_tree.delete(item)
    multi_wafer_header_summary_label.config(text="No wafermaps loaded", fg="gray")

    # Update info label
    multi_wafer_info_label.config(text="No wafermaps loaded", fg="gray")

    # Clear wafer selection list
    update_wafer_selection_list()

    print("Cleared all loaded wafermaps")


def add_multi_wafer_csv_files():
    """Add multiple CSV files to the Multiple Wafermaps tab (appends to existing)"""
    global multi_wafer_stdf_data, multi_wafer_wafer_ids, multi_wafer_test_params
    global multi_wafer_wafer_configs

    csv_paths = filedialog.askopenfilenames(
        title="Select CSV files to add",
        filetypes=[("CSV files", "*.csv"), ("All files", "*.*")],
    )

    if not csv_paths:
        print("No files selected.")
        return

    num_added = 0

    for csv_path in csv_paths:
        try:
            print(f"Loading CSV file: {csv_path}")
            df = pd.read_csv(csv_path)

            # Find coordinate columns
            x_col_candidates = ['x', 'X', 'x_coord', 'X_COORD', 'X_Coordinate', 'x_coordinate', 'DIE_X', 'die_x', 'col', 'COL', 'Column']
            y_col_candidates = ['y', 'Y', 'y_coord', 'Y_COORD', 'Y_Coordinate', 'y_coordinate', 'DIE_Y', 'die_y', 'row', 'ROW', 'Row']

            x_col = None
            y_col = None

            for candidate in x_col_candidates:
                if candidate in df.columns:
                    x_col = candidate
                    break

            for candidate in y_col_candidates:
                if candidate in df.columns:
                    y_col = candidate
                    break

            if x_col is None or y_col is None:
                numeric_cols = df.select_dtypes(include=[np.number]).columns.tolist()
                if len(numeric_cols) >= 2:
                    x_col = numeric_cols[0]
                    y_col = numeric_cols[1]
                else:
                    print(f"ERROR: Could not identify coordinate columns in {csv_path}")
                    continue

            df = df.rename(columns={x_col: 'x', y_col: 'y'})

            # Look for bin column
            bin_col_candidates = ['bin', 'BIN', 'Bin', 'HARD_BIN', 'hard_bin', 'SOFT_BIN', 'soft_bin', 'HB', 'SB']
            for candidate in bin_col_candidates:
                if candidate in df.columns:
                    df = df.rename(columns={candidate: 'bin'})
                    break

            if 'bin' not in df.columns:
                df['bin'] = 1

            wafer_id = os.path.basename(csv_path).replace('.csv', '').replace('.CSV', '')

            # Build test parameters from numeric columns
            numeric_columns = df.select_dtypes(include=[np.number]).columns.tolist()
            exclude_cols = ['x', 'y', 'bin']
            test_columns = [col for col in numeric_columns if col not in exclude_cols]

            for idx, col in enumerate(test_columns):
                test_num = idx + 1
                test_key = f"test_{test_num}"

                if test_key not in multi_wafer_test_params:
                    multi_wafer_test_params[test_key] = col

                df = df.rename(columns={col: test_num})

            # Detect notch orientation from CSV
            notch_orientation = _detect_notch_from_csv(df, csv_path)
            wafer_config = {
                'notch_orientation': notch_orientation,
                'wafer_size': None,
                'die_width': None,
                'die_height': None,
                'pos_x': None,
                'pos_y': None
            }
            print(f"CSV notch orientation for {wafer_id}: {notch_orientation}")

            multi_wafer_stdf_data.append(df)
            multi_wafer_wafer_ids.append(wafer_id)
            multi_wafer_wafer_configs.append(wafer_config)
            num_added += 1
            print(f"Added CSV: {wafer_id} with {len(df)} dies")

        except Exception as e:
            print(f"Error loading CSV {csv_path}: {e}")

    if num_added > 0:
        # Update parameter combobox
        param_options = ["BIN (Bin Number)"]
        for test_key, test_name in sorted(multi_wafer_test_params.items()):
            param_options.append(f"{test_key}: {test_name}")

        multi_wafer_param_combobox["values"] = param_options
        if not multi_wafer_param_combobox.get():
            multi_wafer_param_combobox.current(0)

        multi_wafer_info_label.config(
            text=f"Total: {len(multi_wafer_stdf_data)} wafermaps (added {num_added})",
            fg="green"
        )

        # Update wafer selection list
        update_wafer_selection_list()

        # Update reference combobox for alignment
        update_multi_wafer_ref_combobox()

        # Detect alignment status
        detect_multi_wafer_alignment()

        update_multi_wafer_display()
    else:
        print("No valid CSV files added")


def add_multi_wafer_stdf_files():
    """Load multiple STDF files for the Multiple Wafermaps tab - THREADED VERSION"""
    global multi_wafer_stdf_data, multi_wafer_wafer_ids, multi_wafer_test_params

    stdf_paths = filedialog.askopenfilenames(
        title="Select multiple STDF files for comparison",
        filetypes=[("STDF files", "*.stdf *.std"), ("All files", "*.*")],
    )

    if not stdf_paths:
        print("No files selected.")
        return

    num_files = len(stdf_paths)

    # Create progress dialog
    progress_win = tk.Toplevel(main_win)
    progress_win.title("Loading Multiple Wafermaps")
    progress_win.geometry("500x200")
    progress_win.transient(main_win)
    progress_win.grab_set()

    # Center the dialog
    progress_win.update_idletasks()
    x = main_win.winfo_x() + (main_win.winfo_width() - 500) // 2
    y = main_win.winfo_y() + (main_win.winfo_height() - 200) // 2
    progress_win.geometry(f"+{x}+{y}")

    # Progress UI elements
    title_label = tk.Label(
        progress_win,
        text=f"Loading {num_files} STDF file(s) for comparison...",
        font=("Helvetica", 12, "bold")
    )
    title_label.pack(pady=10)

    progress_var = tk.DoubleVar(value=0)
    progress_bar = ttk.Progressbar(
        progress_win,
        variable=progress_var,
        maximum=num_files,
        length=400,
        mode='determinate'
    )
    progress_bar.pack(pady=10)

    status_label = tk.Label(progress_win, text="Initializing...", font=("Helvetica", 10))
    status_label.pack(pady=5)

    time_label = tk.Label(progress_win, text="", font=("Helvetica", 9), fg="blue")
    time_label.pack(pady=5)

    # Storage for results
    results = []
    results_lock = threading.Lock()
    completed_count = [0]
    start_time = [time.time()]
    cancelled = [False]

    def load_single_file(stdf_path):
        if cancelled[0]:
            return None

        filename = os.path.basename(stdf_path)
        try:
            df, wafer_id, test_params, grouped_params, test_limits_dict, wafer_cfg = read_wafermap_from_stdf(stdf_path)
            if not df.empty:
                return {
                    'df': df,
                    'wafer_id': wafer_id if wafer_id else filename,
                    'test_params': test_params,
                    'grouped_params': grouped_params,
                    'test_limits': test_limits_dict,
                    'wafer_config': wafer_cfg,
                    'filename': filename
                }
        except Exception as e:
            print(f"[Thread] Error loading {filename}: {e}")
        return None

    def on_file_complete(future, stdf_path):
        if cancelled[0]:
            return

        try:
            result = future.result()

            with results_lock:
                if result:
                    results.append(result)
                completed_count[0] += 1
                count = completed_count[0]

            def update_ui():
                if cancelled[0] or not progress_win.winfo_exists():
                    return

                elapsed = time.time() - start_time[0]
                files_per_sec = count / elapsed if elapsed > 0 else 0
                eta = (num_files - count) / files_per_sec if files_per_sec > 0 else 0

                progress_var.set(count)
                status_label.config(text=f"Loaded {count}/{num_files} files")
                time_label.config(text=f"Speed: {files_per_sec:.1f} files/sec | ETA: {eta:.1f}s")

                if count >= num_files:
                    finish_loading()

            main_win.after(0, update_ui)
        except Exception as e:
            print(f"Error in callback: {e}")

    def finish_loading():
        global multi_wafer_stdf_data, multi_wafer_wafer_ids, multi_wafer_test_params
        global multi_wafer_test_params_per_wafer, multi_wafer_test_limits_per_wafer
        global multi_wafer_wafer_configs

        elapsed = time.time() - start_time[0]

        # Store count before adding new files
        existing_count = len(multi_wafer_stdf_data)

        # Append new results to existing data (don't reset)
        for result in results:
            multi_wafer_stdf_data.append(result['df'])
            multi_wafer_wafer_ids.append(result['wafer_id'])

            # Store per-wafer test params and limits for comparison
            multi_wafer_test_params_per_wafer.append(result['test_params'])
            multi_wafer_test_limits_per_wafer.append(result.get('test_limits', {}))
            multi_wafer_wafer_configs.append(result.get('wafer_config', {}))

            if not multi_wafer_test_params:
                multi_wafer_test_params = result['test_params']
            else:
                multi_wafer_test_params.update(result['test_params'])

        num_added = len(multi_wafer_stdf_data) - existing_count

        if multi_wafer_stdf_data:
            param_options = ["BIN (Bin Number)"]
            for test_key, test_name in sorted(multi_wafer_test_params.items()):
                param_options.append(f"{test_key}: {test_name}")

            multi_wafer_param_combobox["values"] = param_options
            if not multi_wafer_param_combobox.get():
                multi_wafer_param_combobox.current(0)

        multi_wafer_info_label.config(
            text=f"Total: {len(multi_wafer_stdf_data)} wafermaps (added {num_added} in {elapsed:.1f}s)",
            fg="green"
        )

        print(f"Added {num_added} wafermaps in {elapsed:.1f}s (Total: {len(multi_wafer_stdf_data)})")

        progress_win.destroy()

        # Update wafer selection list
        update_wafer_selection_list()

        # Update header comparison comboboxes
        update_header_comparison_comboboxes()

        # Update reference combobox for alignment
        update_multi_wafer_ref_combobox()

        # Detect alignment status
        detect_multi_wafer_alignment()

        if multi_wafer_stdf_data:
            update_multi_wafer_display()

    def on_cancel():
        cancelled[0] = True
        progress_win.destroy()

    cancel_btn = tk.Button(progress_win, text="Cancel", command=on_cancel, font=("Helvetica", 10))
    cancel_btn.pack(pady=10)

    # Submit all files to thread pool
    for stdf_path in stdf_paths:
        future = thread_pool.submit(load_single_file, stdf_path)
        future.add_done_callback(lambda f, p=stdf_path: on_file_complete(f, p))

    print(f"Submitted {num_files} files to thread pool for Multiple Wafermaps")


# Global variable for plot data cache
multi_wafer_plot_data_cache = {}
multi_wafer_highlight_rectangles = []  # Store highlight rectangles for clearing
multi_wafer_current_popup = None  # Store current popup reference to close it when new one opens

def show_die_data_popup_all_wafers(x_coord, y_coord):
    """Show a popup window with die data from ALL wafers at the clicked position"""
    global multi_wafer_highlight_rectangles, multi_wafer_current_popup

    # Close previous popup if exists
    if multi_wafer_current_popup is not None:
        try:
            multi_wafer_current_popup.destroy()
        except:
            pass
        multi_wafer_current_popup = None

    # Get the currently selected parameter
    selected = multi_wafer_param_combobox.get()
    if selected.startswith("BIN"):
        selected_param_column = "bin"
    else:
        test_key = selected.split(":")[0].strip()
        if test_key.startswith("test_"):
            selected_param_column = int(test_key.replace("test_", ""))
        else:
            selected_param_column = int(test_key)

    # Clear previous highlights
    for rect in multi_wafer_highlight_rectangles:
        try:
            rect.remove()
        except:
            pass
    multi_wafer_highlight_rectangles = []

    # Highlight the die position on ALL plots
    axes_info_list = multi_wafer_plot_data_cache.get('axes_info', [])
    fig = None

    for info in axes_info_list:
        ax = info['ax']
        fig = ax.figure
        x_min = info['x_min']
        y_min = info['y_min']

        # Calculate grid position
        grid_x = x_coord - x_min
        grid_y = y_coord - y_min

        # Draw highlight rectangle with black border
        rect = plt.Rectangle(
            (grid_x - 0.5, grid_y - 0.5), 1, 1,
            fill=False, edgecolor='black', linewidth=3, zorder=100
        )
        ax.add_patch(rect)
        multi_wafer_highlight_rectangles.append(rect)

    # Redraw canvas
    if fig is not None:
        fig.canvas.draw_idle()

    # Collect data from ALL wafers at this position
    all_wafer_data = []
    for info in axes_info_list:
        df = info['df']
        wafer_id = info['wafer_id']

        # Find the die at this coordinate
        die_data = df[(df['x'] == x_coord) & (df['y'] == y_coord)]

        if len(die_data) > 0:
            all_wafer_data.append({
                'wafer_id': wafer_id,
                'die_row': die_data.iloc[0]
            })

    if not all_wafer_data:
        print(f"No data found at position X={x_coord}, Y={y_coord}")
        return

    # Create popup window (75% of original size: 450x375 instead of 600x500)
    popup = tk.Toplevel(main_win)
    popup.title(f"Die Data - Position: X={x_coord}, Y={y_coord} - {len(all_wafer_data)} Wafers")
    popup.geometry("450x375")
    popup.transient(main_win)

    # Store reference to current popup
    multi_wafer_current_popup = popup

    # Clear highlights when popup is closed
    def on_popup_close():
        global multi_wafer_current_popup
        for rect in multi_wafer_highlight_rectangles:
            try:
                rect.remove()
            except:
                pass
        multi_wafer_highlight_rectangles.clear()
        multi_wafer_current_popup = None
        if fig is not None:
            fig.canvas.draw_idle()
        popup.destroy()

    popup.protocol("WM_DELETE_WINDOW", on_popup_close)

    # Center the popup
    popup.update_idletasks()
    x = main_win.winfo_x() + (main_win.winfo_width() - 450) // 2
    y = main_win.winfo_y() + (main_win.winfo_height() - 375) // 2
    popup.geometry(f"+{x}+{y}")

    # Header frame
    header_frame = tk.Frame(popup, bg="#2196F3", pady=5)
    header_frame.pack(fill=tk.X)

    tk.Label(
        header_frame,
        text=f"Position: X = {x_coord}, Y = {y_coord}",
        font=("Helvetica", 10, "bold"),
        bg="#2196F3",
        fg="white"
    ).pack()

    tk.Label(
        header_frame,
        text=f"Data from {len(all_wafer_data)} Wafer(s)",
        font=("Helvetica", 9),
        bg="#2196F3",
        fg="white"
    ).pack()

    # Legend frame
    legend_frame = tk.Frame(popup, pady=3)
    legend_frame.pack(fill=tk.X)
    tk.Label(legend_frame, text="■ Selected Parameter", fg="red", font=("Helvetica", 8, "bold")).pack(side=tk.LEFT, padx=5)

    # Create Treeview with scrollbars - dynamic columns based on number of wafers
    tree_frame = tk.Frame(popup)
    tree_frame.pack(fill=tk.BOTH, expand=True, padx=5, pady=3)

    tree_scroll_y = ttk.Scrollbar(tree_frame, orient="vertical")
    tree_scroll_x = ttk.Scrollbar(tree_frame, orient="horizontal")

    # Build column list: Parameter + one column per wafer
    columns = ["parameter"]
    for wd in all_wafer_data:
        short_id = str(wd['wafer_id'])[:15] + "..." if len(str(wd['wafer_id'])) > 15 else str(wd['wafer_id'])
        columns.append(short_id)

    tree = ttk.Treeview(
        tree_frame,
        columns=columns,
        show="headings",
        yscrollcommand=tree_scroll_y.set,
        xscrollcommand=tree_scroll_x.set
    )

    # Set headings
    tree.heading("parameter", text="Parameter")
    tree.column("parameter", width=150, anchor="w")

    col_width = max(60, 250 // len(all_wafer_data)) if all_wafer_data else 80
    for i, wd in enumerate(all_wafer_data):
        short_id = str(wd['wafer_id'])[:15] + "..." if len(str(wd['wafer_id'])) > 15 else str(wd['wafer_id'])
        tree.heading(columns[i+1], text=short_id)
        tree.column(columns[i+1], width=col_width, anchor="center")

    tree_scroll_y.config(command=tree.yview)
    tree_scroll_x.config(command=tree.xview)

    tree_scroll_y.pack(side=tk.RIGHT, fill=tk.Y)
    tree_scroll_x.pack(side=tk.BOTTOM, fill=tk.X)
    tree.pack(fill=tk.BOTH, expand=True)

    # Configure tags for highlighting
    tree.tag_configure("selected", background="#FFCDD2", foreground="red")
    tree.tag_configure("normal", background="white")
    tree.tag_configure("coordinate", background="#E3F2FD")

    # Get test params for readable names
    test_params = multi_wafer_test_params

    # Get all unique columns from all wafers
    all_columns = set()
    for wd in all_wafer_data:
        all_columns.update(wd['die_row'].index)

    # Remove x, y from display columns
    display_columns = [c for c in all_columns if c not in ['x', 'y']]

    # Sort columns: bin first, then numeric test numbers
    def sort_key(col):
        if col == 'bin':
            return (0, 0)
        elif isinstance(col, int):
            return (1, col)
        else:
            return (2, str(col))

    display_columns = sorted(display_columns, key=sort_key)

    # Add data rows
    for col in display_columns:
        # Determine parameter name
        if col == 'bin':
            param_name = "BIN"
        else:
            test_key = f"test_{col}"
            if test_key in test_params:
                param_name = f"{col}: {test_params[test_key]}"
            else:
                param_name = f"Test {col}"

        # Determine if this is the selected parameter
        is_selected = False
        if col == selected_param_column:
            is_selected = True
        elif isinstance(col, int) and col == selected_param_column:
            is_selected = True

        # Collect values from all wafers
        row_values = [param_name]
        for wd in all_wafer_data:
            die_row = wd['die_row']
            if col in die_row.index:
                value = die_row[col]
                if pd.notna(value):
                    if isinstance(value, float):
                        row_values.append(f"{value:.4g}")
                    else:
                        row_values.append(str(value))
                else:
                    row_values.append("-")
            else:
                row_values.append("-")

        tag = "selected" if is_selected else "normal"
        tree.insert("", "end", values=row_values, tags=(tag,))

    # Button frame
    btn_frame = tk.Frame(popup, pady=5)
    btn_frame.pack(fill=tk.X)

    def copy_to_clipboard():
        """Copy die data to clipboard"""
        data_str = f"Position: X={x_coord}, Y={y_coord}\n\n"
        # Header row
        header = "Parameter"
        for wd in all_wafer_data:
            header += f"\t{wd['wafer_id']}"
        data_str += header + "\n"

        for item in tree.get_children():
            values = tree.item(item, "values")
            data_str += "\t".join(str(v) for v in values) + "\n"

        popup.clipboard_clear()
        popup.clipboard_append(data_str)
        print("Die data copied to clipboard")

    tk.Button(
        btn_frame,
        text="Copy",
        command=copy_to_clipboard,
        font=("Helvetica", 8),
        bg="#4CAF50",
        fg="white"
    ).pack(side=tk.LEFT, padx=5)

    tk.Button(
        btn_frame,
        text="Close",
        command=on_popup_close,
        font=("Helvetica", 8)
    ).pack(side=tk.RIGHT, padx=5)


def update_multi_wafer_display():
    """Update the multiple wafermaps display - only show selected wafers"""
    global multi_wafer_canvas, multi_wafer_stdf_data, multi_wafer_wafer_ids
    global multi_wafer_plot_data_cache  # Store data for click handler

    # Remove independent selector frame if it exists
    if hasattr(create_independent_param_selectors, 'selector_frame') and create_independent_param_selectors.selector_frame:
        create_independent_param_selectors.selector_frame.destroy()
        create_independent_param_selectors.selector_frame = None

    if not multi_wafer_stdf_data:
        print("No wafermap data loaded for Multiple Wafermaps tab")
        return

    # Get selected wafer indices
    selected_indices = get_selected_wafer_indices()

    if not selected_indices:
        print("No wafers selected for display")
        # Clear existing canvas
        if multi_wafer_canvas:
            multi_wafer_canvas.get_tk_widget().destroy()
            multi_wafer_canvas = None
        return

    selected = multi_wafer_param_combobox.get()

    if not selected:
        print("No parameter selected")
        return

    # Parse parameter selection
    if selected.startswith("BIN"):
        param_column = "bin"
        param_label = "Bin"
    else:
        test_key = selected.split(":")[0].strip()
        if test_key.startswith("test_"):
            param_column = int(test_key.replace("test_", ""))
        else:
            param_column = int(test_key)
        param_label = selected

    # Get selected data and IDs
    selected_data = [multi_wafer_stdf_data[i] for i in selected_indices]
    selected_ids = [multi_wafer_wafer_ids[i] for i in selected_indices]

    # Store cache for click handler
    multi_wafer_plot_data_cache = {
        'selected_data': selected_data,
        'selected_ids': selected_ids,
        'selected_indices': selected_indices,
        'param_column': param_column,
        'param_label': param_label,
        'axes_info': []  # Will store (ax, x_min, y_min, df, wafer_id) for each plot
    }

    # Determine layout based on selected wafers
    num_files = len(selected_data)
    layout_choice = multi_wafer_layout_var.get()

    if layout_choice == "1 Column":
        cols = 1
    elif layout_choice == "2 Columns":
        cols = 2
    elif layout_choice == "3 Columns":
        cols = 3
    elif layout_choice == "4 Columns":
        cols = 4
    else:  # Auto
        if num_files <= 2:
            cols = num_files
        elif num_files <= 4:
            cols = 2
        elif num_files <= 9:
            cols = 3
        else:
            cols = 4

    rows = (num_files + cols - 1) // cols

    # Calculate figure size based on number of wafermaps
    fig_width = min(5 * cols, 20)
    fig_height = min(5 * rows, 16)

    fig, axes = plt.subplots(rows, cols, figsize=(fig_width, fig_height), constrained_layout=True)

    # Flatten axes for easy iteration
    if num_files == 1:
        axes = [axes]
    elif rows == 1 or cols == 1:
        axes = axes.flatten() if hasattr(axes, 'flatten') else [axes]
    else:
        axes = axes.flatten()

    # Choose colormap
    cmap = "tab20" if param_column == "bin" else "viridis"

    # Track min/max for consistent color scale across all plots
    global_min = float('inf')
    global_max = float('-inf')

    # First pass: compute global min/max for selected wafers only
    for df in selected_data:
        if param_column in df.columns:
            mask = df[param_column].notna()
            values = df.loc[mask, param_column].values
            if len(values) > 0:
                global_min = min(global_min, np.nanmin(values))
                global_max = max(global_max, np.nanmax(values))

    # Store imshow references for colorbar slider
    imshow_list = []
    first_cbar = None

    # Second pass: create plots for selected wafers
    for idx, (df, wafer_id) in enumerate(zip(selected_data, selected_ids)):
        ax = axes[idx]

        if param_column not in df.columns:
            ax.set_title(f"{wafer_id}\nParameter not found")
            ax.axis('off')
            continue

        mask = df[param_column].notna()
        plot_data = df[mask]

        if len(plot_data) == 0:
            ax.set_title(f"{wafer_id}\nNo data")
            ax.axis('off')
            continue

        # Compute grid
        grid, x_min, y_min, x_max, y_max, grid_width, grid_height = _compute_grid_fast(
            plot_data, param_column
        )

        # Store axes info for click handler
        multi_wafer_plot_data_cache['axes_info'].append({
            'ax': ax,
            'x_min': x_min,
            'y_min': y_min,
            'x_max': x_max,
            'y_max': y_max,
            'df': df,
            'wafer_id': wafer_id,
            'wafer_idx': selected_indices[idx]
        })

        im = ax.imshow(
            grid,
            cmap=cmap,
            aspect="equal",
            interpolation="nearest",
            origin="upper",
            vmin=global_min,
            vmax=global_max
        )
        imshow_list.append(im)

        # Truncate long wafer IDs
        short_id = str(wafer_id)[:25] + "..." if len(str(wafer_id)) > 25 else str(wafer_id)
        ax.set_title(f"{short_id}", fontsize=9, fontweight="bold")
        ax.set_xlabel("X", fontsize=8)
        ax.set_ylabel("Y", fontsize=8)
        ax.tick_params(axis='both', labelsize=7)

        # Set correct axis tick labels to show real coordinates
        num_x_ticks = min(6, grid_width)
        num_y_ticks = min(6, grid_height)
        x_tick_positions = np.linspace(0, grid_width - 1, num_x_ticks)
        y_tick_positions = np.linspace(0, grid_height - 1, num_y_ticks)
        ax.set_xticks(x_tick_positions)
        ax.set_yticks(y_tick_positions)
        ax.set_xticklabels([f"{int(x_min + pos)}" for pos in x_tick_positions])
        ax.set_yticklabels([f"{int(y_min + pos)}" for pos in y_tick_positions])

        # Show grid if enabled
        if multi_wafer_grid_var.get():
            ax.set_xticks(np.arange(-0.5, grid_width, 1), minor=True)
            ax.set_yticks(np.arange(-0.5, grid_height, 1), minor=True)
            ax.grid(which="minor", color="black", linewidth=0.2)
            ax.tick_params(which="minor", size=0)

        # Draw notch marker if orientation is available for this wafer
        if idx < len(multi_wafer_wafer_configs) and multi_wafer_wafer_configs[selected_indices[idx]]:
            wafer_cfg = multi_wafer_wafer_configs[selected_indices[idx]]
            if wafer_cfg.get('notch_orientation'):
                draw_notch_marker(ax, grid_width, grid_height, wafer_cfg['notch_orientation'],
                                 marker_size=min(grid_width, grid_height) * 0.06)

        # Add colorbar only to first plot
        if idx == 0:
            cbar = fig.colorbar(im, ax=ax, fraction=0.046, pad=0.04)
            cbar.ax.tick_params(labelsize=7)
            first_cbar = cbar

    # Hide unused axes
    for idx in range(num_files, len(axes)):
        axes[idx].axis('off')

    # Add a shared title
    fig.suptitle(f"Multiple Wafermaps: {param_label}", fontsize=12, fontweight="bold")

    # Initialize slider_state for click handler reference
    slider_state = None

    # Add professional sliders to the first colorbar
    if first_cbar is not None and len(imshow_list) > 0:
        data_min = global_min
        data_max = global_max

        # Create triangular arrow markers
        low_marker = first_cbar.ax.scatter(
            [0.5], [data_min],
            marker='^', s=150, c='#D32F2F',
            edgecolors='white', linewidths=1.5,
            zorder=15, clip_on=False
        )

        high_marker = first_cbar.ax.scatter(
            [0.5], [data_max],
            marker='v', s=150, c='#1976D2',
            edgecolors='white', linewidths=1.5,
            zorder=15, clip_on=False
        )

        # Add horizontal limit lines
        low_line, = first_cbar.ax.plot(
            [-0.6, 1.6], [data_min, data_min],
            color='#D32F2F', linewidth=2,
            solid_capstyle='round', zorder=12,
            transform=first_cbar.ax.get_yaxis_transform()
        )

        high_line, = first_cbar.ax.plot(
            [-0.6, 1.6], [data_max, data_max],
            color='#1976D2', linewidth=2,
            solid_capstyle='round', zorder=12,
            transform=first_cbar.ax.get_yaxis_transform()
        )

        # Value labels
        low_text = first_cbar.ax.text(
            1.8, data_min, f'{data_min:.2f}',
            fontsize=6, fontweight='bold',
            color='#D32F2F', va='center', ha='left',
            transform=first_cbar.ax.get_yaxis_transform()
        )

        high_text = first_cbar.ax.text(
            1.8, data_max, f'{data_max:.2f}',
            fontsize=6, fontweight='bold',
            color='#1976D2', va='center', ha='left',
            transform=first_cbar.ax.get_yaxis_transform()
        )

        # Slider state
        slider_state = {
            'low_val': data_min,
            'high_val': data_max,
            'data_min': data_min,
            'data_max': data_max,
            'dragging': None,
            'low_marker': low_marker,
            'high_marker': high_marker,
            'low_line': low_line,
            'high_line': high_line,
            'low_text': low_text,
            'high_text': high_text,
            'cbar_ax': first_cbar.ax
        }

        def on_multi_press(event):
            if event.inaxes == slider_state['cbar_ax']:
                y_click = event.ydata
                if y_click is not None:
                    dist_to_low = abs(y_click - slider_state['low_val'])
                    dist_to_high = abs(y_click - slider_state['high_val'])
                    threshold = 0.1 * (slider_state['data_max'] - slider_state['data_min'])

                    if dist_to_low < dist_to_high and dist_to_low < threshold:
                        slider_state['dragging'] = 'low'
                        low_marker.set_sizes([220])
                    elif dist_to_high < threshold:
                        slider_state['dragging'] = 'high'
                        high_marker.set_sizes([220])
                    fig.canvas.draw_idle()

        def on_multi_motion(event):
            if slider_state['dragging']:
                if event.ydata is not None:
                    y_click = event.ydata
                elif event.y is not None:
                    inv = slider_state['cbar_ax'].transData.inverted()
                    _, y_click = inv.transform((event.x, event.y))
                else:
                    return

                y_click = max(slider_state['data_min'], min(slider_state['data_max'], y_click))
                min_gap = 0.02 * (slider_state['data_max'] - slider_state['data_min'])

                if slider_state['dragging'] == 'low':
                    new_val = min(y_click, slider_state['high_val'] - min_gap)
                    new_val = max(new_val, slider_state['data_min'])
                    slider_state['low_val'] = new_val
                    low_marker.set_offsets([[0.5, new_val]])
                    low_line.set_ydata([new_val, new_val])
                    low_text.set_position((1.8, new_val))
                    low_text.set_text(f'{new_val:.2f}')

                elif slider_state['dragging'] == 'high':
                    new_val = max(y_click, slider_state['low_val'] + min_gap)
                    new_val = min(new_val, slider_state['data_max'])
                    slider_state['high_val'] = new_val
                    high_marker.set_offsets([[0.5, new_val]])
                    high_line.set_ydata([new_val, new_val])
                    high_text.set_position((1.8, new_val))
                    high_text.set_text(f'{new_val:.2f}')

                # Update all heatmaps
                for img in imshow_list:
                    img.set_clim(vmin=slider_state['low_val'], vmax=slider_state['high_val'])

                fig.canvas.draw_idle()

        def on_multi_release(event):
            if slider_state['dragging'] == 'low':
                low_marker.set_sizes([150])
            elif slider_state['dragging'] == 'high':
                high_marker.set_sizes([150])
            slider_state['dragging'] = None
            fig.canvas.draw_idle()

        fig.canvas.mpl_connect('button_press_event', on_multi_press)
        fig.canvas.mpl_connect('motion_notify_event', on_multi_motion)
        fig.canvas.mpl_connect('button_release_event', on_multi_release)

    # Add zoom functionality (mouse wheel and rectangle selection)
    # Store original axis limits for reset
    original_axis_limits = {}
    for info in multi_wafer_plot_data_cache.get('axes_info', []):
        ax = info['ax']
        original_axis_limits[id(ax)] = {
            'xlim': ax.get_xlim(),
            'ylim': ax.get_ylim()
        }

    # Selection and zoom state
    zoom_state = {
        'selected_ax': None,           # Currently selected/active axis for zoom
        'selection_rect': None,        # Visual selection rectangle around selected heatmap
        'start_point': None,           # Start point for rectangle zoom
        'start_pixel': None,           # Pixel position for drag detection
        'zoom_rect': None,             # Rectangle being drawn for zoom
        'is_dragging': False           # Flag to track if user is dragging
    }

    def highlight_selected_axis(ax):
        """Add visual highlight to selected axis"""
        # Remove old selection rectangle from fig.patches
        if zoom_state['selection_rect'] is not None:
            try:
                if zoom_state['selection_rect'] in fig.patches:
                    fig.patches.remove(zoom_state['selection_rect'])
            except:
                pass
            zoom_state['selection_rect'] = None

        if ax is None:
            fig.canvas.draw_idle()
            return

        # Add blue border around selected axis
        bbox = ax.get_position()
        selection_rect = plt.Rectangle(
            (bbox.x0, bbox.y0), bbox.width, bbox.height,
            transform=fig.transFigure,
            fill=False, edgecolor='blue', linewidth=3, linestyle='-',
            zorder=1000
        )
        fig.patches.append(selection_rect)
        zoom_state['selection_rect'] = selection_rect
        fig.canvas.draw_idle()

    def on_multi_wafer_press(event):
        """Handle mouse press - left click selects heatmap or starts zoom drag"""
        if event.inaxes is None:
            return

        # Skip if click is on colorbar
        if slider_state is not None and event.inaxes == slider_state.get('cbar_ax'):
            return

        if event.button == 1:  # Left click
            clicked_ax = event.inaxes

            # Check if click is on a valid heatmap axis
            axes_info_list = multi_wafer_plot_data_cache.get('axes_info', [])
            is_valid_ax = any(info['ax'] == clicked_ax for info in axes_info_list)

            if not is_valid_ax:
                return

            # If no axis selected yet, or clicking on a different axis -> select it
            if zoom_state['selected_ax'] is None or zoom_state['selected_ax'] != clicked_ax:
                zoom_state['selected_ax'] = clicked_ax
                highlight_selected_axis(clicked_ax)
                zoom_state['start_point'] = None
                zoom_state['is_dragging'] = False
            else:
                # Already selected - start rectangle zoom
                zoom_state['start_point'] = (event.xdata, event.ydata)
                zoom_state['start_pixel'] = (event.x, event.y)
                zoom_state['is_dragging'] = False

    def on_multi_wafer_motion(event):
        """Handle mouse motion for rectangle zoom on selected axis"""
        if event.inaxes is None:
            return

        # Only process if we have a start point (dragging for zoom)
        if zoom_state['start_point'] is None:
            return

        # Only allow zoom on selected axis
        if zoom_state['selected_ax'] is None or event.inaxes != zoom_state['selected_ax']:
            return

        # Check if user has moved enough to be considered dragging (5 pixels threshold)
        if zoom_state['start_pixel'] is not None:
            dx = abs(event.x - zoom_state['start_pixel'][0])
            dy = abs(event.y - zoom_state['start_pixel'][1])
            if dx > 5 or dy > 5:
                zoom_state['is_dragging'] = True

        # Only draw rectangle if actually dragging
        if not zoom_state['is_dragging']:
            return

        ax = zoom_state['selected_ax']

        # Remove old zoom rectangle
        if zoom_state['zoom_rect'] is not None:
            try:
                zoom_state['zoom_rect'].remove()
            except:
                pass

        x0, y0 = zoom_state['start_point']
        x1, y1 = event.xdata, event.ydata

        # Draw new zoom rectangle
        zoom_state['zoom_rect'] = plt.Rectangle(
            (min(x0, x1), min(y0, y1)),
            abs(x1 - x0), abs(y1 - y0),
            fill=False, edgecolor='red', linewidth=2, linestyle='--'
        )
        ax.add_patch(zoom_state['zoom_rect'])
        fig.canvas.draw_idle()

    def on_multi_wafer_release(event):
        """Handle mouse release for rectangle zoom"""
        if zoom_state['start_point'] is None:
            return

        was_dragging = zoom_state['is_dragging']
        ax = zoom_state['selected_ax']

        # Remove zoom rectangle
        if zoom_state['zoom_rect'] is not None:
            try:
                zoom_state['zoom_rect'].remove()
            except:
                pass
            zoom_state['zoom_rect'] = None

        # Apply zoom if was dragging on selected axis
        if was_dragging and ax is not None and event.xdata is not None and event.ydata is not None:
            x0, y0 = zoom_state['start_point']
            x1, y1 = event.xdata, event.ydata

            # Check if rectangle is significant
            if abs(x1 - x0) > 0.5 and abs(y1 - y0) > 0.5:
                ax.set_xlim(min(x0, x1), max(x0, x1))
                ax.set_ylim(min(y0, y1), max(y0, y1))
                fig.canvas.draw_idle()

        # Reset drag state
        zoom_state['start_point'] = None
        zoom_state['start_pixel'] = None
        zoom_state['is_dragging'] = False

    def on_multi_wafer_scroll(event):
        """Handle mouse wheel scroll to zoom - only on selected axis"""
        # Only zoom on selected axis
        if zoom_state['selected_ax'] is None:
            return

        if event.inaxes != zoom_state['selected_ax']:
            return

        ax = zoom_state['selected_ax']

        cur_xlim = ax.get_xlim()
        cur_ylim = ax.get_ylim()

        x_data = event.xdata
        y_data = event.ydata

        if x_data is None or y_data is None:
            return

        # Determine zoom direction
        zoom_in = False
        if hasattr(event, 'step') and event.step is not None:
            zoom_in = event.step > 0
        elif hasattr(event, 'button') and event.button in ['up', 'down']:
            zoom_in = (event.button == 'up')
        else:
            return

        zoom_factor = 0.8 if zoom_in else 1.25

        x_range = cur_xlim[1] - cur_xlim[0]
        y_range = cur_ylim[1] - cur_ylim[0]

        new_x_range = x_range * zoom_factor
        new_y_range = y_range * zoom_factor

        rel_x = (x_data - cur_xlim[0]) / x_range
        rel_y = (y_data - cur_ylim[0]) / y_range

        new_xlim = [x_data - rel_x * new_x_range, x_data + (1 - rel_x) * new_x_range]
        new_ylim = [y_data - rel_y * new_y_range, y_data + (1 - rel_y) * new_y_range]

        ax.set_xlim(new_xlim)
        ax.set_ylim(new_ylim)

        fig.canvas.draw_idle()

    # Double-click handler for die data display (ONLY double-click)
    def on_die_double_click(event):
        """Handle DOUBLE click on die to show measurement data"""
        # Only respond to double-click (dblclick)
        if not event.dblclick:
            return

        # Only respond to left button
        if event.button != 1:
            return

        if event.inaxes is None:
            return

        # Skip if click is on colorbar
        if slider_state is not None and event.inaxes == slider_state.get('cbar_ax'):
            return

        # Find which subplot was clicked
        clicked_ax = event.inaxes
        axes_info_list = multi_wafer_plot_data_cache.get('axes_info', [])

        for info in axes_info_list:
            if info['ax'] == clicked_ax:
                # Convert click coordinates to die coordinates
                x_click = int(round(event.xdata)) + info['x_min']
                y_click = int(round(event.ydata)) + info['y_min']

                # Show popup with data from ALL wafers at this position
                show_die_data_popup_all_wafers(x_click, y_click)
                break

    # Function to reset zoom to original
    def reset_multi_wafer_zoom():
        """Reset all axes to original zoom level"""
        for info in multi_wafer_plot_data_cache.get('axes_info', []):
            ax = info['ax']
            ax_id = id(ax)
            if ax_id in original_axis_limits:
                ax.set_xlim(original_axis_limits[ax_id]['xlim'])
                ax.set_ylim(original_axis_limits[ax_id]['ylim'])
        # Also clear selection
        zoom_state['selected_ax'] = None
        highlight_selected_axis(None)
        fig.canvas.draw_idle()

    # Store reset function for button access
    multi_wafer_plot_data_cache['reset_zoom_func'] = reset_multi_wafer_zoom

    # Connect events
    fig.canvas.mpl_connect('button_press_event', on_multi_wafer_press)
    fig.canvas.mpl_connect('motion_notify_event', on_multi_wafer_motion)
    fig.canvas.mpl_connect('button_release_event', on_multi_wafer_release)
    fig.canvas.mpl_connect('scroll_event', on_multi_wafer_scroll)
    fig.canvas.mpl_connect('button_press_event', on_die_double_click)  # Double-click for data

    # Destroy old canvas
    if multi_wafer_canvas:
        multi_wafer_canvas.get_tk_widget().destroy()

    # Create new canvas
    multi_wafer_canvas = FigureCanvasTkAgg(fig, master=multi_wafer_display_frame)
    canvas_widget = multi_wafer_canvas.get_tk_widget()
    canvas_widget.pack(fill=tk.BOTH, expand=True)

    # Bind tkinter native scroll events for Windows compatibility
    def on_multi_wafer_tk_mousewheel(event):
        """Handle Windows native mouse wheel events for zoom - only on selected axis"""
        # Only zoom if an axis is selected
        if zoom_state['selected_ax'] is None:
            return

        ax = zoom_state['selected_ax']

        cur_xlim = ax.get_xlim()
        cur_ylim = ax.get_ylim()

        # On Windows, event.delta is positive for scroll up, negative for scroll down
        zoom_in = event.delta > 0
        zoom_factor = 0.8 if zoom_in else 1.25

        # Calculate center of current view
        x_center = (cur_xlim[0] + cur_xlim[1]) / 2
        y_center = (cur_ylim[0] + cur_ylim[1]) / 2

        # Calculate current ranges
        x_range = cur_xlim[1] - cur_xlim[0]
        y_range = cur_ylim[1] - cur_ylim[0]

        # New ranges after zoom
        new_x_range = x_range * zoom_factor
        new_y_range = y_range * zoom_factor

        # New limits centered on current view center
        new_xlim = [x_center - new_x_range / 2, x_center + new_x_range / 2]
        new_ylim = [y_center - new_y_range / 2, y_center + new_y_range / 2]

        ax.set_xlim(new_xlim)
        ax.set_ylim(new_ylim)

        multi_wafer_canvas.draw_idle()

    canvas_widget.bind("<MouseWheel>", on_multi_wafer_tk_mousewheel)
    canvas_widget.bind("<Button-4>", lambda e: on_multi_wafer_tk_mousewheel(type('Event', (), {'delta': 120, 'x': e.x, 'y': e.y})))
    canvas_widget.bind("<Button-5>", lambda e: on_multi_wafer_tk_mousewheel(type('Event', (), {'delta': -120, 'x': e.x, 'y': e.y})))

    multi_wafer_canvas.draw()

    print(f"Multiple Wafermaps display updated with {num_files} selected wafermaps ({rows}x{cols} layout)")


def update_multi_wafer_independent_display():
    """Update the multiple wafermaps display with independent parameter selection per wafer"""
    global multi_wafer_canvas, multi_wafer_stdf_data, multi_wafer_wafer_ids
    global multi_wafer_plot_data_cache, multi_wafer_independent_params, multi_wafer_independent_comboboxes

    if not multi_wafer_stdf_data:
        print("No wafermap data loaded for Multiple Wafermaps tab")
        return

    # Get selected wafer indices
    selected_indices = get_selected_wafer_indices()

    if not selected_indices:
        print("No wafers selected for display")
        if multi_wafer_canvas:
            multi_wafer_canvas.get_tk_widget().destroy()
            multi_wafer_canvas = None
        return

    # Get selected data and IDs
    selected_data = [multi_wafer_stdf_data[i] for i in selected_indices]
    selected_ids = [multi_wafer_wafer_ids[i] for i in selected_indices]

    # Determine layout based on selected wafers
    num_files = len(selected_data)
    layout_choice = multi_wafer_layout_var.get()

    if layout_choice == "1 Column":
        cols = 1
    elif layout_choice == "2 Columns":
        cols = 2
    elif layout_choice == "3 Columns":
        cols = 3
    elif layout_choice == "4 Columns":
        cols = 4
    else:  # Auto
        if num_files <= 2:
            cols = num_files
        elif num_files <= 4:
            cols = 2
        elif num_files <= 9:
            cols = 3
        else:
            cols = 4

    rows = (num_files + cols - 1) // cols

    # Calculate figure size
    fig_width = min(5 * cols, 20)
    fig_height = min(5.5 * rows, 18)  # Extra height for combobox

    fig, axes = plt.subplots(rows, cols, figsize=(fig_width, fig_height), constrained_layout=True)

    # Flatten axes for easy iteration
    if num_files == 1:
        axes = [axes]
    elif rows == 1 or cols == 1:
        axes = axes.flatten() if hasattr(axes, 'flatten') else [axes]
    else:
        axes = axes.flatten()

    # Clear old combobox widgets
    for widget in multi_wafer_independent_comboboxes.values():
        if hasattr(widget, 'destroy'):
            widget.destroy()
    multi_wafer_independent_comboboxes.clear()

    # Store cache for click handler
    multi_wafer_plot_data_cache = {
        'selected_data': selected_data,
        'selected_ids': selected_ids,
        'selected_indices': selected_indices,
        'axes_info': [],
        'independent_mode': True
    }

    # Get available parameters from param combobox
    available_params = list(multi_wafer_param_combobox['values'])
    default_param = multi_wafer_param_combobox.get() if multi_wafer_param_combobox.get() else (available_params[0] if available_params else "")

    # Create plots for each selected wafer with individual parameter selection
    for idx, (df, wafer_id) in enumerate(zip(selected_data, selected_ids)):
        ax = axes[idx]
        wafer_idx = selected_indices[idx]

        # Get or set default parameter for this wafer
        if wafer_idx not in multi_wafer_independent_params:
            multi_wafer_independent_params[wafer_idx] = default_param

        selected_param = multi_wafer_independent_params.get(wafer_idx, default_param)

        # Parse parameter selection
        if not selected_param:
            ax.set_title(f"{wafer_id}\nNo parameter selected")
            ax.axis('off')
            continue

        if selected_param.startswith("BIN"):
            param_column = "bin"
            param_label = "Bin"
        else:
            test_key = selected_param.split(":")[0].strip()
            if test_key.startswith("test_"):
                param_column = int(test_key.replace("test_", ""))
            else:
                try:
                    param_column = int(test_key)
                except ValueError:
                    param_column = test_key
            param_label = selected_param

        if param_column not in df.columns:
            ax.set_title(f"{wafer_id}\nParameter not found")
            ax.axis('off')
            continue

        mask = df[param_column].notna()
        plot_data = df[mask]

        if len(plot_data) == 0:
            ax.set_title(f"{wafer_id}\nNo data")
            ax.axis('off')
            continue

        # Compute grid
        grid, x_min, y_min, x_max, y_max, grid_width, grid_height = _compute_grid_fast(
            plot_data, param_column
        )

        # Compute individual min/max for this wafer's parameter
        values = df.loc[mask, param_column].values
        vmin = np.nanmin(values)
        vmax = np.nanmax(values)

        # Store axes info for click handler
        multi_wafer_plot_data_cache['axes_info'].append({
            'ax': ax,
            'x_min': x_min,
            'y_min': y_min,
            'x_max': x_max,
            'y_max': y_max,
            'df': df,
            'wafer_id': wafer_id,
            'wafer_idx': wafer_idx,
            'param_column': param_column,
            'param_label': param_label
        })

        # Choose colormap
        cmap = "tab20" if param_column == "bin" else "viridis"

        im = ax.imshow(
            grid,
            cmap=cmap,
            aspect="equal",
            interpolation="nearest",
            origin="upper",
            vmin=vmin,
            vmax=vmax
        )

        # Truncate long wafer IDs
        short_id = str(wafer_id)[:25] + "..." if len(str(wafer_id)) > 25 else str(wafer_id)

        # Truncate parameter label for title
        short_param = param_label[:30] + "..." if len(param_label) > 30 else param_label
        ax.set_title(f"{short_id}\n{short_param}", fontsize=8, fontweight="bold")
        ax.set_xlabel("X", fontsize=7)
        ax.set_ylabel("Y", fontsize=7)
        ax.tick_params(axis='both', labelsize=6)

        # Show grid if enabled
        if multi_wafer_grid_var.get():
            ax.set_xticks(np.arange(-0.5, grid_width, 1), minor=True)
            ax.set_yticks(np.arange(-0.5, grid_height, 1), minor=True)
            ax.grid(which="minor", color="black", linewidth=0.2)
            ax.tick_params(which="minor", size=0)

        # Add individual colorbar
        cbar = fig.colorbar(im, ax=ax, fraction=0.046, pad=0.04)
        cbar.ax.tick_params(labelsize=6)

    # Hide unused axes
    for idx in range(num_files, len(axes)):
        axes[idx].axis('off')

    # Add a shared title
    fig.suptitle("Independent Wafermap Comparison", fontsize=12, fontweight="bold")

    # Destroy old canvas
    if multi_wafer_canvas:
        multi_wafer_canvas.get_tk_widget().destroy()

    # Create new canvas
    multi_wafer_canvas = FigureCanvasTkAgg(fig, master=multi_wafer_display_frame)
    multi_wafer_canvas.get_tk_widget().pack(fill=tk.BOTH, expand=True)

    # Create individual parameter selection comboboxes as overlays
    # We'll place them in a separate frame above the canvas
    create_independent_param_selectors(selected_indices, selected_ids, available_params, cols, rows)

    multi_wafer_canvas.draw()

    print(f"Independent Wafermaps display updated with {num_files} wafermaps")


def create_independent_param_selectors(selected_indices, selected_ids, available_params, cols, rows):
    """Create individual parameter selector comboboxes for each wafer in independent mode"""
    global multi_wafer_independent_comboboxes, multi_wafer_independent_selector_frame

    # Remove old selector frame if exists
    if hasattr(create_independent_param_selectors, 'selector_frame') and create_independent_param_selectors.selector_frame:
        create_independent_param_selectors.selector_frame.destroy()

    # Create a frame for parameter selectors above the canvas
    selector_frame = tk.Frame(multi_wafer_display_frame)
    selector_frame.pack(side=tk.TOP, fill=tk.X, before=multi_wafer_canvas.get_tk_widget())
    create_independent_param_selectors.selector_frame = selector_frame

    # Create a grid of selectors matching the wafermap layout
    num_wafers = len(selected_indices)

    for idx, (wafer_idx, wafer_id) in enumerate(zip(selected_indices, selected_ids)):
        row = idx // cols
        col = idx % cols

        # Create a small frame for each selector
        cell_frame = tk.Frame(selector_frame)
        cell_frame.grid(row=row, column=col, padx=5, pady=2, sticky='ew')

        # Short wafer ID label
        short_id = str(wafer_id)[:15] + "..." if len(str(wafer_id)) > 15 else str(wafer_id)
        label = tk.Label(cell_frame, text=short_id, font=("Helvetica", 8), width=18, anchor='w')
        label.pack(side=tk.LEFT, padx=2)

        # Combobox for parameter selection
        current_param = multi_wafer_independent_params.get(wafer_idx, available_params[0] if available_params else "")

        combo_var = tk.StringVar(value=current_param)
        combo = ttk.Combobox(
            cell_frame,
            textvariable=combo_var,
            values=available_params,
            state="readonly",
            width=35,
            font=("Helvetica", 8)
        )
        combo.pack(side=tk.LEFT, padx=2)

        # Bind selection change
        def on_param_change(event, w_idx=wafer_idx, var=combo_var):
            multi_wafer_independent_params[w_idx] = var.get()
            update_multi_wafer_independent_display()

        combo.bind("<<ComboboxSelected>>", on_param_change)

        multi_wafer_independent_comboboxes[wafer_idx] = combo

    # Configure grid columns to expand evenly
    for col in range(cols):
        selector_frame.grid_columnconfigure(col, weight=1)


def update_multi_wafer_boxplot():
    """Update the boxplot display for selected wafers"""
    global multi_wafer_boxplot_canvas, multi_wafer_stdf_data, multi_wafer_wafer_ids

    if not multi_wafer_stdf_data:
        print("No wafermap data loaded for Boxplot")
        return

    # Get selected wafer indices
    selected_indices = get_selected_wafer_indices()

    if not selected_indices:
        print("No wafers selected for boxplot")
        if multi_wafer_boxplot_canvas:
            multi_wafer_boxplot_canvas.get_tk_widget().destroy()
            multi_wafer_boxplot_canvas = None
        return

    selected = multi_wafer_param_combobox.get()

    if not selected:
        print("No parameter selected for boxplot")
        return

    # Parse parameter selection
    if selected.startswith("BIN"):
        param_column = "bin"
        param_label = "Bin"
    else:
        test_key = selected.split(":")[0].strip()
        if test_key.startswith("test_"):
            param_column = int(test_key.replace("test_", ""))
        else:
            param_column = int(test_key)
        param_label = selected

    # Get selected data and IDs
    selected_data = [multi_wafer_stdf_data[i] for i in selected_indices]
    selected_ids = [multi_wafer_wafer_ids[i] for i in selected_indices]

    # Collect data for boxplot
    boxplot_data = []
    labels = []

    for df, wafer_id in zip(selected_data, selected_ids):
        if param_column in df.columns:
            values = df[param_column].dropna().values
            if len(values) > 0:
                boxplot_data.append(values)
                # Truncate long wafer IDs for labels
                short_id = str(wafer_id)[:20] + "..." if len(str(wafer_id)) > 20 else str(wafer_id)
                labels.append(short_id)

    if not boxplot_data:
        print("No valid data for boxplot")
        return

    # Create figure
    fig, ax = plt.subplots(figsize=(max(8, len(labels) * 0.8), 6), constrained_layout=True)

    # Create boxplot with professional styling
    bp = ax.boxplot(boxplot_data, labels=labels, patch_artist=True, notch=True)

    # Color palette for boxes
    colors = plt.cm.Set3(np.linspace(0, 1, len(boxplot_data)))

    for patch, color in zip(bp['boxes'], colors):
        patch.set_facecolor(color)
        patch.set_alpha(0.7)

    # Style whiskers, caps, and medians
    for whisker in bp['whiskers']:
        whisker.set(color='#555555', linewidth=1.5, linestyle='--')
    for cap in bp['caps']:
        cap.set(color='#555555', linewidth=2)
    for median in bp['medians']:
        median.set(color='#D32F2F', linewidth=2)
    for flier in bp['fliers']:
        flier.set(marker='o', markerfacecolor='#FF5722', markersize=5, alpha=0.5)

    ax.set_ylabel(param_label, fontsize=11)
    ax.set_xlabel("Wafer", fontsize=11)
    ax.set_title(f"Boxplot Comparison: {param_label}", fontsize=14, fontweight="bold")

    # Rotate labels if many wafers
    if len(labels) > 5:
        plt.xticks(rotation=45, ha='right')

    ax.grid(True, alpha=0.3, axis='y')
    ax.set_axisbelow(True)

    # Add statistics annotation if checkbox is enabled
    if multi_wafer_boxplot_stats_var.get():
        for i, data in enumerate(boxplot_data):
            mean_val = np.mean(data)
            median_val = np.median(data)
            # Position text slightly to the right of each boxplot
            x_pos = i + 1 + 0.35
            # Add mean and median as small text next to each box
            ax.annotate(f'Mean: {mean_val:.3g}\nMedian: {median_val:.3g}',
                        xy=(x_pos, mean_val), fontsize=7, color='#333333',
                        ha='left', va='center',
                        bbox=dict(boxstyle='round,pad=0.3', facecolor='#FFFFCC',
                                  edgecolor='#CCCCCC', alpha=0.8))
        # Extend x-axis limits to make room for annotations
        ax.set_xlim(0.4, len(boxplot_data) + 1.2)

    # Destroy old canvas
    if multi_wafer_boxplot_canvas:
        multi_wafer_boxplot_canvas.get_tk_widget().destroy()

    # Create new canvas
    multi_wafer_boxplot_canvas = FigureCanvasTkAgg(fig, master=multi_wafer_boxplot_frame)
    multi_wafer_boxplot_canvas.get_tk_widget().pack(fill=tk.BOTH, expand=True)
    multi_wafer_boxplot_canvas.draw()

    print(f"Boxplot updated with {len(boxplot_data)} wafers")


def update_multi_wafer_distribution():
    """Update the distribution (CDF) display for selected wafers"""
    global multi_wafer_distribution_canvas, multi_wafer_stdf_data, multi_wafer_wafer_ids

    if not multi_wafer_stdf_data:
        print("No wafermap data loaded for Distribution")
        return

    # Get selected wafer indices
    selected_indices = get_selected_wafer_indices()

    if not selected_indices:
        print("No wafers selected for distribution")
        if multi_wafer_distribution_canvas:
            multi_wafer_distribution_canvas.get_tk_widget().destroy()
            multi_wafer_distribution_canvas = None
        return

    selected = multi_wafer_param_combobox.get()

    if not selected:
        print("No parameter selected for distribution")
        return

    # Parse parameter selection
    if selected.startswith("BIN"):
        param_column = "bin"
        param_label = "Bin"
    else:
        test_key = selected.split(":")[0].strip()
        if test_key.startswith("test_"):
            param_column = int(test_key.replace("test_", ""))
        else:
            param_column = int(test_key)
        param_label = selected

    # Get selected data and IDs
    selected_data = [multi_wafer_stdf_data[i] for i in selected_indices]
    selected_ids = [multi_wafer_wafer_ids[i] for i in selected_indices]

    # Create figure with two subplots: PDF (histogram) and CDF
    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(14, 6), constrained_layout=True)

    # Color palette
    colors = plt.cm.tab10(np.linspace(0, 1, len(selected_data)))

    # Plot for each selected wafer
    all_values = []
    for idx, (df, wafer_id, color) in enumerate(zip(selected_data, selected_ids, colors)):
        if param_column not in df.columns:
            continue

        values = df[param_column].dropna().values
        if len(values) == 0:
            continue

        all_values.extend(values)

        # Truncate long wafer IDs
        short_id = str(wafer_id)[:20] + "..." if len(str(wafer_id)) > 20 else str(wafer_id)

        # PDF (Histogram)
        ax1.hist(values, bins=50, alpha=0.5, label=short_id, color=color, density=True, edgecolor='white', linewidth=0.5)

        # CDF
        sorted_vals = np.sort(values)
        cdf = np.arange(1, len(sorted_vals) + 1) / len(sorted_vals)
        ax2.plot(sorted_vals, cdf, label=short_id, color=color, linewidth=2)

    if not all_values:
        print("No valid data for distribution")
        return

    # Style PDF plot
    ax1.set_xlabel(param_label, fontsize=11)
    ax1.set_ylabel("Density", fontsize=11)
    ax1.set_title(f"Probability Density Function (PDF)", fontsize=12, fontweight="bold")
    ax1.legend(loc='upper right', fontsize=8)
    ax1.grid(True, alpha=0.3)
    ax1.set_axisbelow(True)

    # Style CDF plot
    ax2.set_xlabel(param_label, fontsize=11)
    ax2.set_ylabel("Cumulative Probability", fontsize=11)
    ax2.set_title(f"Cumulative Distribution Function (CDF)", fontsize=12, fontweight="bold")
    ax2.legend(loc='lower right', fontsize=8)
    ax2.grid(True, alpha=0.3)
    ax2.set_axisbelow(True)
    ax2.set_ylim(0, 1)

    # Add horizontal lines at key percentiles on CDF
    for perc, label in [(0.25, '25%'), (0.5, '50%'), (0.75, '75%')]:
        ax2.axhline(y=perc, color='#888888', linestyle='--', linewidth=0.8, alpha=0.5)
        ax2.text(ax2.get_xlim()[1], perc, f' {label}', va='center', fontsize=8, color='#666666')

    # Main title
    fig.suptitle(f"Distribution Analysis: {param_label}", fontsize=14, fontweight="bold")

    # Destroy old canvas
    if multi_wafer_distribution_canvas:
        multi_wafer_distribution_canvas.get_tk_widget().destroy()

    # Create new canvas
    multi_wafer_distribution_canvas = FigureCanvasTkAgg(fig, master=multi_wafer_distribution_frame)
    multi_wafer_distribution_canvas.get_tk_widget().pack(fill=tk.BOTH, expand=True)
    multi_wafer_distribution_canvas.draw()

    print(f"Distribution plot updated with {len(selected_data)} wafers")


# ============================================================================
# Tab 4: Diffmap - Difference between reference and comparison wafermaps
# ============================================================================

# Storage for diffmap data
diffmap_reference_data = None
diffmap_reference_id = None
diffmap_compare_data = None
diffmap_compare_id = None
diffmap_result_data = None
diffmap_test_params = {}
diffmap_canvas = None
diffmap_wafer_config = None  # Wafer configuration (notch orientation) for diffmap

# Control frame for diffmap
diffmap_control_frame = tk.Frame(tab_diffmap)
diffmap_control_frame.pack(side=tk.TOP, fill=tk.X, padx=10, pady=10)

# Reference file section
ref_frame = tk.LabelFrame(diffmap_control_frame, text="Reference File", font=("Helvetica", 10, "bold"))
ref_frame.pack(side=tk.LEFT, fill=tk.X, padx=5, pady=5)

diffmap_ref_btn = tk.Button(
    ref_frame,
    text="Select Reference STDF",
    command=lambda: select_diffmap_reference(),
    font=("Helvetica", 10),
    bg="#4CAF50",
    fg="white"
)
diffmap_ref_btn.pack(side=tk.LEFT, padx=5, pady=5)

diffmap_ref_csv_btn = tk.Button(
    ref_frame,
    text="Select Reference CSV",
    command=lambda: select_diffmap_reference_csv(),
    font=("Helvetica", 10),
    bg="#FF9800",
    fg="white"
)
diffmap_ref_csv_btn.pack(side=tk.LEFT, padx=5, pady=5)

diffmap_ref_label = tk.Label(
    ref_frame,
    text="No file selected",
    font=("Helvetica", 9),
    fg="gray",
    width=30,
    anchor="w"
)
diffmap_ref_label.pack(side=tk.LEFT, padx=5, pady=5)

# Comparison file section
comp_frame = tk.LabelFrame(diffmap_control_frame, text="Comparison File", font=("Helvetica", 10, "bold"))
comp_frame.pack(side=tk.LEFT, fill=tk.X, padx=5, pady=5)

diffmap_comp_btn = tk.Button(
    comp_frame,
    text="Select Comparison STDF",
    command=lambda: select_diffmap_comparison(),
    font=("Helvetica", 10),
    bg="#2196F3",
    fg="white"
)
diffmap_comp_btn.pack(side=tk.LEFT, padx=5, pady=5)

diffmap_comp_csv_btn = tk.Button(
    comp_frame,
    text="Select Comparison CSV",
    command=lambda: select_diffmap_comparison_csv(),
    font=("Helvetica", 10),
    bg="#E91E63",
    fg="white"
)
diffmap_comp_csv_btn.pack(side=tk.LEFT, padx=5, pady=5)

diffmap_comp_label = tk.Label(
    comp_frame,
    text="No file selected",
    font=("Helvetica", 9),
    fg="gray",
    width=30,
    anchor="w"
)
diffmap_comp_label.pack(side=tk.LEFT, padx=5, pady=5)

# Calculate button
diffmap_calc_btn = tk.Button(
    diffmap_control_frame,
    text="Calculate Diff",
    command=lambda: calculate_diffmap(),
    font=("Helvetica", 11, "bold"),
    bg="#FF9800",
    fg="white",
    state=tk.DISABLED
)
diffmap_calc_btn.pack(side=tk.LEFT, padx=10, pady=5)

# Status frame for offset and notch orientation warnings - NEW ROW
diffmap_status_row = tk.Frame(tab_diffmap)
diffmap_status_row.pack(side=tk.TOP, fill=tk.X, padx=10, pady=2)

diffmap_status_frame = tk.LabelFrame(diffmap_status_row, text="Data Alignment", font=("Helvetica", 8))
diffmap_status_frame.pack(side=tk.LEFT, padx=5, pady=2)

# X/Y Offset status label
diffmap_offset_label = tk.Label(
    diffmap_status_frame,
    text=" X/Y: -- ",
    font=("Consolas", 9),
    fg="gray",
    padx=8,
    pady=2
)
diffmap_offset_label.pack(side=tk.LEFT, padx=3, pady=2)

# Notch orientation status label
diffmap_notch_label = tk.Label(
    diffmap_status_frame,
    text=" Notch: -- ",
    font=("Consolas", 9),
    fg="gray",
    padx=8,
    pady=2
)
diffmap_notch_label.pack(side=tk.LEFT, padx=3, pady=2)

# Overlay Coordinates Button
diffmap_overlay_btn = tk.Button(
    diffmap_status_row,
    text="Overlay Coordinates",
    command=lambda: overlay_coordinates(),
    font=("Helvetica", 9),
    bg="#9C27B0",
    fg="white",
    state=tk.DISABLED
)
diffmap_overlay_btn.pack(side=tk.LEFT, padx=10, pady=2)

# Info label for overlay status
diffmap_overlay_info = tk.Label(
    diffmap_status_row,
    text="",
    font=("Helvetica", 8),
    fg="gray"
)
diffmap_overlay_info.pack(side=tk.LEFT, padx=5, pady=2)

# Store detected values globally for overlay function
diffmap_detected_x_offset = 0
diffmap_detected_y_offset = 0
diffmap_detected_ref_notch = "Unknown"
diffmap_detected_comp_notch = "Unknown"

# Store original comparison data for reset
diffmap_compare_data_original = None
diffmap_aligned = False

def detect_wafer_offset_and_notch():
    """Detect X/Y offset and notch orientation between reference and comparison data"""
    global diffmap_reference_data, diffmap_compare_data
    global diffmap_detected_x_offset, diffmap_detected_y_offset
    global diffmap_detected_ref_notch, diffmap_detected_comp_notch
    global diffmap_aligned

    if diffmap_reference_data is None or diffmap_compare_data is None:
        diffmap_offset_label.config(text=" X/Y: -- ", fg="gray", bg="#f0f0f0")
        diffmap_notch_label.config(text=" Notch: -- ", fg="gray", bg="#f0f0f0")
        diffmap_overlay_btn.config(state=tk.DISABLED)
        diffmap_overlay_info.config(text="")
        return

    ref_df = diffmap_reference_data
    comp_df = diffmap_compare_data

    # Calculate X/Y ranges for both datasets
    ref_x_min, ref_x_max = ref_df['x'].min(), ref_df['x'].max()
    ref_y_min, ref_y_max = ref_df['y'].min(), ref_df['y'].max()
    comp_x_min, comp_x_max = comp_df['x'].min(), comp_df['x'].max()
    comp_y_min, comp_y_max = comp_df['y'].min(), comp_df['y'].max()

    # Calculate offset
    x_offset = int(comp_x_min - ref_x_min)
    y_offset = int(comp_y_min - ref_y_min)

    diffmap_detected_x_offset = x_offset
    diffmap_detected_y_offset = y_offset

    # Determine offset status
    if x_offset == 0 and y_offset == 0:
        diffmap_offset_label.config(
            text=" X/Y: OK ",
            fg="white",
            bg="#4CAF50"
        )
    else:
        diffmap_offset_label.config(
            text=f" X/Y: ΔX={x_offset:+d} ΔY={y_offset:+d} ",
            fg="white",
            bg="#f44336"
        )

    # Detect notch orientation based on wafer shape
    ref_notch = detect_notch_orientation(ref_df)
    comp_notch = detect_notch_orientation(comp_df)

    diffmap_detected_ref_notch = ref_notch
    diffmap_detected_comp_notch = comp_notch

    if ref_notch == comp_notch:
        diffmap_notch_label.config(
            text=f" Notch: ✓ {ref_notch} ",
            fg="white",
            bg="#4CAF50"
        )
    else:
        diffmap_notch_label.config(
            text=f" Notch: {ref_notch}↔{comp_notch} ",
            fg="white",
            bg="#f44336"
        )

    # Enable overlay button if there's a mismatch
    if x_offset != 0 or y_offset != 0 or ref_notch != comp_notch:
        diffmap_overlay_btn.config(state=tk.NORMAL)
        diffmap_overlay_info.config(text="Click to align comparison data to reference", fg="#9C27B0")
    else:
        diffmap_overlay_btn.config(state=tk.DISABLED)
        diffmap_overlay_info.config(text="Data already aligned", fg="green")

def detect_notch_orientation(df):
    """Detect notch orientation based on die distribution pattern"""
    if df is None or len(df) == 0:
        return "?"

    x_min, x_max = df['x'].min(), df['x'].max()
    y_min, y_max = df['y'].min(), df['y'].max()

    # Define edge regions (outer 15% of each side)
    x_range = x_max - x_min
    y_range = y_max - y_min

    if x_range == 0 or y_range == 0:
        return "?"

    edge_threshold_x = x_range * 0.15
    edge_threshold_y = y_range * 0.15

    # Count dies at each edge
    top_edge = len(df[df['y'] >= (y_max - edge_threshold_y)])
    bottom_edge = len(df[df['y'] <= (y_min + edge_threshold_y)])
    left_edge = len(df[df['x'] <= (x_min + edge_threshold_x)])
    right_edge = len(df[df['x'] >= (x_max - edge_threshold_x)])

    edges = {
        '12h': top_edge,
        '6h': bottom_edge,
        '9h': left_edge,
        '3h': right_edge
    }

    # The notch is typically where there are fewer dies
    min_edge = min(edges, key=edges.get)
    return min_edge

def overlay_coordinates():
    """Adjust comparison data coordinates to match reference data"""
    global diffmap_compare_data, diffmap_detected_x_offset, diffmap_detected_y_offset
    global diffmap_detected_ref_notch, diffmap_detected_comp_notch
    global diffmap_compare_data_original, diffmap_aligned

    if diffmap_compare_data is None or diffmap_reference_data is None:
        print("No data loaded - run Calculate Diffmap first")
        return

    # Check if already aligned
    if diffmap_aligned:
        print("Data already aligned. Load new files to realign.")
        return

    print(f"\n{'='*60}")
    print("Overlaying Coordinates")
    print(f"{'='*60}")

    # Save original data before modifying
    if diffmap_compare_data_original is None:
        diffmap_compare_data_original = diffmap_compare_data.copy()

    # Apply X/Y offset correction
    if diffmap_detected_x_offset != 0 or diffmap_detected_y_offset != 0:
        print(f"Applying offset correction: dX={-diffmap_detected_x_offset}, dY={-diffmap_detected_y_offset}")
        diffmap_compare_data['x'] = diffmap_compare_data['x'] - diffmap_detected_x_offset
        diffmap_compare_data['y'] = diffmap_compare_data['y'] - diffmap_detected_y_offset

    # Apply rotation if notch orientations differ
    if diffmap_detected_ref_notch != diffmap_detected_comp_notch:
        print(f"Applying notch rotation: {diffmap_detected_comp_notch} -> {diffmap_detected_ref_notch}")
        diffmap_compare_data = rotate_wafer_data(
            diffmap_compare_data,
            diffmap_detected_comp_notch,
            diffmap_detected_ref_notch
        )

    # Mark as aligned
    diffmap_aligned = True

    # Re-detect to update status
    detect_wafer_offset_and_notch()

    # Update comparison label
    current_text = diffmap_comp_label.cget("text")
    if "[aligned]" not in current_text:
        diffmap_comp_label.config(text=current_text + " [aligned]", fg="purple")
    diffmap_overlay_info.config(text="OK Coordinates aligned!", fg="green")

    # Disable overlay button after alignment
    diffmap_overlay_btn.config(state=tk.DISABLED)

    print("Overlay complete!")

    # Recalculate diffmap with aligned coordinates (without reloading files)
    recalculate_diffmap_only()


def recalculate_diffmap_only():
    """Recalculate difference map without reloading files - OPTIMIZED"""
    global diffmap_result_data, diffmap_reference_data, diffmap_compare_data

    if diffmap_reference_data is None or diffmap_compare_data is None:
        print("No data to calculate")
        return

    print(f"\n{'='*60}")
    print("Recalculating Difference Map (Reference - Comparison)")
    print(f"{'='*60}")

    calc_start = time.time()

    # Merge dataframes on x, y coordinates
    ref_df = diffmap_reference_data
    comp_df = diffmap_compare_data

    # Create a result dataframe with matching coordinates
    merged = ref_df.merge(comp_df, on=['x', 'y'], suffixes=('_ref', '_comp'), how='inner')

    if len(merged) == 0:
        diffmap_info_label.config(text="No matching die coordinates found!", fg="red")
        print("ERROR: No matching die coordinates between reference and comparison files")
        return

    print(f"Matched {len(merged)} dies between reference and comparison")

    # OPTIMIZED: Build all columns at once to avoid DataFrame fragmentation
    result_data = {
        'x': merged['x'].values,
        'y': merged['y'].values
    }

    # Calculate difference for bin
    if 'bin_ref' in merged.columns and 'bin_comp' in merged.columns:
        result_data['bin'] = merged['bin_ref'].values - merged['bin_comp'].values

    # Calculate difference for all numeric parameters
    diff_count = 0
    for col in ref_df.columns:
        if col in ['x', 'y', 'bin']:
            continue

        ref_col = f"{col}_ref"
        comp_col = f"{col}_comp"

        if ref_col in merged.columns and comp_col in merged.columns:
            try:
                result_data[col] = merged[ref_col].values - merged[comp_col].values
                diff_count += 1
            except Exception as e:
                pass  # Skip silently for speed

    # Create DataFrame all at once (much faster, no fragmentation)
    diffmap_result_data = pd.DataFrame(result_data)

    calc_time = time.time() - calc_start
    print(f"Calculated differences for {diff_count} parameters in {calc_time:.2f}s")
    print(f"Result dataframe: {len(diffmap_result_data)} dies")

    # Update info label
    diffmap_info_label.config(
        text=f"Diff: {len(diffmap_result_data)} dies, {diff_count} params",
        fg="green"
    )

    # Update display
    update_diffmap_display()

def rotate_wafer_data(df, from_notch, to_notch):
    """Rotate wafer data to align notch orientations"""
    # Define rotation angles (clockwise, in 90-degree steps)
    notch_positions = {'12h': 0, '3h': 90, '6h': 180, '9h': 270}

    if from_notch not in notch_positions or to_notch not in notch_positions:
        return df

    from_angle = notch_positions[from_notch]
    to_angle = notch_positions[to_notch]
    rotation = (to_angle - from_angle) % 360

    if rotation == 0:
        return df

    # Get center of the wafer
    x_center = (df['x'].min() + df['x'].max()) / 2
    y_center = (df['y'].min() + df['y'].max()) / 2

    # Create a copy
    rotated_df = df.copy()

    # Translate to origin, rotate, translate back
    x_rel = df['x'] - x_center
    y_rel = df['y'] - y_center

    if rotation == 90:
        # 90° clockwise: (x, y) -> (y, -x)
        rotated_df['x'] = (y_rel + x_center).round().astype(int)
        rotated_df['y'] = (-x_rel + y_center).round().astype(int)
    elif rotation == 180:
        # 180°: (x, y) -> (-x, -y)
        rotated_df['x'] = (-x_rel + x_center).round().astype(int)
        rotated_df['y'] = (-y_rel + y_center).round().astype(int)
    elif rotation == 270:
        # 270° clockwise (90° counter-clockwise): (x, y) -> (-y, x)
        rotated_df['x'] = (-y_rel + x_center).round().astype(int)
        rotated_df['y'] = (x_rel + y_center).round().astype(int)

    print(f"  Rotated {rotation}° clockwise")
    return rotated_df

# Second row of controls
diffmap_control_frame2 = tk.Frame(tab_diffmap)
diffmap_control_frame2.pack(side=tk.TOP, fill=tk.X, padx=10, pady=5)

# View type selection (Diffmap vs Correlation)
diffmap_view_label = tk.Label(
    diffmap_control_frame2,
    text="View:",
    font=("Helvetica", 10)
)
diffmap_view_label.pack(side=tk.LEFT, padx=5)

diffmap_view_var = tk.StringVar(value="Diffmap")
diffmap_view_combobox = ttk.Combobox(
    diffmap_control_frame2,
    textvariable=diffmap_view_var,
    values=["Diffmap", "Correlation Plot"],
    state="readonly",
    width=15,
    font=("Helvetica", 10)
)
diffmap_view_combobox.pack(side=tk.LEFT, padx=5)
diffmap_view_combobox.bind("<<ComboboxSelected>>", lambda e: update_diffmap_display())

# Parameter selection with custom colored dropdown
diffmap_param_label = tk.Label(
    diffmap_control_frame2,
    text="Parameter:",
    font=("Helvetica", 10)
)
diffmap_param_label.pack(side=tk.LEFT, padx=5)

# Custom dropdown frame
diffmap_param_frame = tk.Frame(diffmap_control_frame2)
diffmap_param_frame.pack(side=tk.LEFT, padx=5)

diffmap_param_var = tk.StringVar(value="")
diffmap_param_entry = tk.Entry(
    diffmap_param_frame,
    textvariable=diffmap_param_var,
    state="readonly",
    width=42,
    font=("Helvetica", 10),
    readonlybackground="white",
    cursor="hand2"
)
diffmap_param_entry.pack(side=tk.LEFT)

diffmap_dropdown_btn = tk.Button(
    diffmap_param_frame,
    text="\u25BC",
    font=("Helvetica", 8),
    width=2,
    relief=tk.RAISED
)
diffmap_dropdown_btn.pack(side=tk.LEFT)

# Store options with status
diffmap_param_items = []  # List of (text, is_missing)
diffmap_popup = None

def show_param_dropdown(event=None):
    """Show the custom parameter dropdown"""
    global diffmap_popup

    # Close if already open
    if diffmap_popup is not None:
        try:
            diffmap_popup.destroy()
        except:
            pass
        diffmap_popup = None
        return

    if not diffmap_param_items:
        print("No parameters available yet - load files first")
        return

    # Create popup
    diffmap_popup = tk.Toplevel(main_win)
    diffmap_popup.wm_overrideredirect(True)
    diffmap_popup.attributes("-topmost", True)

    # Position below entry
    x = diffmap_param_entry.winfo_rootx()
    y = diffmap_param_entry.winfo_rooty() + diffmap_param_entry.winfo_height()

    # Frame with border
    frame = tk.Frame(diffmap_popup, bd=2, relief=tk.RAISED, bg="white")
    frame.pack(fill=tk.BOTH, expand=True)

    # Scrollbar
    scrollbar = tk.Scrollbar(frame)
    scrollbar.pack(side=tk.RIGHT, fill=tk.Y)

    # Listbox
    listbox = tk.Listbox(
        frame,
        width=48,
        height=min(12, len(diffmap_param_items)),
        font=("Helvetica", 10),
        yscrollcommand=scrollbar.set,
        selectmode=tk.SINGLE,
        activestyle='dotbox',
        exportselection=False,
        bg="white",
        selectbackground="#0078D7",
        selectforeground="white"
    )
    listbox.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)
    scrollbar.config(command=listbox.yview)

    # Add items with colors
    for i, (text, is_missing) in enumerate(diffmap_param_items):
        listbox.insert(tk.END, text)
        if is_missing:
            listbox.itemconfig(i, bg="#FFCCCC", fg="#CC0000", selectbackground="#FF6666", selectforeground="white")

    # Find and select current item
    current_val = diffmap_param_var.get()
    for i, (text, _) in enumerate(diffmap_param_items):
        if text == current_val:
            listbox.selection_set(i)
            listbox.see(i)
            break

    def do_select():
        global diffmap_popup
        sel = listbox.curselection()
        if sel:
            idx = sel[0]
            text, is_missing = diffmap_param_items[idx]
            diffmap_param_var.set(text)
            if is_missing:
                diffmap_param_entry.config(readonlybackground="#FFCCCC")
            else:
                diffmap_param_entry.config(readonlybackground="white")
        if diffmap_popup:
            try:
                diffmap_popup.destroy()
            except:
                pass
            diffmap_popup = None
        # Trigger heatmap update
        main_win.after(10, lambda: update_diffmap_display())

    def on_click(event):
        # Get clicked item
        idx = listbox.nearest(event.y)
        if idx >= 0:
            listbox.selection_clear(0, tk.END)
            listbox.selection_set(idx)
            listbox.activate(idx)
            main_win.after(50, do_select)

    listbox.bind("<Button-1>", on_click)
    listbox.bind("<Return>", lambda e: do_select())
    listbox.bind("<Escape>", lambda e: close_param_popup())

    diffmap_popup.geometry(f"+{x}+{y}")
    listbox.focus_set()

    # Close when clicking outside
    def check_focus():
        global diffmap_popup
        if diffmap_popup:
            try:
                if not diffmap_popup.focus_get():
                    close_param_popup()
            except:
                pass

    main_win.after(500, check_focus)

def close_param_popup():
    global diffmap_popup
    if diffmap_popup:
        try:
            diffmap_popup.destroy()
        except:
            pass
        diffmap_popup = None

diffmap_param_entry.bind("<Button-1>", show_param_dropdown)
diffmap_dropdown_btn.config(command=show_param_dropdown)

# Compatibility wrapper for existing code
class DiffmapParamComboboxWrapper:
    def get(self):
        return diffmap_param_var.get()
    def current(self, idx):
        if 0 <= idx < len(diffmap_param_items):
            text, is_missing = diffmap_param_items[idx]
            diffmap_param_var.set(text)
            if is_missing:
                diffmap_param_entry.config(readonlybackground="#FFCCCC")
            else:
                diffmap_param_entry.config(readonlybackground="white")
    def __setitem__(self, key, value):
        pass  # Ignore ["values"] = ... calls

diffmap_param_combobox = DiffmapParamComboboxWrapper()

# Refresh button
diffmap_refresh_btn = tk.Button(
    diffmap_control_frame2,
    text="Refresh",
    command=lambda: update_diffmap_display(),
    font=("Helvetica", 10)
)
diffmap_refresh_btn.pack(side=tk.LEFT, padx=5)

# Show grid checkbox
diffmap_grid_var = tk.BooleanVar(value=False)
diffmap_grid_checkbox = tk.Checkbutton(
    diffmap_control_frame2,
    text="Show Grid",
    variable=diffmap_grid_var,
    command=lambda: update_diffmap_display(),
    font=("Helvetica", 10)
)
diffmap_grid_checkbox.pack(side=tk.LEFT, padx=5)

# Colormap selection for diff
diffmap_cmap_label = tk.Label(
    diffmap_control_frame2,
    text="Colormap:",
    font=("Helvetica", 10)
)
diffmap_cmap_label.pack(side=tk.LEFT, padx=5)

diffmap_cmap_var = tk.StringVar(value="RdBu_r")
diffmap_cmap_combobox = ttk.Combobox(
    diffmap_control_frame2,
    textvariable=diffmap_cmap_var,
    values=["RdBu_r", "coolwarm", "seismic", "bwr", "PiYG", "PRGn", "viridis", "plasma"],
    state="readonly",
    width=10,
    font=("Helvetica", 9)
)
diffmap_cmap_combobox.pack(side=tk.LEFT, padx=5)
diffmap_cmap_combobox.bind("<<ComboboxSelected>>", lambda e: update_diffmap_display())

# Symmetric colorscale checkbox
diffmap_symmetric_var = tk.BooleanVar(value=True)
diffmap_symmetric_checkbox = tk.Checkbutton(
    diffmap_control_frame2,
    text="Symmetric Scale",
    variable=diffmap_symmetric_var,
    command=lambda: update_diffmap_display(),
    font=("Helvetica", 10)
)
diffmap_symmetric_checkbox.pack(side=tk.LEFT, padx=5)

# Info label
diffmap_info_label = tk.Label(
    diffmap_control_frame2,
    text="",
    font=("Helvetica", 9),
    fg="gray"
)
diffmap_info_label.pack(side=tk.LEFT, padx=10)

# Main display area with stats panel (same as Wafermap)
diffmap_main_container = tk.Frame(tab_diffmap)
diffmap_main_container.pack(side=tk.BOTTOM, fill=tk.BOTH, expand=True)

# Left panel for statistics
diffmap_stats_panel = tk.Frame(diffmap_main_container, width=320, bg="#f0f0f0")
diffmap_stats_panel.pack(side=tk.LEFT, fill=tk.Y, padx=5, pady=5)
diffmap_stats_panel.pack_propagate(False)

# Stats panel title
diffmap_stats_title = tk.Label(
    diffmap_stats_panel,
    text="Diff Statistics",
    font=("Helvetica", 12, "bold"),
    bg="#f0f0f0"
)
diffmap_stats_title.pack(side=tk.TOP, pady=5)

# Frame for boxplot
diffmap_boxplot_frame = tk.Frame(diffmap_stats_panel, bg="#f0f0f0")
diffmap_boxplot_frame.pack(side=tk.TOP, fill=tk.BOTH, expand=True, padx=5, pady=5)

diffmap_boxplot_label = tk.Label(
    diffmap_boxplot_frame,
    text="Difference Boxplot",
    font=("Helvetica", 10, "bold"),
    bg="#f0f0f0"
)
diffmap_boxplot_label.pack(side=tk.TOP)

# Frame for histogram/distribution
diffmap_hist_frame = tk.Frame(diffmap_stats_panel, bg="#f0f0f0")
diffmap_hist_frame.pack(side=tk.TOP, fill=tk.BOTH, expand=True, padx=5, pady=5)

diffmap_hist_label = tk.Label(
    diffmap_hist_frame,
    text="Difference Distribution",
    font=("Helvetica", 10, "bold"),
    bg="#f0f0f0"
)
diffmap_hist_label.pack(side=tk.TOP)

# Right panel for diffmap display
diffmap_display_frame = tk.Frame(diffmap_main_container)
diffmap_display_frame.pack(side=tk.RIGHT, fill=tk.BOTH, expand=True)

# Canvas references for stats
diffmap_boxplot_canvas = None
diffmap_hist_canvas = None

# Store file paths for deferred loading
diffmap_reference_path = None
diffmap_compare_path = None


def select_diffmap_reference():
    """Select reference STDF file for diffmap (does not load yet)"""
    global diffmap_reference_path

    stdf_path = filedialog.askopenfilename(
        title="Select Reference STDF File",
        filetypes=[("STDF files", "*.stdf *.std"), ("All files", "*.*")],
    )

    if not stdf_path:
        return

    diffmap_reference_path = stdf_path
    filename = os.path.basename(stdf_path)
    short_name = filename[:25] + "..." if len(filename) > 25 else filename
    diffmap_ref_label.config(text=f"{short_name}", fg="darkgreen")

    # Enable calculate button if both files selected
    check_diffmap_ready()
    print(f"Reference selected: {filename}")


def select_diffmap_comparison():
    """Select comparison STDF file for diffmap (does not load yet)"""
    global diffmap_compare_path, diffmap_compare_data_original, diffmap_aligned

    stdf_path = filedialog.askopenfilename(
        title="Select Comparison STDF File",
        filetypes=[("STDF files", "*.stdf *.std"), ("All files", "*.*")],
    )

    if not stdf_path:
        return

    diffmap_compare_path = stdf_path
    # Reset alignment state when new file is selected
    diffmap_compare_data_original = None
    diffmap_aligned = False

    filename = os.path.basename(stdf_path)
    short_name = filename[:25] + "..." if len(filename) > 25 else filename
    diffmap_comp_label.config(text=f"{short_name}", fg="darkblue")

    # Enable calculate button if both files selected
    check_diffmap_ready()
    print(f"Comparison selected: {filename}")


def select_diffmap_reference_csv():
    """Select reference CSV file for diffmap (does not load yet)"""
    global diffmap_reference_path

    csv_path = filedialog.askopenfilename(
        title="Select Reference CSV File",
        filetypes=[("CSV files", "*.csv"), ("All files", "*.*")],
    )

    if not csv_path:
        return

    diffmap_reference_path = csv_path
    filename = os.path.basename(csv_path)
    short_name = filename[:25] + "..." if len(filename) > 25 else filename
    diffmap_ref_label.config(text=f"{short_name}", fg="darkgreen")

    # Enable calculate button if both files selected
    check_diffmap_ready()
    print(f"Reference CSV selected: {filename}")


def select_diffmap_comparison_csv():
    """Select comparison CSV file for diffmap (does not load yet)"""
    global diffmap_compare_path, diffmap_compare_data_original, diffmap_aligned

    csv_path = filedialog.askopenfilename(
        title="Select Comparison CSV File",
        filetypes=[("CSV files", "*.csv"), ("All files", "*.*")],
    )

    if not csv_path:
        return

    diffmap_compare_path = csv_path
    # Reset alignment state when new file is selected
    diffmap_compare_data_original = None
    diffmap_aligned = False

    filename = os.path.basename(csv_path)
    short_name = filename[:25] + "..." if len(filename) > 25 else filename
    diffmap_comp_label.config(text=f"{short_name}", fg="darkblue")

    # Enable calculate button if both files selected
    check_diffmap_ready()
    print(f"Comparison CSV selected: {filename}")


def load_diffmap_reference():
    """Load reference file (STDF or CSV) - called during calculate"""
    global diffmap_reference_data, diffmap_reference_id, diffmap_test_params, diffmap_reference_path
    global diffmap_wafer_config

    if not diffmap_reference_path:
        print("ERROR: No reference path set")
        return False

    filename = os.path.basename(diffmap_reference_path)
    print(f"Loading Reference: {filename}")

    try:
        if diffmap_reference_path.lower().endswith(('.stdf', '.std')):
            df, wafer_id, test_params, grouped_params, test_limits, wafer_cfg = read_wafermap_from_stdf(diffmap_reference_path)
            if df is None or df.empty:
                print(f"ERROR: STDF returned empty dataframe")
                diffmap_ref_label.config(text="Error loading file", fg="red")
                return False
            diffmap_reference_data = df
            diffmap_reference_id = wafer_id if wafer_id else filename
            diffmap_test_params.update(test_params)
            diffmap_wafer_config = wafer_cfg  # Store wafer configuration for notch display
        else:
            # CSV loading
            if not load_csv_as_reference():
                print("ERROR: CSV loading failed")
                return False

        short_name = filename[:25] + "..." if len(filename) > 25 else filename
        diffmap_ref_label.config(text=f"{short_name} ({len(diffmap_reference_data)} dies)", fg="green")
        print(f"Reference loaded: {len(diffmap_reference_data)} dies")
        return True
    except Exception as e:
        print(f"ERROR loading reference: {e}")
        import traceback
        traceback.print_exc()
        return False


def load_csv_as_reference():
    """Helper to load CSV file as reference"""
    global diffmap_reference_data, diffmap_reference_id, diffmap_test_params, diffmap_reference_path
    global diffmap_wafer_config

    try:
        df = pd.read_csv(diffmap_reference_path)
        filename = os.path.basename(diffmap_reference_path)

        # Find coordinate columns
        x_col_candidates = ['x', 'X', 'x_coord', 'X_COORD', 'DIE_X', 'die_x', 'col', 'COL']
        y_col_candidates = ['y', 'Y', 'y_coord', 'Y_COORD', 'DIE_Y', 'die_y', 'row', 'ROW']

        x_col = y_col = None
        for c in x_col_candidates:
            if c in df.columns:
                x_col = c
                break
        for c in y_col_candidates:
            if c in df.columns:
                y_col = c
                break

        if x_col is None or y_col is None:
            numeric_cols = df.select_dtypes(include=[np.number]).columns.tolist()
            if len(numeric_cols) >= 2:
                x_col, y_col = numeric_cols[0], numeric_cols[1]
            else:
                return False

        df = df.rename(columns={x_col: 'x', y_col: 'y'})

        # Look for bin column
        for c in ['bin', 'BIN', 'Bin', 'HARD_BIN', 'SOFT_BIN']:
            if c in df.columns:
                df = df.rename(columns={c: 'bin'})
                break
        if 'bin' not in df.columns:
            df['bin'] = 1

        # Build test parameters
        test_params = {}
        numeric_columns = [c for c in df.select_dtypes(include=[np.number]).columns if c not in ['x', 'y', 'bin']]
        for idx, col in enumerate(numeric_columns):
            test_num = idx + 1
            test_params[f"test_{test_num}"] = col
            df = df.rename(columns={col: test_num})

        # Detect notch orientation from CSV
        notch_orientation = _detect_notch_from_csv(df, diffmap_reference_path)

        # Update diffmap wafer config
        diffmap_wafer_config = {
            'notch_orientation': notch_orientation,
            'wafer_size': None,
            'die_width': None,
            'die_height': None,
            'pos_x': None,
            'pos_y': None
        }
        print(f"Reference CSV notch orientation: {notch_orientation}")

        diffmap_reference_data = df
        diffmap_reference_id = filename.replace('.csv', '').replace('.CSV', '')
        diffmap_test_params.update(test_params)
        return True
    except Exception as e:
        print(f"Error loading CSV: {e}")
        return False


def _detect_notch_from_csv(df, csv_path):
    """Helper function to detect notch orientation from CSV data"""
    notch_orientation = None

    # Check for notch/flat column in CSV
    notch_col_candidates = ['notch', 'Notch', 'NOTCH', 'flat', 'Flat', 'FLAT',
                            'WF_FLAT', 'wf_flat', 'orientation', 'Orientation',
                            'ORIENTATION', 'wafer_flat', 'WAFER_FLAT']

    for candidate in notch_col_candidates:
        if candidate in df.columns:
            notch_values = df[candidate].dropna()
            if len(notch_values) > 0:
                notch_val = str(notch_values.iloc[0]).strip().upper()
                notch_mapping = {
                    'U': 'U', 'UP': 'U', '0': 'U', 'TOP': 'U', 'NORTH': 'U',
                    'D': 'D', 'DOWN': 'D', '180': 'D', 'BOTTOM': 'D', 'SOUTH': 'D',
                    'L': 'L', 'LEFT': 'L', '270': 'L', 'WEST': 'L',
                    'R': 'R', 'RIGHT': 'R', '90': 'R', 'EAST': 'R'
                }
                notch_orientation = notch_mapping.get(notch_val, notch_val[0] if notch_val else None)
                print(f"Detected notch from CSV column '{candidate}': {notch_orientation}")
            break

    # If not found in columns, check in header comments
    if notch_orientation is None:
        try:
            with open(csv_path, 'r') as f:
                for i, line in enumerate(f):
                    if i > 20:
                        break
                    line_upper = line.upper()
                    if 'NOTCH' in line_upper or 'FLAT' in line_upper:
                        for orient in ['UP', 'DOWN', 'LEFT', 'RIGHT', 'NORTH', 'SOUTH', 'EAST', 'WEST']:
                            if orient in line_upper:
                                notch_mapping = {'UP': 'U', 'DOWN': 'D', 'LEFT': 'L', 'RIGHT': 'R',
                                                'NORTH': 'U', 'SOUTH': 'D', 'WEST': 'L', 'EAST': 'R'}
                                notch_orientation = notch_mapping.get(orient, 'D')
                                print(f"Detected notch from CSV header: {notch_orientation}")
                                break
                        if notch_orientation:
                            break
        except:
            pass

    # Default to 'D' (down) if not found
    if notch_orientation is None:
        notch_orientation = 'D'
        print(f"No notch found in CSV, using default: {notch_orientation}")

    return notch_orientation


def load_diffmap_comparison():
    """Load comparison file (STDF or CSV) - called during calculate"""
    global diffmap_compare_data, diffmap_compare_id, diffmap_test_params, diffmap_compare_path
    global diffmap_compare_data_original, diffmap_aligned

    if not diffmap_compare_path:
        print("ERROR: No comparison path set")
        return False

    filename = os.path.basename(diffmap_compare_path)
    print(f"Loading Comparison: {filename}")

    try:
        if diffmap_compare_path.lower().endswith(('.stdf', '.std')):
            df, wafer_id, test_params, grouped_params, test_limits, wafer_cfg = read_wafermap_from_stdf(diffmap_compare_path)
            if df is None or df.empty:
                print(f"ERROR: STDF returned empty dataframe")
                diffmap_comp_label.config(text="Error loading file", fg="red")
                return False
            diffmap_compare_data = df
            diffmap_compare_id = wafer_id if wafer_id else filename
            diffmap_test_params.update(test_params)
        else:
            # CSV loading
            if not load_csv_as_comparison():
                print("ERROR: CSV loading failed")
                return False

        # Reset alignment state
        diffmap_compare_data_original = None
        diffmap_aligned = False

        short_name = filename[:25] + "..." if len(filename) > 25 else filename
        diffmap_comp_label.config(text=f"{short_name} ({len(diffmap_compare_data)} dies)", fg="blue")
        print(f"Comparison loaded: {len(diffmap_compare_data)} dies")
        return True
    except Exception as e:
        print(f"ERROR loading comparison: {e}")
        import traceback
        traceback.print_exc()
        return False


def load_csv_as_comparison():
    """Helper to load CSV file as comparison"""
    global diffmap_compare_data, diffmap_compare_id, diffmap_test_params, diffmap_compare_path

    try:
        df = pd.read_csv(diffmap_compare_path)
        filename = os.path.basename(diffmap_compare_path)

        # Find coordinate columns
        x_col_candidates = ['x', 'X', 'x_coord', 'X_COORD', 'DIE_X', 'die_x', 'col', 'COL']
        y_col_candidates = ['y', 'Y', 'y_coord', 'Y_COORD', 'DIE_Y', 'die_y', 'row', 'ROW']

        x_col = y_col = None
        for c in x_col_candidates:
            if c in df.columns:
                x_col = c
                break
        for c in y_col_candidates:
            if c in df.columns:
                y_col = c
                break

        if x_col is None or y_col is None:
            numeric_cols = df.select_dtypes(include=[np.number]).columns.tolist()
            if len(numeric_cols) >= 2:
                x_col, y_col = numeric_cols[0], numeric_cols[1]
            else:
                return False

        df = df.rename(columns={x_col: 'x', y_col: 'y'})

        # Look for bin column
        for c in ['bin', 'BIN', 'Bin', 'HARD_BIN', 'SOFT_BIN']:
            if c in df.columns:
                df = df.rename(columns={c: 'bin'})
                break
        if 'bin' not in df.columns:
            df['bin'] = 1

        # Build test parameters
        test_params = {}
        numeric_columns = [c for c in df.select_dtypes(include=[np.number]).columns if c not in ['x', 'y', 'bin']]
        for idx, col in enumerate(numeric_columns):
            test_num = idx + 1
            test_params[f"test_{test_num}"] = col
            df = df.rename(columns={col: test_num})

        diffmap_compare_data = df
        diffmap_compare_id = filename.replace('.csv', '').replace('.CSV', '')
        diffmap_test_params.update(test_params)
        return True
    except Exception as e:
        print(f"Error loading CSV: {e}")
        return False



def update_diffmap_params():
    """Update parameter dropdown for diffmap - highlights missing params in red"""
    global diffmap_test_params, diffmap_reference_data, diffmap_compare_data, diffmap_param_items
    global diffmap_result_data

    diffmap_param_items = []  # Clear list

    # Add BIN option - check if it has valid data
    bin_has_data = False
    if diffmap_result_data is not None and 'bin' in diffmap_result_data.columns:
        bin_has_data = diffmap_result_data['bin'].notna().sum() > 0
    diffmap_param_items.append(("BIN (Bin Number)", not bin_has_data))

    for test_key, test_name in sorted(diffmap_test_params.items()):
        # Get the column key used in diffmap_result_data
        if test_key.startswith("test_"):
            param_num = int(test_key.replace("test_", ""))
        else:
            try:
                param_num = int(test_key)
            except:
                param_num = test_key

        # Check if parameter has valid (non-NaN) data in the result
        has_valid_data = False
        if diffmap_result_data is not None and param_num in diffmap_result_data.columns:
            has_valid_data = diffmap_result_data[param_num].notna().sum() > 0

        display_text = f"{test_key}: {test_name}"
        is_missing = not has_valid_data

        diffmap_param_items.append((display_text, is_missing))

    # Set first item if nothing selected
    if diffmap_param_items and not diffmap_param_var.get():
        diffmap_param_combobox.current(0)


def check_diffmap_ready():
    """Check if both files are selected and enable calculate button"""
    if diffmap_reference_path is not None and diffmap_compare_path is not None:
        diffmap_calc_btn.config(state=tk.NORMAL)
    else:
        diffmap_calc_btn.config(state=tk.DISABLED)
        # Reset status labels
        diffmap_offset_label.config(text=" X/Y: -- ", fg="gray", bg="#f0f0f0")
        diffmap_notch_label.config(text=" Notch: -- ", fg="gray", bg="#f0f0f0")


def calculate_diffmap():
    """Load files and calculate difference map: Reference - Comparison - OPTIMIZED"""
    global diffmap_result_data, diffmap_reference_data, diffmap_compare_data
    global diffmap_test_params

    if diffmap_reference_path is None or diffmap_compare_path is None:
        print("Both reference and comparison files must be selected")
        return

    # Clear old test params
    diffmap_test_params = {}

    print(f"\n{'='*60}")
    print("Loading files...")
    print(f"{'='*60}")

    # Disable button during loading
    diffmap_calc_btn.config(state=tk.DISABLED, text="Loading...")
    diffmap_info_label.config(text="Loading reference file...", fg="blue")
    main_win.update_idletasks()

    start_time = time.time()

    # Load reference file
    ref_start = time.time()
    if not load_diffmap_reference():
        print("Failed to load reference file")
        diffmap_calc_btn.config(state=tk.NORMAL, text="Calculate Diffmap")
        diffmap_info_label.config(text="Failed to load reference file", fg="red")
        return
    ref_time = time.time() - ref_start
    print(f"Reference loaded in {ref_time:.1f}s")

    diffmap_info_label.config(text="Loading comparison file...", fg="blue")
    main_win.update_idletasks()

    # Load comparison file
    comp_start = time.time()
    if not load_diffmap_comparison():
        print("Failed to load comparison file")
        diffmap_calc_btn.config(state=tk.NORMAL, text="Calculate Diffmap")
        diffmap_info_label.config(text="Failed to load comparison file", fg="red")
        return
    comp_time = time.time() - comp_start
    print(f"Comparison loaded in {comp_time:.1f}s")

    load_time = time.time() - start_time
    print(f"Total loading completed in {load_time:.1f}s")

    # Re-enable button
    diffmap_calc_btn.config(state=tk.NORMAL, text="Calculate Diffmap")

    # Update parameter combobox after loading
    update_diffmap_params()

    # Check for offset and notch alignment
    detect_wafer_offset_and_notch()

    print(f"\n{'='*60}")
    print("Calculating Difference Map (Reference - Comparison)")
    print(f"{'='*60}")

    calc_start = time.time()
    diffmap_info_label.config(text="Calculating differences...", fg="blue")
    main_win.update_idletasks()

    # Merge dataframes on x, y coordinates (no copy needed)
    ref_df = diffmap_reference_data
    comp_df = diffmap_compare_data

    # Create a result dataframe with matching coordinates
    merged = ref_df.merge(comp_df, on=['x', 'y'], suffixes=('_ref', '_comp'), how='inner')

    if len(merged) == 0:
        diffmap_info_label.config(text="No matching die coordinates found!", fg="red")
        print("ERROR: No matching die coordinates between reference and comparison files")
        return

    print(f"Matched {len(merged)} dies between reference and comparison")

    # OPTIMIZED: Build all columns at once to avoid DataFrame fragmentation
    result_data = {
        'x': merged['x'].values,
        'y': merged['y'].values
    }

    # Calculate difference for bin
    if 'bin_ref' in merged.columns and 'bin_comp' in merged.columns:
        result_data['bin'] = merged['bin_ref'].values - merged['bin_comp'].values

    # Calculate difference for all numeric parameters
    diff_count = 0
    for col in ref_df.columns:
        if col in ['x', 'y', 'bin']:
            continue

        ref_col = f"{col}_ref"
        comp_col = f"{col}_comp"

        if ref_col in merged.columns and comp_col in merged.columns:
            try:
                result_data[col] = merged[ref_col].values - merged[comp_col].values
                diff_count += 1
            except:
                pass  # Skip silently for speed

    # Create DataFrame all at once (much faster, no fragmentation warnings)
    diffmap_result_data = pd.DataFrame(result_data)

    calc_time = time.time() - calc_start
    total_time = time.time() - start_time

    # Display load time and calculation info in the UI
    diffmap_info_label.config(
        text=f"✓ {len(diffmap_result_data)} dies, {diff_count} params | Load: {load_time:.1f}s | Total: {total_time:.1f}s",
        fg="green"
    )

    print(f"Difference map calculated: {len(diffmap_result_data)} dies, {diff_count} parameters")
    print(f"Calculation time: {calc_time:.2f}s")
    print(f"Total time: {total_time:.2f}s")
    print(f"{'='*60}\n")

    # Update display
    update_diffmap_display()


def update_diffmap_display():
    """Update the diffmap display based on selected view type"""
    global diffmap_canvas, diffmap_result_data, diffmap_boxplot_canvas, diffmap_hist_canvas

    view_type = diffmap_view_var.get()

    if view_type == "Correlation Plot":
        update_correlation_plot_display()
    else:
        update_diffmap_heatmap_display()


def update_correlation_plot_display():
    """Display correlation plot between reference and comparison data"""
    global diffmap_canvas, diffmap_reference_data, diffmap_compare_data
    global diffmap_boxplot_canvas, diffmap_hist_canvas

    if diffmap_reference_data is None or diffmap_compare_data is None:
        print("Both reference and comparison data needed for correlation plot")
        return

    selected = diffmap_param_combobox.get()

    if not selected:
        print("No parameter selected")
        return

    # Parse parameter selection
    if selected.startswith("BIN"):
        param_column = "bin"
        param_label = "Bin"
    else:
        test_key = selected.split(":")[0].strip()
        if test_key.startswith("test_"):
            param_column = int(test_key.replace("test_", ""))
        else:
            param_column = int(test_key)
        param_label = selected.split(":")[-1].strip() if ":" in selected else selected

    # Check if parameter exists in both datasets
    if param_column not in diffmap_reference_data.columns:
        print(f"Parameter {param_column} not found in reference data")
        return
    if param_column not in diffmap_compare_data.columns:
        print(f"Parameter {param_column} not found in comparison data")
        return

    # Merge on x, y coordinates
    merged = pd.merge(
        diffmap_reference_data[['x', 'y', param_column]],
        diffmap_compare_data[['x', 'y', param_column]],
        on=['x', 'y'],
        suffixes=('_ref', '_comp')
    )

    ref_col = f"{param_column}_ref"
    comp_col = f"{param_column}_comp"

    # Remove NaN values
    merged = merged.dropna(subset=[ref_col, comp_col])

    if len(merged) == 0:
        print("No matching data points for correlation")
        return

    ref_values = merged[ref_col].values
    comp_values = merged[comp_col].values

    # Calculate correlation
    correlation = np.corrcoef(ref_values, comp_values)[0, 1]

    # Create figure
    fig, ax = plt.subplots(figsize=(10, 9))

    # Scatter plot
    scatter = ax.scatter(
        ref_values, comp_values,
        c='#3498DB', alpha=0.6, s=30, edgecolors='white', linewidth=0.5
    )

    # Add unity line (y=x)
    min_val = min(ref_values.min(), comp_values.min())
    max_val = max(ref_values.max(), comp_values.max())
    margin = (max_val - min_val) * 0.05
    ax.plot([min_val - margin, max_val + margin], [min_val - margin, max_val + margin],
            'r--', linewidth=2, label='Unity (y=x)', alpha=0.8)

    # Linear regression fit using numpy (no scipy needed)
    # Calculate slope and intercept using least squares
    n = len(ref_values)
    mean_x = np.mean(ref_values)
    mean_y = np.mean(comp_values)

    # Calculate slope and intercept
    numerator = np.sum((ref_values - mean_x) * (comp_values - mean_y))
    denominator = np.sum((ref_values - mean_x) ** 2)
    slope = numerator / denominator if denominator != 0 else 0
    intercept = mean_y - slope * mean_x

    # Calculate R-squared
    y_pred = slope * ref_values + intercept
    ss_res = np.sum((comp_values - y_pred) ** 2)
    ss_tot = np.sum((comp_values - mean_y) ** 2)
    r_squared = 1 - (ss_res / ss_tot) if ss_tot != 0 else 0

    # Standard error of slope
    if n > 2 and denominator != 0:
        mse = ss_res / (n - 2)
        std_err = np.sqrt(mse / denominator)
    else:
        std_err = 0

    fit_line = slope * np.array([min_val, max_val]) + intercept
    ax.plot([min_val, max_val], fit_line, 'g-', linewidth=2,
            label=f'Fit: y={slope:.4f}x+{intercept:.4f}', alpha=0.8)

    # Title and labels - smaller font to fit window
    # Truncate title to fit within plot width
    param_short = param_label[:40] + "..." if len(param_label) > 40 else param_label
    ref_short = diffmap_reference_id[:20] + "..." if len(diffmap_reference_id) > 20 else diffmap_reference_id
    comp_short = diffmap_compare_id[:20] + "..." if len(diffmap_compare_id) > 20 else diffmap_compare_id
    title = f"Correlation: {param_short}\nRef: {ref_short} vs Comp: {comp_short}"
    ax.set_title(title, fontsize=6, fontweight="bold")
    ax.set_xlabel(f"Reference", fontsize=15)
    ax.set_ylabel(f"Comparison", fontsize=15)

    # Set equal aspect and limits
    ax.set_xlim(min_val - margin, max_val + margin)
    ax.set_ylim(min_val - margin, max_val + margin)
    ax.set_aspect('equal')

    # Add grid
    ax.grid(True, alpha=0.3, linestyle='-')

    # Statistics text box
    stats_text = (
        f"Correlation Statistics:\n"
        f"R² = {r_squared:.4f}\n"
        f"R = {correlation:.4f}\n"
        f"Slope = {slope:.4f}\n"
        f"Intercept = {intercept:.4f}\n"
        f"Std Error = {std_err:.4f}\n"
        f"N = {len(merged)}"
    )

    ax.text(
        0.02, 0.98, stats_text,
        transform=ax.transAxes,
        fontsize=9,
        verticalalignment='top',
        bbox=dict(boxstyle='round', facecolor='wheat', alpha=0.9)
    )

    ax.legend(loc='lower right', fontsize=9)
    fig.tight_layout()

    # Destroy old canvas
    if diffmap_canvas:
        diffmap_canvas.get_tk_widget().destroy()

    # Create new canvas
    diffmap_canvas = FigureCanvasTkAgg(fig, master=diffmap_display_frame)
    diffmap_canvas.get_tk_widget().pack(fill=tk.BOTH, expand=True)
    diffmap_canvas.draw()

    # Update statistics plots for correlation
    update_correlation_stats(ref_values, comp_values, param_label)

    print(f"Correlation plot updated for parameter: {param_label}, R²={r_squared:.4f}")


def update_correlation_stats(ref_values, comp_values, param_label):
    """Update the statistics plots for correlation view"""
    global diffmap_boxplot_canvas, diffmap_hist_canvas

    # Clear existing
    for widget in diffmap_boxplot_frame.winfo_children():
        if widget != diffmap_boxplot_label:
            widget.destroy()

    for widget in diffmap_hist_frame.winfo_children():
        if widget != diffmap_hist_label:
            widget.destroy()

    # Create comparison boxplot
    fig_box, ax_box = plt.subplots(figsize=(2.8, 2.5))
    fig_box.patch.set_facecolor('white')

    bp = ax_box.boxplot(
        [ref_values, comp_values],
        tick_labels=["Ref", "Comp"],
        vert=True,
        patch_artist=True,
        showmeans=True,
        widths=0.6,
        meanprops=dict(marker="D", markerfacecolor="#E74C3C", markeredgecolor="white", markersize=5),
        medianprops=dict(color="#2C3E50", linewidth=2),
        whiskerprops=dict(color="#2C3E50", linewidth=1.5),
        capprops=dict(color="#2C3E50", linewidth=1.5),
        flierprops=dict(marker='o', markerfacecolor='#95A5A6', markersize=3, alpha=0.6),
        boxprops=dict(linewidth=1.5)
    )

    colors = ['#3498DB', '#E74C3C']
    for patch, color in zip(bp['boxes'], colors):
        patch.set_facecolor(color)
        patch.set_alpha(0.75)

    ax_box.set_title("Ref vs Comp", fontsize=8, fontweight="bold", color='#2C3E50')
    ax_box.tick_params(axis='both', which='major', labelsize=6)
    ax_box.set_ylabel("Value", fontsize=7)
    ax_box.grid(True, alpha=0.4)

    fig_box.tight_layout()

    diffmap_boxplot_canvas = FigureCanvasTkAgg(fig_box, master=diffmap_boxplot_frame)
    diffmap_boxplot_canvas.get_tk_widget().pack(fill=tk.BOTH, expand=True)
    diffmap_boxplot_canvas.draw()

    # Create residual histogram
    residuals = comp_values - ref_values
    fig_hist, ax_hist = plt.subplots(figsize=(2.8, 2.5))
    fig_hist.patch.set_facecolor('white')

    n, bins, patches = ax_hist.hist(residuals, bins=30, density=True, alpha=0.7, edgecolor='white')

    for patch, left_edge in zip(patches, bins[:-1]):
        if left_edge < 0:
            patch.set_facecolor('#3498DB')
        else:
            patch.set_facecolor('#E74C3C')

    ax_hist.axvline(x=0, color='#2C3E50', linewidth=1.5, linestyle='-')
    ax_hist.axvline(x=np.mean(residuals), color='#E74C3C', linewidth=1.5, linestyle='--')

    ax_hist.set_title("Residuals", fontsize=8, fontweight="bold", color='#2C3E50')
    ax_hist.set_xlabel("Comp - Ref", fontsize=7)
    ax_hist.set_ylabel("Density", fontsize=7)
    ax_hist.tick_params(axis='both', which='major', labelsize=6)
    ax_hist.grid(True, alpha=0.4)

    fig_hist.tight_layout()

    diffmap_hist_canvas = FigureCanvasTkAgg(fig_hist, master=diffmap_hist_frame)
    diffmap_hist_canvas.get_tk_widget().pack(fill=tk.BOTH, expand=True)
    diffmap_hist_canvas.draw()


# Store current diffmap die popup window reference
_diffmap_die_popup = None


def show_diffmap_die_popup(x_coord, y_coord, selected_param_column):
    """Show a popup window with all measurement parameters for the clicked die in diffmap
    Design consistent with show_die_data_popup_all_wafers from Multiwafer tab"""
    global _diffmap_die_popup, diffmap_reference_data, diffmap_compare_data, diffmap_result_data
    global diffmap_test_params

    # Close previous popup if exists
    if _diffmap_die_popup is not None:
        try:
            _diffmap_die_popup.destroy()
        except:
            pass
        _diffmap_die_popup = None

    # Get data for this die from both reference and compare files
    ref_die = None
    comp_die = None
    diff_die = None

    if diffmap_reference_data is not None:
        ref_data = diffmap_reference_data[(diffmap_reference_data['x'] == x_coord) &
                                           (diffmap_reference_data['y'] == y_coord)]
        if len(ref_data) > 0:
            ref_die = ref_data.iloc[0]

    if diffmap_compare_data is not None:
        comp_data = diffmap_compare_data[(diffmap_compare_data['x'] == x_coord) &
                                          (diffmap_compare_data['y'] == y_coord)]
        if len(comp_data) > 0:
            comp_die = comp_data.iloc[0]

    if diffmap_result_data is not None:
        diff_data = diffmap_result_data[(diffmap_result_data['x'] == x_coord) &
                                         (diffmap_result_data['y'] == y_coord)]
        if len(diff_data) > 0:
            diff_die = diff_data.iloc[0]

    if ref_die is None and comp_die is None:
        print(f"No data found at position X={x_coord}, Y={y_coord}")
        return

    # Count available data sources
    data_count = sum([1 for d in [ref_die, comp_die] if d is not None])

    # Create popup window (75% of original size like Multiwafer: 450x375)
    popup = tk.Toplevel(main_win)
    popup.title(f"Diffmap Die Data - Position: X={x_coord}, Y={y_coord}")
    popup.geometry("450x375")
    popup.transient(main_win)

    # Store reference to current popup
    _diffmap_die_popup = popup

    def on_popup_close():
        global _diffmap_die_popup
        _diffmap_die_popup = None
        popup.destroy()

    popup.protocol("WM_DELETE_WINDOW", on_popup_close)

    # Center the popup
    popup.update_idletasks()
    x = main_win.winfo_x() + (main_win.winfo_width() - 450) // 2
    y = main_win.winfo_y() + (main_win.winfo_height() - 375) // 2
    popup.geometry(f"+{x}+{y}")

    # Header frame (same blue color as Multiwafer)
    header_frame = tk.Frame(popup, bg="#2196F3", pady=5)
    header_frame.pack(fill=tk.X)

    tk.Label(
        header_frame,
        text=f"Position: X = {x_coord}, Y = {y_coord}",
        font=("Helvetica", 10, "bold"),
        bg="#2196F3",
        fg="white"
    ).pack()

    tk.Label(
        header_frame,
        text="Reference | Compare | Difference",
        font=("Helvetica", 9),
        bg="#2196F3",
        fg="white"
    ).pack()

    # Legend frame
    legend_frame = tk.Frame(popup, pady=3)
    legend_frame.pack(fill=tk.X)
    tk.Label(legend_frame, text="■ Selected Parameter", fg="red", font=("Helvetica", 8, "bold")).pack(side=tk.LEFT, padx=5)

    # Create Treeview with scrollbars
    tree_frame = tk.Frame(popup)
    tree_frame.pack(fill=tk.BOTH, expand=True, padx=5, pady=3)

    tree_scroll_y = ttk.Scrollbar(tree_frame, orient="vertical")
    tree_scroll_x = ttk.Scrollbar(tree_frame, orient="horizontal")

    columns = ["parameter", "reference", "compare", "diff"]

    tree = ttk.Treeview(
        tree_frame,
        columns=columns,
        show="headings",
        yscrollcommand=tree_scroll_y.set,
        xscrollcommand=tree_scroll_x.set
    )

    tree.heading("parameter", text="Parameter")
    tree.heading("reference", text="Reference")
    tree.heading("compare", text="Compare")
    tree.heading("diff", text="Diff (C-R)")
    tree.column("parameter", width=120, anchor="w")
    tree.column("reference", width=90, anchor="center")
    tree.column("compare", width=90, anchor="center")
    tree.column("diff", width=90, anchor="center")

    tree_scroll_y.config(command=tree.yview)
    tree_scroll_x.config(command=tree.xview)

    tree_scroll_y.pack(side=tk.RIGHT, fill=tk.Y)
    tree_scroll_x.pack(side=tk.BOTTOM, fill=tk.X)
    tree.pack(fill=tk.BOTH, expand=True)

    # Configure tags for highlighting (same as Multiwafer)
    tree.tag_configure("selected", background="#FFCDD2", foreground="red")
    tree.tag_configure("normal", background="white")
    tree.tag_configure("coordinate", background="#E3F2FD")

    # Get all unique columns
    all_columns = set()
    if ref_die is not None:
        all_columns.update(ref_die.index)
    if comp_die is not None:
        all_columns.update(comp_die.index)

    # Remove x, y from display
    all_columns = [c for c in all_columns if c not in ['x', 'y']]

    # Sort columns: bin first, then numeric test numbers
    def sort_key(col):
        if col == 'bin':
            return (0, 0)
        elif isinstance(col, int):
            return (1, col)
        else:
            return (2, str(col))

    all_columns = sorted(all_columns, key=sort_key)

    # Add data rows
    for col in all_columns:
        # Determine parameter name
        if col == 'bin':
            param_name = "BIN"
        else:
            test_key = f"test_{col}"
            if test_key in diffmap_test_params:
                param_name = f"{col}: {diffmap_test_params[test_key]}"
            else:
                param_name = f"Test {col}"

        # Determine if this is the selected parameter
        is_selected = (col == selected_param_column)

        # Get values
        ref_val = "-"
        comp_val = "-"
        diff_val = "-"

        if ref_die is not None and col in ref_die.index:
            val = ref_die[col]
            if pd.notna(val):
                ref_val = f"{val:.4g}" if isinstance(val, float) else str(val)

        if comp_die is not None and col in comp_die.index:
            val = comp_die[col]
            if pd.notna(val):
                comp_val = f"{val:.4g}" if isinstance(val, float) else str(val)

        if diff_die is not None and col in diff_die.index:
            val = diff_die[col]
            if pd.notna(val):
                diff_val = f"{val:.4g}" if isinstance(val, float) else str(val)

        tag = "selected" if is_selected else "normal"
        tree.insert("", "end", values=[param_name, ref_val, comp_val, diff_val], tags=(tag,))

    # Button frame
    btn_frame = tk.Frame(popup, pady=5)
    btn_frame.pack(fill=tk.X)

    def copy_to_clipboard():
        """Copy die data to clipboard"""
        data_str = f"Position: X={x_coord}, Y={y_coord}\n\n"
        data_str += "Parameter\tReference\tCompare\tDiff\n"

        for item in tree.get_children():
            values = tree.item(item, "values")
            data_str += "\t".join(str(v) for v in values) + "\n"

        popup.clipboard_clear()
        popup.clipboard_append(data_str)
        print("Diffmap die data copied to clipboard")

    tk.Button(
        btn_frame,
        text="Copy",
        command=copy_to_clipboard,
        font=("Helvetica", 8),
        bg="#4CAF50",
        fg="white"
    ).pack(side=tk.LEFT, padx=5)

    tk.Button(
        btn_frame,
        text="Close",
        command=on_popup_close,
        font=("Helvetica", 8)
    ).pack(side=tk.RIGHT, padx=5)


def update_diffmap_heatmap_display():
    """Update the diffmap heatmap display"""
    global diffmap_canvas, diffmap_result_data, diffmap_boxplot_canvas, diffmap_hist_canvas

    try:
        if diffmap_result_data is None or diffmap_result_data.empty:
            return

        selected = diffmap_param_combobox.get()

        if not selected:
            return

        # Parse parameter selection
        if selected.startswith("BIN"):
            param_column = "bin"
            param_label = "Bin Difference"
        else:
            test_key = selected.split(":")[0].strip()
            if test_key.startswith("test_"):
                param_column = int(test_key.replace("test_", ""))
            else:
                try:
                    param_column = int(test_key)
                except ValueError:
                    return
            param_label = selected.split(":")[-1].strip() if ":" in selected else selected
            param_label = f"Diff: {param_label}"

        if param_column not in diffmap_result_data.columns:
            return

        # Get plot data
        mask = diffmap_result_data[param_column].notna()
        plot_data = diffmap_result_data[mask]

        if len(plot_data) == 0:
            # Show message in display area instead of leaving old heatmap
            if diffmap_canvas is not None:
                try:
                    diffmap_canvas.get_tk_widget().destroy()
                    plt.close(diffmap_canvas.figure)
                except:
                    pass
                diffmap_canvas = None
            plt.close('all')

            # Create a simple figure with "No Data" message
            fig, ax = plt.subplots(figsize=(10, 9))
            ax.text(0.5, 0.5, f"No valid data for:\n{selected}",
                   ha='center', va='center', fontsize=14, color='red',
                   transform=ax.transAxes)
            ax.set_xlim(0, 1)
            ax.set_ylim(0, 1)
            ax.axis('off')

            diffmap_canvas = FigureCanvasTkAgg(fig, master=diffmap_display_frame)
            diffmap_canvas.get_tk_widget().pack(fill=tk.BOTH, expand=True)
            diffmap_canvas.draw()
            return

        # Close old matplotlib figures first
        if diffmap_canvas is not None:
            try:
                old_widget = diffmap_canvas.get_tk_widget()
                old_fig = diffmap_canvas.figure
                old_widget.destroy()
                plt.close(old_fig)
            except:
                pass
            diffmap_canvas = None

        # Close any orphaned figures
        plt.close('all')

        # Compute grid
        grid, x_min, y_min, x_max, y_max, grid_width, grid_height = _compute_grid_fast(
            plot_data, param_column
        )

        # Create NEW figure
        fig, ax = plt.subplots(figsize=(10, 9))

        # Get colormap
        cmap = diffmap_cmap_var.get()

        # Calculate color limits
        data_min = float(np.nanmin(grid))
        data_max = float(np.nanmax(grid))

        if diffmap_symmetric_var.get():
            abs_max = max(abs(data_min), abs(data_max))
            vmin, vmax = -abs_max, abs_max
        else:
            vmin, vmax = data_min, data_max

        # Prevent identical vmin/vmax
        if vmin == vmax:
            vmin = vmin - 0.1
            vmax = vmax + 0.1

        im = ax.imshow(
            grid,
            cmap=cmap,
            aspect="equal",
            interpolation="nearest",
            origin="upper",
            vmin=vmin,
            vmax=vmax
        )

        # Invert Y-axis so Y increases downward (0 at top)
        ax.invert_yaxis()

        # Title
        param_short = param_label[:40] + "..." if len(param_label) > 40 else param_label
        ref_short = diffmap_reference_id[:20] + "..." if len(diffmap_reference_id) > 20 else diffmap_reference_id
        comp_short = diffmap_compare_id[:20] + "..." if len(diffmap_compare_id) > 20 else diffmap_compare_id
        title = f"Diff: {param_short}\nRef: {ref_short} vs Comp: {comp_short}"
        ax.set_title(title, fontsize=8, fontweight="bold")
        ax.set_xlabel("X Coordinate", fontsize=10)
        ax.set_ylabel("Y Coordinate", fontsize=10)

        # Set correct axis tick labels to show real coordinates
        num_x_ticks = min(10, grid_width)
        num_y_ticks = min(10, grid_height)
        x_tick_positions = np.linspace(0, grid_width - 1, num_x_ticks)
        y_tick_positions = np.linspace(0, grid_height - 1, num_y_ticks)
        ax.set_xticks(x_tick_positions)
        ax.set_yticks(y_tick_positions)
        ax.set_xticklabels([f"{int(x_min + pos)}" for pos in x_tick_positions])
        ax.set_yticklabels([f"{int(y_min + pos)}" for pos in y_tick_positions])

        # Show grid if enabled
        if diffmap_grid_var.get():
            ax.set_xticks(np.arange(-0.5, grid_width, 1), minor=True)
            ax.set_yticks(np.arange(-0.5, grid_height, 1), minor=True)
            ax.grid(which="minor", color="black", linewidth=0.3)
            ax.tick_params(which="minor", size=0)

        # Draw notch marker if orientation is available
        if diffmap_wafer_config and diffmap_wafer_config.get('notch_orientation'):
            draw_notch_marker(ax, grid_width, grid_height, diffmap_wafer_config['notch_orientation'])

        # Colorbar
        cbar = fig.colorbar(im, ax=ax, fraction=0.024, pad=0.01)
        cbar.ax.tick_params(labelsize=9)
        cbar.set_label(param_label, fontsize=10)

        # Statistics text box
        diff_values = plot_data[param_column].values
        stats_text = (
            f"Diff Statistics:\n"
            f"Min: {np.nanmin(diff_values):.3f}\n"
            f"Max: {np.nanmax(diff_values):.3f}\n"
            f"Mean: {np.nanmean(diff_values):.3f}\n"
            f"Std: {np.nanstd(diff_values):.3f}\n"
            f"Dies: {len(plot_data)}"
        )

        ax.text(
            0.02, 0.98, stats_text,
            transform=ax.transAxes,
            fontsize=9,
            verticalalignment='top',
            bbox=dict(boxstyle='round', facecolor='wheat', alpha=0.8)
        )

        fig.tight_layout()

        # Store grid info for click handler
        _diffmap_click_info = {
            'x_min': x_min,
            'y_min': y_min,
            'param_column': param_column
        }

        def on_diffmap_die_click(event):
            """Handle click on diffmap to show die data"""
            if event.inaxes is None:
                return
            if event.xdata is None or event.ydata is None:
                return

            # Convert click position to actual coordinates
            actual_x = int(round(event.xdata + _diffmap_click_info['x_min']))
            actual_y = int(round(event.ydata + _diffmap_click_info['y_min']))

            # Show popup with diffmap data
            show_diffmap_die_popup(actual_x, actual_y, _diffmap_click_info['param_column'])

        # Create new canvas
        diffmap_canvas = FigureCanvasTkAgg(fig, master=diffmap_display_frame)
        diffmap_canvas.get_tk_widget().pack(fill=tk.BOTH, expand=True)
        diffmap_canvas.mpl_connect('button_press_event', on_diffmap_die_click)
        diffmap_canvas.draw()

        # Update statistics plots
        update_diffmap_stats(diff_values, param_label)

    except Exception as e:
        import traceback
        traceback.print_exc()


def update_diffmap_stats(diff_values, param_label):
    """Update the statistics plots for diffmap"""
    global diffmap_boxplot_canvas, diffmap_hist_canvas

    # Clear existing
    for widget in diffmap_boxplot_frame.winfo_children():
        if widget != diffmap_boxplot_label:
            widget.destroy()

    for widget in diffmap_hist_frame.winfo_children():
        if widget != diffmap_hist_label:
            widget.destroy()

    # Create boxplot - professional styling
    fig_box, ax_box = plt.subplots(figsize=(2.8, 2.5))
    fig_box.patch.set_facecolor('white')

    bp = ax_box.boxplot(
        [diff_values],
        tick_labels=["Diff"],
        vert=True,
        patch_artist=True,
        showmeans=True,
        widths=0.6,
        meanprops=dict(marker="D", markerfacecolor="#E74C3C", markeredgecolor="white", markersize=5, markeredgewidth=1),
        medianprops=dict(color="#2C3E50", linewidth=2),
        whiskerprops=dict(color="#2C3E50", linewidth=1.5, linestyle='-'),
        capprops=dict(color="#2C3E50", linewidth=1.5),
        flierprops=dict(marker='o', markerfacecolor='#95A5A6', markeredgecolor='#7F8C8D', markersize=3, alpha=0.6),
        boxprops=dict(linewidth=1.5)
    )

    # Color based on mean (green if near zero, red/blue if positive/negative)
    mean_val = np.mean(diff_values)
    if abs(mean_val) < np.std(diff_values) * 0.1:
        box_color = '#2ECC71'  # Green - near zero
    elif mean_val > 0:
        box_color = '#E74C3C'  # Red - positive diff
    else:
        box_color = '#3498DB'  # Blue - negative diff

    for patch in bp['boxes']:
        patch.set_facecolor(box_color)
        patch.set_alpha(0.75)
        patch.set_edgecolor('#2C3E50')

    # Calculate statistics
    stats_max = np.max(diff_values)
    stats_min = np.min(diff_values)
    stats_q1 = np.percentile(diff_values, 25)
    stats_q3 = np.percentile(diff_values, 75)
    stats_median = np.median(diff_values)

    # Add statistics box in top left corner
    stats_text = (
        f"Max: {stats_max:.3g}\n"
        f"Min: {stats_min:.3g}\n"
        f"Q1: {stats_q1:.3g}\n"
        f"Q3: {stats_q3:.3g}\n"
        f"Mean: {mean_val:.3g}\n"
        f"Median: {stats_median:.3g}"
    )
    ax_box.text(0.02, 0.98, stats_text,
        transform=ax_box.transAxes,
        fontsize=5,
        fontweight='normal',
        fontfamily='monospace',
        verticalalignment='top',
        horizontalalignment='left',
        bbox=dict(boxstyle='round,pad=0.3', facecolor='white', edgecolor='#BDC3C7', alpha=0.9, linewidth=1)
    )

    # Add zero reference line
    ax_box.axhline(y=0, color='#2C3E50', linewidth=1, linestyle='--', alpha=0.5)

    ax_box.set_title("Difference", fontsize=8, fontweight="bold", color='#2C3E50')
    ax_box.tick_params(axis='both', which='major', labelsize=6, colors='#2C3E50')
    ax_box.set_ylabel("Δ Value", fontsize=7, color='#2C3E50')
    ax_box.set_facecolor('#FAFAFA')
    ax_box.grid(True, alpha=0.4, linestyle='-', linewidth=0.5, color='#BDC3C7')
    ax_box.spines['top'].set_visible(False)
    ax_box.spines['right'].set_visible(False)
    ax_box.spines['left'].set_color('#BDC3C7')
    ax_box.spines['bottom'].set_color('#BDC3C7')

    fig_box.tight_layout()

    diffmap_boxplot_canvas = FigureCanvasTkAgg(fig_box, master=diffmap_boxplot_frame)
    diffmap_boxplot_canvas.get_tk_widget().pack(fill=tk.BOTH, expand=True)
    diffmap_boxplot_canvas.draw()

    # Create histogram - professional styling
    fig_hist, ax_hist = plt.subplots(figsize=(2.8, 2.5))
    fig_hist.patch.set_facecolor('white')

    # Histogram with color based on positive/negative - professional colors
    n, bins, patches = ax_hist.hist(diff_values, bins=30, density=True, alpha=0.7, edgecolor='white', linewidth=0.5)

    # Color bars based on value with professional colors
    for patch, left_edge in zip(patches, bins[:-1]):
        if left_edge < 0:
            patch.set_facecolor('#3498DB')  # Professional blue for negative
        else:
            patch.set_facecolor('#E74C3C')  # Professional red for positive

    # Add zero line
    ax_hist.axvline(x=0, color='#2C3E50', linewidth=1.5, linestyle='-', alpha=0.8)

    # Add mean line
    ax_hist.axvline(x=mean_val, color='#E74C3C', linewidth=1.5, linestyle='--', alpha=0.8, label=f'Mean={mean_val:.3g}')

    ax_hist.set_title("Distribution", fontsize=8, fontweight="bold", color='#2C3E50')
    ax_hist.set_xlabel("Δ Value", fontsize=7, color='#2C3E50')
    ax_hist.set_ylabel("Density", fontsize=7, color='#2C3E50')
    ax_hist.tick_params(axis='both', which='major', labelsize=6, colors='#2C3E50')
    ax_hist.set_facecolor('#FAFAFA')
    ax_hist.grid(True, alpha=0.4, linestyle='-', linewidth=0.5, color='#BDC3C7')
    ax_hist.spines['top'].set_visible(False)
    ax_hist.spines['right'].set_visible(False)
    ax_hist.spines['left'].set_color('#BDC3C7')
    ax_hist.spines['bottom'].set_color('#BDC3C7')
    ax_hist.legend(fontsize=5, loc='upper right', framealpha=0.9, edgecolor='#BDC3C7')

    fig_hist.tight_layout()

    diffmap_hist_canvas = FigureCanvasTkAgg(fig_hist, master=diffmap_hist_frame)
    diffmap_hist_canvas.get_tk_widget().pack(fill=tk.BOTH, expand=True)
    diffmap_hist_canvas.draw()


# ============================================================================
# Tab 5: Gage R&R Analysis - Multi-file Comparison
# ============================================================================

# Gage R&R Global Variables
grr_file_data = []  # List of loaded file data (dict with 'data', 'type', 'path', 'wafer_id')
grr_selected_dies = []  # List of selected die coordinates for comparison
grr_num_files_var = None
grr_num_dies_var = None
grr_file_canvases = []  # Canvases for wafermap displays
grr_file_frames = []  # Frames for each file display
grr_viz_param_var = None  # Variable for visualization parameter selection

# Main container for Gage R&R tab
grr_main_frame = tk.Frame(tab_grr, bg='#f5f5f5')
grr_main_frame.pack(fill=tk.BOTH, expand=True)

# Header
grr_header = tk.Label(
    grr_main_frame,
    text="Gage R&R Analysis - Multi-File Comparison",
    font=("Helvetica", 14, "bold"),
    bg='#f5f5f5',
    fg='#2C3E50'
)
grr_header.pack(pady=10)

# Top Control Panel
grr_control_frame = tk.Frame(grr_main_frame, bg='#f5f5f5')
grr_control_frame.pack(fill=tk.X, padx=10, pady=5)

# Left side controls - Row 1
grr_left_controls = tk.Frame(grr_control_frame, bg='#f5f5f5')
grr_left_controls.pack(side=tk.LEFT, padx=10)

# File Type Selection
tk.Label(grr_left_controls, text="File Type:", font=("Helvetica", 10, "bold"), bg='#f5f5f5').grid(row=0, column=0, padx=5, pady=5, sticky='w')
grr_file_type_var = tk.StringVar(value="CSV")
grr_file_type_combo = ttk.Combobox(grr_left_controls, textvariable=grr_file_type_var, values=["STDF", "CSV"], state="readonly", width=8)
grr_file_type_combo.grid(row=0, column=1, padx=5, pady=5)

# Number of files dropdown
tk.Label(grr_left_controls, text="Number of Files:", font=("Helvetica", 10, "bold"), bg='#f5f5f5').grid(row=0, column=2, padx=15, pady=5, sticky='w')
grr_num_files_var = tk.StringVar(value="3")
grr_num_files_combo = ttk.Combobox(grr_left_controls, textvariable=grr_num_files_var, values=["2", "3", "4", "5", "6"], state="readonly", width=5)
grr_num_files_combo.grid(row=0, column=3, padx=5, pady=5)

# Number of dies dropdown
tk.Label(grr_left_controls, text="Number of Dies:", font=("Helvetica", 10, "bold"), bg='#f5f5f5').grid(row=0, column=4, padx=15, pady=5, sticky='w')
grr_num_dies_var = tk.StringVar(value="5")
grr_num_dies_combo = ttk.Combobox(grr_left_controls, textvariable=grr_num_dies_var, values=["3", "5", "10", "15", "20", "25", "30"], state="readonly", width=5)
grr_num_dies_combo.grid(row=0, column=5, padx=5, pady=5)

# Row 2 - Parameter selection and visualization
# Parameter selection button (opens multi-select dialog)
tk.Label(grr_left_controls, text="Parameters:", font=("Helvetica", 10, "bold"), bg='#f5f5f5').grid(row=1, column=0, padx=5, pady=5, sticky='w')
grr_selected_params = []  # List of selected parameters
grr_available_params = ["BIN"]  # Available parameters from loaded files
grr_analysis_results = {}  # Store results for PPT export
grr_param_label = tk.Label(grr_left_controls, text="0 selected", font=("Helvetica", 9), bg='#E8F4F8', fg='#2C3E50', width=15, relief='groove')
grr_param_label.grid(row=1, column=1, padx=5, pady=5)

def open_param_selector():
    """Open parameter selection dialog"""
    global grr_selected_params, grr_available_params

    # Create dialog window
    param_dialog = tk.Toplevel(main_win)
    param_dialog.title("Select Parameters for Gage R&R")
    param_dialog.geometry("500x600")
    param_dialog.transient(main_win)
    param_dialog.grab_set()

    # Instructions
    tk.Label(param_dialog, text="Select parameters for Gage R&R analysis:",
             font=("Helvetica", 10, "bold")).pack(pady=10)
    tk.Label(param_dialog, text="(Hold Ctrl to select multiple, Shift for range)",
             font=("Helvetica", 8), fg='gray').pack()

    # Search/filter frame
    filter_frame = tk.Frame(param_dialog)
    filter_frame.pack(fill=tk.X, padx=10, pady=5)
    tk.Label(filter_frame, text="Filter:", font=("Helvetica", 9)).pack(side=tk.LEFT, padx=5)
    filter_var = tk.StringVar()
    filter_entry = tk.Entry(filter_frame, textvariable=filter_var, width=40)
    filter_entry.pack(side=tk.LEFT, padx=5, fill=tk.X, expand=True)

    # Listbox frame with scrollbars
    list_frame = tk.Frame(param_dialog)
    list_frame.pack(fill=tk.BOTH, expand=True, padx=10, pady=5)

    scrollbar_y = ttk.Scrollbar(list_frame, orient="vertical")
    scrollbar_y.pack(side=tk.RIGHT, fill=tk.Y)

    scrollbar_x = ttk.Scrollbar(list_frame, orient="horizontal")
    scrollbar_x.pack(side=tk.BOTTOM, fill=tk.X)

    param_listbox = tk.Listbox(list_frame, selectmode=tk.EXTENDED,
                                yscrollcommand=scrollbar_y.set,
                                xscrollcommand=scrollbar_x.set,
                                font=("Consolas", 9), height=20)
    param_listbox.pack(fill=tk.BOTH, expand=True)

    scrollbar_y.config(command=param_listbox.yview)
    scrollbar_x.config(command=param_listbox.xview)

    # Populate listbox
    def populate_listbox(filter_text=""):
        param_listbox.delete(0, tk.END)
        filter_lower = filter_text.lower()
        for param in grr_available_params:
            if filter_lower in param.lower():
                param_listbox.insert(tk.END, param)
        # Re-select previously selected items
        for i, item in enumerate(param_listbox.get(0, tk.END)):
            if item in grr_selected_params:
                param_listbox.selection_set(i)

    populate_listbox()

    # Filter entry callback
    def on_filter_change(*args):
        populate_listbox(filter_var.get())
    filter_var.trace('w', on_filter_change)

    # Selection info label
    selection_info = tk.Label(param_dialog, text=f"{len(grr_selected_params)} parameter(s) selected",
                              font=("Helvetica", 9), fg='#3498DB')
    selection_info.pack(pady=5)

    def update_selection_info(event=None):
        count = len(param_listbox.curselection())
        selection_info.config(text=f"{count} parameter(s) selected")

    param_listbox.bind('<<ListboxSelect>>', update_selection_info)

    # Quick select buttons frame
    quick_frame = tk.Frame(param_dialog)
    quick_frame.pack(fill=tk.X, padx=10, pady=5)

    def select_all():
        param_listbox.select_set(0, tk.END)
        update_selection_info()

    def clear_all():
        param_listbox.selection_clear(0, tk.END)
        update_selection_info()

    tk.Button(quick_frame, text="Select All", command=select_all,
              font=("Helvetica", 8), bg='#3498DB', fg='white').pack(side=tk.LEFT, padx=5)
    tk.Button(quick_frame, text="Clear All", command=clear_all,
              font=("Helvetica", 8), bg='#95A5A6', fg='white').pack(side=tk.LEFT, padx=5)

    # OK/Cancel buttons
    button_frame = tk.Frame(param_dialog)
    button_frame.pack(fill=tk.X, padx=10, pady=10)

    def on_ok():
        global grr_selected_params
        grr_selected_params = [param_listbox.get(i) for i in param_listbox.curselection()]
        grr_param_label.config(text=f"{len(grr_selected_params)} selected")
        param_dialog.destroy()

    def on_cancel():
        param_dialog.destroy()

    tk.Button(button_frame, text="OK", command=on_ok,
              font=("Helvetica", 10, "bold"), bg='#27AE60', fg='white', width=10).pack(side=tk.LEFT, padx=10)
    tk.Button(button_frame, text="Cancel", command=on_cancel,
              font=("Helvetica", 10), bg='#E74C3C', fg='white', width=10).pack(side=tk.LEFT, padx=10)

grr_param_btn = tk.Button(grr_left_controls, text="Select...", command=open_param_selector,
                          font=("Helvetica", 9), bg='#3498DB', fg='white')
grr_param_btn.grid(row=1, column=2, padx=5, pady=5)

# Visualization parameter selection (for wafermap display color)
tk.Label(grr_left_controls, text="Viz Param:", font=("Helvetica", 10, "bold"), bg='#f5f5f5').grid(row=1, column=3, padx=15, pady=5, sticky='w')
grr_viz_param_var = tk.StringVar(value="None")
grr_viz_param_combo = ttk.Combobox(grr_left_controls, textvariable=grr_viz_param_var, values=["None"], state="readonly", width=25)
grr_viz_param_combo.grid(row=1, column=4, columnspan=2, padx=5, pady=5, sticky='w')

def on_viz_param_change(event=None):
    """Update wafermap visualization when parameter changes"""
    update_grr_file_displays()

grr_viz_param_combo.bind("<<ComboboxSelected>>", on_viz_param_change)

# Right side controls - buttons
grr_right_controls = tk.Frame(grr_control_frame, bg='#f5f5f5')
grr_right_controls.pack(side=tk.RIGHT, padx=10)

# Status label
grr_status_var = tk.StringVar(value="Ready to load files")
grr_status_label = tk.Label(grr_right_controls, textvariable=grr_status_var, font=("Helvetica", 9), bg='#f5f5f5', fg='#666666')
grr_status_label.pack(side=tk.RIGHT, padx=20)

# Separator
ttk.Separator(grr_main_frame, orient='horizontal').pack(fill=tk.X, padx=10, pady=5)

# Info label for instructions
grr_info_label = tk.Label(
    grr_main_frame,
    text="📋 Load files (CSV, TXT, STDF), then click on wafermaps to select dies for comparison. Selected dies are highlighted with black frames.",
    font=("Helvetica", 9),
    bg='#f5f5f5',
    fg='#3498DB'
)
grr_info_label.pack(pady=5)

# Selected dies display frame
grr_selected_frame = tk.Frame(grr_main_frame, bg='#E8F4F8', relief='groove', bd=1)
grr_selected_frame.pack(fill=tk.X, padx=10, pady=5)

tk.Label(grr_selected_frame, text="Selected Dies:", font=("Helvetica", 10, "bold"), bg='#E8F4F8').pack(side=tk.LEFT, padx=10, pady=5)
grr_selected_dies_label = tk.Label(grr_selected_frame, text="None selected", font=("Helvetica", 9), bg='#E8F4F8', fg='#2C3E50')
grr_selected_dies_label.pack(side=tk.LEFT, padx=10, pady=5)

# Clear selection button
def clear_grr_selection():
    global grr_selected_dies
    grr_selected_dies = []
    grr_selected_dies_label.config(text="None selected")
    update_grr_wafermaps()
    grr_status_var.set(f"Selection cleared")

grr_clear_btn = tk.Button(grr_selected_frame, text="Clear Selection", command=clear_grr_selection, font=("Helvetica", 8), bg='#E74C3C', fg='white')
grr_clear_btn.pack(side=tk.RIGHT, padx=10, pady=5)

# ============================================================================
# Main Notebook for Gage R&R - Two tabs: Maps/Selection and Results/Analysis
# ============================================================================
grr_main_notebook = ttk.Notebook(grr_main_frame, style="GRRMain.TNotebook")
grr_main_notebook.pack(fill=tk.BOTH, expand=True, padx=5, pady=2)

# Style for the main notebook tabs - reduced size by 50%
style = ttk.Style()
style.configure("GRRMain.TNotebook.Tab", font=("Helvetica", 8), padding=[7, 2])

# ============================================================================
# TAB 1: Wafermap Selection Tab
# ============================================================================
grr_maps_tab = tk.Frame(grr_main_notebook, bg='white')
grr_main_notebook.add(grr_maps_tab, text="🗺️ Wafermap Selection")

# Wafermap display area
grr_content_frame = tk.LabelFrame(grr_maps_tab, text="📁 Loaded Files - Wafermaps (click to select dies)", bg='white', font=("Helvetica", 10, "bold"))
grr_content_frame.pack(fill=tk.BOTH, expand=True, padx=5, pady=5)

# Create canvas with scrollbars inside the frame
grr_wafermap_inner = tk.Frame(grr_content_frame, bg='white')
grr_wafermap_inner.pack(fill=tk.BOTH, expand=True)

grr_content_canvas = tk.Canvas(grr_wafermap_inner, bg='white', highlightthickness=1, highlightbackground='#BDC3C7', height=500)
grr_content_scrollbar_x = ttk.Scrollbar(grr_wafermap_inner, orient="horizontal", command=grr_content_canvas.xview)
grr_content_scrollbar_y = ttk.Scrollbar(grr_wafermap_inner, orient="vertical", command=grr_content_canvas.yview)
grr_content_canvas.configure(xscrollcommand=grr_content_scrollbar_x.set, yscrollcommand=grr_content_scrollbar_y.set)

# Pack scrollbars and canvas
grr_content_scrollbar_y.pack(side=tk.RIGHT, fill=tk.Y)
grr_content_scrollbar_x.pack(side=tk.BOTTOM, fill=tk.X)
grr_content_canvas.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)

# Create container frame inside canvas
grr_files_container = tk.Frame(grr_content_canvas, bg='white')
grr_canvas_window = grr_content_canvas.create_window((0, 0), window=grr_files_container, anchor='nw')

def on_grr_container_configure(event):
    grr_content_canvas.configure(scrollregion=grr_content_canvas.bbox("all"))

def on_grr_canvas_configure(event):
    # Make the inner frame at least as wide as the canvas
    grr_content_canvas.itemconfig(grr_canvas_window, height=max(500, grr_files_container.winfo_reqheight()))

grr_files_container.bind("<Configure>", on_grr_container_configure)
grr_content_canvas.bind("<Configure>", on_grr_canvas_configure)

# ============================================================================
# TAB 2: Results and Analysis Tab
# ============================================================================
grr_results_tab = tk.Frame(grr_main_notebook, bg='#f5f5f5')
grr_main_notebook.add(grr_results_tab, text="📊 Results & Analysis")

# Buttons frame at top of results tab
grr_buttons_frame = tk.Frame(grr_results_tab, bg='#f5f5f5')
grr_buttons_frame.pack(fill=tk.X, pady=5, padx=5)

# Results display frame - contains notebook with sub-tabs
grr_results_display_frame = tk.Frame(grr_results_tab, bg='white', relief='groove', bd=1)
grr_results_display_frame.pack(fill=tk.BOTH, expand=True, padx=2, pady=2)

# Results title
grr_results_title = tk.Label(
    grr_results_display_frame,
    text="📊 Gage R&R Analysis Results",
    font=("Helvetica", 9, "bold"),
    bg='white',
    fg='#2C3E50'
)
grr_results_title.pack(pady=2)

# Notebook for results sub-tabs (Summary, Table, Die Statistics, Graphs)
grr_results_notebook = ttk.Notebook(grr_results_display_frame, style="GRRSub.TNotebook")
grr_results_notebook.pack(fill=tk.BOTH, expand=True, padx=2, pady=2)

# Style for sub-notebook tabs - reduced size for more graph space
style.configure("GRRSub.TNotebook.Tab", font=("Helvetica", 8), padding=[5, 2])

# Tab 1: Summary Text
grr_summary_tab = tk.Frame(grr_results_notebook, bg='white')
grr_results_notebook.add(grr_summary_tab, text="Summary")

# Results text widget for detailed display
grr_results_text = tk.Text(
    grr_summary_tab,
    height=15,
    width=120,
    font=("Consolas", 9),
    bg='#FAFAFA',
    fg='#2C3E50',
    relief='flat',
    padx=10,
    pady=10
)
grr_results_text.pack(fill=tk.BOTH, expand=True, padx=5, pady=5)
grr_results_text.insert('1.0', "No analysis performed yet.\n\nInstructions:\n1. Load files using 'Load Files' button\n2. Click on wafermaps to select dies for comparison\n3. Click 'Run Gage R&R Analysis' to calculate results")
grr_results_text.config(state='disabled')

# Tab 2: Data Table
grr_table_tab = tk.Frame(grr_results_notebook, bg='white')
grr_results_notebook.add(grr_table_tab, text="Data Table")

# Table controls frame
grr_table_controls = tk.Frame(grr_table_tab, bg='white')
grr_table_controls.pack(fill=tk.X, padx=5, pady=5)

tk.Label(grr_table_controls, text="Select Parameter:", font=("Helvetica", 9, "bold"), bg='white').pack(side=tk.LEFT, padx=5)
grr_table_param_var = tk.StringVar(value="All Parameters")
grr_table_param_combo = ttk.Combobox(grr_table_controls, textvariable=grr_table_param_var, values=["All Parameters"], state="readonly", width=40)
grr_table_param_combo.pack(side=tk.LEFT, padx=5)

# Export table button
def export_grr_table():
    """Export the Gage R&R table to CSV"""
    if not hasattr(grr_data_table, 'get_children') or not grr_data_table.get_children():
        grr_status_var.set("No data to export")
        return

    save_path = filedialog.asksaveasfilename(
        title="Export Gage R&R Table",
        defaultextension=".csv",
        filetypes=[("CSV files", "*.csv"), ("All files", "*.*")]
    )

    if save_path:
        with open(save_path, 'w') as f:
            # Write headers
            headers = [grr_data_table.heading(col)['text'] for col in grr_data_table['columns']]
            f.write(','.join(headers) + '\n')

            # Write data
            for item in grr_data_table.get_children():
                values = grr_data_table.item(item)['values']
                f.write(','.join(map(str, values)) + '\n')

        grr_status_var.set(f"Table exported to {os.path.basename(save_path)}")

grr_export_table_btn = tk.Button(grr_table_controls, text="📥 Export Table", command=export_grr_table,
                                  font=("Helvetica", 9), bg='#3498DB', fg='white')
grr_export_table_btn.pack(side=tk.RIGHT, padx=5)

# Treeview table for data display
grr_table_frame = tk.Frame(grr_table_tab, bg='white')
grr_table_frame.pack(fill=tk.BOTH, expand=True, padx=5, pady=5)

# Scrollbars for table
grr_table_scroll_y = ttk.Scrollbar(grr_table_frame, orient="vertical")
grr_table_scroll_y.pack(side=tk.RIGHT, fill=tk.Y)

grr_table_scroll_x = ttk.Scrollbar(grr_table_frame, orient="horizontal")
grr_table_scroll_x.pack(side=tk.BOTTOM, fill=tk.X)

# Create Treeview
grr_data_table = ttk.Treeview(
    grr_table_frame,
    yscrollcommand=grr_table_scroll_y.set,
    xscrollcommand=grr_table_scroll_x.set,
    show='headings'
)
grr_data_table.pack(fill=tk.BOTH, expand=True)

grr_table_scroll_y.config(command=grr_data_table.yview)
grr_table_scroll_x.config(command=grr_data_table.xview)

# Style for the table
style = ttk.Style()
style.configure("Treeview", font=("Consolas", 9), rowheight=25)
style.configure("Treeview.Heading", font=("Helvetica", 9, "bold"))

# Tab 3: Statistics per Die
grr_die_stats_tab = tk.Frame(grr_results_notebook, bg='white')
grr_results_notebook.add(grr_die_stats_tab, text="Die Statistics")

# Die statistics table
grr_die_table_frame = tk.Frame(grr_die_stats_tab, bg='white')
grr_die_table_frame.pack(fill=tk.BOTH, expand=True, padx=5, pady=5)

grr_die_scroll_y = ttk.Scrollbar(grr_die_table_frame, orient="vertical")
grr_die_scroll_y.pack(side=tk.RIGHT, fill=tk.Y)

grr_die_scroll_x = ttk.Scrollbar(grr_die_table_frame, orient="horizontal")
grr_die_scroll_x.pack(side=tk.BOTTOM, fill=tk.X)

grr_die_table = ttk.Treeview(
    grr_die_table_frame,
    yscrollcommand=grr_die_scroll_y.set,
    xscrollcommand=grr_die_scroll_x.set,
    show='headings'
)
grr_die_table.pack(fill=tk.BOTH, expand=True)

grr_die_scroll_y.config(command=grr_die_table.yview)
grr_die_scroll_x.config(command=grr_die_table.xview)

# Tab 4: Graph Plotting for selected parameters
grr_graph_tab = tk.Frame(grr_results_notebook, bg='white')
grr_results_notebook.add(grr_graph_tab, text="📈 Graphs")

# Graph controls frame
grr_graph_controls = tk.Frame(grr_graph_tab, bg='white')
grr_graph_controls.pack(fill=tk.X, padx=5, pady=5)

tk.Label(grr_graph_controls, text="Select Parameter:", font=("Helvetica", 9, "bold"), bg='white').pack(side=tk.LEFT, padx=5)
grr_graph_param_var = tk.StringVar(value="Select parameter...")
grr_graph_param_combo = ttk.Combobox(grr_graph_controls, textvariable=grr_graph_param_var, values=[], state="readonly", width=40)
grr_graph_param_combo.pack(side=tk.LEFT, padx=5)

# Plot button
def plot_grr_graph():
    """Plot graph for selected parameter showing values across runs/files"""
    global grr_file_data, grr_selected_dies, grr_graph_canvas

    param = grr_graph_param_var.get()
    if param == "Select parameter..." or not param:
        grr_status_var.set("Select a parameter to plot")
        return

    if len(grr_file_data) < 1:
        grr_status_var.set("Load files first")
        return

    if len(grr_selected_dies) < 1:
        grr_status_var.set("Select at least 1 die")
        return

    # Clear existing plot
    for widget in grr_graph_frame.winfo_children():
        widget.destroy()

    # Collect data for plotting
    file_names = [os.path.basename(f['path'])[:20] for f in grr_file_data]
    all_values = []  # list of values per die

    for die_coord in grr_selected_dies:
        x, y = die_coord
        die_values = []

        for file_info in grr_file_data:
            if file_info['type'] == 'stdf' or file_info['type'] == 'csv_wafermap':
                df = file_info['data']
                param_col = 'bin' if param == 'BIN' else param

                if param_col in df.columns:
                    die_data = df[(df['x'] == x) & (df['y'] == y)]
                    if len(die_data) > 0:
                        die_values.append(die_data[param_col].values[0])
                    else:
                        die_values.append(np.nan)
                else:
                    die_values.append(np.nan)
            else:
                data = file_info['data']
                if hasattr(data, 'shape') and x < data.shape[0] and y < data.shape[1]:
                    die_values.append(data[x, y])
                else:
                    die_values.append(np.nan)

        all_values.append(die_values)

    # Create figure with two subplots
    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(12, 4))
    fig.patch.set_facecolor('white')

    # Plot 1: Line plot showing values across files/runs for each die
    x_pos = np.arange(len(file_names))
    colors = plt.cm.tab10(np.linspace(0, 1, len(grr_selected_dies)))

    for i, (die_coord, values) in enumerate(zip(grr_selected_dies, all_values)):
        valid_values = [v for v in values if not np.isnan(v)]
        valid_idx = [j for j, v in enumerate(values) if not np.isnan(v)]

        if valid_values:
            ax1.plot(valid_idx, valid_values, 'o-', label=f"Die ({die_coord[0]},{die_coord[1]})",
                    color=colors[i], linewidth=2, markersize=8)
            ax1.axhline(y=np.mean(valid_values), color=colors[i], linestyle='--', alpha=0.5)

    ax1.set_xticks(x_pos)
    ax1.set_xticklabels([f"Run {i+1}" for i in range(len(file_names))], fontsize=8, rotation=45, ha='right')
    ax1.set_xlabel("File/Run", fontsize=10)
    ax1.set_ylabel("Value", fontsize=10)
    ax1.set_title(f"Values Across Runs\n{param[:50]}", fontsize=10, fontweight='bold')
    ax1.legend(fontsize=7, loc='best')
    ax1.grid(True, alpha=0.3)
    ax1.set_facecolor('#FAFAFA')

    # Calculate overall stats for reference lines
    all_flat = [v for values in all_values for v in values if not np.isnan(v)]
    if all_flat:
        overall_mean = np.mean(all_flat)
        overall_std = np.std(all_flat)

        # Add reference lines (limits)
        ax1.axhline(y=overall_mean, color='green', linestyle='-', linewidth=2, alpha=0.7, label='Mean')
        ax1.axhline(y=overall_mean + 3*overall_std, color='red', linestyle='--', linewidth=1.5, alpha=0.7, label='+3σ')
        ax1.axhline(y=overall_mean - 3*overall_std, color='red', linestyle='--', linewidth=1.5, alpha=0.7, label='-3σ')

        # Add statistics text
        stats_text = f"Mean: {overall_mean:.4g}\nStd: {overall_std:.4g}\n+3σ: {overall_mean + 3*overall_std:.4g}\n-3σ: {overall_mean - 3*overall_std:.4g}"
        ax1.text(0.02, 0.98, stats_text, transform=ax1.transAxes, fontsize=7,
                verticalalignment='top', bbox=dict(boxstyle='round', facecolor='white', alpha=0.9))

    # Plot 2: Box plot comparing values per file/run
    plot_data = []
    for j in range(len(grr_file_data)):
        run_values = [all_values[i][j] for i in range(len(all_values)) if not np.isnan(all_values[i][j])]
        plot_data.append(run_values)

    if any(len(d) > 0 for d in plot_data):
        bp = ax2.boxplot(plot_data, labels=[f"Run {i+1}" for i in range(len(plot_data))],
                        patch_artist=True, showmeans=True,
                        meanprops=dict(marker='D', markerfacecolor='red', markeredgecolor='white', markersize=6))

        # Color boxes
        colors_box = plt.cm.Pastel1(np.linspace(0, 1, len(plot_data)))
        for patch, color in zip(bp['boxes'], colors_box):
            patch.set_facecolor(color)
            patch.set_alpha(0.7)

    ax2.set_xlabel("File/Run", fontsize=10)
    ax2.set_ylabel("Value", fontsize=10)
    ax2.set_title(f"Distribution per Run\n{param[:50]}", fontsize=10, fontweight='bold')
    ax2.grid(True, alpha=0.3, axis='y')
    ax2.set_facecolor('#FAFAFA')
    ax2.tick_params(axis='x', labelrotation=45)

    fig.tight_layout()

    # Create canvas
    grr_graph_canvas = FigureCanvasTkAgg(fig, master=grr_graph_frame)
    grr_graph_canvas.get_tk_widget().pack(fill=tk.BOTH, expand=True)
    grr_graph_canvas.draw()

    grr_status_var.set(f"Graph plotted for {param[:40]}")

grr_plot_btn = tk.Button(grr_graph_controls, text="📊 Plot Graph", command=plot_grr_graph,
                          font=("Helvetica", 9, "bold"), bg='#27AE60', fg='white')
grr_plot_btn.pack(side=tk.LEFT, padx=10)

# Graph display frame
grr_graph_frame = tk.Frame(grr_graph_tab, bg='white')
grr_graph_frame.pack(fill=tk.BOTH, expand=True, padx=5, pady=5)

# Placeholder for graph
grr_graph_placeholder = tk.Label(grr_graph_frame,
    text="Select a parameter and click 'Plot Graph' to visualize values across runs.\n\nThe graph will show:\n• Line plot: Values for each die across runs/files\n• Box plot: Distribution comparison per run\n• Reference lines: Mean and ±3σ limits",
    font=("Helvetica", 10), bg='white', fg='#666666', justify='center')
grr_graph_placeholder.pack(pady=50)

grr_graph_canvas = None

# Function to update the data table with multiple parameters
def update_grr_multi_param_table(all_param_results, file_names, die_coords):
    """Update the Gage R&R data table with all parameters and their values"""
    # Clear existing data
    for item in grr_data_table.get_children():
        grr_data_table.delete(item)

    # Define columns: Parameter, Die, File values, Statistics
    columns = ['Parameter', 'Die (X,Y)'] + [f'File {i+1}' for i, name in enumerate(file_names)]
    columns += ['Mean', 'Std Dev', 'Range', 'CV%']

    grr_data_table['columns'] = columns

    # Configure columns
    for col in columns:
        grr_data_table.heading(col, text=col)
        if col == 'Parameter':
            width = 200
        elif col == 'Die (X,Y)':
            width = 80
        else:
            width = 100
        grr_data_table.column(col, width=width, anchor='center')

    # Add data rows for each parameter and die
    for param, res in all_param_results.items():
        comparison_data = res['comparison_data']

        for i, die_coord in enumerate(die_coords):
            row_data = [param[:50], f"({die_coord[0]}, {die_coord[1]})"]

            # Values from each file
            values = []
            for j, file_data in enumerate(comparison_data):
                if i < len(file_data):
                    val = file_data[i]
                    row_data.append(f"{val:.4g}" if not np.isnan(val) else "N/A")
                    if not np.isnan(val):
                        values.append(val)
                else:
                    row_data.append("N/A")

            # Calculate statistics for this die across files
            if values:
                mean_val = np.mean(values)
                std_val = np.std(values)
                range_val = np.max(values) - np.min(values)
                cv_val = (std_val / mean_val * 100) if mean_val != 0 else 0
                row_data.extend([f"{mean_val:.4g}", f"{std_val:.4g}", f"{range_val:.4g}", f"{cv_val:.2f}%"])
            else:
                row_data.extend(["N/A", "N/A", "N/A", "N/A"])

            grr_data_table.insert('', 'end', values=row_data)

def update_grr_summary_table(all_param_results):
    """Update the summary statistics table for all parameters"""
    # Clear existing data
    for item in grr_die_table.get_children():
        grr_die_table.delete(item)

    # Define columns for Gage R&R summary
    columns = ['Parameter', 'Mean', 'Std Dev', 'Min', 'Max', 'Range',
               'Repeatability', 'Reproducibility', 'GRR Total', '%GRR', 'ndc', 'Assessment']
    grr_die_table['columns'] = columns

    # Configure columns
    col_widths = {'Parameter': 200, 'Mean': 90, 'Std Dev': 90, 'Min': 90, 'Max': 90,
                  'Range': 90, 'Repeatability': 100, 'Reproducibility': 110,
                  'GRR Total': 90, '%GRR': 80, 'ndc': 60, 'Assessment': 100}

    for col in columns:
        grr_die_table.heading(col, text=col)
        grr_die_table.column(col, width=col_widths.get(col, 80), anchor='center')

    # Add row for each parameter
    for param, res in all_param_results.items():
        row_data = [
            param[:50],
            f"{res['overall_mean']:.4g}",
            f"{res['overall_std']:.4g}",
            f"{res['overall_min']:.4g}",
            f"{res['overall_max']:.4g}",
            f"{res['range']:.4g}",
            f"{res['repeatability']:.4g}",
            f"{res['reproducibility']:.4g}",
            f"{res['grr_total']:.4g}",
            f"{res['grr_percent']:.2f}%",
            f"{res['ndc']:.2f}",
            res['assessment']
        ]
        grr_die_table.insert('', 'end', values=row_data)

# Keep old functions for backward compatibility
def update_grr_data_table(comparison_data, file_names, die_coords, param_name):
    """Update the Gage R&R data table with comparison values (single parameter)"""
    pass  # No longer used, kept for compatibility

def update_grr_die_stats_table(comparison_data, file_names, die_coords):
    """Update the die statistics table (single parameter)"""
    pass  # No longer used, kept for compatibility

# Analysis button
def run_grr_analysis():
    """Run Gage R&R analysis on selected dies for all selected parameters"""
    global grr_file_data, grr_selected_dies, grr_selected_params, grr_analysis_results

    if len(grr_file_data) < 2:
        grr_status_var.set("Error: Load at least 2 files for comparison")
        return

    if len(grr_selected_dies) < 1:
        grr_status_var.set("Error: Select at least 1 die for comparison")
        return

    if len(grr_selected_params) < 1:
        grr_status_var.set("Error: Select at least 1 parameter for analysis")
        return

    # Store all results per parameter
    all_param_results = {}

    for param in grr_selected_params:
        # Extract data for selected dies from all files for this parameter
        comparison_data = []

        for file_info in grr_file_data:
            file_values = []
            for die_coord in grr_selected_dies:
                x, y = die_coord

                if file_info['type'] == 'stdf' or file_info['type'] == 'csv_wafermap':
                    df = file_info['data']
                    param_col = 'bin' if param == 'BIN' else param

                    if param_col in df.columns:
                        die_data = df[(df['x'] == x) & (df['y'] == y)]
                        if len(die_data) > 0:
                            file_values.append(die_data[param_col].values[0])
                        else:
                            file_values.append(np.nan)
                    else:
                        file_values.append(np.nan)
                else:
                    # For CSV/TXT matrix files
                    data = file_info['data']
                    if hasattr(data, 'shape') and x < data.shape[0] and y < data.shape[1]:
                        file_values.append(data[x, y])
                    else:
                        file_values.append(np.nan)

            comparison_data.append(file_values)

        # Calculate basic Gage R&R metrics
        comparison_array = np.array(comparison_data, dtype=float)

        # Remove columns (dies) where any file has NaN
        valid_cols = ~np.isnan(comparison_array).any(axis=0)
        valid_data = comparison_array[:, valid_cols]

        if valid_data.size > 0 and valid_data.shape[1] > 0:
            # Calculate repeatability and reproducibility
            overall_mean = np.nanmean(valid_data)
            overall_std = np.nanstd(valid_data)

            # Between-file variation (reproducibility)
            file_means = np.nanmean(valid_data, axis=1)
            reproducibility = np.std(file_means)

            # Within-file variation (repeatability)
            repeatability = np.mean([np.nanstd(row) for row in valid_data])

            # Total GRR
            grr_total = np.sqrt(repeatability**2 + reproducibility**2)

            # Calculate additional metrics
            total_variation = overall_std
            grr_percent = (grr_total / total_variation * 100) if total_variation > 0 else 0
            repeatability_percent = (repeatability / total_variation * 100) if total_variation > 0 else 0
            reproducibility_percent = (reproducibility / total_variation * 100) if total_variation > 0 else 0

            # Number of distinct categories (ndc)
            ndc = 1.41 * (total_variation / grr_total) if grr_total > 0 else 0

            # Assessment
            if grr_percent < 10:
                assessment = "EXCELLENT"
            elif grr_percent < 30:
                assessment = "ACCEPTABLE"
            else:
                assessment = "UNACCEPTABLE"

            all_param_results[param] = {
                'comparison_data': comparison_data,
                'valid_data': valid_data,
                'overall_mean': overall_mean,
                'overall_std': overall_std,
                'overall_min': np.nanmin(valid_data),
                'overall_max': np.nanmax(valid_data),
                'range': np.nanmax(valid_data) - np.nanmin(valid_data),
                'repeatability': repeatability,
                'reproducibility': reproducibility,
                'grr_total': grr_total,
                'total_variation': total_variation,
                'grr_percent': grr_percent,
                'repeatability_percent': repeatability_percent,
                'reproducibility_percent': reproducibility_percent,
                'ndc': ndc,
                'assessment': assessment,
                'file_means': file_means
            }

    if not all_param_results:
        grr_status_var.set("Error: No valid data found for any parameter")
        return

    # Build detailed results text for all parameters
    results_text = "=" * 100 + "\n"
    results_text += "                         GAGE R&R ANALYSIS RESULTS - MULTI-PARAMETER\n"
    results_text += "=" * 100 + "\n\n"

    results_text += "STUDY INFORMATION:\n"
    results_text += "-" * 50 + "\n"
    results_text += f"  Number of Files (Operators/Trials):  {len(grr_file_data)}\n"
    results_text += f"  Number of Dies (Parts):              {len(grr_selected_dies)}\n"
    results_text += f"  Number of Parameters:                {len(grr_selected_params)}\n"
    results_text += f"  Selected Die Coordinates:            {grr_selected_dies}\n\n"

    # Summary table for all parameters
    results_text += "PARAMETER SUMMARY TABLE:\n"
    results_text += "-" * 100 + "\n"
    results_text += f"{'Parameter':<40} {'Mean':>12} {'Std Dev':>12} {'%GRR':>10} {'ndc':>8} {'Assessment':>15}\n"
    results_text += "-" * 100 + "\n"

    for param, res in all_param_results.items():
        param_short = param[:38] if len(param) > 38 else param
        results_text += f"{param_short:<40} {res['overall_mean']:>12.4g} {res['overall_std']:>12.4g} {res['grr_percent']:>9.2f}% {res['ndc']:>8.2f} {res['assessment']:>15}\n"

    results_text += "-" * 100 + "\n\n"

    # Detailed results for each parameter
    for param, res in all_param_results.items():
        results_text += f"\n{'='*80}\n"
        results_text += f"PARAMETER: {param}\n"
        results_text += f"{'='*80}\n\n"

        results_text += "DESCRIPTIVE STATISTICS:\n"
        results_text += f"  Overall Mean:     {res['overall_mean']:.6g}\n"
        results_text += f"  Overall Std Dev:  {res['overall_std']:.6g}\n"
        results_text += f"  Min:              {res['overall_min']:.6g}\n"
        results_text += f"  Max:              {res['overall_max']:.6g}\n"
        results_text += f"  Range:            {res['range']:.6g}\n\n"

        results_text += "VARIANCE COMPONENTS:\n"
        results_text += f"  Repeatability:    {res['repeatability']:.6g} ({res['repeatability_percent']:.2f}%)\n"
        results_text += f"  Reproducibility:  {res['reproducibility']:.6g} ({res['reproducibility_percent']:.2f}%)\n"
        results_text += f"  Total Gage R&R:   {res['grr_total']:.6g} ({res['grr_percent']:.2f}%)\n\n"

        results_text += "ASSESSMENT:\n"
        results_text += f"  %GRR: {res['grr_percent']:.2f}%  |  ndc: {res['ndc']:.2f}  |  {res['assessment']}\n"

    results_text += "\n" + "=" * 100

    # Update results display
    grr_results_text.config(state='normal')
    grr_results_text.delete('1.0', tk.END)
    grr_results_text.insert('1.0', results_text)
    grr_results_text.config(state='disabled')

    # Update status bar
    grr_status_var.set(f"Analysis complete: {len(all_param_results)} parameters analyzed")

    # Update data tables with all parameters
    file_names = [os.path.basename(f['path']) for f in grr_file_data]
    update_grr_multi_param_table(all_param_results, file_names, grr_selected_dies)
    update_grr_summary_table(all_param_results)

    # Update graph parameter combo with analyzed parameters
    grr_graph_param_combo['values'] = list(all_param_results.keys())
    if all_param_results:
        grr_graph_param_var.set(list(all_param_results.keys())[0])

    # Store results globally for PPT export
    grr_analysis_results.clear()
    grr_analysis_results.update(all_param_results)
    grr_analysis_results['_meta'] = {
        'file_names': file_names,
        'die_coords': list(grr_selected_dies),
        'num_files': len(grr_file_data),
        'num_dies': len(grr_selected_dies),
        'num_params': len(grr_selected_params)
    }

    # Switch to Data Table tab to show results
    grr_results_notebook.select(1)

grr_analyze_btn = tk.Button(
    grr_buttons_frame,
    text="🔬 Run Gage R&R Analysis",
    command=run_grr_analysis,
    font=("Helvetica", 10, "bold"),
    bg='#27AE60',
    fg='white',
    padx=20,
    pady=5
)
grr_analyze_btn.pack(side=tk.LEFT, padx=10)

# Load files button
def load_grr_files():
    """Load files for Gage R&R comparison based on selected file type"""
    global grr_file_data, grr_selected_dies, grr_available_params, grr_selected_params

    num_files = int(grr_num_files_var.get())
    file_type = grr_file_type_var.get()

    # Set file filter based on selected type
    if file_type == "STDF":
        filetypes = [
            ("STDF files", "*.stdf"),
            ("All files", "*.*")
        ]
        title = f"Select {num_files} STDF files for Gage R&R comparison"
    else:  # CSV
        filetypes = [
            ("CSV files", "*.csv"),
            ("All files", "*.*")
        ]
        title = f"Select {num_files} CSV wafermap files for Gage R&R comparison"

    file_paths = filedialog.askopenfilenames(
        title=title,
        filetypes=filetypes
    )

    if not file_paths:
        return

    if len(file_paths) != num_files:
        grr_status_var.set(f"Warning: Selected {len(file_paths)} files, expected {num_files}")

    # Clear existing data
    grr_file_data = []
    grr_selected_dies = []
    grr_selected_params = []
    grr_available_params = ["BIN"]
    grr_selected_dies_label.config(text="None selected")
    grr_param_label.config(text="0 selected")

    # Load each file
    for path in file_paths:
        file_ext = os.path.splitext(path)[1].lower()
        file_info = {'path': path, 'type': None, 'data': None, 'wafer_id': None}

        try:
            if file_ext == '.stdf':
                # Load STDF file using existing parser
                df, params, limits, groups, wafer_id = parse_stdf_file(path)
                file_info['type'] = 'stdf'
                file_info['data'] = df
                file_info['wafer_id'] = wafer_id
                file_info['params'] = params

                # Update available parameters with STDF parameters
                if params:
                    for k, v in params.items():
                        param_name = f"{k}: {v}"
                        if param_name not in grr_available_params:
                            grr_available_params.append(param_name)

            elif file_ext == '.csv':
                # Load CSV wafermap file - header row + data with comma delimiter
                # Format: x, y coordinates + measurement parameters
                try:
                    import pandas as pd
                    # Load CSV with comma delimiter and header
                    df_csv = pd.read_csv(path, delimiter=',')

                    print(f"CSV columns: {list(df_csv.columns)}")

                    # Look for x and y coordinate columns (case-insensitive)
                    x_col = None
                    y_col = None
                    for col in df_csv.columns:
                        col_lower = col.lower().strip()
                        if col_lower in ['x', 'x_coord', 'x_coordinate', 'xcoord', 'die_x']:
                            x_col = col
                        elif col_lower in ['y', 'y_coord', 'y_coordinate', 'ycoord', 'die_y']:
                            y_col = col

                    if x_col and y_col:
                        # This is a wafermap format with x, y coordinates
                        # Convert to numeric
                        df_csv[x_col] = pd.to_numeric(df_csv[x_col], errors='coerce')
                        df_csv[y_col] = pd.to_numeric(df_csv[y_col], errors='coerce')

                        # Rename columns to standard names
                        df_csv = df_csv.rename(columns={x_col: 'x', y_col: 'y'})

                        # Get numeric parameter columns (excluding x, y)
                        param_cols = []
                        for col in df_csv.columns:
                            if col not in ['x', 'y']:
                                df_csv[col] = pd.to_numeric(df_csv[col], errors='coerce')
                                if df_csv[col].notna().any():
                                    param_cols.append(col)

                        # Store as DataFrame (similar to STDF format)
                        file_info['type'] = 'csv_wafermap'
                        file_info['data'] = df_csv
                        file_info['wafer_id'] = os.path.basename(path)
                        file_info['params'] = {col: col for col in param_cols}

                        # Update available parameters list
                        if param_cols:
                            for col in param_cols:
                                if col not in grr_available_params:
                                    grr_available_params.append(col)

                        print(f"CSV wafermap loaded: {path}, dies={len(df_csv)}, params={len(param_cols)}")
                        print(f"  X range: {df_csv['x'].min()} to {df_csv['x'].max()}")
                        print(f"  Y range: {df_csv['y'].min()} to {df_csv['y'].max()}")
                    else:
                        # No x, y columns found - treat as regular data matrix
                        numeric_cols = df_csv.select_dtypes(include=[np.number]).columns
                        if len(numeric_cols) > 0:
                            data = df_csv[numeric_cols].values
                        else:
                            data = df_csv.apply(pd.to_numeric, errors='coerce').values

                        file_info['type'] = 'csv'
                        file_info['data'] = data
                        file_info['wafer_id'] = os.path.basename(path)

                        valid = data[~np.isnan(data)]
                        if len(valid) > 0:
                            print(f"CSV matrix loaded: {path}, shape={data.shape}")

                except Exception as e:
                    print(f"CSV loading error: {e}")
                    # Fallback to numpy
                    try:
                        data = np.genfromtxt(path, delimiter=',', skip_header=1)
                        file_info['type'] = 'csv'
                        file_info['data'] = data
                        file_info['wafer_id'] = os.path.basename(path)
                    except Exception as e2:
                        print(f"Fallback also failed: {e2}")

            elif file_ext == '.txt':
                # Load TXT file (assume space or tab delimited)
                try:
                    data = np.loadtxt(path)
                except:
                    try:
                        data = np.genfromtxt(path, delimiter='\t')
                    except:
                        data = np.genfromtxt(path, delimiter=' ')

                # Ensure data is 2D
                if data.ndim == 1:
                    n = len(data)
                    cols = int(np.ceil(np.sqrt(n)))
                    rows = int(np.ceil(n / cols))
                    padded = np.full(rows * cols, np.nan)
                    padded[:n] = data
                    data = padded.reshape(rows, cols)

                file_info['type'] = 'txt'
                file_info['data'] = data
                file_info['wafer_id'] = os.path.basename(path)

                # Calculate basic stats for status
                valid = data[~np.isnan(data)]
                if len(valid) > 0:
                    print(f"TXT loaded: {path}, shape={data.shape}, min={np.min(valid):.3g}, max={np.max(valid):.3g}")

            grr_file_data.append(file_info)

        except Exception as e:
            print(f"Error loading {path}: {e}")
            grr_status_var.set(f"Error loading {os.path.basename(path)}: {str(e)[:50]}")

    grr_status_var.set(f"Loaded {len(grr_file_data)} files successfully")

    # Update visualization parameter combobox
    viz_params = ["None"] + [p for p in grr_available_params if p != "BIN"]
    grr_viz_param_combo['values'] = viz_params
    grr_viz_param_var.set("None")

    # Update display
    update_grr_file_displays()

grr_load_btn = tk.Button(
    grr_buttons_frame,
    text="📁 Load Files",
    command=load_grr_files,
    font=("Helvetica", 10, "bold"),
    bg='#3498DB',
    fg='white',
    padx=20,
    pady=5
)
grr_load_btn.pack(side=tk.LEFT, padx=10)

# Export results button
def export_grr_results():
    """Export Gage R&R results to CSV"""
    if len(grr_file_data) < 2 or len(grr_selected_dies) < 1:
        grr_status_var.set("Error: Run analysis first before exporting")
        return

    save_path = filedialog.asksaveasfilename(
        title="Save Gage R&R Results",
        defaultextension=".csv",
        filetypes=[("CSV files", "*.csv"), ("All files", "*.*")]
    )

    if save_path:
        # Create export data
        export_data = []
        header = ["Die_X", "Die_Y"] + [os.path.basename(f['path']) for f in grr_file_data]
        export_data.append(header)

        for die_coord in grr_selected_dies:
            x, y = die_coord
            row = [x, y]

            for file_info in grr_file_data:
                if file_info['type'] == 'stdf' or file_info['type'] == 'csv_wafermap':
                    df = file_info['data']
                    param = grr_param_var.get()
                    param_col = 'bin' if param == 'BIN' else param

                    die_data = df[(df['x'] == x) & (df['y'] == y)]
                    if len(die_data) > 0 and param_col in df.columns:
                        row.append(die_data[param_col].values[0])
                    else:
                        row.append('')
                else:
                    data = file_info['data']
                    if hasattr(data, 'shape') and x < data.shape[0] and y < data.shape[1]:
                        row.append(data[x, y])
                    else:
                        row.append('')

            export_data.append(row)

        # Write to CSV
        with open(save_path, 'w') as f:
            for row in export_data:
                f.write(','.join(map(str, row)) + '\n')

        grr_status_var.set(f"Results exported to {os.path.basename(save_path)}")

grr_export_btn = tk.Button(
    grr_buttons_frame,
    text="💾 Export Results",
    command=export_grr_results,
    font=("Helvetica", 10, "bold"),
    bg='#9B59B6',
    fg='white',
    padx=20,
    pady=5
)
grr_export_btn.pack(side=tk.LEFT, padx=10)


def update_grr_file_displays():
    """Update the file display panels with loaded data"""
    global grr_file_canvases, grr_file_frames

    # Clear existing displays
    for widget in grr_files_container.winfo_children():
        widget.destroy()

    grr_file_canvases = []
    grr_file_frames = []

    if not grr_file_data:
        # Show placeholder
        placeholder = tk.Label(grr_files_container, text="No files loaded. Click 'Load Files' to begin.",
                              font=("Helvetica", 11), bg='white', fg='#666666')
        placeholder.pack(pady=30, padx=30)
        return

    # Create a display panel for each file
    for idx, file_info in enumerate(grr_file_data):
        # File frame with fixed size - LARGER for better visualization
        file_frame = tk.Frame(grr_files_container, bg='white', relief='raised', bd=2, width=400, height=380)
        file_frame.pack(side=tk.LEFT, padx=8, pady=8)
        file_frame.pack_propagate(False)
        grr_file_frames.append(file_frame)

        # File title
        file_name = os.path.basename(file_info['path'])
        title_text = f"File {idx + 1}: {file_name[:20]}..." if len(file_name) > 20 else f"File {idx + 1}: {file_name}"
        title_label = tk.Label(file_frame, text=title_text, font=("Helvetica", 8, "bold"), bg='white', fg='#2C3E50')
        title_label.pack(pady=3)

        # Create wafermap/data display
        if file_info['type'] == 'stdf' or file_info['type'] == 'csv_wafermap':
            create_grr_stdf_display(file_frame, file_info, idx)
        else:
            create_grr_data_display(file_frame, file_info, idx)

    # Update scroll region
    grr_files_container.update_idletasks()
    grr_content_canvas.configure(scrollregion=grr_content_canvas.bbox("all"))


def create_grr_stdf_display(parent_frame, file_info, file_idx):
    """Create wafermap display for STDF file with clickable dies"""
    global grr_file_canvases

    df = file_info['data']

    if df is None or (hasattr(df, 'empty') and df.empty):
        tk.Label(parent_frame, text="No data", font=("Helvetica", 9), bg='white').pack()
        return

    # Check if we have x and y columns
    if 'x' not in df.columns or 'y' not in df.columns:
        tk.Label(parent_frame, text="Missing x/y", font=("Helvetica", 9), bg='white').pack()
        return

    # Get data with valid x,y coordinates only
    plot_data = df[df['x'].notna() & df['y'].notna()]

    if len(plot_data) == 0:
        tk.Label(parent_frame, text="No valid points", font=("Helvetica", 9), bg='white').pack()
        return

    # Create LARGER figure
    fig, ax = plt.subplots(figsize=(4.2, 3.6))
    fig.patch.set_facecolor('white')

    # Check if visualization parameter is selected
    viz_param = grr_viz_param_var.get() if grr_viz_param_var else "None"

    if viz_param and viz_param != "None" and viz_param in df.columns:
        # Color by the selected visualization parameter
        color_data = plot_data[viz_param].values
        valid_mask = pd.notna(color_data)

        if valid_mask.sum() > 0:
            scatter = ax.scatter(
                plot_data["x"].values[valid_mask],
                plot_data["y"].values[valid_mask],
                c=color_data[valid_mask],
                s=40,
                cmap='viridis',
                edgecolors="gray",
                linewidth=0.3
            )
            # Add colorbar
            cbar = plt.colorbar(scatter, ax=ax, shrink=0.8, pad=0.02)
            cbar.ax.tick_params(labelsize=6)
            # Short label
            param_short = viz_param[:15] + "..." if len(viz_param) > 15 else viz_param
            cbar.set_label(param_short, fontsize=6)
    else:
        # Simply plot all dies as blue points
        ax.scatter(
            plot_data["x"].values,
            plot_data["y"].values,
            c='#3498DB',
            s=40,
            edgecolors="gray",
            linewidth=0.3
        )

    # Highlight selected dies
    for die_coord in grr_selected_dies:
        x, y = die_coord
        if ((plot_data['x'] == x) & (plot_data['y'] == y)).any():
            ax.scatter([x], [y], c='red', s=60, edgecolors='black', linewidth=2)

    ax.set_xlabel("X", fontsize=8)
    ax.set_ylabel("Y", fontsize=8)
    ax.set_title(f"Click to select", fontsize=8, color='#3498DB')
    ax.set_aspect("equal")
    ax.tick_params(axis='both', labelsize=6)
    fig.tight_layout()

    # Create canvas
    canvas = FigureCanvasTkAgg(fig, master=parent_frame)
    canvas.get_tk_widget().pack(fill=tk.BOTH, expand=True, padx=2, pady=2)
    canvas.draw()

    grr_file_canvases.append({'canvas': canvas, 'fig': fig, 'ax': ax, 'file_idx': file_idx, 'data': df})

    # Click handler
    def on_click(event):
        if event.inaxes == ax and event.button == 1 and event.xdata and event.ydata:
            x_click = int(round(event.xdata))
            y_click = int(round(event.ydata))
            if ((df['x'] == x_click) & (df['y'] == y_click)).any():
                toggle_grr_die_selection(x_click, y_click)

    canvas.mpl_connect('button_press_event', on_click)


def create_grr_data_display(parent_frame, file_info, file_idx):
    """Create heatmap display for CSV/TXT file with clickable cells and statistics"""
    global grr_file_canvases

    data = file_info['data']

    if data is None or (hasattr(data, '__len__') and len(data) == 0):
        tk.Label(parent_frame, text="No data", font=("Helvetica", 9), bg='white').pack()
        return

    # Convert to numpy array if needed
    if not isinstance(data, np.ndarray):
        data = np.array(data)

    # Handle 1D data by reshaping
    if data.ndim == 1:
        # Try to reshape to square-ish shape
        n = len(data)
        cols = int(np.ceil(np.sqrt(n)))
        rows = int(np.ceil(n / cols))
        padded = np.full(rows * cols, np.nan)
        padded[:n] = data
        data = padded.reshape(rows, cols)
        file_info['data'] = data  # Update file info with reshaped data

    # Create figure for heatmap - professional style
    fig, ax = plt.subplots(figsize=(4.5, 4))
    fig.patch.set_facecolor('white')

    # Calculate statistics for display
    valid_data = data[~np.isnan(data)] if np.any(np.isnan(data)) else data.flatten()
    stats_max = np.max(valid_data) if len(valid_data) > 0 else 0
    stats_min = np.min(valid_data) if len(valid_data) > 0 else 0
    stats_mean = np.mean(valid_data) if len(valid_data) > 0 else 0
    stats_std = np.std(valid_data) if len(valid_data) > 0 else 0

    # Create heatmap with professional colormap
    im = ax.imshow(data, cmap='viridis', aspect='auto', interpolation='nearest')

    # Highlight selected cells with black frame
    for die_coord in grr_selected_dies:
        x, y = die_coord
        if x < data.shape[0] and y < data.shape[1]:
            rect = plt.Rectangle((y - 0.5, x - 0.5), 1, 1, fill=False, edgecolor='black', linewidth=3)
            ax.add_patch(rect)
            rect2 = plt.Rectangle((y - 0.5, x - 0.5), 1, 1, fill=False, edgecolor='yellow', linewidth=1.5)
            ax.add_patch(rect2)

    ax.set_xlabel("Column", fontsize=9, color='#2C3E50')
    ax.set_ylabel("Row", fontsize=9, color='#2C3E50')
    ax.set_title("Click to select cells", fontsize=9, color='#3498DB', fontweight='bold')
    ax.tick_params(axis='both', labelsize=7, colors='#2C3E50')
    ax.set_facecolor('#FAFAFA')

    # Add colorbar with professional styling
    cbar = plt.colorbar(im, ax=ax, shrink=0.85, pad=0.02)
    cbar.ax.tick_params(labelsize=7)
    cbar.set_label('Value', fontsize=8)

    # Add statistics text box
    stats_text = (
        f"Max: {stats_max:.3g}\n"
        f"Min: {stats_min:.3g}\n"
        f"Mean: {stats_mean:.3g}\n"
        f"Std: {stats_std:.3g}\n"
        f"Shape: {data.shape[0]}x{data.shape[1]}"
    )
    ax.text(0.02, 0.98, stats_text,
        transform=ax.transAxes,
        fontsize=6,
        fontweight='normal',
        fontfamily='monospace',
        verticalalignment='top',
        horizontalalignment='left',
        bbox=dict(boxstyle='round,pad=0.3', facecolor='white', edgecolor='#BDC3C7', alpha=0.9, linewidth=1)
    )

    fig.tight_layout()

    # Create canvas
    canvas = FigureCanvasTkAgg(fig, master=parent_frame)
    canvas.get_tk_widget().pack(fill=tk.BOTH, expand=True, padx=5, pady=5)
    canvas.draw()

    grr_file_canvases.append({'canvas': canvas, 'fig': fig, 'ax': ax, 'file_idx': file_idx, 'data': data})

    # Add click event handler
    def on_click(event):
        if event.inaxes == ax and event.button == 1:
            x_click = int(round(event.ydata))  # Row
            y_click = int(round(event.xdata))  # Column

            if 0 <= x_click < data.shape[0] and 0 <= y_click < data.shape[1]:
                toggle_grr_die_selection(x_click, y_click)

    canvas.mpl_connect('button_press_event', on_click)


def toggle_grr_die_selection(x, y):
    """Toggle die selection for Gage R&R comparison"""
    global grr_selected_dies

    max_dies = int(grr_num_dies_var.get())
    die_coord = (x, y)

    if die_coord in grr_selected_dies:
        # Deselect
        grr_selected_dies.remove(die_coord)
        grr_status_var.set(f"Deselected die ({x}, {y})")
    else:
        # Check if max dies reached
        if len(grr_selected_dies) >= max_dies:
            grr_status_var.set(f"Max {max_dies} dies allowed. Deselect a die first.")
            return

        # Select
        grr_selected_dies.append(die_coord)
        grr_status_var.set(f"Selected die ({x}, {y}) - Total: {len(grr_selected_dies)}/{max_dies}")

    # Update label
    if grr_selected_dies:
        dies_text = ", ".join([f"({x},{y})" for x, y in grr_selected_dies])
        grr_selected_dies_label.config(text=f"{len(grr_selected_dies)} dies: {dies_text[:80]}{'...' if len(dies_text) > 80 else ''}")
    else:
        grr_selected_dies_label.config(text="None selected")

    # Update all displays
    update_grr_wafermaps()


def update_grr_wafermaps():
    """Update all wafermap displays to show current selection"""
    global grr_file_canvases

    for canvas_info in grr_file_canvases:
        canvas = canvas_info['canvas']
        fig = canvas_info['fig']
        ax = canvas_info['ax']
        file_idx = canvas_info['file_idx']
        df = canvas_info['data']

        if file_idx < len(grr_file_data):
            # Clear and redraw
            ax.clear()

            if 'x' in df.columns and 'y' in df.columns:
                plot_data = df[df['x'].notna() & df['y'].notna()]

                if len(plot_data) > 0:
                    # Plot all dies as blue
                    ax.scatter(
                        plot_data["x"].values,
                        plot_data["y"].values,
                        c='#3498DB',
                        s=30,
                        edgecolors="gray",
                        linewidth=0.3
                    )

                    # Highlight selected dies in red
                    for die_coord in grr_selected_dies:
                        x, y = die_coord
                        if ((plot_data['x'] == x) & (plot_data['y'] == y)).any():
                            ax.scatter([x], [y], c='red', s=50, edgecolors='black', linewidth=2)

                    ax.set_xlabel("X", fontsize=7)
                    ax.set_ylabel("Y", fontsize=7)
                    ax.set_title(f"Click to select", fontsize=7, color='#3498DB')
                    ax.set_aspect("equal")
                    ax.tick_params(axis='both', labelsize=5)

            fig.tight_layout()
            canvas.draw()


# ============================================================================
# Tab 6: Create Presentation - PowerPoint Export
# ============================================================================

# Presentation configuration variables
pptx_title_var = tk.StringVar(value="Wafermap Analysis Report")
pptx_include_wafermap_var = tk.BooleanVar(value=False)
pptx_include_multi_wafer_var = tk.BooleanVar(value=False)
pptx_include_diffmap_var = tk.BooleanVar(value=False)
pptx_include_stats_var = tk.BooleanVar(value=False)
pptx_selected_wafermaps = []

# Header - compact
pptx_header = tk.Label(
    tab_presentation,
    text="📑 Create PowerPoint Presentation",
    font=("Helvetica", 11, "bold")
)
pptx_header.pack(pady=3)

# Title entry frame - compact
title_main_frame = tk.Frame(tab_presentation)
title_main_frame.pack(fill=tk.X, padx=20, pady=3)

title_label = tk.Label(title_main_frame, text="Title:", font=("Helvetica", 9, "bold"))
title_label.pack(side=tk.LEFT, padx=3)

title_entry = tk.Entry(title_main_frame, textvariable=pptx_title_var, width=40, font=("Helvetica", 9))
title_entry.pack(side=tk.LEFT, padx=3, fill=tk.X, expand=True)

# ============ BUTTONS AT TOP - VISIBLE! ============
pptx_top_btn_frame = tk.Frame(tab_presentation, bg="#E8F5E9", relief="ridge", bd=2)
pptx_top_btn_frame.pack(fill=tk.X, padx=20, pady=5)

create_pptx_btn = tk.Button(
    pptx_top_btn_frame,
    text="🚀 CREATE PRESENTATION",
    command=lambda: create_powerpoint_presentation(),
    font=("Helvetica", 12, "bold"),
    bg="#4CAF50",
    fg="white",
    padx=20,
    pady=8
)
create_pptx_btn.pack(side=tk.LEFT, padx=10, pady=5)

# Quick Google Drive button
def open_google_drive_upload_top():
    import webbrowser
    webbrowser.open("https://drive.google.com/drive/my-drive")

open_drive_btn_top = tk.Button(
    pptx_top_btn_frame,
    text="📁 Open Google Drive",
    command=open_google_drive_upload_top,
    font=("Helvetica", 9),
    bg="#FFA726",
    fg="white",
    padx=10,
    pady=5
)
open_drive_btn_top.pack(side=tk.LEFT, padx=5, pady=5)

# Progress bar at top
pptx_progress_var = tk.DoubleVar(value=0)
pptx_progress = ttk.Progressbar(
    pptx_top_btn_frame,
    variable=pptx_progress_var,
    maximum=100,
    length=150
)
pptx_progress.pack(side=tk.LEFT, padx=10, pady=5)

pptx_status_label = tk.Label(
    pptx_top_btn_frame,
    text="Ready",
    font=("Helvetica", 8),
    fg="gray",
    bg="#E8F5E9"
)
pptx_status_label.pack(side=tk.LEFT, padx=5)

# Create sub-notebook for content selection - smaller
pptx_sub_notebook = ttk.Notebook(tab_presentation)
pptx_sub_notebook.pack(fill=tk.BOTH, expand=True, padx=10, pady=3)

# Style for smaller sub-tabs
style.configure("SubTab.TNotebook.Tab", font=("Helvetica", 9))

# ============ Sub-Tab 1: Wafermap ============
pptx_wafermap_tab = ttk.Frame(pptx_sub_notebook)
pptx_sub_notebook.add(pptx_wafermap_tab, text="Wafermap")

# Wafermap sub-tab content
wafermap_content_frame = tk.Frame(pptx_wafermap_tab)
wafermap_content_frame.pack(fill=tk.BOTH, expand=True, padx=20, pady=10)

# Select All checkbox for Wafermap
pptx_wafermap_select_all_var = tk.BooleanVar(value=False)

def toggle_wafermap_all():
    """Toggle all wafermap checkboxes"""
    state = pptx_wafermap_select_all_var.get()
    pptx_wafermap_binmap_var.set(state)
    pptx_wafermap_heatmap_var.set(state)
    pptx_wafermap_stats_var.set(state)
    pptx_wafermap_boxplot_var.set(state)
    pptx_wafermap_histogram_var.set(state)

wafermap_select_all_check = tk.Checkbutton(
    wafermap_content_frame,
    text="Include Wafermap Tab in Report",
    variable=pptx_wafermap_select_all_var,
    command=toggle_wafermap_all,
    font=("Helvetica", 10, "bold")
)
wafermap_select_all_check.pack(anchor="w", pady=10)

ttk.Separator(wafermap_content_frame, orient="horizontal").pack(fill=tk.X, pady=5)

# Individual Wafermap options
pptx_wafermap_binmap_var = tk.BooleanVar(value=False)
pptx_wafermap_heatmap_var = tk.BooleanVar(value=False)
pptx_wafermap_stats_var = tk.BooleanVar(value=False)
pptx_wafermap_boxplot_var = tk.BooleanVar(value=False)
pptx_wafermap_histogram_var = tk.BooleanVar(value=False)

wafermap_options_label = tk.Label(wafermap_content_frame, text="Select content to include:", font=("Helvetica", 9))
wafermap_options_label.pack(anchor="w", pady=5)

tk.Checkbutton(wafermap_content_frame, text="Bin Map (Die binning visualization)",
               variable=pptx_wafermap_binmap_var, font=("Helvetica", 9)).pack(anchor="w", padx=20, pady=2)
tk.Checkbutton(wafermap_content_frame, text="Parameter Heatmap (Selected parameter)",
               variable=pptx_wafermap_heatmap_var, font=("Helvetica", 9)).pack(anchor="w", padx=20, pady=2)
tk.Checkbutton(wafermap_content_frame, text="Statistics Summary (Min, Max, Mean, Median)",
               variable=pptx_wafermap_stats_var, font=("Helvetica", 9)).pack(anchor="w", padx=20, pady=2)
tk.Checkbutton(wafermap_content_frame, text="Boxplot Analysis",
               variable=pptx_wafermap_boxplot_var, font=("Helvetica", 9)).pack(anchor="w", padx=20, pady=2)
tk.Checkbutton(wafermap_content_frame, text="Histogram / Distribution",
               variable=pptx_wafermap_histogram_var, font=("Helvetica", 9)).pack(anchor="w", padx=20, pady=2)

# Parameter group selection for wafermap
ttk.Separator(wafermap_content_frame, orient="horizontal").pack(fill=tk.X, pady=10)

wafermap_group_frame = tk.Frame(wafermap_content_frame)
wafermap_group_frame.pack(fill=tk.X, pady=5)

wafermap_group_label = tk.Label(wafermap_group_frame, text="Parameter Group for Boxplots:", font=("Helvetica", 9))
wafermap_group_label.pack(side=tk.LEFT, padx=5)

pptx_wafermap_group_var = tk.StringVar(value="All Groups")
pptx_wafermap_group_combobox = ttk.Combobox(
    wafermap_group_frame,
    textvariable=pptx_wafermap_group_var,
    state="readonly",
    width=25,
    font=("Helvetica", 9)
)
pptx_wafermap_group_combobox["values"] = ["All Groups"]
pptx_wafermap_group_combobox.pack(side=tk.LEFT, padx=5)

# ============ Sub-Tab 2: Multiple Wafermaps ============
pptx_multi_wafer_tab = ttk.Frame(pptx_sub_notebook)
pptx_sub_notebook.add(pptx_multi_wafer_tab, text="Multiple Wafermaps")

# Multiple Wafermaps sub-tab content
multi_wafer_content_frame = tk.Frame(pptx_multi_wafer_tab)
multi_wafer_content_frame.pack(fill=tk.BOTH, expand=True, padx=20, pady=10)

# Select All checkbox for Multiple Wafermaps
pptx_multi_select_all_var = tk.BooleanVar(value=False)

def toggle_multi_all():
    """Toggle all multiple wafermap checkboxes"""
    state = pptx_multi_select_all_var.get()
    pptx_multi_overview_var.set(state)
    pptx_multi_comparison_var.set(state)
    pptx_multi_yield_var.set(state)
    pptx_multi_individual_var.set(state)

multi_select_all_check = tk.Checkbutton(
    multi_wafer_content_frame,
    text="Include Multiple Wafermaps Tab in Report",
    variable=pptx_multi_select_all_var,
    command=toggle_multi_all,
    font=("Helvetica", 10, "bold")
)
multi_select_all_check.pack(anchor="w", pady=10)

ttk.Separator(multi_wafer_content_frame, orient="horizontal").pack(fill=tk.X, pady=5)

# Individual Multiple Wafermap options
pptx_multi_overview_var = tk.BooleanVar(value=False)
pptx_multi_comparison_var = tk.BooleanVar(value=False)
pptx_multi_yield_var = tk.BooleanVar(value=False)
pptx_multi_individual_var = tk.BooleanVar(value=False)

multi_options_label = tk.Label(multi_wafer_content_frame, text="Select content to include:", font=("Helvetica", 9))
multi_options_label.pack(anchor="w", pady=5)

tk.Checkbutton(multi_wafer_content_frame, text="Overview Slide (All wafermaps grid)",
               variable=pptx_multi_overview_var, font=("Helvetica", 9)).pack(anchor="w", padx=20, pady=2)
tk.Checkbutton(multi_wafer_content_frame, text="Wafer-to-Wafer Comparison (Boxplots)",
               variable=pptx_multi_comparison_var, font=("Helvetica", 9)).pack(anchor="w", padx=20, pady=2)
tk.Checkbutton(multi_wafer_content_frame, text="Yield Summary Table",
               variable=pptx_multi_yield_var, font=("Helvetica", 9)).pack(anchor="w", padx=20, pady=2)
tk.Checkbutton(multi_wafer_content_frame, text="Individual Wafermap Slides",
               variable=pptx_multi_individual_var, font=("Helvetica", 9)).pack(anchor="w", padx=20, pady=2)

# Wafer selection listbox
ttk.Separator(multi_wafer_content_frame, orient="horizontal").pack(fill=tk.X, pady=10)

wafer_selection_label = tk.Label(multi_wafer_content_frame, text="Select individual wafermaps to include:", font=("Helvetica", 9))
wafer_selection_label.pack(anchor="w", pady=5)

# Listbox for wafermap selection with scrollbar
wafer_listbox_frame = tk.Frame(multi_wafer_content_frame)
wafer_listbox_frame.pack(fill=tk.BOTH, expand=True, padx=10, pady=5)

wafer_scrollbar = tk.Scrollbar(wafer_listbox_frame)
wafer_scrollbar.pack(side=tk.RIGHT, fill=tk.Y)

wafer_listbox = tk.Listbox(
    wafer_listbox_frame,
    selectmode=tk.MULTIPLE,
    font=("Helvetica", 8),
    height=8,
    yscrollcommand=wafer_scrollbar.set
)
wafer_listbox.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)
wafer_scrollbar.config(command=wafer_listbox.yview)

# Buttons for selection
selection_btn_frame = tk.Frame(multi_wafer_content_frame)
selection_btn_frame.pack(fill=tk.X, padx=10, pady=5)

select_all_btn = tk.Button(
    selection_btn_frame,
    text="Select All",
    command=lambda: wafer_listbox.select_set(0, tk.END),
    font=("Helvetica", 8)
)
select_all_btn.pack(side=tk.LEFT, padx=5)

deselect_all_btn = tk.Button(
    selection_btn_frame,
    text="Deselect All",
    command=lambda: wafer_listbox.selection_clear(0, tk.END),
    font=("Helvetica", 8)
)
deselect_all_btn.pack(side=tk.LEFT, padx=5)

refresh_list_btn = tk.Button(
    selection_btn_frame,
    text="Refresh List",
    command=lambda: update_wafer_listbox(),
    font=("Helvetica", 8)
)
refresh_list_btn.pack(side=tk.LEFT, padx=5)

# ============ Sub-Tab 3: Diffmap ============
pptx_diffmap_tab = ttk.Frame(pptx_sub_notebook)
pptx_sub_notebook.add(pptx_diffmap_tab, text="Diffmap")

# Diffmap sub-tab content
diffmap_content_frame = tk.Frame(pptx_diffmap_tab)
diffmap_content_frame.pack(fill=tk.BOTH, expand=True, padx=20, pady=10)

# Select All checkbox for Diffmap
pptx_diffmap_select_all_var = tk.BooleanVar(value=False)

def toggle_diffmap_all():
    """Toggle all diffmap checkboxes"""
    state = pptx_diffmap_select_all_var.get()
    pptx_diffmap_heatmap_var.set(state)
    pptx_diffmap_histogram_var.set(state)
    pptx_diffmap_stats_var.set(state)
    pptx_diffmap_reference_var.set(state)

diffmap_select_all_check = tk.Checkbutton(
    diffmap_content_frame,
    text="Include Diffmap Tab in Report",
    variable=pptx_diffmap_select_all_var,
    command=toggle_diffmap_all,
    font=("Helvetica", 10, "bold")
)
diffmap_select_all_check.pack(anchor="w", pady=10)

ttk.Separator(diffmap_content_frame, orient="horizontal").pack(fill=tk.X, pady=5)

# Individual Diffmap options
pptx_diffmap_heatmap_var = tk.BooleanVar(value=False)
pptx_diffmap_histogram_var = tk.BooleanVar(value=False)
pptx_diffmap_stats_var = tk.BooleanVar(value=False)
pptx_diffmap_reference_var = tk.BooleanVar(value=False)

diffmap_options_label = tk.Label(diffmap_content_frame, text="Select content to include:", font=("Helvetica", 9))
diffmap_options_label.pack(anchor="w", pady=5)

tk.Checkbutton(diffmap_content_frame, text="Difference Heatmap",
               variable=pptx_diffmap_heatmap_var, font=("Helvetica", 9)).pack(anchor="w", padx=20, pady=2)
tk.Checkbutton(diffmap_content_frame, text="Difference Distribution Histogram",
               variable=pptx_diffmap_histogram_var, font=("Helvetica", 9)).pack(anchor="w", padx=20, pady=2)
tk.Checkbutton(diffmap_content_frame, text="Statistics Summary (Δ Mean, Δ Std, etc.)",
               variable=pptx_diffmap_stats_var, font=("Helvetica", 9)).pack(anchor="w", padx=20, pady=2)
tk.Checkbutton(diffmap_content_frame, text="Reference and Comparison Wafermaps",
               variable=pptx_diffmap_reference_var, font=("Helvetica", 9)).pack(anchor="w", padx=20, pady=2)

# ============ Sub-Tab 4: Gage R&R ============
pptx_grr_tab = ttk.Frame(pptx_sub_notebook)
pptx_sub_notebook.add(pptx_grr_tab, text="Gage R&R")

# Gage R&R sub-tab content
grr_content_frame = tk.Frame(pptx_grr_tab)
grr_content_frame.pack(fill=tk.BOTH, expand=True, padx=20, pady=10)

# Select All checkbox for Gage R&R
pptx_grr_select_all_var = tk.BooleanVar(value=False)

def toggle_grr_all():
    """Toggle all Gage R&R checkboxes"""
    state = pptx_grr_select_all_var.get()
    pptx_grr_summary_var.set(state)
    pptx_grr_variance_var.set(state)
    pptx_grr_charts_var.set(state)
    pptx_grr_anova_var.set(state)

grr_select_all_check = tk.Checkbutton(
    grr_content_frame,
    text="Include Gage R&R Tab in Report",
    variable=pptx_grr_select_all_var,
    command=toggle_grr_all,
    font=("Helvetica", 10, "bold")
)
grr_select_all_check.pack(anchor="w", pady=10)

ttk.Separator(grr_content_frame, orient="horizontal").pack(fill=tk.X, pady=5)

# Individual Gage R&R options
pptx_grr_summary_var = tk.BooleanVar(value=False)
pptx_grr_variance_var = tk.BooleanVar(value=False)
pptx_grr_charts_var = tk.BooleanVar(value=False)
pptx_grr_anova_var = tk.BooleanVar(value=False)

grr_options_label = tk.Label(grr_content_frame, text="Select content to include:", font=("Helvetica", 9))
grr_options_label.pack(anchor="w", pady=5)

tk.Checkbutton(grr_content_frame, text="Gage R&R Summary Table",
               variable=pptx_grr_summary_var, font=("Helvetica", 9)).pack(anchor="w", padx=20, pady=2)
tk.Checkbutton(grr_content_frame, text="Variance Components Chart",
               variable=pptx_grr_variance_var, font=("Helvetica", 9)).pack(anchor="w", padx=20, pady=2)
tk.Checkbutton(grr_content_frame, text="Measurement by Part/Operator Charts",
               variable=pptx_grr_charts_var, font=("Helvetica", 9)).pack(anchor="w", padx=20, pady=2)
tk.Checkbutton(grr_content_frame, text="ANOVA Table",
               variable=pptx_grr_anova_var, font=("Helvetica", 9)).pack(anchor="w", padx=20, pady=2)

# ============ Save Options Frame ============
pptx_save_options_frame = tk.LabelFrame(tab_presentation, text="Save Options", font=("Helvetica", 10, "bold"))
pptx_save_options_frame.pack(fill=tk.X, padx=40, pady=10)

# Save method selection
pptx_save_method_var = tk.StringVar(value="local")

save_method_frame = tk.Frame(pptx_save_options_frame)
save_method_frame.pack(fill=tk.X, padx=10, pady=5)

tk.Label(save_method_frame, text="Save Method:", font=("Helvetica", 9, "bold")).pack(side=tk.LEFT, padx=5)

tk.Radiobutton(save_method_frame, text="💾 Save Locally (.pptx)", variable=pptx_save_method_var,
               value="local", font=("Helvetica", 9)).pack(side=tk.LEFT, padx=10)
tk.Radiobutton(save_method_frame, text="☁️ Upload to Google Drive", variable=pptx_save_method_var,
               value="google_drive", font=("Helvetica", 9)).pack(side=tk.LEFT, padx=10)
tk.Radiobutton(save_method_frame, text="📊 Google Slides (Convert)", variable=pptx_save_method_var,
               value="google_slides", font=("Helvetica", 9)).pack(side=tk.LEFT, padx=10)

# Google options frame (shown when Google options selected)
google_options_frame = tk.Frame(pptx_save_options_frame)
google_options_frame.pack(fill=tk.X, padx=10, pady=5)

# Gemini beautification option
pptx_use_gemini_var = tk.BooleanVar(value=False)
gemini_check = tk.Checkbutton(google_options_frame, text="✨ Beautify with Gemini AI",
                              variable=pptx_use_gemini_var, font=("Helvetica", 9))
gemini_check.pack(side=tk.LEFT, padx=20)

# Theme selection for Gemini
tk.Label(google_options_frame, text="Theme:", font=("Helvetica", 9)).pack(side=tk.LEFT, padx=5)
pptx_theme_var = tk.StringVar(value="professional")
theme_combo = ttk.Combobox(google_options_frame, textvariable=pptx_theme_var,
                           values=["professional", "modern", "vibrant"], state="readonly", width=12)
theme_combo.pack(side=tk.LEFT, padx=5)

# Gemini API key entry (optional)
tk.Label(google_options_frame, text="Gemini API Key:", font=("Helvetica", 9)).pack(side=tk.LEFT, padx=10)
pptx_gemini_key_var = tk.StringVar(value="")
gemini_key_entry = tk.Entry(google_options_frame, textvariable=pptx_gemini_key_var, width=30, show="*", font=("Helvetica", 8))
gemini_key_entry.pack(side=tk.LEFT, padx=5)

# Status info for Google API
google_status_frame = tk.Frame(pptx_save_options_frame)
google_status_frame.pack(fill=tk.X, padx=10, pady=5)

google_api_status = "✅ Available" if GOOGLE_API_AVAILABLE else "❌ Not installed (pip install google-api-python-client google-auth-oauthlib)"
gemini_api_status = "✅ Available" if GEMINI_AVAILABLE else "❌ Not installed (pip install google-generativeai)"

tk.Label(google_status_frame, text=f"Google API: {google_api_status}", font=("Helvetica", 8),
         fg="green" if GOOGLE_API_AVAILABLE else "red").pack(side=tk.LEFT, padx=10)
tk.Label(google_status_frame, text=f"Gemini AI: {gemini_api_status}", font=("Helvetica", 8),
         fg="green" if GEMINI_AVAILABLE else "red").pack(side=tk.LEFT, padx=10)

# Setup credentials button
def setup_google_credentials():
    """Open dialog to setup Google credentials"""
    import webbrowser
    msg = """To use Google Drive/Slides integration:

1. Go to Google Cloud Console (https://console.cloud.google.com)
2. Create a new project or select existing
3. Enable 'Google Drive API' and 'Google Slides API'
4. Create OAuth 2.0 credentials (Desktop App)
5. Download the credentials JSON file
6. Save it as 'credentials.json' in the application folder:
   """ + os.path.dirname(__file__) + """

Click 'Open Console' to go to Google Cloud Console."""

    result = tk.messagebox.askquestion("Google API Setup", msg, icon='info')
    if result == 'yes':
        webbrowser.open("https://console.cloud.google.com/apis/credentials")

setup_btn = tk.Button(google_status_frame, text="🔧 Setup Google Credentials", command=setup_google_credentials,
                      font=("Helvetica", 8), bg="#2196F3", fg="white")
setup_btn.pack(side=tk.RIGHT, padx=10)

# ============ Status and Create Button Frame ============
# (Buttons moved to top of tab for visibility)

# Legacy variable references for backward compatibility
# pptx_status_frame already created at top

# Helper function to update group comboboxes in presentation tab
def update_pptx_group_combobox():
    """Update the PowerPoint group combobox with available groups"""
    if grouped_parameters:
        group_names = ["All Groups"] + sorted(grouped_parameters.keys())
        pptx_wafermap_group_combobox["values"] = group_names
        pptx_wafermap_group_combobox.current(0)

# Legacy variable mappings for backward compatibility with create_powerpoint_presentation
pptx_include_boxplots_var = tk.BooleanVar(value=False)
pptx_boxplot_group_var = pptx_wafermap_group_var


def update_wafer_listbox():
    """Update the wafermap listbox with available wafermaps"""
    wafer_listbox.delete(0, tk.END)

    # Add wafermaps from multiple_stdf_data
    if multiple_stdf_data and multiple_wafer_ids:
        for idx, wafer_id in enumerate(multiple_wafer_ids):
            display_name = f"Wafer {idx+1}: {wafer_id}"
            wafer_listbox.insert(tk.END, display_name)

    # Add single wafermap if available
    if current_stdf_data is not None and not current_stdf_data.empty:
        wafer_listbox.insert(tk.END, f"Single Wafermap: {current_wafer_id}")

    if wafer_listbox.size() == 0:
        wafer_listbox.insert(tk.END, "(No wafermaps loaded)")

    print(f"Wafer listbox updated with {wafer_listbox.size()} items")


def figure_to_image_bytes(fig, dpi=150):
    """Convert a matplotlib figure to PNG bytes"""
    buf = BytesIO()
    fig.savefig(buf, format='png', dpi=dpi, bbox_inches='tight', facecolor='white', edgecolor='none')
    buf.seek(0)
    return buf


def create_powerpoint_presentation():
    """Create a professional PowerPoint presentation from the selected data"""
    if not PPTX_AVAILABLE:
        pptx_status_label.config(text="Error: python-pptx not installed. Run: pip install python-pptx", fg="red")
        return

    # Check if at least one tab is selected for inclusion
    include_wafermap = pptx_wafermap_select_all_var.get()
    include_multi_wafer = pptx_multi_select_all_var.get()
    include_diffmap = pptx_diffmap_select_all_var.get()
    include_grr = pptx_grr_select_all_var.get()

    if not any([include_wafermap, include_multi_wafer, include_diffmap, include_grr]):
        pptx_status_label.config(text="Error: Please select at least one tab to include in the report", fg="red")
        return

    # Get save method
    save_method = pptx_save_method_var.get()
    use_gemini = pptx_use_gemini_var.get()
    theme_style = pptx_theme_var.get()

    # Validate Google API availability for Google options
    if save_method in ["google_drive", "google_slides"]:
        if not GOOGLE_API_AVAILABLE:
            pptx_status_label.config(text="Error: Google API not available. Install with: pip install google-api-python-client google-auth-oauthlib", fg="red")
            return

    # Get save path for local save, or use temp file for Google upload
    save_path = None
    if save_method == "local":
        save_path = filedialog.asksaveasfilename(
            title="Save PowerPoint Presentation",
            defaultextension=".pptx",
            filetypes=[("PowerPoint files", "*.pptx"), ("All files", "*.*")],
            initialfile=f"{pptx_title_var.get().replace(' ', '_')}.pptx"
        )
        if not save_path:
            return
    else:
        # Use temp file for Google upload
        import tempfile
        temp_dir = tempfile.gettempdir()
        save_path = os.path.join(temp_dir, f"{pptx_title_var.get().replace(' ', '_')}.pptx")

    pptx_status_label.config(text="Generating PowerPoint...", fg="blue")
    pptx_progress_var.set(0)
    main_win.update()

    try:
        # Create presentation
        prs = Presentation()
        prs.slide_width = Inches(13.333)  # 16:9 aspect ratio
        prs.slide_height = Inches(7.5)

        # Slide layouts
        title_slide_layout = prs.slide_layouts[6]  # Blank layout
        content_slide_layout = prs.slide_layouts[6]  # Blank layout

        # Calculate total steps based on selected tabs
        total_steps = 1  # Title slide
        if include_wafermap:
            total_steps += 1
        if include_multi_wafer:
            total_steps += 1
        if include_diffmap:
            total_steps += 1
        if include_grr:
            total_steps += 1
        current_step = 0

        # ============ Title Slide ============
        slide = prs.slides.add_slide(title_slide_layout)

        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = pptx_title_var.get()
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        title_para.alignment = PP_ALIGN.CENTER

        # Add subtitle with date
        from datetime import datetime
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(0.8))
        subtitle_frame = subtitle_box.text_frame
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.text = f"Generated on {datetime.now().strftime('%Y-%m-%d %H:%M')}"
        subtitle_para.font.size = Pt(20)
        subtitle_para.font.color.rgb = RGBColor(128, 128, 128)
        subtitle_para.alignment = PP_ALIGN.CENTER

        current_step += 1
        pptx_progress_var.set((current_step / total_steps) * 100)
        main_win.update()

        # ============ Wafermap Slide with Bin Distribution ============
        # Determine which data to use - prefer multiple_stdf_data, then current_stdf_data
        wafer_data_for_slide = None
        wafer_id_for_slide = None

        if include_wafermap:
            if multiple_stdf_data and len(multiple_stdf_data) > 0:
                # Use first wafermap from multiple data
                wafer_data_for_slide = multiple_stdf_data[0]
                wafer_id_for_slide = multiple_wafer_ids[0] if multiple_wafer_ids else "Wafer 1"
            elif current_stdf_data is not None and not current_stdf_data.empty:
                wafer_data_for_slide = current_stdf_data
                wafer_id_for_slide = current_wafer_id

        if wafer_data_for_slide is not None and not wafer_data_for_slide.empty:
            slide = prs.slides.add_slide(content_slide_layout)

            # Slide title
            title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(12.733), Inches(0.6))
            title_frame = title_box.text_frame
            title_para = title_frame.paragraphs[0]
            title_para.text = f"Wafermap Analysis: {wafer_id_for_slide}"
            title_para.font.size = Pt(28)
            title_para.font.bold = True

            # LEFT SIDE: Wafermap in binning mode
            fig_wafer, ax_wafer = plt.subplots(figsize=(7, 7))

            # Always use bin for wafermap display
            param_column = "bin"

            if param_column in wafer_data_for_slide.columns:
                plot_data = wafer_data_for_slide.dropna(subset=[param_column])
                if len(plot_data) > 0:
                    # Get unique bins and create color mapping
                    unique_bins = np.sort(plot_data[param_column].unique())

                    # Create custom colormap: green for bin 0 (pass), red shades for fails
                    from matplotlib.colors import ListedColormap
                    bin_colors = []
                    for b in unique_bins:
                        if b == 0:
                            bin_colors.append('#4CAF50')  # Green for pass (bin 0)
                        else:
                            # Different red shades for different fail bins
                            red_shade = max(0.3, 1.0 - (b * 0.1))
                            bin_colors.append(plt.cm.Reds(red_shade))

                    # Create scatter plot with discrete colors
                    sc = ax_wafer.scatter(
                        plot_data["x"],
                        plot_data["y"],
                        c=plot_data[param_column],
                        cmap="tab20",
                        s=100,
                        edgecolors="black",
                        linewidth=0.5
                    )
                    ax_wafer.set_xlabel("X Coordinate", fontsize=12)
                    ax_wafer.set_ylabel("Y Coordinate", fontsize=12)
                    ax_wafer.set_title("Wafermap - Binning", fontsize=14, fontweight="bold")
                    ax_wafer.set_aspect("equal")

                    # Add grid lines for each die
                    x_min, x_max = plot_data["x"].min(), plot_data["x"].max()
                    y_min, y_max = plot_data["y"].min(), plot_data["y"].max()

                    # Set minor ticks at die boundaries
                    ax_wafer.set_xticks(np.arange(x_min - 0.5, x_max + 1.5, 1), minor=True)
                    ax_wafer.set_yticks(np.arange(y_min - 0.5, y_max + 1.5, 1), minor=True)
                    ax_wafer.grid(which='minor', color='black', linewidth=0.3, alpha=0.5)
                    ax_wafer.grid(which='major', color='gray', linewidth=0.5, alpha=0.3)
                    ax_wafer.tick_params(which='minor', size=0)

                    # Add colorbar with bin labels
                    cbar = plt.colorbar(sc, ax=ax_wafer)
                    cbar.set_label("Bin", fontsize=11)

            fig_wafer.tight_layout()

            # Add wafermap to left side of slide - use high DPI for better quality
            img_stream_wafer = figure_to_image_bytes(fig_wafer, dpi=300)
            slide.shapes.add_picture(img_stream_wafer, Inches(0.3), Inches(0.9), width=Inches(6.0), height=Inches(5.8))
            plt.close(fig_wafer)

            # RIGHT SIDE: Bin distribution bar chart
            if "bin" in wafer_data_for_slide.columns:
                bins_data = wafer_data_for_slide["bin"].dropna().values

                if len(bins_data) > 0:
                    unique_bins, bin_counts = np.unique(bins_data, return_counts=True)
                    total_dies = len(bins_data)
                    bin_percentages = (bin_counts / total_dies) * 100

                    # Pass bin is 0
                    pass_bin = 0

                    fig_bins, ax_bins = plt.subplots(figsize=(6, 7))

                    # Create color list
                    colors = []
                    for b in unique_bins:
                        if b == pass_bin:
                            colors.append('#4CAF50')  # Green for pass
                        else:
                            colors.append('#F44336')  # Red for fail

                    x_pos = np.arange(len(unique_bins))
                    bars = ax_bins.bar(x_pos, bin_percentages, color=colors, edgecolor='black', linewidth=1)

                    # Add percentage labels on bars
                    for bar, pct, count in zip(bars, bin_percentages, bin_counts):
                        height = bar.get_height()
                        ax_bins.annotate(f'{pct:.1f}%\n({int(count)})',
                            xy=(bar.get_x() + bar.get_width() / 2, height),
                            xytext=(0, 5),
                            textcoords="offset points",
                            ha='center', va='bottom',
                            fontsize=10, fontweight='bold')

                    ax_bins.set_xticks(x_pos)
                    ax_bins.set_xticklabels([f'Bin {int(b)}' for b in unique_bins], fontsize=11, rotation=45, ha='right')
                    ax_bins.set_ylabel("Percentage (%)", fontsize=12)
                    ax_bins.set_title("Bin Distribution", fontsize=14, fontweight="bold")
                    ax_bins.tick_params(axis='both', which='major', labelsize=10)
                    ax_bins.grid(True, alpha=0.3, linestyle='--', linewidth=0.5, axis='y')
                    ax_bins.set_ylim(0, max(bin_percentages) * 1.25)

                    # Add yield summary box
                    pass_count = bin_counts[unique_bins == pass_bin].sum() if pass_bin in unique_bins else 0
                    fail_count = total_dies - pass_count
                    yield_pct = (pass_count / total_dies) * 100

                    summary_text = f"Summary:\n"
                    summary_text += f"Total Dies: {total_dies}\n"
                    summary_text += f"Pass (Bin 0): {pass_count} ({yield_pct:.1f}%)\n"
                    summary_text += f"Fail: {fail_count} ({100-yield_pct:.1f}%)\n"
                    summary_text += f"Yield: {yield_pct:.1f}%"

                    ax_bins.text(0.98, 0.98, summary_text,
                        transform=ax_bins.transAxes, fontsize=11, fontweight='bold',
                        va='top', ha='right',
                        bbox=dict(boxstyle='round', facecolor='lightyellow', alpha=0.9, edgecolor='black'))

                    fig_bins.tight_layout()

                    # Add bin chart to right side of slide - use high DPI for better quality
                    img_stream_bins = figure_to_image_bytes(fig_bins, dpi=300)
                    slide.shapes.add_picture(img_stream_bins, Inches(6.5), Inches(0.9), width=Inches(6.0), height=Inches(5.8))
                    plt.close(fig_bins)

        current_step += 1
        pptx_progress_var.set((current_step / total_steps) * 100)
        main_win.update()

        # ============ Multiple Wafermaps Slide ============
        if include_multi_wafer and multiple_stdf_data and len(multiple_stdf_data) > 0:
            # Get selected wafermaps from listbox
            selected_indices = list(wafer_listbox.curselection())

            # Filter to only include wafermaps (not the single wafermap option)
            wafermap_indices = [i for i in selected_indices if i < len(multiple_wafer_ids)]

            if not wafermap_indices:
                wafermap_indices = list(range(len(multiple_wafer_ids)))

            num_wafermaps = len(wafermap_indices)

            if num_wafermaps > 0:
                slide = prs.slides.add_slide(content_slide_layout)

                # Slide title
                title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(12.733), Inches(0.6))
                title_frame = title_box.text_frame
                title_para = title_frame.paragraphs[0]
                title_para.text = f"Multiple Wafermaps Overview ({num_wafermaps} wafers)"
                title_para.font.size = Pt(28)
                title_para.font.bold = True

                # Create multi-wafermap figure with better sizing
                cols = min(4, num_wafermaps)
                rows = (num_wafermaps + cols - 1) // cols

                # Larger figure size for better quality
                fig_width = min(14, 3.5 * cols)
                fig_height = min(7, 3.5 * rows)

                fig, axes = plt.subplots(rows, cols, figsize=(fig_width, fig_height))

                if num_wafermaps == 1:
                    axes = [axes]
                elif rows == 1 or cols == 1:
                    axes = axes.flatten() if hasattr(axes, 'flatten') else [axes]
                else:
                    axes = axes.flatten()

                # Get parameter
                selected_param = heatmap_param_combobox.get() if heatmap_param_combobox.get() else "bin"
                if selected_param.startswith("BIN"):
                    param_column = "bin"
                else:
                    test_key = selected_param.split(":")[0].strip()
                    if test_key.startswith("test_"):
                        param_column = int(test_key.replace("test_", ""))
                    else:
                        param_column = int(test_key)

                cmap = "tab20" if param_column == "bin" else "viridis"

                for plot_idx, wafer_idx in enumerate(wafermap_indices):
                    ax = axes[plot_idx]
                    df = multiple_stdf_data[wafer_idx]
                    wafer_id = multiple_wafer_ids[wafer_idx]

                    if param_column in df.columns:
                        plot_data = df.dropna(subset=[param_column])
                        if len(plot_data) > 0:
                            sc = ax.scatter(
                                plot_data["x"],
                                plot_data["y"],
                                c=plot_data[param_column],
                                cmap=cmap,
                                s=20,
                                edgecolors="black",
                                linewidth=0.2
                            )
                            short_id = str(wafer_id)[:20] + "..." if len(str(wafer_id)) > 20 else str(wafer_id)
                            ax.set_title(short_id, fontsize=8, fontweight="bold")
                            ax.set_aspect("equal")
                            ax.tick_params(axis='both', labelsize=6)
                    else:
                        ax.set_title(f"{wafer_id}\nNo data", fontsize=8)
                        ax.axis('off')

                # Hide unused axes
                for idx in range(num_wafermaps, len(axes)):
                    axes[idx].axis('off')

                fig.tight_layout()

                # Add to slide - higher DPI for better quality and proper sizing
                img_stream = figure_to_image_bytes(fig, dpi=250)
                # Center the image on the slide with proper dimensions
                slide.shapes.add_picture(img_stream, Inches(0.5), Inches(0.9), width=Inches(12.3), height=Inches(6.3))
                plt.close(fig)

        current_step += 1
        pptx_progress_var.set((current_step / total_steps) * 100)
        main_win.update()

        # ============ Diffmap Slide ============
        if include_diffmap and diffmap_result_data is not None:
            slide = prs.slides.add_slide(content_slide_layout)

            # Slide title
            title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(12.733), Inches(0.6))
            title_frame = title_box.text_frame
            title_para = title_frame.paragraphs[0]
            title_para.text = f"Diffmap: {diffmap_reference_id} vs {diffmap_compare_id}"
            title_para.font.size = Pt(28)
            title_para.font.bold = True

            # Get selected parameter
            selected_param = diffmap_param_combobox.get() if diffmap_param_combobox.get() else "bin"
            if selected_param.startswith("BIN"):
                param_column = "bin"
                param_label = "Bin"
            else:
                test_key = selected_param.split(":")[0].strip()
                if test_key.startswith("test_"):
                    param_column = int(test_key.replace("test_", ""))
                else:
                    param_column = int(test_key)
                param_label = selected_param

            diff_column = f"diff_{param_column}"

            if diff_column in diffmap_result_data.columns:
                # Create diffmap figure
                fig, ax = plt.subplots(figsize=(8, 8))

                plot_data = diffmap_result_data.dropna(subset=[diff_column])

                if len(plot_data) > 0:
                    diff_values = plot_data[diff_column].values
                    max_abs = max(abs(diff_values.min()), abs(diff_values.max()))

                    sc = ax.scatter(
                        plot_data["x"],
                        plot_data["y"],
                        c=diff_values,
                        cmap="RdBu_r",
                        s=80,
                        edgecolors="black",
                        linewidth=0.5,
                        vmin=-max_abs,
                        vmax=max_abs
                    )
                    ax.set_xlabel("X Coordinate", fontsize=12)
                    ax.set_ylabel("Y Coordinate", fontsize=12)
                    ax.set_title(f"Difference: {param_label}", fontsize=14)
                    ax.set_aspect("equal")
                    ax.grid(True, alpha=0.3)
                    plt.colorbar(sc, ax=ax, label=f"Δ {param_label}")

                fig.tight_layout()

                # Add to slide
                img_stream = figure_to_image_bytes(fig)
                slide.shapes.add_picture(img_stream, Inches(2.5), Inches(0.9), width=Inches(8))
                plt.close(fig)

                # Add statistics
                if pptx_include_stats_var.get():
                    diff_vals = diffmap_result_data[diff_column].dropna()
                    stats_text = f"Diff Statistics:\n"
                    stats_text += f"  Count: {len(diff_vals)}\n"
                    stats_text += f"  Mean: {diff_vals.mean():.4f}\n"
                    stats_text += f"  Std: {diff_vals.std():.4f}\n"
                    stats_text += f"  Min: {diff_vals.min():.4f}\n"
                    stats_text += f"  Max: {diff_vals.max():.4f}"

                    stats_box = slide.shapes.add_textbox(Inches(10.5), Inches(1.5), Inches(2.5), Inches(2))
                    stats_frame = stats_box.text_frame
                    stats_para = stats_frame.paragraphs[0]
                    stats_para.text = stats_text
                    stats_para.font.size = Pt(11)
                    stats_para.font.name = "Consolas"

        current_step += 1
        pptx_progress_var.set((current_step / total_steps) * 100)
        main_win.update()

        # ============ Boxplot Slides for Selected Group (part of Wafermap tab) ============
        if include_wafermap and pptx_wafermap_boxplot_var.get():
            # Determine which data to use
            if multiple_stdf_data and len(multiple_stdf_data) > 0:
                data_sources = multiple_stdf_data
                wafer_labels = multiple_wafer_ids
            elif current_stdf_data is not None and not current_stdf_data.empty:
                data_sources = [current_stdf_data]
                wafer_labels = [current_wafer_id if current_wafer_id else "Wafer"]
            else:
                data_sources = []
                wafer_labels = []

            if data_sources and grouped_parameters:
                # Get selected group
                selected_group = pptx_boxplot_group_var.get()

                # Get parameters for the selected group
                if selected_group == "All Groups":
                    # Get all parameters from all groups
                    params_to_plot = []
                    for group_name, params in grouped_parameters.items():
                        for test_num, short_name, full_name in params:
                            params_to_plot.append((test_num, short_name, full_name, group_name))
                else:
                    # Get parameters from selected group only
                    params_to_plot = []
                    if selected_group in grouped_parameters:
                        for test_num, short_name, full_name in grouped_parameters[selected_group]:
                            params_to_plot.append((test_num, short_name, full_name, selected_group))

                if params_to_plot:
                    # Create boxplot slides - 5 boxplots per slide
                    boxplots_per_slide = 5
                    num_slides = (len(params_to_plot) + boxplots_per_slide - 1) // boxplots_per_slide

                    for slide_num in range(num_slides):
                        slide = prs.slides.add_slide(content_slide_layout)

                        # Get parameters for this slide
                        start_idx = slide_num * boxplots_per_slide
                        end_idx = min(start_idx + boxplots_per_slide, len(params_to_plot))
                        slide_params = params_to_plot[start_idx:end_idx]

                        # Slide title
                        title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.1), Inches(12.733), Inches(0.5))
                        title_frame = title_box.text_frame
                        title_para = title_frame.paragraphs[0]
                        title_para.text = f"Parameter Boxplots: {selected_group} (Slide {slide_num + 1}/{num_slides})"
                        title_para.font.size = Pt(22)
                        title_para.font.bold = True

                        # Create figure with subplots for boxplots - professional styling
                        num_params = len(slide_params)
                        fig, axes = plt.subplots(1, num_params, figsize=(2.8 * num_params, 5.5))

                        # Set professional style
                        plt.style.use('seaborn-v0_8-whitegrid')

                        if num_params == 1:
                            axes = [axes]

                        for idx, (test_num, short_name, full_name, group_name) in enumerate(slide_params):
                            ax = axes[idx]
                            param_column = test_num

                            # Collect data from all sources
                            all_data = []
                            labels = []
                            for df, label in zip(data_sources, wafer_labels):
                                if param_column in df.columns:
                                    values = df[param_column].dropna().values
                                    if len(values) > 0:
                                        all_data.append(values)
                                        short_label = str(label)[:10] + "..." if len(str(label)) > 10 else str(label)
                                        labels.append(short_label)

                            if all_data:
                                # Professional boxplot styling
                                bp = ax.boxplot(
                                    all_data,
                                    tick_labels=labels if len(labels) <= 3 else [f"W{i+1}" for i in range(len(labels))],
                                    vert=True,
                                    patch_artist=True,
                                    showmeans=True,
                                    notch=False,
                                    widths=0.6,
                                    meanprops=dict(marker="D", markerfacecolor="#E74C3C", markeredgecolor="white", markersize=6, markeredgewidth=1),
                                    medianprops=dict(color="#2C3E50", linewidth=2),
                                    whiskerprops=dict(color="#2C3E50", linewidth=1.5, linestyle='-'),
                                    capprops=dict(color="#2C3E50", linewidth=1.5),
                                    flierprops=dict(marker='o', markerfacecolor='#95A5A6', markeredgecolor='#7F8C8D', markersize=4, alpha=0.6),
                                    boxprops=dict(linewidth=1.5)
                                )

                                # Professional color palette
                                professional_colors = ['#3498DB', '#2ECC71', '#9B59B6', '#E67E22', '#1ABC9C', '#E74C3C', '#34495E', '#F39C12']
                                for box_idx, patch in enumerate(bp["boxes"]):
                                    color = professional_colors[box_idx % len(professional_colors)]
                                    patch.set_facecolor(color)
                                    patch.set_alpha(0.75)
                                    patch.set_edgecolor('#2C3E50')

                                # Calculate and display statistics in top left corner
                                # Combine all data for overall statistics
                                combined_data = np.concatenate(all_data)
                                stats_max = np.max(combined_data)
                                stats_min = np.min(combined_data)
                                stats_q1 = np.percentile(combined_data, 25)
                                stats_q3 = np.percentile(combined_data, 75)
                                stats_mean = np.mean(combined_data)
                                stats_median = np.median(combined_data)

                                # Format statistics text
                                stats_text = (
                                    f"Max: {stats_max:.3g}\n"
                                    f"Min: {stats_min:.3g}\n"
                                    f"Q1: {stats_q1:.3g}\n"
                                    f"Q3: {stats_q3:.3g}\n"
                                    f"Mean: {stats_mean:.3g}\n"
                                    f"Median: {stats_median:.3g}"
                                )

                                # Add statistics box in top left corner
                                ax.text(0.02, 0.98, stats_text,
                                    transform=ax.transAxes,
                                    fontsize=6,
                                    fontweight='normal',
                                    fontfamily='monospace',
                                    verticalalignment='top',
                                    horizontalalignment='left',
                                    bbox=dict(
                                        boxstyle='round,pad=0.3',
                                        facecolor='white',
                                        edgecolor='#BDC3C7',
                                        alpha=0.9,
                                        linewidth=1
                                    )
                                )

                            # Truncate title if too long
                            display_name = short_name if len(short_name) <= 20 else short_name[:17] + "..."
                            ax.set_title(display_name, fontsize=9, fontweight="bold", color='#2C3E50', pad=8)
                            ax.tick_params(axis='both', which='major', labelsize=7, colors='#2C3E50')
                            ax.set_facecolor('#FAFAFA')
                            ax.grid(True, alpha=0.4, linestyle='-', linewidth=0.5, color='#BDC3C7')
                            ax.spines['top'].set_visible(False)
                            ax.spines['right'].set_visible(False)
                            ax.spines['left'].set_color('#BDC3C7')
                            ax.spines['bottom'].set_color('#BDC3C7')

                            # Rotate x labels if needed
                            if len(labels) > 1:
                                ax.set_xticklabels(ax.get_xticklabels(), rotation=45, ha='right')

                        fig.tight_layout()
                        fig.patch.set_facecolor('white')

                        # Add figure to slide
                        img_stream = figure_to_image_bytes(fig, dpi=150)
                        # Center the image on the slide
                        img_width = min(12.5, 2.5 * num_params)
                        x_pos = (13.333 - img_width) / 2
                        slide.shapes.add_picture(img_stream, Inches(x_pos), Inches(0.7), width=Inches(img_width))
                        plt.close(fig)

                    print(f"Added {num_slides} boxplot slide(s) for {len(params_to_plot)} parameters from '{selected_group}'")

        # ============ Gage R&R Slides ============
        if include_grr and grr_analysis_results and '_meta' in grr_analysis_results:
            meta = grr_analysis_results['_meta']

            # Filter out metadata key to get only parameter results
            param_results = {k: v for k, v in grr_analysis_results.items() if k != '_meta'}

            if param_results:
                # Calculate overall assessment summary for context
                excellent_count = sum(1 for r in param_results.values() if r['assessment'] == 'EXCELLENT')
                acceptable_count = sum(1 for r in param_results.values() if r['assessment'] == 'ACCEPTABLE')
                unacceptable_count = sum(1 for r in param_results.values() if r['assessment'] == 'UNACCEPTABLE')
                total_params = len(param_results)

                # Find worst parameters for context
                worst_params = [(p, r['grr_percent']) for p, r in param_results.items() if r['assessment'] == 'UNACCEPTABLE']
                worst_params.sort(key=lambda x: x[1], reverse=True)

                # Generate conclusion text
                if unacceptable_count == 0:
                    overall_status = "PASS"
                    status_color = '#4CAF50'
                    if excellent_count == total_params:
                        conclusion_text = "✓ All parameters show EXCELLENT measurement system capability. No action required."
                    else:
                        conclusion_text = f"✓ Measurement system is acceptable. {excellent_count} excellent, {acceptable_count} acceptable parameters."
                else:
                    overall_status = "FAIL"
                    status_color = '#F44336'
                    failed_names = ", ".join([p[:25] for p, _ in worst_params[:3]])
                    conclusion_text = f"✗ Gage R&R FAILED for {unacceptable_count} parameter(s) due to excessive measurement variation (>30%).\n"
                    conclusion_text += f"   Failing parameters: {failed_names}"
                    if len(worst_params) > 3:
                        conclusion_text += f" (+{len(worst_params)-3} more)"
                    conclusion_text += "\n   Recommendation: Review measurement equipment, operator training, or test methodology."

                # ==================== SLIDE 1: Overview with Wafermaps ====================
                slide = prs.slides.add_slide(content_slide_layout)

                # Title
                title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(12.733), Inches(0.5))
                title_frame = title_box.text_frame
                title_para = title_frame.paragraphs[0]
                title_para.text = "Gage R&R Study Overview"
                title_para.font.size = Pt(28)
                title_para.font.bold = True

                # LEFT: Wafermap visualization
                if grr_file_data and len(grr_file_data) > 0:
                    num_files = min(len(grr_file_data), 4)
                    fig_wafers, axes = plt.subplots(1, num_files, figsize=(2.5 * num_files, 2.5))
                    if num_files == 1:
                        axes = [axes]

                    for idx, (ax, file_info) in enumerate(zip(axes, grr_file_data[:4])):
                        df = file_info.get('data')
                        file_name = os.path.basename(file_info.get('path', f'File {idx+1}'))[:15]

                        if df is not None and hasattr(df, 'columns') and 'x' in df.columns and 'y' in df.columns:
                            plot_data = df[df['x'].notna() & df['y'].notna()]
                            if len(plot_data) > 0:
                                ax.scatter(plot_data['x'], plot_data['y'], c='#B3E5FC', s=15, edgecolors='gray', linewidth=0.2)
                                for die_coord in meta.get('die_coords', []):
                                    x, y = die_coord
                                    if ((plot_data['x'] == x) & (plot_data['y'] == y)).any():
                                        ax.scatter([x], [y], c='red', s=25, edgecolors='black', linewidth=1)
                                ax.set_aspect('equal')
                        ax.set_title(file_name, fontsize=7, fontweight='bold')
                        ax.tick_params(axis='both', labelsize=5)

                    fig_wafers.suptitle(f'Selected Dies: {len(meta.get("die_coords", []))} locations (red)', fontsize=8, y=0.02)
                    fig_wafers.tight_layout()
                    img_stream = figure_to_image_bytes(fig_wafers, dpi=200)
                    slide.shapes.add_picture(img_stream, Inches(0.3), Inches(0.75), width=Inches(6.0))
                    plt.close(fig_wafers)

                # RIGHT: Study info and parameters
                info_box = slide.shapes.add_textbox(Inches(6.5), Inches(0.75), Inches(6.5), Inches(2.8))
                tf = info_box.text_frame
                tf.word_wrap = True

                p = tf.paragraphs[0]
                p.text = "Study Configuration"
                p.font.size = Pt(14)
                p.font.bold = True

                p = tf.add_paragraph()
                p.text = f"• Files/Trials: {meta['num_files']}   • Dies/Parts: {meta['num_dies']}   • Parameters: {meta['num_params']}"
                p.font.size = Pt(10)

                p = tf.add_paragraph()
                p.text = f"\nParameters Analyzed:"
                p.font.size = Pt(12)
                p.font.bold = True

                for i, param in enumerate(list(param_results.keys())[:6]):
                    p = tf.add_paragraph()
                    p.text = f"  {i+1}. {param[:45]}"
                    p.font.size = Pt(9)
                    p.font.name = "Consolas"
                if len(param_results) > 6:
                    p = tf.add_paragraph()
                    p.text = f"  ... +{len(param_results)-6} more"
                    p.font.size = Pt(9)

                # BOTTOM: Conclusion box
                conclusion_box = slide.shapes.add_textbox(Inches(0.3), Inches(3.8), Inches(12.7), Inches(1.5))
                tf = conclusion_box.text_frame
                tf.word_wrap = True

                p = tf.paragraphs[0]
                p.text = f"Overall Result: {overall_status}"
                p.font.size = Pt(16)
                p.font.bold = True
                p.font.color.rgb = RGBColor(244, 67, 54) if overall_status == "FAIL" else RGBColor(76, 175, 80)

                p = tf.add_paragraph()
                p.text = conclusion_text
                p.font.size = Pt(11)

                # Assessment summary
                p = tf.add_paragraph()
                p.text = f"\nAssessment: {excellent_count} Excellent | {acceptable_count} Acceptable | {unacceptable_count} Unacceptable"
                p.font.size = Pt(10)
                p.font.color.rgb = RGBColor(100, 100, 100)

                # ==================== SLIDE 2: GRR Summary Table ====================
                if pptx_grr_summary_var.get():
                    slide = prs.slides.add_slide(content_slide_layout)

                    # Slide title
                    title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(12.733), Inches(0.6))
                    title_frame = title_box.text_frame
                    title_para = title_frame.paragraphs[0]
                    title_para.text = "Gage R&R Analysis Summary"
                    title_para.font.size = Pt(28)
                    title_para.font.bold = True

                    # Study info
                    info_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.9), Inches(5), Inches(1.2))
                    info_frame = info_box.text_frame
                    info_text = f"Study Information:\n"
                    info_text += f"  • Files (Operators/Trials): {meta['num_files']}\n"
                    info_text += f"  • Dies (Parts): {meta['num_dies']}\n"
                    info_text += f"  • Parameters Analyzed: {meta['num_params']}"
                    info_para = info_frame.paragraphs[0]
                    info_para.text = info_text
                    info_para.font.size = Pt(12)
                    info_para.font.name = "Consolas"

                    # Create summary table figure
                    fig_summary, ax_summary = plt.subplots(figsize=(12, 5))
                    ax_summary.axis('off')

                    # Prepare table data
                    table_data = [['Parameter', 'Mean', 'Std Dev', '%GRR', 'ndc', 'Assessment']]
                    cell_colors = [['#E3F2FD'] * 6]

                    for param, res in param_results.items():
                        param_short = param[:35] if len(param) > 35 else param
                        assessment = res['assessment']

                        # Color based on assessment
                        if assessment == 'EXCELLENT':
                            row_color = '#E8F5E9'  # Light green
                        elif assessment == 'ACCEPTABLE':
                            row_color = '#FFF8E1'  # Light yellow
                        else:
                            row_color = '#FFEBEE'  # Light red

                        table_data.append([
                            param_short,
                            f"{res['overall_mean']:.4g}",
                            f"{res['overall_std']:.4g}",
                            f"{res['grr_percent']:.2f}%",
                            f"{res['ndc']:.2f}",
                            assessment
                        ])
                        cell_colors.append([row_color] * 6)

                    table = ax_summary.table(
                        cellText=table_data,
                        cellColours=cell_colors,
                        loc='center',
                        cellLoc='center'
                    )
                    table.auto_set_font_size(False)
                    table.set_fontsize(9)
                    table.scale(1.2, 1.8)

                    # Style header row
                    for j in range(6):
                        table[(0, j)].set_text_props(fontweight='bold')

                    fig_summary.tight_layout()

                    img_stream = figure_to_image_bytes(fig_summary, dpi=200)
                    slide.shapes.add_picture(img_stream, Inches(0.5), Inches(2.2), width=Inches(12.3))
                    plt.close(fig_summary)

                # Slide 2: Variance Components Chart
                if pptx_grr_variance_var.get() and len(param_results) > 0:
                    slide = prs.slides.add_slide(content_slide_layout)

                    # Slide title
                    title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(12.733), Inches(0.5))
                    title_frame = title_box.text_frame
                    title_para = title_frame.paragraphs[0]
                    title_para.text = "Gage R&R - Variance Components Analysis"
                    title_para.font.size = Pt(26)
                    title_para.font.bold = True

                    # Create variance chart
                    fig_var, ax_var = plt.subplots(figsize=(10, 4.5))

                    params = list(param_results.keys())
                    x_pos = np.arange(len(params))

                    repeatability = [param_results[p]['repeatability_percent'] for p in params]
                    reproducibility = [param_results[p]['reproducibility_percent'] for p in params]

                    width = 0.35
                    bars1 = ax_var.bar(x_pos - width/2, repeatability, width, label='Repeatability %', color='#3498DB', alpha=0.8)
                    bars2 = ax_var.bar(x_pos + width/2, reproducibility, width, label='Reproducibility %', color='#E74C3C', alpha=0.8)

                    ax_var.set_ylabel('% of Total Variation', fontsize=11)
                    ax_var.set_xlabel('Parameter', fontsize=11)
                    ax_var.set_title('Variance Components by Parameter', fontsize=14, fontweight='bold')
                    ax_var.set_xticks(x_pos)

                    # Truncate parameter names for x-axis
                    param_labels = [p[:15] + '...' if len(p) > 15 else p for p in params]
                    ax_var.set_xticklabels(param_labels, rotation=45, ha='right', fontsize=8)
                    ax_var.legend(loc='upper right')
                    ax_var.grid(True, alpha=0.3, axis='y')

                    # Add reference lines
                    ax_var.axhline(y=10, color='green', linestyle='--', linewidth=1, alpha=0.7, label='10% (Excellent)')
                    ax_var.axhline(y=30, color='orange', linestyle='--', linewidth=1, alpha=0.7, label='30% (Acceptable)')

                    fig_var.tight_layout()

                    img_stream = figure_to_image_bytes(fig_var, dpi=200)
                    slide.shapes.add_picture(img_stream, Inches(0.5), Inches(0.7), width=Inches(12.3), height=Inches(4.5))
                    plt.close(fig_var)

                    # Context box at bottom
                    context_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.4), Inches(12.3), Inches(1.8))
                    tf = context_box.text_frame
                    tf.word_wrap = True

                    # Find which component is dominant
                    avg_repeat = np.mean(repeatability)
                    avg_reprod = np.mean(reproducibility)

                    p = tf.paragraphs[0]
                    p.text = "Interpretation:"
                    p.font.size = Pt(12)
                    p.font.bold = True

                    p = tf.add_paragraph()
                    p.text = "• Repeatability (blue): Variation when same operator measures same part multiple times (equipment variation)"
                    p.font.size = Pt(10)

                    p = tf.add_paragraph()
                    p.text = "• Reproducibility (red): Variation between different operators/trials measuring same part (operator variation)"
                    p.font.size = Pt(10)

                    p = tf.add_paragraph()
                    if avg_repeat > avg_reprod:
                        p.text = f"→ Analysis: Repeatability ({avg_repeat:.1f}%) > Reproducibility ({avg_reprod:.1f}%). Equipment/method variation is dominant. Consider equipment calibration or method refinement."
                    else:
                        p.text = f"→ Analysis: Reproducibility ({avg_reprod:.1f}%) > Repeatability ({avg_repeat:.1f}%). Operator/trial variation is dominant. Consider operator training or standardization."
                    p.font.size = Pt(10)
                    p.font.italic = True
                    p.font.color.rgb = RGBColor(0, 100, 150)

                # Slide 3: %GRR and ndc Chart
                if pptx_grr_charts_var.get() and len(param_results) > 0:
                    slide = prs.slides.add_slide(content_slide_layout)

                    # Slide title
                    title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(12.733), Inches(0.6))
                    title_frame = title_box.text_frame
                    title_para = title_frame.paragraphs[0]
                    title_para.text = "Gage R&R - %GRR and ndc Analysis"
                    title_para.font.size = Pt(28)
                    title_para.font.bold = True

                    # Create figure with two subplots
                    fig_charts, (ax_grr, ax_ndc) = plt.subplots(1, 2, figsize=(12, 5))

                    params = list(param_results.keys())
                    grr_values = [param_results[p]['grr_percent'] for p in params]
                    ndc_values = [param_results[p]['ndc'] for p in params]

                    # %GRR bar chart with color coding
                    colors_grr = []
                    for grr in grr_values:
                        if grr < 10:
                            colors_grr.append('#4CAF50')  # Green - Excellent
                        elif grr < 30:
                            colors_grr.append('#FFC107')  # Yellow - Acceptable
                        else:
                            colors_grr.append('#F44336')  # Red - Unacceptable

                    x_pos = np.arange(len(params))
                    bars_grr = ax_grr.bar(x_pos, grr_values, color=colors_grr, alpha=0.8, edgecolor='black', linewidth=0.5)

                    # Add reference lines
                    ax_grr.axhline(y=10, color='green', linestyle='--', linewidth=1.5, label='10% (Excellent)')
                    ax_grr.axhline(y=30, color='red', linestyle='--', linewidth=1.5, label='30% (Acceptable limit)')

                    ax_grr.set_ylabel('%GRR', fontsize=11)
                    ax_grr.set_xlabel('Parameter', fontsize=11)
                    ax_grr.set_title('%GRR by Parameter', fontsize=12, fontweight='bold')
                    ax_grr.set_xticks(x_pos)
                    param_labels = [p[:12] + '...' if len(p) > 12 else p for p in params]
                    ax_grr.set_xticklabels(param_labels, rotation=45, ha='right', fontsize=7)
                    ax_grr.legend(loc='upper right', fontsize=8)
                    ax_grr.grid(True, alpha=0.3, axis='y')

                    # Add value labels on bars
                    for bar, val in zip(bars_grr, grr_values):
                        ax_grr.annotate(f'{val:.1f}%', xy=(bar.get_x() + bar.get_width()/2, bar.get_height()),
                                       xytext=(0, 3), textcoords='offset points', ha='center', va='bottom', fontsize=7)

                    # ndc bar chart
                    colors_ndc = ['#4CAF50' if n >= 5 else '#FFC107' if n >= 2 else '#F44336' for n in ndc_values]
                    bars_ndc = ax_ndc.bar(x_pos, ndc_values, color=colors_ndc, alpha=0.8, edgecolor='black', linewidth=0.5)

                    ax_ndc.axhline(y=5, color='green', linestyle='--', linewidth=1.5, label='ndc ≥ 5 (Good)')
                    ax_ndc.axhline(y=2, color='orange', linestyle='--', linewidth=1.5, label='ndc ≥ 2 (Minimum)')

                    ax_ndc.set_ylabel('ndc (Number of Distinct Categories)', fontsize=11)
                    ax_ndc.set_xlabel('Parameter', fontsize=11)
                    ax_ndc.set_title('ndc by Parameter', fontsize=12, fontweight='bold')
                    ax_ndc.set_xticks(x_pos)
                    ax_ndc.set_xticklabels(param_labels, rotation=45, ha='right', fontsize=7)
                    ax_ndc.legend(loc='upper right', fontsize=8)
                    ax_ndc.grid(True, alpha=0.3, axis='y')

                    # Add value labels on bars
                    for bar, val in zip(bars_ndc, ndc_values):
                        ax_ndc.annotate(f'{val:.2f}', xy=(bar.get_x() + bar.get_width()/2, bar.get_height()),
                                       xytext=(0, 3), textcoords='offset points', ha='center', va='bottom', fontsize=7)

                    fig_charts.tight_layout()

                    img_stream = figure_to_image_bytes(fig_charts, dpi=200)
                    slide.shapes.add_picture(img_stream, Inches(0.5), Inches(0.7), width=Inches(12.3), height=Inches(4.3))
                    plt.close(fig_charts)

                    # Context box at bottom
                    context_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(12.3), Inches(2.0))
                    tf = context_box.text_frame
                    tf.word_wrap = True

                    p = tf.paragraphs[0]
                    p.text = "Acceptance Criteria & Interpretation:"
                    p.font.size = Pt(12)
                    p.font.bold = True

                    p = tf.add_paragraph()
                    p.text = "• %GRR < 10%: EXCELLENT - Measurement system is excellent for this application"
                    p.font.size = Pt(9)
                    p.font.color.rgb = RGBColor(76, 175, 80)

                    p = tf.add_paragraph()
                    p.text = "• %GRR 10-30%: ACCEPTABLE - May be acceptable depending on application criticality"
                    p.font.size = Pt(9)
                    p.font.color.rgb = RGBColor(255, 152, 0)

                    p = tf.add_paragraph()
                    p.text = "• %GRR > 30%: UNACCEPTABLE - Measurement system needs improvement"
                    p.font.size = Pt(9)
                    p.font.color.rgb = RGBColor(244, 67, 54)

                    p = tf.add_paragraph()
                    p.text = f"\n• ndc (Number of Distinct Categories): Indicates how many distinct categories the measurement can reliably distinguish. ndc ≥ 5 is recommended."
                    p.font.size = Pt(9)

                    # Add specific findings
                    low_ndc_params = [p for p, r in param_results.items() if r['ndc'] < 5]
                    if low_ndc_params:
                        p = tf.add_paragraph()
                        p.text = f"⚠ Warning: {len(low_ndc_params)} parameter(s) have ndc < 5: {', '.join([p[:20] for p in low_ndc_params[:3]])}"
                        p.font.size = Pt(9)
                        p.font.color.rgb = RGBColor(255, 152, 0)

                print(f"Added Gage R&R slides for {len(param_results)} parameters")

        current_step += 1
        pptx_progress_var.set((current_step / total_steps) * 100)
        main_win.update()

        # Save presentation locally first
        prs.save(save_path)
        print(f"PowerPoint presentation saved to: {save_path}")

        # Handle different save methods
        if save_method == "local":
            pptx_progress_var.set(100)
            pptx_status_label.config(text=f"✅ PowerPoint saved to: {os.path.basename(save_path)}", fg="green")

        elif save_method in ["google_drive", "google_slides"]:
            # Upload to Google Drive
            pptx_status_label.config(text="☁️ Uploading to Google Drive...", fg="blue")
            main_win.update()

            convert_to_slides = (save_method == "google_slides")
            file_name = f"{pptx_title_var.get().replace(' ', '_')}.pptx"

            file_id, result = upload_to_google_drive(save_path, file_name, convert_to_slides)

            if file_id:
                # Store URL for "Open Last" button
                open_last_presentation.url = result

                # Apply Gemini beautification if requested
                if use_gemini and convert_to_slides:
                    pptx_status_label.config(text="✨ Applying Gemini AI beautification...", fg="blue")
                    main_win.update()

                    # Apply theme
                    theme_success, theme_msg = apply_gemini_theme_to_slides(file_id, theme_style)

                    # Get Gemini suggestions
                    api_key = pptx_gemini_key_var.get() or None
                    suggestions, error = beautify_slides_with_gemini(file_id, api_key)

                    if suggestions:
                        # Show suggestions in a dialog
                        suggestion_window = tk.Toplevel(main_win)
                        suggestion_window.title("✨ Gemini AI Suggestions")
                        suggestion_window.geometry("600x400")

                        tk.Label(suggestion_window, text="Gemini AI Beautification Suggestions:",
                                font=("Helvetica", 11, "bold")).pack(pady=10)

                        text_widget = tk.Text(suggestion_window, wrap=tk.WORD, font=("Helvetica", 10))
                        text_widget.pack(fill=tk.BOTH, expand=True, padx=10, pady=5)
                        text_widget.insert('1.0', suggestions)
                        text_widget.config(state='disabled')

                        tk.Button(suggestion_window, text="Open Google Slides",
                                 command=lambda: __import__('webbrowser').open(result),
                                 bg="#4CAF50", fg="white").pack(pady=10)

                        pptx_status_label.config(
                            text=f"✅ Uploaded to Google {'Slides' if convert_to_slides else 'Drive'} with AI suggestions!",
                            fg="green"
                        )
                    elif error:
                        pptx_status_label.config(
                            text=f"✅ Uploaded but Gemini error: {error[:50]}",
                            fg="orange"
                        )
                    else:
                        pptx_status_label.config(
                            text=f"✅ Uploaded to Google {'Slides' if convert_to_slides else 'Drive'} with {theme_style} theme!",
                            fg="green"
                        )
                else:
                    pptx_status_label.config(
                        text=f"✅ Uploaded to Google {'Slides' if convert_to_slides else 'Drive'}!",
                        fg="green"
                    )

                # Ask to open in browser
                if tk.messagebox.askyesno("Success", f"Presentation uploaded successfully!\n\nOpen in browser?"):
                    import webbrowser
                    webbrowser.open(result)

            else:
                pptx_status_label.config(text=f"❌ Upload failed: {result}", fg="red")

            # Clean up temp file if not local save
            if save_method != "local" and os.path.exists(save_path):
                try:
                    os.remove(save_path)
                except:
                    pass

        pptx_progress_var.set(100)

    except Exception as e:
        pptx_status_label.config(text=f"Error: {str(e)}", fg="red")
        print(f"Error creating PowerPoint: {e}")
        import traceback
        traceback.print_exc()


# Initialize wafer listbox when tab is selected
def on_presentation_tab_selected(event):
    """Update wafer listbox when presentation tab is selected"""
    selected_tab = event.widget.select()
    tab_text = event.widget.tab(selected_tab, "text")
    if tab_text == "Create Presentation":
        update_wafer_listbox()
        update_pptx_group_combobox()


notebook.bind("<<NotebookTabChanged>>", on_presentation_tab_selected)

# Initialize file source buttons (show STDF buttons by default)
update_source_buttons()

main_win.mainloop()
