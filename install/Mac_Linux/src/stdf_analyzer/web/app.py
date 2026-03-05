"""
STDF Wafermap Analyzer - Web Application (Full Version)

A comprehensive Streamlit-based web interface for analyzing semiconductor test data.
Includes: Wafermap, Heatmap with Limits, Distribution, Boxplot, Statistics, and PowerPoint Export.

Usage:
    streamlit run app.py
"""

import os
import sys
import io
import tempfile
from pathlib import Path
from datetime import datetime
from typing import List, Optional, Tuple, Dict, Any

import streamlit as st
import pandas as pd
import numpy as np

# Add the src directory to the path for imports
src_path = Path(__file__).parent.parent.parent.parent
if str(src_path) not in sys.path:
    sys.path.insert(0, str(src_path))

from src.stdf_analyzer.core.binning import BinningLookup, BIN_COLORS
from src.stdf_analyzer.core.stdf_parser import parse_stdf_file, parse_csv_file, STDFData, TestLimits
from src.stdf_analyzer.core.wafermap import WafermapGenerator

# ============================================================================
# Configuration
# ============================================================================

DEFAULT_SHARE_PATHS = [
    r"C:\Users\szenklarz\Desktop\VS_Folder\AM Data",
    r"C:\Users\szenklarz\Desktop\VS_Folder\Tooltest",
]

SUPPORTED_EXTENSIONS = ['.stdf', '.csv']

PARAM_GROUP_PATTERNS = {
    'DC_LKG': ['DC_LKG', 'LEAKAGE', 'LKG_', 'ILKG', 'IDDQ'],
    'DC_CONT': ['DC_CONT', 'CONTINUITY', 'CONT_', 'OPEN', 'SHORT'],
    'DC_VTH': ['DC_VTH', 'VTH_', 'THRESHOLD', 'VT_'],
    'POWER': ['POWER', 'PWR_', 'IDD', 'ICC', 'CURRENT'],
    'ANALOG': ['ANALOG', 'ADC', 'DAC', 'VREF'],
}

# ============================================================================
# Page Configuration
# ============================================================================

st.set_page_config(
    page_title="STDF Wafermap Analyzer",
    page_icon="🔬",
    layout="wide",
    initial_sidebar_state="expanded"
)

st.markdown("""
<style>
    .main-header { font-size: 2.2rem; font-weight: bold; color: #1f77b4; }
    .stat-good { color: #28a745; font-weight: bold; }
    .stat-bad { color: #dc3545; font-weight: bold; }
    div[data-testid="stMetricValue"] { font-size: 1.5rem; }
</style>
""", unsafe_allow_html=True)

# ============================================================================
# Session State
# ============================================================================

def init_session_state():
    defaults = {
        'loaded_files': {},
        'current_file': None,
        'binning_lookup': BinningLookup(),
        'current_path': None,
        'data_source': "network",
    }
    for key, value in defaults.items():
        if key not in st.session_state:
            st.session_state[key] = value

# ============================================================================
# Utility Functions
# ============================================================================

def extract_param_group(param_name: str) -> str:
    if not param_name:
        return "Other"
    param_upper = param_name.upper()
    for group, patterns in PARAM_GROUP_PATTERNS.items():
        for pattern in patterns:
            if pattern in param_upper:
                return group
    return "Other"

def get_grouped_parameters(data: STDFData) -> Dict[str, List[str]]:
    groups = {}
    for param_key, param_name in data.test_parameters.items():
        group = extract_param_group(param_name)
        if group not in groups:
            groups[group] = []
        groups[group].append((param_key, param_name))
    return groups

def is_valid_path(path: str) -> bool:
    try:
        return os.path.exists(path)
    except:
        return False

def get_folder_contents(path: str):
    folders, files = [], []
    try:
        for item in os.listdir(path):
            item_path = os.path.join(path, item)
            try:
                if os.path.isdir(item_path):
                    folders.append(item)
                elif os.path.isfile(item_path):
                    ext = os.path.splitext(item)[1].lower()
                    if ext in SUPPORTED_EXTENSIONS:
                        stat = os.stat(item_path)
                        files.append({
                            'name': item, 'path': item_path,
                            'size': stat.st_size, 'extension': ext
                        })
            except:
                continue
    except Exception as e:
        st.error(f"Fehler: {e}")
    return sorted(folders), sorted(files, key=lambda x: x['name'])

def format_size(size):
    for unit in ['B', 'KB', 'MB', 'GB']:
        if size < 1024:
            return f"{size:.1f} {unit}"
        size /= 1024

def load_file(file_path: str):
    ext = os.path.splitext(file_path)[1].lower()
    try:
        if ext == '.stdf':
            return parse_stdf_file(file_path, verbose=False)
        elif ext == '.csv':
            return parse_csv_file(file_path)
    except Exception as e:
        st.error(f"Fehler: {e}")
    return None

def load_uploaded(uploaded):
    ext = Path(uploaded.name).suffix.lower()
    with tempfile.NamedTemporaryFile(delete=False, suffix=ext) as tmp:
        tmp.write(uploaded.getvalue())
        tmp_path = tmp.name
    try:
        return load_file(tmp_path)
    finally:
        os.unlink(tmp_path)

# ============================================================================
# Visualization Functions
# ============================================================================

def create_heatmap_with_limits(data: STDFData, param_key: str):
    import plotly.graph_objects as go

    df = data.dataframe
    if param_key not in df.columns:
        return None

    values = df[param_key].dropna()
    x_coords = df.loc[values.index, 'x']
    y_coords = df.loc[values.index, 'y']

    limits = data.test_limits.get(param_key)

    if limits and (limits.lo_limit is not None or limits.hi_limit is not None):
        colors = []
        for v in values:
            fail = False
            if limits.lo_limit is not None and v < limits.lo_limit:
                fail = True
            if limits.hi_limit is not None and v > limits.hi_limit:
                fail = True
            colors.append('#dc3545' if fail else '#28a745')

        fig = go.Figure(go.Scatter(
            x=x_coords, y=y_coords, mode='markers',
            marker=dict(color=colors, size=12, line=dict(width=0.5, color='black')),
            text=[f"Value: {v:.4f}" if isinstance(v, (int, float)) else f"Value: {v}" for v in values],
            hovertemplate="X: %{x}<br>Y: %{y}<br>%{text}<extra></extra>"
        ))
    else:
        fig.add_trace(go.Scatter(
            x=x_coords, y=y_coords, mode='markers',
            marker=dict(color=values, colorscale='Viridis', size=12,
                       colorbar=dict(title=param_key), line=dict(width=0.5, color='black')),
            text=[f"Value: {v:.4f}" if isinstance(v, (int, float)) else f"Value: {v}" for v in values],
            hovertemplate="X: %{x}<br>Y: %{y}<br>%{text}<extra></extra>"
        ))

    title = f"Heatmap - {data.wafer_id}<br>{param_key}"
    if limits:
        parts = []
        if limits.lo_limit is not None:
            parts.append(f"LO: {limits.lo_limit}")
        if limits.hi_limit is not None:
            parts.append(f"HI: {limits.hi_limit}")
        if parts:
            title += f"<br><sub>Limits: {', '.join(parts)} {limits.units}</sub>"

    fig.update_layout(title=title, xaxis_title="X", yaxis_title="Y",
                     yaxis=dict(scaleanchor="x"), width=650, height=650)
    return fig

def create_distribution(data: STDFData, param_key: str):
    import plotly.graph_objects as go
    from scipy import stats

    df = data.dataframe
    if param_key not in df.columns:
        return None

    values = df[param_key].dropna()
    if len(values) == 0:
        return None

    mean, std = values.mean(), values.std()
    limits = data.test_limits.get(param_key)

    fig = go.Figure()
    fig.add_trace(go.Histogram(x=values, nbinsx=50, name='Distribution', opacity=0.7))

    if std > 0:
        x_range = np.linspace(values.min(), values.max(), 100)
        y_normal = stats.norm.pdf(x_range, mean, std) * len(values) * (values.max() - values.min()) / 50
        fig.add_trace(go.Scatter(x=x_range, y=y_normal, mode='lines', name='Normal', line=dict(color='red', width=2)))

    if limits:
        if limits.lo_limit is not None:
            fig.add_vline(x=limits.lo_limit, line_dash="dash", line_color="red", annotation_text=f"LO")
        if limits.hi_limit is not None:
            fig.add_vline(x=limits.hi_limit, line_dash="dash", line_color="red", annotation_text=f"HI")

    fig.add_vline(x=mean, line_dash="dot", line_color="green", annotation_text="μ")
    fig.add_vline(x=mean - 3*std, line_dash="dot", line_color="orange", annotation_text="-3σ")
    fig.add_vline(x=mean + 3*std, line_dash="dot", line_color="orange", annotation_text="+3σ")

    fig.update_layout(
        title=f"Distribution - {param_key}<br><sub>μ={mean:.4f}, σ={std:.4f}, n={len(values)}</sub>",
        xaxis_title=param_key, yaxis_title="Count", width=800, height=450
    )
    return fig

def create_boxplot(loaded_files: Dict[str, STDFData], param_key: str):
    import plotly.graph_objects as go

    fig = go.Figure()
    for fname, data in loaded_files.items():
        df = data.dataframe
        if param_key not in df.columns:
            continue
        values = df[param_key].dropna()
        if len(values) > 0:
            fig.add_trace(go.Box(y=values, name=data.wafer_id or fname, boxpoints='outliers'))

    fig.update_layout(title=f"Boxplot - {param_key}", yaxis_title=param_key, width=850, height=450)
    return fig

def create_stats_table(data: STDFData, params: List[str] = None) -> pd.DataFrame:
    df = data.dataframe
    if params is None:
        params = [c for c in df.columns if c not in ['x', 'y', 'bin'] and df[c].dtype in ['float64', 'int64']]

    rows = []
    for p in params:
        if p not in df.columns:
            continue
        values = df[p].dropna()
        if len(values) == 0:
            continue

        limits = data.test_limits.get(p)
        lo = limits.lo_limit if limits else None
        hi = limits.hi_limit if limits else None
        units = limits.units if limits else ""

        pass_count = len(values)
        if lo is not None:
            pass_count = sum(values >= lo)
        if hi is not None:
            pass_count = sum((values <= hi) & (values >= (lo or float('-inf'))))

        rows.append({
            'Parameter': p, 'Units': units, 'Count': len(values),
            'Mean': f"{values.mean():.4f}", 'Std': f"{values.std():.4f}",
            'Min': f"{values.min():.4f}", 'Max': f"{values.max():.4f}",
            'LO': lo or "-", 'HI': hi or "-",
            'Pass': pass_count, 'Fail': len(values) - pass_count,
            'Yield': f"{pass_count/len(values)*100:.1f}%"
        })
    return pd.DataFrame(rows)

# ============================================================================
# PowerPoint Export
# ============================================================================

def create_pptx(loaded_files: Dict[str, STDFData], params: List[str]) -> bytes:
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        st.error("pip install python-pptx")
        return None

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Title
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(12), Inches(2))
    p = box.text_frame.paragraphs[0]
    p.text = "STDF Wafermap Analysis Report"
    p.font.size = Pt(44)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    box2 = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(12), Inches(1))
    p2 = box2.text_frame.paragraphs[0]
    p2.text = datetime.now().strftime('%Y-%m-%d %H:%M')
    p2.font.size = Pt(18)
    p2.alignment = PP_ALIGN.CENTER

    for fname, data in loaded_files.items():
        gen = WafermapGenerator(data.dataframe, data.wafer_id)

        # Wafermap slide
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
        box.text_frame.paragraphs[0].text = f"Wafermap - {data.wafer_id or fname}"
        box.text_frame.paragraphs[0].font.size = Pt(28)
        box.text_frame.paragraphs[0].font.bold = True

        fig = gen.create_matplotlib_figure(parameter="bin", figsize=(7, 7))
        buf = io.BytesIO()
        fig.savefig(buf, format='png', dpi=120, bbox_inches='tight')
        buf.seek(0)
        slide.shapes.add_picture(buf, Inches(0.5), Inches(1.2), width=Inches(6))

        info = slide.shapes.add_textbox(Inches(7), Inches(2), Inches(5), Inches(4))
        tf = info.text_frame
        tf.paragraphs[0].text = f"Yield: {gen.get_yield():.1f}%"
        tf.paragraphs[0].font.size = Pt(24)
        tf.paragraphs[0].font.bold = True

        for _, row in gen.get_bin_summary().iterrows():
            p = tf.add_paragraph()
            p.text = f"Bin {int(row['Bin'])}: {row['Count']} ({row['Percentage']}%)"
            p.font.size = Pt(14)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()

# ============================================================================
# UI Components
# ============================================================================

def render_sidebar():
    with st.sidebar:
        st.markdown("## 📊 Datenquelle")
        source = st.radio("Laden von:", ["🌐 Netzwerk", "📤 Upload"], horizontal=True)
        st.session_state.data_source = "network" if "Netzwerk" in source else "upload"

        if st.session_state.loaded_files:
            st.markdown("---")
            st.markdown("## 📋 Geladene Dateien")
            for fname in list(st.session_state.loaded_files.keys()):
                c1, c2 = st.columns([4, 1])
                is_sel = fname == st.session_state.current_file
                with c1:
                    if st.button(f"{'▶' if is_sel else '○'} {fname[:16]}", key=f"s_{fname}", use_container_width=True):
                        st.session_state.current_file = fname
                        st.rerun()
                with c2:
                    if st.button("🗑️", key=f"d_{fname}"):
                        del st.session_state.loaded_files[fname]
                        if st.session_state.current_file == fname:
                            st.session_state.current_file = None
                        st.rerun()

def render_file_browser():
    st.markdown("### 📂 Datei-Browser")

    valid = [p for p in DEFAULT_SHARE_PATHS if is_valid_path(p)]
    if not valid:
        path = st.text_input("Pfad eingeben:")
        if path and st.button("Öffnen") and is_valid_path(path):
            st.session_state.current_path = path
            st.rerun()
        return

    if st.session_state.current_path is None:
        st.session_state.current_path = valid[0]

    current = st.session_state.current_path
    st.caption(f"📍 {current}")

    parent = os.path.dirname(current)
    if parent and parent != current:
        if st.button("⬆️ Hoch"):
            st.session_state.current_path = parent
            st.rerun()

    folders, files = get_folder_contents(current)

    c1, c2 = st.columns([1, 2])
    with c1:
        st.markdown("**📁 Ordner**")
        for f in folders[:15]:
            if st.button(f"📁 {f}", key=f"fo_{f}", use_container_width=True):
                st.session_state.current_path = os.path.join(current, f)
                st.rerun()

    with c2:
        st.markdown("**📄 Dateien**")
        filt = st.text_input("🔍 Filter")
        filtered = [f for f in files if not filt or filt.lower() in f['name'].lower()]
        for f in filtered[:20]:
            cc1, cc2, cc3 = st.columns([3, 1, 1])
            cc1.write(f"📊 {f['name']}")
            cc2.caption(format_size(f['size']))
            if cc3.button("📥", key=f"l_{f['name']}"):
                with st.spinner("Laden..."):
                    data = load_file(f['path'])
                    if data:
                        st.session_state.loaded_files[f['name']] = data
                        st.session_state.current_file = f['name']
                        st.rerun()

def render_upload():
    st.markdown("### 📤 Upload")
    files = st.file_uploader("STDF/CSV", type=['stdf', 'csv'], accept_multiple_files=True)
    for f in files or []:
        if f.name not in st.session_state.loaded_files:
            with st.spinner(f"Lade {f.name}..."):
                data = load_uploaded(f)
                if data:
                    st.session_state.loaded_files[f.name] = data
                    st.session_state.current_file = f.name

def render_wafermap_tab():
    if not st.session_state.current_file:
        st.info("👆 Datei laden")
        return

    data = st.session_state.loaded_files[st.session_state.current_file]
    gen = WafermapGenerator(data.dataframe, data.wafer_id)

    c1, c2, c3, c4 = st.columns(4)
    c1.metric("Wafer", data.wafer_id or "N/A")
    c2.metric("Dies", f"{gen.die_count:,}")
    c3.metric("Yield", f"{gen.get_yield():.1f}%")
    c4.metric("Params", data.parameter_count)

    cc1, cc2 = st.columns([1, 3])
    with cc1:
        params = ["bin"] + list(data.test_parameters.keys())
        param = st.selectbox("Parameter", params)
        size = st.slider("Größe", 5, 25, 10)
    with cc2:
        fig = gen.create_plotly_figure(parameter=param, marker_size=size)
        st.plotly_chart(fig, width='stretch')

    if param == "bin":
        st.dataframe(gen.get_bin_summary(), width='stretch')

def render_heatmap_tab():
    if not st.session_state.current_file:
        st.info("👆 Datei laden")
        return

    data = st.session_state.loaded_files[st.session_state.current_file]
    groups = get_grouped_parameters(data)

    cc1, cc2 = st.columns([1, 3])
    with cc1:
        grp = st.selectbox("Gruppe", ["All"] + list(groups.keys()))
        params = list(data.test_parameters.keys()) if grp == "All" else [p[0] for p in groups.get(grp, [])]
        param = st.selectbox("Parameter", params) if params else None

        if param:
            lim = data.test_limits.get(param)
            if lim:
                st.write(f"**LO:** {lim.lo_limit or '-'}")
                st.write(f"**HI:** {lim.hi_limit or '-'}")
                st.write(f"**Units:** {lim.units or '-'}")
    with cc2:
        if param:
            fig = create_heatmap_with_limits(data, param)
            if fig:
                st.plotly_chart(fig, use_container_width=True)

def render_distribution_tab():
    if not st.session_state.current_file:
        st.info("👆 Datei laden")
        return

    data = st.session_state.loaded_files[st.session_state.current_file]

    cc1, cc2 = st.columns([1, 3])
    with cc1:
        params = list(data.test_parameters.keys())
        param = st.selectbox("Parameter", params, key="dist_p") if params else None
    with cc2:
        if param:
            fig = create_distribution(data, param)
            if fig:
                st.plotly_chart(fig, use_container_width=True)

            vals = data.dataframe[param].dropna()
            if len(vals) > 0:
                mc = st.columns(5)
                mc[0].metric("Mean", f"{vals.mean():.4f}")
                mc[1].metric("Std", f"{vals.std():.4f}")
                mc[2].metric("Min", f"{vals.min():.4f}")
                mc[3].metric("Max", f"{vals.max():.4f}")
                mc[4].metric("n", len(vals))

def render_boxplot_tab():
    if not st.session_state.loaded_files:
        st.info("👆 Dateien laden")
        return

    all_params = set()
    for d in st.session_state.loaded_files.values():
        all_params.update(d.test_parameters.keys())

    cc1, cc2 = st.columns([1, 3])
    with cc1:
        param = st.selectbox("Parameter", sorted(all_params), key="box_p")
        files = st.multiselect("Wafer", list(st.session_state.loaded_files.keys()),
                               default=list(st.session_state.loaded_files.keys()))
    with cc2:
        if files and param:
            sel = {k: v for k, v in st.session_state.loaded_files.items() if k in files}
            fig = create_boxplot(sel, param)
            if fig:
                st.plotly_chart(fig, use_container_width=True)

def render_stats_tab():
    if not st.session_state.current_file:
        st.info("👆 Datei laden")
        return

    data = st.session_state.loaded_files[st.session_state.current_file]
    st.markdown(f"### Statistik: {data.wafer_id}")

    groups = get_grouped_parameters(data)
    grp = st.selectbox("Gruppe", ["All"] + list(groups.keys()), key="stat_grp")
    params = list(data.test_parameters.keys()) if grp == "All" else [p[0] for p in groups.get(grp, [])]

    if params:
        df = create_stats_table(data, params)
        st.dataframe(df, height=500, width='stretch')
        st.download_button("📥 CSV", df.to_csv(index=False), f"stats_{data.wafer_id}.csv")

def render_export_tab():
    st.markdown("### 📊 PowerPoint Export")

    if not st.session_state.loaded_files:
        st.info("👆 Dateien laden")
        return

    files = st.multiselect("Wafer", list(st.session_state.loaded_files.keys()),
                           default=list(st.session_state.loaded_files.keys()), key="exp_f")

    if files:
        data = st.session_state.loaded_files[files[0]]
        params = st.multiselect("Parameter", list(data.test_parameters.keys())[:30],
                                default=list(data.test_parameters.keys())[:3])

    if st.button("🎨 Report erstellen", type="primary"):
        if not files:
            st.error("Wafer auswählen!")
            return
        with st.spinner("Erstelle PowerPoint..."):
            sel = {k: v for k, v in st.session_state.loaded_files.items() if k in files}
            pptx = create_pptx(sel, params)
            if pptx:
                st.download_button("📥 Download PPTX", pptx,
                                  f"Report_{datetime.now().strftime('%Y%m%d_%H%M')}.pptx",
                                  "application/vnd.openxmlformats-officedocument.presentationml.presentation")

def render_data_tab():
    if not st.session_state.current_file:
        st.info("👆 Datei laden")
        return

    data = st.session_state.loaded_files[st.session_state.current_file]
    st.markdown(f"### Rohdaten: {data.wafer_id} ({len(data.dataframe)} Zeilen)")

    cols = st.multiselect("Spalten", data.dataframe.columns.tolist(),
                         default=['x', 'y', 'bin'] if 'bin' in data.dataframe.columns else data.dataframe.columns[:5].tolist())
    if cols:
        st.dataframe(data.dataframe[cols], height=400, width='stretch')

    st.download_button("📥 CSV Export", data.dataframe.to_csv(index=False),
                      f"{data.wafer_id or 'data'}.csv", "text/csv")

# ============================================================================
# Main Application
# ============================================================================

def main():
    init_session_state()

    st.markdown('<p class="main-header">🔬 STDF Wafermap Analyzer</p>', unsafe_allow_html=True)
    st.caption("Halbleiter-Testdaten Analyse und Visualisierung")

    render_sidebar()

    # Data loading section
    if st.session_state.data_source == "network":
        if not st.session_state.loaded_files:
            render_file_browser()
        else:
            tabs = st.tabs(["📂 Dateien", "📊 Wafermap", "🎨 Heatmap", "📈 Distribution",
                           "📦 Boxplot", "📋 Statistik", "🎬 Export", "📄 Rohdaten"])
            with tabs[0]: render_file_browser()
            with tabs[1]: render_wafermap_tab()
            with tabs[2]: render_heatmap_tab()
            with tabs[3]: render_distribution_tab()
            with tabs[4]: render_boxplot_tab()
            with tabs[5]: render_stats_tab()
            with tabs[6]: render_export_tab()
            with tabs[7]: render_data_tab()
    else:
        if not st.session_state.loaded_files:
            render_upload()
        else:
            tabs = st.tabs(["📤 Upload", "📊 Wafermap", "🎨 Heatmap", "📈 Distribution",
                           "📦 Boxplot", "📋 Statistik", "🎬 Export", "📄 Rohdaten"])
            with tabs[0]: render_upload()
            with tabs[1]: render_wafermap_tab()
            with tabs[2]: render_heatmap_tab()
            with tabs[3]: render_distribution_tab()
            with tabs[4]: render_boxplot_tab()
            with tabs[5]: render_stats_tab()
            with tabs[6]: render_export_tab()
            with tabs[7]: render_data_tab()

    st.markdown("---")
    st.caption("STDF Wafermap Analyzer v1.1.0 | Krzysztof Szenklarz")

if __name__ == "__main__":
    main()
